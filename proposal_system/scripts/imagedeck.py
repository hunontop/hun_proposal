"""imagedeck — W28 이미지 렌더 트랙 (D8~D13).

storyline.json(내용 SSOT) + wireframe.json(배치 계약, 선택) + 스킨(디자인 계약)을
장별 이미지 생성 프롬프트로 결정론 조립하고(bundle), Codex가 그린 PNG를 수거·실측 검증하며
(collect), 승인된 이미지를 HTML 크롬(발주처·제안사 로고)과 조합한다(compose).

세 함수는 상태머신을 모른다 — 기존 run 디렉터리에 직접 실행·검증 가능한 순수 엔진이다.
프롬프트 계약은 <개발 원본 전용 경로> 원본을 이식했다.

콘솔 출력 규율(이 프로젝트): cp949-안전 문자만 — em-dash 금지, 하이픈 사용.
"""

from __future__ import annotations

import datetime as dt
import hashlib
import json
import re
import shutil
import struct
from pathlib import Path
from typing import Any

import sys as _sys
_sys.path.insert(0, str(Path(__file__).resolve().parent))
import design_contract  # noqa: E402  (sibling module — W31 R2·R5 run별 디자인 계약)
import gates  # noqa: E402  (sibling module — W31 리허설 마찰2 관문 다이얼, KC 패킷 ③ 프로파일 읽기)
import design_knowledge_cards  # noqa: E402  (sibling module — δ패킷: A6 선택 카드 결정론 운반)
import knowledge_ledger  # noqa: E402  (sibling module — ε패킷: master_design/imagedeck_review 지식 소비)

SCHEMA_VERSION = 1

PROMPTS_DIR = "imagedeck_prompts"
# W31 R6: 사람 편집물(오버레이) 정본 위치 — JOURNEY.md R6. bundle()은 이 폴더를 절대 채우지
# 않는다(빈 폴더 + 안내 README만) - NN.md는 사람 몫이고, 있으면 재번들 때마다 프롬프트 말미에
# 병합 주입한다(재번들에도 항상 살아남는다).
PROMPTS_LOCAL_DIR = "imagedeck_prompts_local"
SLIDES_DIR = "imagedeck/slides"
MANIFEST_NAME = "imagedeck_manifest.json"
COLLECT_NAME = "imagedeck_collect.json"
COLLECT_MD = "imagedeck/collect_report.md"
COMPOSE_HTML = "deck.images.html"
# DF6(DECK_FIRST_DESIGN.md §2-⑦ 경로 B): 마무리 타임 장별 오버라이드 채널 — run 워크스페이스
# 루트에 두는 사람 편집 파일(회귀 방지 목적 - 존재하지 않으면 조립 결과가 이전과 바이트 동일).
DECK_OVERRIDES_NAME = "deck_overrides.json"
# 허용 키만 통과(콘텐츠 우회로 방지 - 값 자체는 storyline 재동결로만 바꾼다).
# style_variant/background_image = 안전 분류 ⓐ(재조립만·본문 px 불변),
# chrome_override/deck_class = 안전 분류 ⓑ(구조 변형·본문 생성 px가 바뀔 수 있음).
_DECK_OVERRIDE_KEYS_A = ("style_variant", "background_image")
_DECK_OVERRIDE_KEYS_B = ("chrome_override", "deck_class")
_DECK_OVERRIDE_ALLOWED_KEYS = set(_DECK_OVERRIDE_KEYS_A) | set(_DECK_OVERRIDE_KEYS_B)
# DF4(DECK_FIRST_DESIGN.md §2-③·§3 DF4 행): 덱 프리뷰(틀+배경, 본문 비움) PNG 저장 폴더명 —
# `imagedeck_refs/deck_preview/<class>.png`. resolve_slide_refs의 4계층 조회(slide>global>
# deck_preview>seed)와 render_deck_preview(렌더 본체)가 공유하는 단일 정의.
DECK_PREVIEW_DIRNAME = "deck_preview"

# W31 R9: review_badges 저점수 판정(밋밋·발산추천 — 둘 다 thin_score>=임계치) → 이미지 단계에
# "배경이미지 생성 권장·디자인지식 적극 적용" 신호 주입. 충실(thin 아님)은 과잉 장식 방지를
# 위해 아무것도 주입하지 않는다(app/review_badges.py VERDICT_* 상수와 문자열 계약 — import는
# 하지 않는다: imagedeck은 review_badges를 몰라도 되고, 소비는 design_brief.json 텍스트만으로
# 충분하다 — 판정 로직 자체는 손대지 않는다는 과제 지시와도 부합).
_LOW_SCORE_VERDICTS = ("밋밋", "발산추천")

# 금지 스타일(lecture 디자인 계약 §금지) — collect/review가 참조.
FORBIDDEN_STYLE = ["olive", "green", "photorealism", "3D", "watermark", "dark full background"]

# W29 하이브리드 크롬: storyline template_id -> 장 클래스 기본 매핑.
# slide_classes를 선언한 스킨(inkline 등)에서만 활성 - 미선언 스킨(quartz)은 전 장 content(하위호환).
_TEMPLATE_CLASS_MAP = {
    "cover_slide": "cover", "cover": "cover",
    "toc": "toc", "index": "toc", "table_of_contents": "toc",
    "agenda": "toc",  # 실 storyline 어휘(W29 파일럿 실측 - 목차 장)
    "divider": "divider", "section_divider": "divider", "chapter": "divider",
}


class ImagedeckError(Exception):
    """imagedeck 입력·계약 위반."""


# ---------------------------------------------------------------------------
# 로드 / 해석
# ---------------------------------------------------------------------------

def _now() -> str:
    return dt.datetime.now().isoformat(timespec="seconds")


def _load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def resolve_skin(name_or_path: str, skins_dir: Path) -> Path:
    """스킨 이름(skins/<name>.json) 또는 절대/상대 경로를 파일로 해석."""
    p = Path(name_or_path)
    if p.suffix == ".json" and p.exists():
        return p
    cand = skins_dir / f"{Path(name_or_path).stem}.json"
    if cand.exists():
        return cand
    if p.exists():
        return p
    raise ImagedeckError(f"스킨을 찾을 수 없다: {name_or_path} (본 곳: {cand})")


def _slug(text: str, n: int) -> str:
    """제목 -> 파일명 슬러그. 한글·영숫자 보존, 공백·기호는 하이픈.

    em-dash(—)·en-dash(–)도 하이픈으로 - 비cp949 문자가 파일명에 남으면 Windows 콘솔
    출력·구식 도구에서 사고가 난다(W29 파일럿 실측).
    W28 마찰 L1(2026-07-24 강의 덱 run 실측): Windows 예약문자(`? * " < > |`)가 제목에 있으면
    `out_name`이 그 OS에서 생성 불가 - `--produce`가 조용히 그 장을 누락시킨다. 하이픈군에 합류.
    """
    text = (text or "").strip()
    text = re.sub(r"[\[\]()（）]", "", text)
    text = re.sub(r'[\s/\\:·，,.—–?*"<>|]+', "-", text)
    text = re.sub(r"-{2,}", "-", text).strip("-")
    text = text[:40] if text else f"slide{n}"
    return text or f"slide{n}"


def canvas_dims(skin: dict) -> dict:
    """D12 캔버스 역산. export(최종 크기)에서 크롬 밴드를 빼 생성 캔버스를 구한다.

    W31 마찰15 ⓒ: 이 오류는 대개 design_contract.json이 중립 병합 없이 부분 스킨을 그대로
    대체해 구조 키가 빠졌을 때 나온다(2026-07-22 이전 계약, 또는 --skin으로 부분 스킨 파일을
    직접 지정한 경우). 동작 로직은 그대로(여전히 ImagedeckError만 던진다) - 메시지만 원인·조치를
    안내하도록 보강했다. design_contract.build()는 이제 동결 시점에 이 키들을 검증해 대부분
    막지만, 폴백 경로(--skin 직접 지정)는 이 검증을 거치지 않으므로 방어선이 남아 있다.
    """
    export = skin.get("export") or skin.get("canvas") or {}
    w = int(export.get("width") or export.get("w") or 0)
    h = int(export.get("height") or export.get("h") or 0)
    if not (w and h):
        raise ImagedeckError(
            "스킨에 export.width/height(또는 canvas)가 없다 - 캔버스 역산 불가. "
            "원인: 부분 스킨(colors/brand뿐인 기관 스킨 등)을 design_contract 없이 --skin으로 "
            "직접 지정했을 가능성이 높다. 조치: (a) design_contract.json이 있으면 그걸 쓰게 "
            "--skin을 빼거나, (b) `go`로 계약을 동결(중립 위 딥머지)한 뒤 다시 시도하거나, "
            "(c) design_contract.json이 낡았다면 `go --refreeze-contract`로 재동결하라(마찰15)."
        )
    chrome = skin.get("chrome") or {}
    header_h = int(chrome.get("header_h") or 0)
    footer_h = int(chrome.get("footer_h") or 0)
    gen_h = h - header_h - footer_h
    if gen_h <= 0:
        raise ImagedeckError(
            f"크롬 밴드 합({header_h}+{footer_h})이 캔버스 높이({h}) 이상이다 - 스킨 chrome 확인."
        )
    return {
        "export": {"w": w, "h": h},
        "chrome": {"header_h": header_h, "footer_h": footer_h},
        "gen_canvas": {"w": w, "h": gen_h},
    }


def slide_class(slide: dict, skin: dict) -> str:
    """장 클래스 판정. 우선순위: slide.deck_class(명시) > template_id 매핑 > content.

    slide_classes 미선언 스킨 = 전 장 content(하위호환 게이트 - 기존 quartz 동작 불변).
    """
    classes = skin.get("slide_classes") or {}
    if not classes:
        return "content"
    explicit = slide.get("deck_class")
    if explicit in classes:
        return str(explicit)
    mapped = _TEMPLATE_CLASS_MAP.get(str(slide.get("template_id") or ""))
    if mapped in classes:
        return mapped
    return "content"


def resolve_slide_layout(skin: dict, dims: dict, slide: dict) -> dict:
    """장별 레이아웃 계약(W29): 클래스·크롬 높이·이미지 범위·기대 px.

    - image=body: 생성 px = W x (H - header_h - footer_h)  (하이브리드 본문)
    - image=full: 생성 px = 최종 캔버스 전체(lecture 원본 방식·비주얼 히어로)
    - image=none: 프롬프트를 만들지 않는다(HTML 전용 장 - cover/toc/divider)
    크롬 높이 우선순위: storyline slide.chrome_override > 클래스 spec > 스킨 chrome 전역(푸터 가변 요구).
    """
    classes = skin.get("slide_classes") or {}
    cls = slide_class(slide, skin)
    spec = classes.get(cls) or {}
    ch = skin.get("chrome") or {}
    header_h = int(spec.get("header_h", ch.get("header_h") or 0))
    footer_h = int(spec.get("footer_h", ch.get("footer_h") or 0))
    if spec.get("header") is False:
        header_h = 0
    footer_flag = spec.get("footer", True)
    overlay_footer = footer_flag == "overlay"
    if footer_flag is False or overlay_footer:
        footer_h = 0  # overlay는 높이를 예약하지 않는다(이미지 위에 겹침)
    if classes:  # 장별 오버라이드는 slide_classes 선언 스킨에서만(레거시 스킨 불변)
        ov = slide.get("chrome_override") or {}
        header_h = int(ov.get("header_h", header_h))
        footer_h = int(ov.get("footer_h", footer_h))
    # 바깥 프레임 띠(inset 밴드): 오버레이가 아니라 안쪽 전체(크롬+본문)가 줄어든다 -
    # 이미지 px도 함께 역산되어 콘텐츠 가림이 없다(D12 원칙 유지). 클래스별 opt-out = spec.frame=false.
    frame_w = 0
    if classes:
        frame_w = int(((ch.get("frame") or {}).get("width")) or 0)
        if spec.get("frame") is False:
            frame_w = 0
    # 본문 여백(body_margin): HTML이 고정 예약(페이지간 좌우 여백 일관성) - 이미지는 edge-to-edge.
    body_m = {}
    if classes:
        body_m = dict(ch.get("body_margin") or {})
        if "body_margin" in spec:
            body_m = dict(spec.get("body_margin") or {})
    bx = int(body_m.get("x", body_m.get("left", 0)) or 0)
    btop = int(body_m.get("top", 0) or 0)
    bbot = int(body_m.get("bottom", 0) or 0)
    image_mode = str(spec.get("image", "body")) if classes else "body"
    exp = dims["export"]
    if image_mode == "none":
        expected = None
    elif image_mode == "full":
        expected = {"w": exp["w"] - 2 * frame_w, "h": exp["h"] - 2 * frame_w}
    else:
        gen_w = exp["w"] - 2 * frame_w - 2 * bx
        gen_h = exp["h"] - header_h - footer_h - 2 * frame_w - btop - bbot
        if gen_h <= 0 or gen_w <= 0:
            raise ImagedeckError(
                f"장 {slide.get('n')}: 크롬({header_h}+{footer_h})+프레임({frame_w}x2)"
                f"+본문여백(x{bx}/t{btop}/b{bbot})이 캔버스({exp['w']}x{exp['h']})를 초과.")
        expected = {"w": gen_w, "h": gen_h}
    return {"cls": cls, "image": image_mode,
            "chrome": {"header_h": header_h, "footer_h": footer_h},
            "frame_w": frame_w,
            "body_margin": {"x": bx, "top": btop, "bottom": bbot},
            "expected_px": expected, "overlay_footer": overlay_footer}


def _wf_by_id(wireframe: dict | None) -> dict:
    if not wireframe:
        return {}
    return {str(s.get("slide_id")): s for s in (wireframe.get("slides") or [])}


def resolve_wireframe(wf_by_id: dict, n: Any, mode: str) -> tuple[dict | None, str]:
    """D13 wireframe 적용 모드. (적용할 wireframe 슬라이드 or None, 판정 라벨)."""
    slide = wf_by_id.get(str(n))
    if mode == "off":
        return None, "off"
    if mode == "on":
        return (slide, "on") if slide else (None, "on_no_match")
    # auto (기본): 같은 장 wireframe이 있으면 적용, 없으면 off 규칙.
    return (slide, "auto_on") if slide else (None, "auto_off")


def _flags_of(slide: dict) -> list[str]:
    """내용 안전: 이미지 안에 보존해야 할 flag 문자열([예시]·검토요망 등)을 수집."""
    flags: list[str] = []
    if slide.get("example"):
        flags.append("[예시]")
    fl = slide.get("flag")
    if fl:
        flags.append(str(fl))
    blob = json.dumps(slide, ensure_ascii=False)
    for token in ("[검토요망]", "[예시]"):
        if token in blob and token not in flags:
            flags.append(token)
    return flags


def _missing_binds(wf_slide: dict | None, fields: dict) -> tuple[list[str], list[str]]:
    """wireframe slot의 binds 키가 storyline fields에 있는지 검사. (있는 키, 없는 키)."""
    have, missing = [], []
    if not wf_slide:
        return have, missing
    for slot in wf_slide.get("slots") or []:
        b = slot.get("binds")
        if not b:
            continue
        for key in ([b] if isinstance(b, str) else list(b)):
            if key in (fields or {}):
                have.append(key)
            else:
                missing.append(key)
    return have, missing


# ---------------------------------------------------------------------------
# W31 γ패킷(리허설 마찰23) — 결정론 사전 분할(오버플로 종단 처리)
#
# 1차 탈출구(β1, 유지): 하한(minimum_body_size) 내 폰트 축소 허용(_freedom_line) — 강제 금지 폐지.
# 2차: 하한으로도 수용 불가가 "실측 확정"된 장만 bundle이 사전 분할한다. 판정은 결정론
# (글자수 vs 하한 폰트 기준 용량 상수) — LLM 판단에 맡기지 않는다(마찰23의 교훈: 생성기는
# "줄이지도 말고 넘치지도 말라"를 받으면 분할하지 않고 조용히 넘치는 쪽을 택했다).
# ---------------------------------------------------------------------------

# 글자당 실면적 근사 계수 — 한글(CJK) 글자 1개는 정사각형에 가깝다(가로≈세로≈font_size)에
# 줄간격(약 1.5배)·자간/카드 내부 여백 등 안전 오버헤드(약 1.3배)를 곱한 경험적 상수(1.5*1.3≈1.95).
# app/render/design_checks.py의 임계처럼 실 코퍼스 분포로 교정된 값이 아니라 사전(a priori) 보수치다
# — 사람 평가와 대조되기 전까지 "안 걸림"은 "잘 맞다"가 아니라 "이 근사가 못 걸렀다"는 뜻이다.
CHAR_AREA_FACTOR = 1.95
# 근소 초과(용량의 15% 이내)는 아직 "확정"이 아니다 — 오탐(분할 남발) 방지를 위한 여유.
SPLIT_CONFIDENCE_MARGIN = 1.15


def _slide_text_chars(slide: dict) -> int:
    """storyline 슬라이드 콘텐츠 총 글자수(message+bullets+fields 문자열 리프).

    app/review_badges._filled_leaves와 자매 함수이지만 저건 "채워진 개수"를, 이건 "총 글자수"를
    센다(용도가 다르다 — 여기는 용량 초과 판정용). 예시 표기([예시]/예시 데이터) 리프는 세지
    않는다(그쪽 관례와 동일 — 데모 데이터가 분할 판정을 왜곡하지 않게)."""
    parts: list[str] = []
    msg = slide.get("message")
    if isinstance(msg, str):
        parts.append(msg)
    for b in (slide.get("bullets") or []):
        if isinstance(b, str):
            parts.append(b)

    def _walk(v: Any) -> None:
        if isinstance(v, str):
            if "[예시]" not in v and "예시 데이터" not in v:
                parts.append(v)
        elif isinstance(v, dict):
            if v.get("is_example") is True:
                return
            for k, x in v.items():
                if isinstance(k, str) and (k.startswith("_") or k == "is_example"):
                    continue
                _walk(x)
        elif isinstance(v, (list, tuple)):
            for x in v:
                _walk(x)

    _walk(slide.get("fields") or {})
    return sum(len(p) for p in parts if p and p.strip())


def _capacity_chars(body_w: int, body_h: int, min_font_px: int) -> int:
    """body 캔버스(px) x 하한 폰트에서 결정론적으로 수용 가능한 글자수 상한(근사)."""
    if min_font_px <= 0 or body_w <= 0 or body_h <= 0:
        return 0
    area = body_w * body_h
    return int(area / ((min_font_px ** 2) * CHAR_AREA_FACTOR))


def _needs_split(slide: dict, layout: dict, chrome_skin: dict) -> tuple[bool, dict]:
    """마찰23 2차 탈출구 판정 — 하한 폰트로도 수용 불가가 실측 확정됐는지(결정론).

    반환: (필요 여부, detail — chars/capacity_at_min_font/min_font_px, manifest·뷰에 그대로 기록)."""
    exp = layout.get("expected_px") or {}
    min_font = int((chrome_skin.get("typography") or {}).get("minimum_body_size", 22))
    capacity = _capacity_chars(int(exp.get("w") or 0), int(exp.get("h") or 0), min_font)
    chars = _slide_text_chars(slide)
    detail = {"chars": chars, "capacity_at_min_font": capacity, "min_font_px": min_font}
    need = capacity > 0 and chars > capacity * SPLIT_CONFIDENCE_MARGIN
    return need, detail


def _split_list_value(value: Any) -> "tuple[Any, Any] | None":
    """리스트 값을 절반으로 항목 단위 분할. 항목 2개 미만이면 분할 불가(None)."""
    if not isinstance(value, list) or len(value) < 2:
        return None
    mid = (len(value) + 1) // 2
    return value[:mid], value[mid:]


def _split_slide_content(slide: dict) -> "tuple[dict, dict, list[str]] | None":
    """W31 γ패킷(마찰23): 슬라이드 콘텐츠를 A/B로 의미 단위 분할.

    분할 대상: bullets(문자열 리스트)와 fields 중 **리스트값**(항목 2개 이상 — dict 리스트든
    문자열 리스트든 항목 단위로 반씩 나눈다). 스칼라·단일 dict 필드(예: 지표 이름 하나, 결론
    문장 하나)는 분할할 수 없는 **맥락**이므로 A/B 양쪽에 그대로 남긴다(정보 손실 방지 — 양쪽
    다 "무엇에 대한 이야기인지"는 알아야 한다). 분할 가능한 것이 하나도 없으면(fields가 전부
    스칼라/단일 dict이고 bullets도 항목 2개 미만) None을 반환해 "분할 포기"를 호출부에
    알린다(과제 지시: 분할 불가 단일 필드면 분할 포기+경고)."""
    fields = slide.get("fields") if isinstance(slide.get("fields"), dict) else {}
    bullets = slide.get("bullets") if isinstance(slide.get("bullets"), list) else []
    notes: list[str] = []
    did_split = False

    a_fields, b_fields = dict(fields), dict(fields)
    for key, value in fields.items():
        parts = _split_list_value(value)
        if parts is None:
            continue
        a_fields[key], b_fields[key] = parts
        did_split = True
        notes.append(f"fields.{key}({len(value)}건->{len(parts[0])}/{len(parts[1])})")

    a_bullets, b_bullets = list(bullets), list(bullets)
    bparts = _split_list_value(bullets)
    if bparts is not None:
        a_bullets, b_bullets = bparts
        did_split = True
        notes.append(f"bullets({len(bullets)}건->{len(bparts[0])}/{len(bparts[1])})")

    if not did_split:
        return None

    slide_a = dict(slide)
    slide_a["fields"] = a_fields
    slide_a["bullets"] = a_bullets
    slide_b = dict(slide)
    slide_b["fields"] = b_fields
    slide_b["bullets"] = b_bullets
    return slide_a, slide_b, notes


# ---------------------------------------------------------------------------
# W31 R6 — 프롬프트 오버레이(사람 편집물, imagedeck_prompts_local/NN.md)
# ---------------------------------------------------------------------------

_LOCAL_OVERLAY_README = """\
# 프롬프트 오버레이 (imagedeck_prompts_local/) — R6

이 폴더는 **사람 편집물의 정본**이다. `imagedeck --bundle`은 이 폴더 안의 NN.md를 절대 만들지도
지우지도 않는다(이 README만 최초 1회 만든다).

## 쓰는 법

- `imagedeck_prompts/`에 생성된 프롬프트와 **같은 파일명**으로 파일을 만든다.
  예: `05.md`(장 5), A/B 실험 장은 `05A.md`/`05B.md`.
- 그 안에 이 장 이미지 생성에 추가로 넣고 싶은 지시를 자유 형식으로 적는다.
- 다음 `imagedeck --bundle` 때 `imagedeck_prompts/05.md` 말미에
  "## 사람 추가 지시(오버레이 — imagedeck_prompts_local/05.md)" 구획으로 자동 병합된다.
- 정본이 이 폴더 쪽이므로, 몇 번을 재번들해도(스킨을 바꾸든 wireframe 모드를 바꾸든) 오버레이는
  항상 다시 붙는다 — 지우려면 이 폴더의 파일을 지우거나 비워라.

## B2 관문(08_프롬프트확인)과의 관계

프롬프트·레퍼런스를 확인했는데 기대와 다르면, 프롬프트 자체를 고치는 것이 아니라 **여기에
지시를 추가**한다(JOURNEY.md B2). 장별 wireframe 적용 여부(on/off)는 오버레이가 아니라
`imagedeck --bundle --wireframe-mode`(전역) 또는 storyline의 장별 설정이 담당한다.
"""


def _ensure_prompts_local_dir(run: Path) -> None:
    """R6: 오버레이 폴더를 열어둔다(빈 폴더 + 안내 README만 — NN.md는 사람 몫, 절대 생성 안 함)."""
    d = run / PROMPTS_LOCAL_DIR
    d.mkdir(parents=True, exist_ok=True)
    readme = d / "README.md"
    if not readme.is_file():
        readme.write_text(_LOCAL_OVERLAY_README, encoding="utf-8")


def _load_overlay(run: Path, fname: str) -> str | None:
    """`imagedeck_prompts_local/<fname>`이 있고 비어있지 않으면 그 내용을 반환."""
    p = run / PROMPTS_LOCAL_DIR / fname
    if not p.is_file():
        return None
    try:
        text = p.read_text(encoding="utf-8").strip()
    except OSError:
        return None
    return text or None


def _merge_overlay(body: str, run: Path, fname: str) -> tuple[str, bool]:
    """오버레이가 있으면 프롬프트 말미에 명확한 구획으로 병합 주입(R6). 없으면 원본 그대로."""
    overlay = _load_overlay(run, fname)
    if overlay is None:
        return body, False
    block = (
        f"\n\n## 사람 추가 지시(오버레이 — {PROMPTS_LOCAL_DIR}/{fname})\n\n"
        f"{overlay}\n"
    )
    return body + block, True


# ---------------------------------------------------------------------------
# W31 R9 — review_badges 장별 채점 → 이미지 단계 디자인 지표 신호
# ---------------------------------------------------------------------------

def _design_verdicts(run: Path) -> dict[str, str]:
    """design_brief.json(page_rhythm.slides)에서 slide_id -> review_badges verdict.

    채점은 render 시점(design_brief 생성)에만 생기고 bundle 시점에 없을 수 있다(design_brief.json
    자체가 없는 파일럿·직접 bundle 호출 run) - 그럴 땐 빈 dict로 우아하게 생략한다(강제 아님).
    """
    brief_path = run / "design_brief.json"
    if not brief_path.is_file():
        return {}
    try:
        brief = _load_json(brief_path)
    except (OSError, ValueError):
        return {}
    rhythm = ((brief or {}).get("page_rhythm") or {}).get("slides") or []
    out: dict[str, str] = {}
    for r in rhythm:
        if not isinstance(r, dict) or r.get("slide_id") is None:
            continue
        verdict = r.get("verdict")
        if verdict:
            out[str(r["slide_id"])] = str(verdict)
    return out


def _design_signal_block(verdict: str | None) -> str | None:
    """저점수(밋밋·발산추천) 장에만 배경이미지·디자인지식 적용 권장 블록(축A — 보완). 충실/미채점은 None(생략)."""
    if verdict not in _LOW_SCORE_VERDICTS:
        return None
    return (
        f"\n\n## 디자인 지표(review_badges 채점 — R9, 축A 보완)\n\n"
        f"- 이 장의 채점: **{verdict}**(내용이 얇다고 판정됨 — A5 회의 라운드 신호 재활용).\n"
        f"- 배경 이미지 생성을 적극 권장한다. 디자인지식(무드/배경 레퍼런스)을 적극 적용해 "
        f"이 장의 시각적 임팩트를 보강하라(내용이 얇을수록 디자인으로 보완).\n"
    )


# ---------------------------------------------------------------------------
# W31 R9(리허설 마찰17, 2026-07-21 확정) — emphasis(hero) 장 축B 강조 신호
#
# 축A(위 _design_signal_block)는 "내용이 얇으니 디자인으로 보완"이고, 축B는 그 반대다 — 강조
# 장은 A5에서 **의도적으로** 메시지를 줄인 것이라 review_badges가 thin이어도(사실 verdict는
# review_badges.compute_review_badges가 emphasis 장을 강제로 충실 처리해 애초에 밋밋으로
# 안 잡힌다) "보완"이 아니라 "여백을 지켜라"가 맞는 신호다. 두 축은 문구·목적이 다르므로
# 섞지 않는다(사용자 확정 — 05 매뉴얼 고지 참고).
# ---------------------------------------------------------------------------

_EMPHASIS_VALUE = "hero"


def _emphasis_signal_block(slide: dict) -> str | None:
    """emphasis="hero" 장에만 강조(축B) 블록. 그 외는 None(생략 — 과잉 장식 방지, 축A와 대칭)."""
    if slide.get("emphasis") != _EMPHASIS_VALUE:
        return None
    return (
        "\n\n## 디자인 지표 — 강조 축B (A5 확정, R9)\n\n"
        "- 이 장은 A5(내용 동결) 회의에서 **디자인 강조(hero) 장으로 확정**됐다(review_badges "
        "채점과 무관 — 의도된 저밀도이므로 thin 판정 대상이 아니다).\n"
        "- 의도적 저밀도를 유지하라(메시지를 한 줄로) — 공유뇌 디자인지식 `챕터-간지-여백` 카드의 "
        "조작적 정의: \"여백=전환 신호, 채우지 말 것\".\n"
        "- 무드/풀블리드 배경 이미지 생성을 권장한다. 텍스트 존은 한 모서리에만 배치하라.\n"
        "- 배경은 \"있는 듯 없는 듯\" 하게 하라 — 공유뇌 디자인지식 `강조색-하나-나머지-무채색` "
        "카드의 조작적 정의: 배경은 절제하고 본문 가독을 최우선하라.\n"
    )


# ---------------------------------------------------------------------------
# bundle — 장별 프롬프트 조립 (a)
# ---------------------------------------------------------------------------

# W32 마찰36⑸: 이미지 프롬프트의 storyline 덤프에서 거르는 작업용 키. 종전에는 슬라이드를 통째로
# json.dumps 해 내부 메모가 생성기에 유출됐다(내부 어휘의 청중 노출 — 마찰29·31과 동일 계열,
# 이번은 storyline 쪽 경로. 실측 선례 = lecture-검정강의안 프롬프트 39장 전부에 visual·evidence 유출).
# evidence·render 계열=세션 작업 메모 · deck_class=compose 분류(기계 소비) · supports_axis=축 추적
# (마찰29가 장표에서 걷어낸 그 내부 표기의 필드판) · visual(구명)·form_intent=뼈대 단계가 소비하는
# 형태 층 입력 · art_note=아래 명시 구획으로 따로 실린다(중복 방지).
_PROMPT_WORKING_KEYS = frozenset({
    "visual", "form_intent", "art_note", "evidence", "render", "render_note",
    "deck_class", "supports_axis",
})

_PROMPT_TEMPLATE = """\
# 이미지 장표 프롬프트 - {label} (장 {n}{variant_tag})

> 결정론 조립(imagedeck bundle). 이 파일 = Codex 이미지 생산 입력.
> 생성 캔버스: **{gen_w} x {gen_h} px** (최종 {exp_w} x {exp_h}에서 크롬 header {header_h}/footer {footer_h} 역산 - D12).
> wireframe 모드: **{wf_mode}** (판정: {wf_res}).

```text
Use case: infographic-diagram
Asset type: Korean proposal slide image

Source of truth (storyline slide - 정본, 변경 금지):
{storyline_json}

Wireframe mode:
{wf_mode}

Layout contract (wireframe slide):
{wireframe_json}

{design_section}
{art_note_block}
{reference_block}

Primary request:
Create one new slide from the storyline content.
Preserve all Korean text, dates, numbers, filenames, and these flags verbatim: {flags}

Slide scope:
{class_note}

Layout:
{layout_instruction}

Illustration:
Choose simple line illustrations that directly explain each slot.

Constraints:
{constraints_block}
```

## 바인딩 점검
{binds_block}

## 저장 경로
- `{slides_dir}/{out_name}` (반환 px = {gen_w} x {gen_h} 정확히. 불일치 = collect 불합격)
"""

# W28 마찰 L5(2026-07-24 강의 덱 run 실측 - 38장 중 11장 결함: 이중 렌더링·중간 절단·카드 밖
# 넘침·잔여 획·한글 tofu). W31 마찰19의 "자유는 색·배치이지 글자 깨짐 허용이 아니다"를 계약으로
# 명문화 — 완전/축소 두 제약 블록 모두에 합류(중복 문자열 금지, 공통 상수 하나로 관리).
_TEXT_INTEGRITY_CONSTRAINTS = """\
- no duplicate rendering of the same sentence or phrase (같은 문장·구절 두 번 그리지 말 것)
- no mid-sentence truncation (문장 중간 절단 금지)
- all text must fit fully inside its card/frame (텍스트는 카드/프레임 안에 완전히)
- no stray strokes or orphan fragments (잔여 획·고아 조각 금지)
- Korean glyphs required — no tofu boxes or question-mark substitution (한글 글리프 필수, 네모·물음표 대체 금지)"""

# DF1(프롬프트 다이어트, CONTEXT/DECK_FIRST_DESIGN.md §2-①·§3): 완전/축소 두 하드 제약 블록을
# 하나로 수렴한다 — 디자인 정본은 이제 HTML 틀(chrome_contract)뿐이므로, "차용 스킨의 세부
# 스펙(테마 취향)을 프롬프트에 그대로 실을지"를 가르던 두 갈래 자체가 무의미해졌다. 남기는 것은
# "행동 규칙"뿐: 정확 px·정직성(창작 금지)·바인딩 누락 금지·오버플로 대응(A/B 분할, 폰트 축소
# 금지)·사진풍/3D/워터마크 금지·텍스트 무결성(_TEXT_INTEGRITY_CONSTRAINTS). 종전 _FULL_CONSTRAINTS
# 전용이던 "no olive or green"·"no dark full background"는 색 취향(테마)이라 뺐다 — 색은 이제
# 프롬프트가 아니라 HTML 틀의 몫이다(§2-①). 종전 _REDUCED_CONSTRAINTS의 "preserve all Korean
# text..." 줄은 _PROMPT_TEMPLATE의 Primary request 문단과 중복이라 통합 시 생략(바이트만 다르고
# 의미는 이미 존재).
_CONSTRAINTS = """\
- exact {gen_w} x {gen_h} px (생성 캔버스 - 크롬은 HTML이 조합, 이미지에 그리지 말 것)
- no invented facts, numbers, or achievements
- no omitted bound fields{missing_note}
- no font shrinking (넘치면 A/B 장 분할: {slugA}/{slugB})
- no photorealism, no 3D, no watermark
""" + _TEXT_INTEGRITY_CONSTRAINTS


# ---------------------------------------------------------------------------
# DF1(프롬프트 다이어트): design_section도 완전/축소 두 갈래를 하나로 수렴한다 — 이웃 크롬
# 브리핑(_chrome_brief, 마찰19 "너는 이 디자인의 일부다") + 브랜드 고지(있으면) + 분량/넘침
# 행동 한계(있으면) + 색·배경 자유 명시. 종전 완전 스킨 전용이던 "계약 전체 JSON 그대로 주입"
# (_skin_summary 전량 - colors·color_roles·typography·layout·components 등 테마 정의)은 폐지
# 한다 — 그 정의들은 이제 HTML 틀(chrome_contract)의 몫이다(DECK_FIRST_DESIGN.md §2-①).
# ---------------------------------------------------------------------------

def _chrome_brief(chrome_skin: dict, layout: dict) -> str:
    """완성 슬라이드에서 이 이미지 위/아래에 HTML로 얹히는 크롬(배지·제목·푸터)을 알려준다 -
    "너는 이 디자인의 일부다"가 핵심(밋밋 수렴 진단의 반대축: 금지 아니라 맥락 제공, 마찰19).

    DF1 앵커 결정(DECK_FIRST_DESIGN.md §3 DF1 행): 이 브리핑은 테마 사양이 아니라 맥락 제공이라
    과도기 유지한다 - DF4(덱 프리뷰 실물 레퍼런스)가 착지하면 텍스트 브리핑 대신 실물 PNG
    레퍼런스로 격상할 수 있으니, DF4 완료 시 이 함수의 존치 여부를 재검토할 것."""
    typo = chrome_skin.get("typography") or {}
    sizes = typo.get("sizes") or {}
    badge_px = sizes.get("section_badge", 22)
    title_px = sizes.get("title", 44)
    ch = layout["chrome"]
    if layout["image"] == "full":
        return ("이 장은 전체 캔버스를 그린다(deck_class=full_image) - HTML 크롬은 하단 반투명 "
                "오버레이 푸터(프로젝트명·페이지) 하나뿐이다. 표지·제목 영역까지 네가 자유롭게 "
                "디자인하는 히어로 장이다.")
    if ch["header_h"] > 0 or ch["footer_h"] > 0:
        text = (f"완성 슬라이드에서 네 이미지 위에는 섹션 배지({badge_px}px)·제목({title_px}px 볼드)이"
                + (", 아래에는 푸터 띠가 HTML로 얹힌다." if ch["footer_h"] > 0 else " HTML로 얹힌다.")
                + " 너는 이 디자인의 일부다 - 이 타이포·구조와 어울리게 본문을 꾸며라.")
        return text
    return "이 장은 HTML 크롬 없이 캔버스 전체를 그린다 - 자유롭게 전체 구도를 디자인하라."


def _brand_notice(chrome_skin: dict) -> str | None:
    """브랜드 색 고지(있을 때만) - 쓰라/마라 강제 없이 참고만(사용자 확정: 강제 → 가이드)."""
    colors = chrome_skin.get("colors") or {}
    hexes = [colors[k] for k in ("navy", "orange") if colors.get(k)]
    if not hexes:
        return None
    return "이 기관의 브랜드 색: " + "·".join(f"#{c}" for c in hexes) + " (쓰라/마라 강제 없음 - 고지만)."


def _freedom_line(chrome_skin: dict) -> str:
    """색·배경·아이콘 자유 명시 + 글씨 크기 하한 내 축소 허용(사용자 확정: 강제 → 가이드 전환)."""
    min_size = (chrome_skin.get("typography") or {}).get("minimum_body_size", 22)
    return (f"색·서브 컬러·배경 연출·컬러 아이콘은 자유다(위 크롬과 시각적으로 조화를 이루는 "
            f"선에서). 글씨 크기는 필요하면 하한 {min_size}px까지 축소를 허용한다(전면 금지 아님).")


def _behavior_limits_summary(image_skin: dict) -> str | None:
    """DF1(DECK_FIRST_DESIGN.md §2-①): content_limits/overflow_policy는 "디자인"이 아니라
    "행동 한계"라 프롬프트에 남긴다 - colors·color_roles·typography·layout·components 등
    테마 정의는 여기서 다루지 않는다(그건 이제 HTML 틀의 몫). 두 키가 둘 다 없으면 None(생략)."""
    keep = {k: image_skin[k] for k in ("content_limits", "overflow_policy") if image_skin.get(k)}
    if not keep:
        return None
    return json.dumps(keep, ensure_ascii=False, indent=2)


def _design_section(chrome_skin: dict, image_skin: dict, layout: dict,
                     art_look: str | None = None) -> str:
    """DF1(프롬프트 다이어트, DECK_FIRST_DESIGN.md §2-①·§3): 완전/축소 두 갈래를 하나의 문법으로
    수렴 - 이웃 크롬 브리핑 + 브랜드 고지(있으면) + 확정 룩(있으면) + 분량/넘침 행동 한계(있으면)
    + 색·배경 자유 명시. full_skin 값과 무관하게 항상 이 하나의 섹션을 쓴다(더 이상 "차용 스킨의
    세부 스펙 전량을 그대로 주입"하는 갈래가 없다 - 테마는 HTML 틀이 그린다)."""
    lines = ["Art direction (본문 이미지는 틀에 끼워지는 본문 - 색·타이포·배치 정본은 HTML 크롬이다):",
             "- " + _chrome_brief(chrome_skin, layout)]
    notice = _brand_notice(chrome_skin)
    if notice:
        lines.append("- " + notice)
    # W31 R10(β2): 마스터 시안(07_테마확정)에서 확정한 look 서술이 있으면 장별 프롬프트에도
    # 이어붙인다(시리즈 일관성 — 없으면 생략, 기존 프롬프트와 바이트 동일).
    if art_look:
        lines.append(f"- 확정 룩(마스터 시안, 07_테마확정): {art_look}")
    limits = _behavior_limits_summary(image_skin)
    if limits:
        lines.append("- 분량/넘침 한계(행동 규칙 - 디자인 아님):\n" + limits)
    lines.append("- " + _freedom_line(chrome_skin))
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# W31 마찰20(β): 레퍼런스 3계층 조회(장별 > 전체 > 시드) + 문단 유무 분기.
# ---------------------------------------------------------------------------

# design-assets/references/seed/ (repo 루트 기준) — <개발 원본 전용 경로>
# original_style_reference.png·approved_quartz_reference.png 시드 동봉(사용자 확정 2026-07-21).
_REPO_ROOT = Path(__file__).resolve().parents[2]
SEED_REFS_DIR = _REPO_ROOT / "design-assets" / "references" / "seed"

_REF_EXTS = ("*.png", "*.jpg", "*.jpeg")


def _slide_ref_dirname(n: Any) -> str:
    return f"{int(n):02d}" if str(n).isdigit() else str(n)


def _scan_ref_images(d: Path) -> list[str]:
    if not d.is_dir():
        return []
    found: list[str] = []
    for pat in _REF_EXTS:
        found.extend(sorted(str(p) for p in d.glob(pat)))
    return found


def resolve_slide_refs(run: Path, n: Any, cls: str | None = None) -> tuple[list[str], str]:
    """마찰20 -> DF4(4계층 승격, DECK_FIRST_DESIGN.md §2-③·§3): 장별
    (`imagedeck_refs/slides/<NN>/`) > 전체(`imagedeck_refs/global/`) >
    덱 프리뷰(`imagedeck_refs/deck_preview/<cls>.png`, 자동) > 시드
    (`design-assets/references/seed/`) 순으로 조회. 반환: (경로 목록, 출처 태그).

    deck_preview 층은 슬라이드의 장 클래스(`cls`, `resolve_slide_layout`의 `layout["cls"]`)에
    맞는 파일 1장만 본다 - 다른 클래스의 틀 프리뷰를 섞어 첨부하면 오히려 혼선이다. `cls` 미지정
    (하위호환 폴백 호출부)이면 이 층 자체를 건너뛴다.
    """
    run = Path(run)
    slide_dir = run / "imagedeck_refs" / "slides" / _slide_ref_dirname(n)
    found = _scan_ref_images(slide_dir)
    if found:
        return found, "slide"
    global_dir = run / "imagedeck_refs" / "global"
    found = _scan_ref_images(global_dir)
    if found:
        return found, "global"
    if cls:
        preview_png = run / "imagedeck_refs" / DECK_PREVIEW_DIRNAME / f"{cls}.png"
        if preview_png.is_file():
            return [str(preview_png)], "deck_preview"
    found = _scan_ref_images(SEED_REFS_DIR)
    if found:
        return found, "seed"
    return [], "none"


_REF_ROLE_LABELS = (
    "Image 1: strict card-infographic design-language reference",
    "Image 2: approved color and series reference",
)


def _extend_with_knowledge_images(lines: list[str], refs: list[str], knowledge_images: list[dict]) -> None:
    """δ패킷: A6(뼈대 결정기)가 이 장에 선택한 디자인지식 카드의 실물 이미지를 "구조 레퍼런스"
    역할로 같은 Reference roles 목록에 이어붙인다 - β1 3계층(slide/global/deck_preview/seed)과
    같은 지위(별도 문단이 아니라 같은 목록의 뒷번호)다. knowledge_images가 비면 아무 것도 하지
    않는다(기존 프롬프트와 바이트 동일 보존)."""
    base = len(refs)
    for i, img in enumerate(knowledge_images or []):
        role = design_knowledge_cards.structure_reference_role(img)
        lines.append(f"- Image {base + i + 1} ({role}): {img['path']}")


def _reference_block(refs: list[str], source: str, knowledge_images: "list[dict] | None" = None) -> str:
    """마찰20: 레퍼런스가 있으면 역할 명시(시드면 '기본값 - 교체 가능' 표기), 하나도 없으면
    (시드 폴더까지 비면) 문단 자체를 빼고 자립 지시로 대체한다.

    DF4: source="deck_preview"는 별도 역할 문구 - "끼워질 완성 틀"임을 명시하고 본문이 비어
    있음·그 틀에 어울리게 그리라는 지시를 함께 준다(기존 slide/global/seed 문법은 불변).

    δ패킷: knowledge_images(design_knowledge_cards.carry_knowledge(...)["images"])가 있으면
    "구조 레퍼런스" 역할로 이어붙인다. refs가 비어 있어도 knowledge_images만으로 목록을 연다
    (자립 지시로 도망가지 않는다 - 카드가 실제로 이미지를 골라줬으면 그것이 레퍼런스다).
    """
    knowledge_images = knowledge_images or []
    if not refs and not knowledge_images:
        return ("자립 지시(레퍼런스 없음): 위 art direction(크롬 브리핑)이 시각 언어의 기준이다.")
    if source == "deck_preview":
        lines = ["Reference roles:"]
        for i, r in enumerate(refs):
            lines.append(
                f"- Image {i + 1} (deck preview - 이 이미지가 끼워질 완성 틀): 크롬(헤더/푸터)·배경·"
                f"장식까지 다 그려진 실제 슬라이드다. 본문 자리만 비어 있다(placeholder) - 그 색·타이포·"
                f"분위기를 기준 삼아, 그 틀 안에 자연스럽게 어울리는 본문을 그려라: {r}"
            )
        _extend_with_knowledge_images(lines, refs, knowledge_images)
        return "\n".join(lines)
    lines = ["Reference roles:"]
    tag = " (기본값 - 교체 가능)" if source == "seed" else ""
    for i, r in enumerate(refs):
        role = _REF_ROLE_LABELS[i] if i < len(_REF_ROLE_LABELS) else f"Image {i + 1}: additional reference"
        lines.append(f"- {role}: {r}{tag}")
    _extend_with_knowledge_images(lines, refs, knowledge_images)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# W31 R10 v2(β2, 리허설 마찰26) — 덱 마스터 디자인 공정
#
# `imagedeck --master-bundle`: 복합 입력함(발주처 브랜드·자사 아이덴티티·주제·레퍼런스·디자인지식
# pull) 브리핑을 자기완결 프롬프트로 조립한다. **내용(storyline) 유무와 무관하게 동작** —
# 디자인 선행 루트(start → 조사 → 마스터 시안 → 내용)를 가능하게 한다(JOURNEY R10 v2).
# `imagedeck --master-apply`: 확정된 master_design.json을 검증하고 design_contract.json에
# art_direction·density를 기록(재동결 문법 — prev 보존)한 뒤, 확정 시안 이미지를
# `imagedeck_refs/global/`에 시리즈 레퍼런스로 등록한다(이후 bundle의 3계층 조회가 자동 동봉).
# ---------------------------------------------------------------------------

MASTER_PROMPT_NAME = "master_design_prompt.md"
MASTER_DESIGN_NAME = "master_design.json"
DENSITY_VALUES = ("standard", "spacious", "dense")
_MASTER_AXIS_LABELS = ("발주처축", "자사축", "주제축")


def _master_institution_block(run: Path, skins_dir: Path) -> str:
    """① 발주처 브랜드 — institution_research.json + 등록 스킨 + run/assets/client 자산 목록."""
    render_dir = _REPO_ROOT / "app" / "render"
    if str(render_dir) not in _sys.path:
        _sys.path.insert(0, str(render_dir))
    res = None
    ir_mod = None
    try:
        import institution_research as ir_mod  # type: ignore
        res = ir_mod.load(run)
    except Exception:
        pass
    lines = ["### (1) 발주처 브랜드 — institution_research.json + 등록 스킨 + 자산"]
    if res is None:
        lines.append(
            "(institution_research.json 없음 — `research --bundle/--apply`로 먼저 조사하면 "
            "브랜드색·로고·직인용 후보가 여기 채워진다. 없어도 진행 가능하다 — 자사/주제 축으로 대체.)"
        )
    else:
        lines.append(ir_mod.render_for_prompt(res))
        applied = (res.get("_applied_skin") or {}).get("skin_id")
        if applied:
            lines.append(f"- 등록 스킨: {Path(skins_dir) / f'{applied}.json'} (render --skins {applied})")
    client_assets = Path(run) / "assets" / "client"
    if client_assets.is_dir():
        files = sorted(p.name for p in client_assets.glob("*")
                       if p.is_file() and p.name.lower() != "readme.md")
        if files:
            lines.append(f"- 발주처 자산({client_assets}): {', '.join(files)}")
    return "\n".join(lines)


def _master_company_selection(run: Path) -> dict | None:
    try:
        import company  # sibling
        return company.load_selection(run)
    except Exception:
        return None


def _master_company_block(run: Path) -> str:
    """② 자사 아이덴티티 — 선택된 companies/<id>/의 assets·profile 요지."""
    lines = ["### (2) 자사 아이덴티티 — 선택된 제안사 프로필의 assets·요지"]
    sel = _master_company_selection(run)
    if not sel or not sel.get("company_id"):
        lines.append(
            "(선택된 제안사 프로필 없음 — `start --company <id>`로 연결하거나 `company --list`로 "
            "창고를 확인하라. 없어도 진행 가능하다 — 발주처/주제 축으로 대체.)"
        )
        return "\n".join(lines)
    import company  # sibling(위에서 이미 import 성공했으므로 안전)
    cid = sel["company_id"]
    profile = company.load(cid)
    if profile is None:
        lines.append(f"(제안사 프로필 선택됨: {cid} — profile.json 아직 없음. `company --bundle --id {cid}`)")
        return "\n".join(lines)
    lines.append(company.render_for_prompt(profile))
    assets = company.assets_dir(cid)
    if assets.is_dir():
        files = sorted(p.name for p in assets.glob("*")
                       if p.is_file() and p.name.lower() != "readme.md")
        if files:
            lines.append(f"- 자사 자산({assets}): {', '.join(files)}")
    return "\n".join(lines)


def _master_subject_title(run: Path) -> str | None:
    """분석카드 헤드라인(H1의 em-dash 뒤) > brief.md 첫 줄 순으로 공고명을 추정."""
    analysis_dir = Path(run) / "analysis"
    if analysis_dir.is_dir():
        cards = sorted(analysis_dir.glob("*_분석카드.md"))
        if cards:
            try:
                first_line = cards[0].read_text(encoding="utf-8", errors="replace").splitlines()[0]
            except OSError:
                first_line = ""
            if "—" in first_line:
                title = first_line.split("—", 1)[1].strip()
                if title:
                    return title
    brief_path = Path(run) / "brief.md"
    if brief_path.is_file():
        for ln in brief_path.read_text(encoding="utf-8", errors="replace").splitlines():
            if ln.strip():
                return ln.strip().lstrip("#").strip()
    return None


def _master_subject_block(run: Path) -> str:
    """③ 주제 — 공고명 + (있으면) message_map 핵심 메시지."""
    run = Path(run)
    lines = ["### (3) 주제 — 공고명·핵심 메시지"]
    title = _master_subject_title(run)
    lines.append(f"- 공고명/제목: {title or f'(미상 — run명 {run.name}으로 대체)'}")
    gov = None
    try:
        import message_map  # sibling
        mm = message_map.load(run)
        if mm is not None:
            gov = message_map.governing_text(mm)
    except Exception:
        mm = None
    if gov:
        lines.append(f"- 핵심 메시지(governing, message_map): {gov}")
    else:
        lines.append("- 핵심 메시지: (아직 message_map 없음 — 있으면 다음 번들부터 자동 동봉)")
    return "\n".join(lines)


def resolve_master_refs(run: Path) -> tuple[list[str], str]:
    """마스터 시안용 레퍼런스 — 장별 계층 없이 전체(run) > 시드 순(resolve_slide_refs의 자매)."""
    run = Path(run)
    global_dir = run / "imagedeck_refs" / "global"
    found = _scan_ref_images(global_dir)
    if found:
        return found, "global"
    found = _scan_ref_images(SEED_REFS_DIR)
    if found:
        return found, "seed"
    return [], "none"


def _master_reference_block(run: Path) -> str:
    """④ 레퍼런스 — 3계층(run 전체 > 시드) 경로 목록."""
    lines = ["### (4) 레퍼런스 — 3계층(run 전체 > 시드) 경로 목록"]
    refs, source = resolve_master_refs(run)
    if not refs:
        lines.append("(레퍼런스 없음 — 자립 지시로 진행하되, 확정 시안을 이후 시리즈 레퍼런스로 등록한다.)")
    else:
        tag = " (기본값 - 교체 가능)" if source == "seed" else ""
        lines.append(f"- 출처: {source}{tag}")
        lines.extend(f"  - {r}" for r in refs)
    return "\n".join(lines)


# ⑤ 디자인지식 pull — wireframe.py/design_spec.py의 기존 pull 프로토콜 문구를 그대로 재사용한다
# (공유뇌 ref/디자인지식/ 능동 조회, 자동 주입 아님 — MANUAL.md §9.5·§9.6과 동일 관례).
_MASTER_KNOWLEDGE_BLOCK = """### (5) 디자인지식 pull — 능동 조회(자동 주입 아님)
공유뇌 `ref/디자인지식/`을 능동 조회하라(wireframe·디자인 고도화 루프와 같은 기존 pull
프로토콜 — 필드가 요구할 때 세션이 직접 검색해 인용한다, 자동 배선 아님). 특히 배경·여백·
강조색 판단의 근거 카드(예: `챕터-간지-여백`·`강조색-하나-나머지-무채색`)를 찾아 무드·배경
방향을 정하는 근거로 삼아라. 카드가 없으면 지어내지 말고 없다고 명시하라(catalog_gap과 같은
정직화 원칙). 필요하면 Claude 디자인 스킬(예: design-critique)을 활용해 시안을 스스로
비평·정련하라."""


_MASTER_PROMPT_TEMPLATE = """# 마스터 디자인 시안 프롬프트 (imagedeck --master-bundle, W31 R10 v2)

> 자기완결 브리핑 — 이 파일 하나만으로 마스터 시안(공통 배경 + 크롬 조합 + 대표 장 1~2개 실물)을
> 만들 수 있어야 한다. **내용(storyline) 유무와 무관하게 동작한다** — storyline이 아직 없으면
> 대표 장 시안은 표지 + 샘플 본문(가상 제목·불릿)으로 만든다(내용 의존 금지 — 디자인 선행 루트).

## 입력함(복합 브리핑)

{institution_block}

{company_block}

{subject_block}

{reference_block}

{knowledge_block}

## 지시

1. 위 다섯 요소를 종합해 시각적 축을 판단하라: 발주처 브랜드가 강하면 **발주처축**, 자사
   아이덴티티가 강하면 **자사축**, 둘 다 약하고 주제가 뚜렷하면 **주제축**.
2. 축이 하나로 수렴하지 않고 갈리면(예: 발주처 색과 자사 색이 충돌) **복수 후보안**을 제시하라
   ({axis_labels} 중 갈리는 축 전부 — 후보마다 look 서술 한 문단 + 대표 장 시안 1장). 최종 선택은
   회의(대시보드/채팅)에서 사람이 한다 — 세션이 임의로 하나를 확정하지 않는다.
3. 산출물:
   - **전 장 공유 배경 PNG**(선택 — 필요하면 생성하라. 이미지 생성 능력으로, 코드 드로잉 금지):
     캔버스 전체 크기(위 캔버스 px와 동일), 그 위에 헤더·본문·푸터 텍스트가 얹혀도 읽히도록
     **저대비·여백 중심**(진한 무늬·강한 대비 지양 — 본문을 방해하지 않는 것이 최우선). 확정하면
     `background`에 경로를 적는다(모든 장이 공유 — chrome_contract.chrome.frame.image로 동결).
   - **코너 장식 후보**(선택 — 필요하면 1개 이상): 이미지 생성 능력으로만 그린다(코드 드로잉
     금지). **투명 배경 선호**(PNG 알파 채널 — 코너에 자연스럽게 얹히도록). 투명 배경이 어려우면
     위 배경색과 맞춘 사각 이미지로 폴백해도 된다(폴백 여부는 사람이 판단). 확정하면
     `decor_slots`에 위치별로 적는다(chrome_contract.decor_slots로 동결 — 코너 장식 등 배경 외
     자산).
   - 크롬 조합(헤더·푸터·프레임이 배경/색과 어울리는지) 시안.
   - 대표 장 1~2개 실물 시안 이미지(표지 + 샘플 본문 1장 — 내용 확정 전이면 가상 제목/불릿,
     확정 후면 실제 스토리라인 1장을 써도 된다).
   - `{master_design_name}`(스키마는 아래) — 후보가 여럿이면 최종 선택 1건만 기록한다(선택
     근거는 chosen_axis·sources에 남긴다).

## `{master_design_name}` 스키마

```json
{{
  "look": "무드·색·배경 방향을 문장으로 서술 (예: '따뜻한 아이보리 바탕에 손그림풍 플랫 선화, 강조색 하나만 채도 높게')",
  "density": "standard | spacious | dense",
  "assets": ["대표 장 시안 이미지 경로(run 상대 또는 절대), ..."],
  "background": "전 장 공유 배경 PNG 경로(선택 — run 상대 또는 절대, 캔버스 전체 크기)",
  "decor_slots": [
    {{"id": "corner", "image": "코너 장식 PNG 경로(run 상대 또는 절대)", "anchor": "top-right",
      "offset_x": 0, "offset_y": 0, "width": 120, "opacity": 1.0}}
  ],
  "chosen_axis": "발주처축 | 자사축 | 주제축 | null",
  "sources": ["이 판단에 쓴 입력 요소(예: institution_research.json, companies/<id>/profile.json, seed 레퍼런스)"],
  "knowledge_used": {{"cards": ["반영한 지식 카드 슬러그", "..."],
                       "web": [{{"url": "https://...", "purpose": "용도 한 줄"}}]}}
}}
```

- `density`: standard(기본 분량 밴드 유지) · spacious(여백형 — 밴드보다 적은 글자 지향) ·
  dense(밀집형 — 밴드보다 많은 정보 지향). standard가 아니고 스토리라인이 이미 있으면 상류에
  "분량 밴드 재조정 → A5 부분 재생성 권장" 경고가 표면화된다(차단 아님).
- `assets`는 실재하는 파일이어야 한다(`--master-apply`가 존재 검증).
- `background`·`decor_slots`는 **둘 다 선택**이다(생략하면 기존 크롬 그대로 — 계약 바이트 동일).
  지정하면 `--master-apply`가 디자인 타임 1회 생성 → 사람 검수를 거친 것으로 보고 그대로
  `chrome_contract`에 동결한다(이후 모든 장의 compose가 이 배경·장식을 그린다 — 재생성 없이
  재사용). `decor_slots[].anchor`는 `top-left`/`top-right`/`bottom-left`/`bottom-right` 중
  하나, `image`·`width`는 필수(둘 다 파일 실재 검증). `image`·`background`는 실재하는 파일이어야
  한다(`--master-apply`가 검증).

## 적용

시안이 정해지면 저장 후:
`python proposal_system/scripts/proposal_pipeline.py imagedeck --master-apply --file <{master_design_name} 경로> --run {run_name}`
"""


def build_master_prompt(run: Path, *, skins_dir: Path) -> str:
    run = Path(run)
    profile = gates.load_config(run)["profile"]  # ε패킷: 지식 pull+보고 의무(config 표 소비)
    knowledge_block = _MASTER_KNOWLEDGE_BLOCK + "\n\n" + knowledge_ledger.handoff_block(
        run, "master_design", profile
    )
    return _MASTER_PROMPT_TEMPLATE.format(
        institution_block=_master_institution_block(run, skins_dir),
        company_block=_master_company_block(run),
        subject_block=_master_subject_block(run),
        reference_block=_master_reference_block(run),
        knowledge_block=knowledge_block,
        axis_labels="/".join(_MASTER_AXIS_LABELS),
        master_design_name=MASTER_DESIGN_NAME,
        run_name=run.name,
    )


def master_bundle(run: Path, *, skins_dir: Path) -> dict:
    """`imagedeck --master-bundle`: 복합 입력함 브리핑을 run 루트에 쓴다(내용 유무 무관).

    반환값은 CLI 출력·state 기록용 요약 — 프롬프트 재파싱 없이 5요소 존재·레퍼런스 계층·
    내용(storyline) 유무를 알려준다.
    """
    run = Path(run)
    prompt = build_master_prompt(run, skins_dir=skins_dir)
    out = run / MASTER_PROMPT_NAME
    out.write_text(prompt, encoding="utf-8")
    refs, refs_source = resolve_master_refs(run)
    sel = _master_company_selection(run)
    return {
        "prompt": str(out),
        "institution_present": (run / "institution_research.json").is_file(),
        "company_present": bool(sel and sel.get("company_id")),
        "refs_source": refs_source,
        "refs_count": len(refs),
        "content_present": (run / "storyline.json").is_file(),
    }


def validate_master_design(doc: Any, run: Path) -> dict:
    """`imagedeck --master-apply` 검증 — 스키마·density 값·자산 실재.

    반환 {"errors": [...], "warnings": [...]} — research/company 모듈과 같은 문법(오류=적용
    중단·SSOT 안전, 경고=표면화만).
    """
    errors: list[str] = []
    warnings: list[str] = []
    if not isinstance(doc, dict):
        return {"errors": ["master_design.json이 객체가 아니다"], "warnings": []}
    look = doc.get("look")
    if not look or not isinstance(look, str) or not look.strip():
        errors.append("look 없음 — 무드·색·배경 방향 서술은 필수다")
    density = doc.get("density", "standard")
    if density not in DENSITY_VALUES:
        errors.append(f"density={density!r} — {DENSITY_VALUES} 중 하나여야 한다")
    assets = doc.get("assets") or []
    if not isinstance(assets, list):
        errors.append("assets가 배열이 아니다")
    else:
        if not assets:
            warnings.append("assets 비어있음 — 대표 장 시안 이미지 없이 확정한다(권장하지 않음)")
        for a in assets:
            p = Path(a)
            if not p.is_absolute():
                p = Path(run) / a
            if not p.is_file():
                errors.append(f"assets 파일 없음: {a}")
    chosen_axis = doc.get("chosen_axis")
    if chosen_axis is not None and chosen_axis not in _MASTER_AXIS_LABELS:
        warnings.append(
            f"chosen_axis={chosen_axis!r} — 관례 어휘{_MASTER_AXIS_LABELS} 밖(자유 서술 허용, 표면화만)"
        )
    # DF3(2026-07-24): 전 장 공유 배경 PNG — 선택. 지정 시 assets와 동일 수준으로 실재 검증한다
    # (--master-apply가 chrome_contract.chrome.frame.image로 동결하므로 파일이 없으면 동결이
    # 깨진 채로 확정된다 — 오류로 막는다. 미지정이면 아래 검사 자체가 없다 = 하위호환).
    background = doc.get("background")
    if background is not None:
        if not isinstance(background, str) or not background.strip():
            errors.append("background가 문자열이 아니다")
        else:
            p = Path(background)
            if not p.is_absolute():
                p = Path(run) / background
            if not p.is_file():
                errors.append(f"background 파일 없음: {background}")
    # DF3: 코너 장식 등 배경 외 자산 — 선택. DF2 계약 스키마(id/image/anchor/offset_x/offset_y/
    # width/opacity)와 같은 필드를 요구하되, compose 쪽 _resolve_decor_slots(관대 - 무효 슬롯만
    # 건너뛰고 warning)과 달리 여기(동결 시점)는 엄격하다 — 무효 슬롯이 있으면 적용 자체를 막는다
    # (SSOT 안전, assets 검증과 동일 수준).
    decor_slots = doc.get("decor_slots")
    if decor_slots is not None:
        if not isinstance(decor_slots, list):
            errors.append("decor_slots가 배열이 아니다")
        else:
            for i, slot in enumerate(decor_slots):
                if not isinstance(slot, dict):
                    errors.append(f"decor_slots[{i}]가 객체가 아니다")
                    continue
                slot_id = str(slot.get("id") or f"decor{i}")
                anchor = slot.get("anchor")
                if anchor not in _DECOR_ANCHORS:
                    errors.append(
                        f"decor_slots '{slot_id}': anchor 미지원({anchor!r}) — "
                        f"{list(_DECOR_ANCHORS)} 중 하나여야 한다"
                    )
                image = slot.get("image")
                if not image:
                    errors.append(f"decor_slots '{slot_id}': image 경로 없음")
                else:
                    p = Path(str(image))
                    if not p.is_absolute():
                        p = Path(run) / image
                    if not p.is_file():
                        errors.append(f"decor_slots '{slot_id}': 이미지 없음({image})")
                if not slot.get("width"):
                    errors.append(f"decor_slots '{slot_id}': width 없음")
    return {"errors": errors, "warnings": warnings}


def register_master_refs(run: Path, assets: list[str]) -> list[str]:
    """확정 시안 이미지를 `imagedeck_refs/global/`에 등록(β1의 3계층 조회가 자동 동봉, 마찰20).

    이름 충돌은 `<stem>_masterN<suffix>`로 결정론 회피(덮어쓰지 않는다). 반환 = run 상대경로 목록.
    """
    import shutil
    run = Path(run)
    dest_dir = run / "imagedeck_refs" / "global"
    dest_dir.mkdir(parents=True, exist_ok=True)
    registered: list[str] = []
    for a in assets or []:
        src = Path(a)
        if not src.is_absolute():
            src = run / a
        if not src.is_file():
            continue
        dest = dest_dir / src.name
        if dest.exists() and dest.resolve() != src.resolve():
            stem, suf = src.stem, src.suffix
            i = 1
            while (dest_dir / f"{stem}_master{i}{suf}").exists():
                i += 1
            dest = dest_dir / f"{stem}_master{i}{suf}"
        if dest.resolve() != src.resolve():
            shutil.copy2(src, dest)
        registered.append(str(dest.relative_to(run)))
    return registered


# DF3(2026-07-24, CONTEXT/DECK_FIRST_DESIGN.md §2-②·§3): 확정 배경·장식은 `imagedeck_refs/`(참고
# 이미지 3~4계층 조회 대상)와 성격이 다르다 — compose가 매 조립마다 직접 그리는 실제 조립 입력
# 이다(참고용이 아니다). 섞이면 resolve_master_refs/resolve_slide_refs가 이 파일들까지 "레퍼런스"
# 로 오인해 매 장 생성 프롬프트에 끼워 넣는 혼선이 생긴다 — 그래서 별도 고정 위치를 둔다.
MASTER_ASSETS_DIR = "imagedeck/design_assets"


def register_master_assets(run: Path, *, background: str | None,
                            decor_slots: list[dict] | None) -> tuple[str | None, list[dict]]:
    """DF3: 확정 배경 PNG·장식 자산을 run 로컬 고정 위치(`imagedeck/design_assets/`)로 복사해 계약에
    run 상대경로로 동결한다 — 마스터 시안 작성 중 만든 파일이 임시 위치에 있어도(run 밖 절대경로
    등) 계약이 그 임시 위치에 의존하지 않도록 한다(register_master_refs와 같은 동기, 별개 폴더 —
    위 docstring 참고). 이름 충돌은 `<stem>_mN<suffix>`로 결정론 회피.

    호출 전 존재 검증은 `validate_master_design`이 이미 끝냈다는 전제(여기선 재확인하지 않고
    조용히 건너뛴다 — 이중 오류 표면화 방지). 반환: (배경 run상대경로|None, decor_slots 배열
    (image 경로가 등록 위치로 치환된 사본 — 원본 slot dict는 변형하지 않는다))."""
    import shutil
    run = Path(run)
    dest_dir = run / MASTER_ASSETS_DIR
    dest_dir.mkdir(parents=True, exist_ok=True)

    def _copy_in(rel: str) -> str:
        src = Path(rel)
        if not src.is_absolute():
            src = run / rel
        if not src.is_file():
            return rel  # 검증을 이미 통과했어야 하지만, 방어적으로 원본 경로 유지(조용히 실패 안 함)
        dest = dest_dir / src.name
        if dest.exists() and dest.resolve() != src.resolve():
            stem, suf = src.stem, src.suffix
            i = 1
            while (dest_dir / f"{stem}_m{i}{suf}").exists():
                i += 1
            dest = dest_dir / f"{stem}_m{i}{suf}"
        if dest.resolve() != src.resolve():
            shutil.copy2(src, dest)
        return str(dest.relative_to(run)).replace("\\", "/")

    bg_rel = _copy_in(background) if background else None
    resolved_decor: list[dict] = []
    for slot in decor_slots or []:
        s = dict(slot)
        if s.get("image"):
            s["image"] = _copy_in(str(s["image"]))
        resolved_decor.append(s)
    return bg_rel, resolved_decor


def _class_note(layout: dict) -> str:
    """장 클래스별 이미지 범위 규칙(W29 하이브리드 크롬). 프롬프트 Slide scope 블록."""
    ch = layout["chrome"]
    if layout["image"] == "full":
        return ("FULL canvas image (deck_class=full_image, lecture style) - draw the complete "
                "slide including the title area. HTML adds only a translucent footer overlay "
                "(project name, page number) on top - do not reserve space for it.")
    if ch["header_h"] > 0 or ch["footer_h"] > 0:
        bm = layout.get("body_margin") or {}
        margin_note = ""
        if any(bm.get(k) for k in ("x", "top", "bottom")):
            margin_note = (f" Draw content EDGE-TO-EDGE on the canvas - do NOT add outer "
                           f"margins/padding; the HTML slot reserves consistent margins "
                           f"(left/right {bm.get('x', 0)}px, top {bm.get('top', 0)}px, "
                           f"bottom {bm.get('bottom', 0)}px) for every page.")
        return ("BODY ONLY (hybrid chrome). Do NOT draw: section badge, slide title, subtitle, "
                "page number, logos, chrome bands, or flag ribbons ([예시]/[검토요망] banners) - "
                "the HTML chrome renders all of those. In-body text stays verbatim, including "
                "[예시] prefixes inside cards. The bottom summary band (outputs) IS body content - "
                "draw it when bound." + margin_note)
    return "Full canvas (no HTML chrome) - draw the complete slide."


def _layout_instruction(wf_slide: dict | None, slide: dict) -> str:
    if wf_slide:
        frame = wf_slide.get("frame", "?")
        slots = wf_slide.get("slots") or []
        principles = wf_slide.get("principles") or []
        lines = [f"Frame: {frame}"]
        for slot in slots:
            piece = slot.get("piece", "?")
            size = slot.get("size", "?")
            bind = slot.get("binds") or slot.get("data") or "-"
            lines.append(f"- slot: piece={piece} size={size} binds/data={bind}")
        if principles:
            lines.append("Principles: " + " | ".join(str(p) for p in principles))
        return "\n".join(lines)
    # wireframe 없음 - storyline template_id로 배치 추론(lecture Template 보조 매핑).
    tid = slide.get("template_id") or "executive_summary"
    return (f"No wireframe. Infer layout from storyline template_id='{tid}' and fields.\n"
            f"Template mapping: executive_summary=주장+3~4근거카드 | portfolio_cases=사례카드+지표 | "
            f"strategy_pillars=동일크기 3카드 | process_steps=번호단계+화살표 | "
            f"data_interpretation=수치비교+해석 | org_roles=조직연결.")


def _one_prompt(slide: dict, wf_slide: dict | None, wf_mode: str, wf_res: str,
                image_skin: dict, chrome_skin: dict, full_skin: bool, dims: dict,
                refs: list[str], refs_source: str, n: Any, variant: str | None,
                slug: str, layout: dict, art_look: str | None = None) -> tuple[str, dict]:
    fields = slide.get("fields") or {}
    have, missing = _missing_binds(wf_slide, fields)
    flags = _flags_of(slide)
    gen = layout["expected_px"]     # W29: 장별 기대 px(클래스·크롬 오버라이드 반영)
    exp = dims["export"]
    ch = layout["chrome"]
    variant_tag = f" [{variant}]" if variant else ""
    out_name = f"{int(n):02d}{variant or ''}_{slug}.png" if str(n).isdigit() else f"{n}{variant or ''}_{slug}.png"

    missing_note = ""
    if missing:
        missing_note = f"\n- [검토요망] 바인딩 대상 없음(빈 칸 금지, 그대로 표시): {', '.join(missing)}"
    binds_lines = []
    for k in have:
        binds_lines.append(f"- OK  {k}: {json.dumps(fields.get(k), ensure_ascii=False)[:80]}")
    for k in missing:
        binds_lines.append(f"- [검토요망]  {k}: storyline fields에 없음")
    binds_block = "\n".join(binds_lines) if binds_lines else "- (binds 없음 - storyline 직접 출력)"

    # DF1(프롬프트 다이어트, DECK_FIRST_DESIGN.md §3): full_skin 값과 무관하게 하나의 문법으로
    # 수렴 — 종전 "제약 수준 = 차용 수준"(완전 스킨만 세부 스펙 전량 주입) 분기는 폐지됐다. 테마
    # 정본은 이제 HTML 틀(chrome_contract)뿐이라 프롬프트에 실을 것은 행동 규칙뿐이다. full_skin은
    # meta/manifest 기록용으로만 남는다(하위호환 소비처: design_contract.is_full_skin 등).
    design_section = _design_section(chrome_skin, image_skin, layout, art_look)
    constraints_block = _CONSTRAINTS.format(
        gen_w=gen["w"], gen_h=gen["h"], missing_note=missing_note,
        slugA=f"{n}A_{slug}.png", slugB=f"{n}B_{slug}.png")

    # δ패킷: A6(뼈대 결정기)가 이 장에 고른 knowledge_cards를 결정론 운반(본문 발췌+실물 이미지).
    # wf_slide가 없거나 knowledge_cards가 비어 있으면 carry_knowledge가 즉시 빈 결과를 돌려줘
    # vault/knowhow에 손도 대지 않는다(카드 미인용 run의 프롬프트 바이트 보존).
    knowledge_carry = design_knowledge_cards.carry_knowledge(
        (wf_slide or {}).get("knowledge_cards") if wf_slide else None
    )

    # W32 마찰36⑷⑸: 작업용 키를 걸러 덤프하고(유출 차단 — _PROMPT_WORKING_KEYS 주석 참조),
    # art_note는 명시 구획으로 전달한다. 작업용 키가 하나도 없는 슬라이드는 종전과 바이트 동일
    # (supports_axis 등이 있던 기존 덱은 그 키가 빠진 만큼 달라진다 — 유출 차단이 목적이라 의도된 변화).
    slide_public = {k: v for k, v in slide.items() if k not in _PROMPT_WORKING_KEYS}
    _art_note = str(slide.get("art_note") or "").strip()
    art_note_block = (
        f"\nArt-direction note (storyline art_note — 의도 참고. 색·형태 결정값은 design contract가 우선):\n"
        f"{_art_note}\n" if _art_note else ""
    )
    body = _PROMPT_TEMPLATE.format(
        label=slide.get("title") or f"장 {n}",
        n=n, variant_tag=variant_tag,
        gen_w=gen["w"], gen_h=gen["h"], exp_w=exp["w"], exp_h=exp["h"],
        header_h=ch["header_h"], footer_h=ch["footer_h"],
        wf_mode=wf_mode, wf_res=wf_res,
        art_note_block=art_note_block,
        storyline_json=json.dumps(slide_public, ensure_ascii=False, indent=2),
        wireframe_json=(json.dumps(wf_slide, ensure_ascii=False, indent=2) if wf_slide else "none"),
        design_section=design_section,
        reference_block=_reference_block(refs, refs_source, knowledge_images=knowledge_carry.get("images")),
        flags=(", ".join(flags) if flags else "(없음)"),
        class_note=_class_note(layout),
        layout_instruction=_layout_instruction(wf_slide, slide),
        constraints_block=constraints_block,
        binds_block=binds_block,
        slides_dir=SLIDES_DIR, out_name=out_name,
    )
    knowledge_block = design_knowledge_cards.render_prompt_block(knowledge_carry)
    if knowledge_block:
        body += "\n\n" + knowledge_block + "\n"

    meta = {
        "n": n, "slug": slug, "variant": variant,
        "out_name": out_name,
        "expected_px": {"w": gen["w"], "h": gen["h"]},
        "deck_class": layout["cls"],
        "image_scope": layout["image"],
        "chrome": dict(layout["chrome"]),
        "frame_w": layout["frame_w"],
        "body_margin": dict(layout["body_margin"]),
        "wireframe_applied": wf_slide is not None,
        "wireframe_resolution": wf_res,
        "binds": have, "missing_binds": missing,
        "flags": flags,
        "full_skin": full_skin,
        "references": refs,
        "references_source": refs_source,
        # δ패킷: 디자인지표 뷰(journey 08)가 요약 1줄로 표면화한다(§5) — 카드 인용이 없어도
        # 항상 기록(0건과 미측정을 구분하는 기존 매니페스트 관례, 마찰23 overflow_splits와 동일).
        "knowledge_carried": {
            "cards": len([c for c in knowledge_carry["cards"] if c.get("found")]),
            "images": len(knowledge_carry["images"]),
            "missing": knowledge_carry["missing"],
        },
    }
    return body, meta


def _load_deck_overrides(run: Path) -> dict:
    """DF6: run/deck_overrides.json 로드. 없으면 빈 dict(하위호환 - 조립 결과 바이트 동일)."""
    path = run / DECK_OVERRIDES_NAME
    if not path.exists():
        return {}
    try:
        data = _load_json(path)
    except (OSError, json.JSONDecodeError) as e:
        raise ImagedeckError(f"{DECK_OVERRIDES_NAME} 파싱 실패: {e}") from e
    return data or {}


def _apply_deck_overrides(slides: list[dict], run: Path) -> tuple[list[dict], list[dict]]:
    """DF6(DECK_FIRST_DESIGN.md §2-⑦ 경로 B) 공용 헬퍼 - bundle과 compose(HTML/PPTX) 모두
    storyline 로드 직후 이 함수를 통과해야 한다(단일 헬퍼 공유). 마무리 타임에 특정 장만
    style_variant/chrome_override/deck_class/fields.background_image를 바꾸는 채널 -
    storyline 재동결(go) 없이 "오버라이드 파일 수정 -> 재조립" 루프로 끝낸다.

    콘텐츠 키(title/bullets/fields의 다른 값 등)는 거부한다 - 콘텐츠 변경은 storyline
    재동결의 몫이고, 이 채널이 그 우회로가 되면 덱 우선 원칙(산출물이 아니라 정본을 고친다)이
    무너진다.

    반환: (오버라이드가 얹힌 새 slides 리스트 - 오버라이드가 없으면 원본 리스트를 그대로 반환해
           바이트 동일을 보장, 적용 기록 리스트 [{"n", "category"("a"|"b"), "keys"}, ...]).
    """
    overrides = _load_deck_overrides(run).get("slides") or {}
    if not overrides:
        return slides, []
    remaining = {str(k): v for k, v in overrides.items()}
    applied: list[dict] = []
    out: list[dict] = []
    for slide in slides:
        n = slide.get("n")
        ov = remaining.pop(str(n), None)
        if not ov:
            out.append(slide)
            continue
        bad = [k for k in ov if k not in _DECK_OVERRIDE_ALLOWED_KEYS]
        if bad:
            raise ImagedeckError(
                f"{DECK_OVERRIDES_NAME} 장 {n}: 허용되지 않는 키 {bad} - "
                f"허용 키는 {sorted(_DECK_OVERRIDE_ALLOWED_KEYS)}뿐이다(콘텐츠 필드는 storyline "
                "재동결(go)로 바꾸라 - 이 채널은 마무리 타임 디자인 미세조정 전용이다)."
            )
        new_slide = dict(slide)
        keys_applied: list[str] = []
        category = "a"
        if "style_variant" in ov:
            new_slide["style_variant"] = ov["style_variant"]
            keys_applied.append("style_variant")
        if "background_image" in ov:
            fields = dict(new_slide.get("fields") or {})
            fields["background_image"] = ov["background_image"]
            new_slide["fields"] = fields
            keys_applied.append("fields.background_image")
        if "chrome_override" in ov:
            merged = dict(new_slide.get("chrome_override") or {})
            merged.update(ov["chrome_override"] or {})
            new_slide["chrome_override"] = merged
            keys_applied.append("chrome_override")
            category = "b"
        if "deck_class" in ov:
            new_slide["deck_class"] = ov["deck_class"]
            keys_applied.append("deck_class")
            category = "b"
        out.append(new_slide)
        applied.append({"n": n, "category": category, "keys": keys_applied})
    if remaining:
        raise ImagedeckError(
            f"{DECK_OVERRIDES_NAME}: storyline에 없는 장 번호 {sorted(remaining)} - 오타를 "
            "확인하라(장 번호는 storyline.json의 slides[].n과 일치해야 한다)."
        )
    return out, applied


def bundle(run: Path, skin_path: "Path | None" = None, wireframe_mode: str = "auto",
           refs: list[str] | None = None, ab_slides: set | None = None) -> dict:
    """(a) 장별 프롬프트 번들 + manifest 조립. 정지점: Codex 단발 위임 안내.

    W31 R2·R5: `run/design_contract.json`이 있으면 그것이 유일한 디자인 정본이다(`skin_path`는
    무시된다) — `chrome_contract`로 캔버스/레이아웃을 역산한다. DF1(DECK_FIRST_DESIGN.md §2-①)
    이후 `image_contract`에서 프롬프트에 실제로 넘어가는 것은 `content_limits`/`overflow_policy`
    (행동 한계)뿐이다 — colors/color_roles/typography/layout/components 등 테마 정의는 더 이상
    프롬프트로 요약 주입되지 않는다(정본이 `chrome_contract` HTML 틀로 일원화됐다). 계약이 없는
    run(파일럿·레거시, design_contract 도입 전에 만들어졌거나 go 편입 없이 이 함수를 직접 호출하는
    경우)은 `skin_path` 파일 하나를 chrome/image 양쪽에 그대로 쓴다(폴백 — deprecated, 회귀 방지
    목적으로만 유지 — 이 경로도 DF1 수렴 프롬프트 문법을 그대로 쓴다).
    """
    run = Path(run)
    storyline_path = run / "storyline.json"
    if not storyline_path.exists():
        raise ImagedeckError(f"storyline.json 없음: {storyline_path} (먼저 go로 스토리라인 확정)")
    storyline = _load_json(storyline_path)
    slides = storyline.get("slides") or []
    if not slides:
        raise ImagedeckError("storyline.slides가 비어 있다.")
    # DF6(DECK_FIRST_DESIGN.md §2-⑦ 경로 B): storyline 로드 직후 - compose(HTML/PPTX)와 공유하는
    # 단일 헬퍼. chrome_override/deck_class(안전 분류 ⓑ)는 여기서 얹혀야 resolve_slide_layout의
    # px 역산에 반영되고, 그 결과가 manifest -> compose로 흘러 "재조립만으로는 안 바뀜(재번들
    # 필요) -> collect가 옛 px 불일치를 검출"이 성립한다.
    slides, deck_overrides = _apply_deck_overrides(slides, run)

    wf_path = run / "wireframe.json"
    wireframe = _load_json(wf_path) if wf_path.exists() else None
    wf_by_id = _wf_by_id(wireframe)

    contract = design_contract.load(run)
    if contract is not None:
        chrome_skin = contract.get("chrome_contract") or {}
        image_skin = contract.get("image_contract") or {}
        skin_label = (contract.get("meta") or {}).get("source") or "design_contract"
        skin_path_out: str | None = None
        # DF1: 프롬프트 문법은 더 이상 이 값으로 갈리지 않는다(_design_section이 항상 하나) —
        # full_skin은 계약 meta/manifest 기록·하위호환 소비처(is_full_skin 호출부·describe_for_view
        # 등, journey_folders/proposal_pipeline의 안내문)를 위해서만 계속 계산해 둔다.
        full_skin = design_contract.is_full_skin(contract)
        # W31 R10(β2): 마스터 시안(imagedeck --master-apply)이 계약에 기록한 look 서술 — 있으면
        # 장별 프롬프트에도 이어붙여 시리즈 일관성을 강화한다(없으면 None, 바이트 동일).
        art_look = (contract.get("art_direction") or {}).get("look")
    else:
        if skin_path is None:
            raise ImagedeckError(
                "design_contract.json도 skin_path도 없다 - go(B1)가 계약을 동결하거나 --skin을 지정하라."
            )
        chrome_skin = _load_json(skin_path)
        image_skin = chrome_skin  # 폴백(구 경로) — chrome 정보 포함 dict를 그대로 image_skin으로도 쓴다.
        skin_label = Path(skin_path).stem
        skin_path_out = str(skin_path)
        full_skin = True  # 폴백(파일럿·레거시) — meta/manifest 기록용 값(DF1로 프롬프트 문법엔 더 영향 없음).
        art_look = None
    dims = canvas_dims(chrome_skin)
    cli_refs = refs or []  # --ref 명시(있으면 전 장 최우선 — 마찰20, 종전 --ref 동작 보존)
    ab_slides = ab_slides or set()

    prompts_dir = run / PROMPTS_DIR
    prompts_dir.mkdir(parents=True, exist_ok=True)
    # 이전 번들 잔재 제거(스킨/모드 변경 후 stale 프롬프트 방지).
    for old in prompts_dir.glob("*.md"):
        old.unlink()
    _ensure_prompts_local_dir(run)  # R6: 오버레이 수납처(빈 폴더+README만, NN.md는 사람 몫).
    verdicts = _design_verdicts(run)  # R9: {slide_id(str): verdict} — design_brief.json 없으면 {}.

    manifest_slides = []
    overflow_splits: list[dict] = []  # W31 γ패킷(마찰23): 사전 분할이 실제 발생한 장의 사유 기록
    hasher = hashlib.sha256()
    for slide in slides:
        n = slide.get("n")
        slug = _slug(slide.get("title", ""), n)
        layout = resolve_slide_layout(chrome_skin, dims, slide)
        if layout["image"] == "none":
            # HTML 전용 장(cover/toc/divider) - 프롬프트를 만들지 않고 compose가 렌더한다(W29).
            # 오버레이·디자인 신호는 프롬프트가 없어 주입할 곳이 없다 - 지표 뷰 일관성을 위해
            # design_verdict만 기록한다(항상 False - 주입/병합 없음).
            manifest_slides.append({
                "n": n, "slug": slug, "variant": None,
                "render": "html", "deck_class": layout["cls"],
                "image_scope": "none", "chrome": dict(layout["chrome"]),
                "frame_w": layout["frame_w"],
                "flags": _flags_of(slide),
                "design_verdict": verdicts.get(str(n)),
                "design_signal_injected": False,
                "emphasis_signal_injected": False,
                "overlay_merged": False,
                "knowledge_carried": None,  # δ패킷: 프롬프트 자체가 없어 운반 대상 없음(html 전용 장)
            })
            continue
        is_manual_ab = n in ab_slides or str(n) in {str(x) for x in ab_slides}
        # slide_variants: [(variant_label, 프롬프트에 쓸 slide dict, wireframe 모드), ...]
        slide_variants: list[tuple[str | None, dict, str]]
        content_split = False
        split_notes: list[str] = []
        split_skip_reason: str | None = None
        if is_manual_ab:
            slide_variants = [("A", slide, "off"), ("B", slide, "on")]  # A=off, B=on (lecture A/B 실험)
        else:
            # W31 γ패킷(마찰22·23): 수동 A/B 실험 대상이 아닌 장만 자동 오버플로 분할 판정 대상.
            need_split, split_detail = _needs_split(slide, layout, chrome_skin)
            if need_split:
                split = _split_slide_content(slide)
                if split is not None:
                    slide_a, slide_b, split_notes = split
                    content_split = True
                    slide_variants = [("A", slide_a, wireframe_mode), ("B", slide_b, wireframe_mode)]
                    overflow_splits.append({
                        "n": n, "reason": "; ".join(split_notes), **split_detail,
                    })
                else:
                    split_skip_reason = (
                        f"분량({split_detail['chars']}자)이 하한 폰트 용량 추정치"
                        f"({split_detail['capacity_at_min_font']}자)을 초과했지만 분할 가능한 "
                        "리스트형 fields/bullets가 없어 분할을 포기했다(단일 필드/스칼라뿐)."
                    )
                    slide_variants = [(None, slide, wireframe_mode)]
            else:
                slide_variants = [(None, slide, wireframe_mode)]
        for variant, v_slide, mode in slide_variants:
            wf_slide, wf_res = resolve_wireframe(wf_by_id, n, mode)
            # W31 마찰20(β) -> DF4: --ref 명시가 있으면 전 장 최우선(종전 동작), 없으면
            # 장별>전체>덱 프리뷰>시드 4계층 조회(resolve_slide_refs). deck_preview 층은 이 장의
            # 클래스(layout["cls"])에 맞는 프리뷰 1장만 본다.
            if cli_refs:
                slide_refs, refs_source = cli_refs, "cli"
            else:
                slide_refs, refs_source = resolve_slide_refs(run, n, layout["cls"])
            body, meta = _one_prompt(v_slide, wf_slide, mode, wf_res, image_skin, chrome_skin,
                                     full_skin, dims, slide_refs, refs_source, n, variant, slug, layout,
                                     art_look=art_look)
            fname = f"{int(n):02d}{variant or ''}.md" if str(n).isdigit() else f"{n}{variant or ''}.md"
            if content_split:
                # W31 γ패킷(마찰23): 프롬프트 본문에도 분할 사실과 분량을 명시(생성기가 "이건 2장 중
                # 1장"임을 알아야 한다 — 나머지 절반을 다시 그리거나 누락시키지 않게).
                other = "B" if variant == "A" else "A"
                body += (
                    f"\n\n## 사전 분할 안내(bundle 결정론 분할 — 마찰23)\n\n"
                    f"이 장은 분량 초과({split_notes})로 2장(A/B)으로 사전 분할됐다. 이 프롬프트는 "
                    f"{variant}편(전체 2장 중 1장)이다 - 나머지 절반은 별도 프롬프트({other}편)에서 "
                    "생성한다. 여기 담긴 콘텐츠만 그려라(나머지를 추측해 채우지 말 것).\n"
                )
            # W31 R9: 저점수(밋밋·발산추천) 장에만 배경이미지·디자인지식 적용 권장 신호 주입(축A).
            # 충실 장은 과잉 장식 방지를 위해 아무것도 얹지 않는다. 채점 데이터 자체가 없으면
            # (design_brief.json 없음 - 파일럿·직접 bundle 호출) 조용히 생략(강제 아님).
            verdict = verdicts.get(str(n))
            signal_block = _design_signal_block(verdict)
            if signal_block:
                body += signal_block
            meta["design_verdict"] = verdict
            meta["design_signal_injected"] = bool(signal_block)

            # W31 R9(마찰17): emphasis="hero" 장에는 별도로 축B(강조) 신호를 얹는다 — 축A와
            # 배타적이지 않게 코드 구조를 두되(값이 없으면 None), 실제로는 emphasis 장의 verdict가
            # review_badges에서 이미 충실로 강제되어 축A가 뜰 일이 없다(문구 혼재 방지, 이중 확인).
            emphasis_block = _emphasis_signal_block(v_slide)
            if emphasis_block:
                body += emphasis_block
            meta["emphasis_signal_injected"] = bool(emphasis_block)

            # W31 R6: 사람 오버레이(imagedeck_prompts_local/<fname>) — 있으면 말미에 병합.
            # 재번들해도 오버레이 원본은 사람 쪽 폴더에 있으니 항상 다시 붙는다.
            body, overlay_merged = _merge_overlay(body, run, fname)
            meta["overlay_merged"] = overlay_merged

            if content_split:
                meta["content_split"] = True
                meta["split_reason"] = "; ".join(split_notes)
            if split_skip_reason:
                meta["overflow_split_skipped"] = True
                meta["overflow_split_skip_reason"] = split_skip_reason

            (prompts_dir / fname).write_text(body, encoding="utf-8")
            hasher.update(body.encode("utf-8"))
            meta["prompt_file"] = f"{PROMPTS_DIR}/{fname}"
            manifest_slides.append(meta)

    project = (storyline.get("meta") or {}).get("project") or run.name
    manifest = {
        "schema_version": SCHEMA_VERSION,
        "generated_at": _now(),
        "skin": skin_label,
        "skin_path": skin_path_out,
        # W31 R2: 이 번들이 design_contract.json을 정본으로 썼는지(compose가 재조회 시 참고).
        "design_contract": contract is not None,
        "wireframe_mode": wireframe_mode,
        "wireframe_present": wireframe is not None,
        "export": dims["export"],
        "chrome": dims["chrome"],
        "gen_canvas": dims["gen_canvas"],
        "classes_enabled": bool(chrome_skin.get("slide_classes")),
        # W31 마찰20(β): 전역 --ref 명시(있으면 전 장 최우선). 장별 실제 조회 결과는 각
        # slides[].references/references_source(장별>전체>시드 3계층)를 본다.
        "references": cli_refs,
        "full_skin": full_skin,
        "project": project,
        "bundle_hash": "sha256:" + hasher.hexdigest(),
        # W31 γ패킷(마찰23): 이번 bundle에서 실제 사전 분할이 발생한 장(사유 포함) — 빈 리스트도
        # 항상 남긴다(다른 게이트 신호 필드와 같은 관례 — "0건"과 "미측정"을 구분).
        "overflow_splits": overflow_splits,
        # DF6: run/deck_overrides.json이 이번 bundle에 실제로 얹은 장 번호(빈 리스트도 항상
        # 기록 — "0건"과 "미측정"을 구분하는 기존 관례). 장별 category/keys 상세는
        # deck_overrides(CLI 콘솔 ⓐ/ⓑ 안내용 - manifest에도 같이 남겨 사후 추적 가능).
        "deck_overrides_applied": [d["n"] for d in deck_overrides],
        "deck_overrides": deck_overrides,
        "slides": manifest_slides,
    }
    (run / MANIFEST_NAME).write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (run / SLIDES_DIR).mkdir(parents=True, exist_ok=True)
    return manifest


# ---------------------------------------------------------------------------
# produce — Codex 단발 위임 생산 드라이버 (b, W29 정식 승격)
# ---------------------------------------------------------------------------

_PRODUCE_INSTRUCTION = """

## 실행 지시 (Codex 단발 위임)
- **이미지 생성 능력으로 그려라.** 코드 드로잉 금지 - SVG/HTML/CSS/PIL/matplotlib로 도형을
  조립해 래스터하는 방식은 평면적인 벡터 룩이 되어 불합격이다(W30 파일럿 실측).
- Reference roles의 레퍼런스 이미지 파일을 **먼저 열어 보고** 그 디자인 언어(일러스트 질감·
  아이콘 디테일·카드 구성·색 시리즈)를 따르라. 아이콘·일러스트는 단순 기호가 아니라
  디테일 있는 삽화로.
- PNG를 **정확히 {w}x{h} px**로 `{rel_out}` (현재 작업 디렉터리 기준)에 저장하라.
  생성 px가 다르면 저장 전에 맞춰라(내용 훼손 없는 리샘플 허용 - 최종 px가 계약).
- 텍스트는 정본 그대로(오탈자 = 불합격), 디자인 계약(색·타이포·금지 스타일) 준수.
- 완료 후 저장 경로만 한 줄 출력하라.
"""


def produce(run: Path, runner, only: "set | None" = None, progress=None) -> dict:
    """(b) 이미지 생산: manifest의 미생산 이미지 장을 runner(Codex)에 순차 단발 위임.

    W29 파일럿의 임시 드라이버를 정식 승격. runner(prompt, meta) 계약 =
    codex_runner.make_codex_runner (cwd는 run 디렉터리여야 상대 저장 경로가 맞다).
    장마다 종료 즉시 px 실측(빠른 실패 감지). 이미 있는 파일은 skip(재실행 안전).
    progress(문자열) 콜백으로 진행을 알린다(CLI가 출력 담당 - 엔진은 print하지 않는다).
    반환: {"generated": [...], "skipped": [...], "failed": [...]} (out_name 기준).
    """
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(manifest_path)
    slides_dir = run / SLIDES_DIR
    slides_dir.mkdir(parents=True, exist_ok=True)
    only_s = {str(x) for x in only} if only else None
    notify = progress or (lambda _msg: None)

    report = {"generated": [], "skipped": [], "failed": []}
    for s in manifest.get("slides", []):
        if s.get("render") == "html":
            continue
        if only_s and str(s["n"]) not in only_s:
            continue
        out = slides_dir / s["out_name"]
        if out.exists():
            report["skipped"].append(s["out_name"])
            continue
        exp = s["expected_px"]
        prompt = (run / s["prompt_file"]).read_text(encoding="utf-8")
        prompt += _PRODUCE_INSTRUCTION.format(
            w=exp["w"], h=exp["h"], rel_out=f"{SLIDES_DIR}/{s['out_name']}")
        notify(f"[gen ] 장 {s['n']} {s['out_name']} ({exp['w']}x{exp['h']}) ...")
        runner(prompt, {"slide": s["n"], "out": str(out), "expected_px": exp})
        dims = png_dims(out) if out.exists() else None
        if dims == (exp["w"], exp["h"]):
            report["generated"].append(s["out_name"])
            notify(f"[ok  ] 장 {s['n']} 실측={dims}")
        else:
            # W32 마찰34: 종전에는 불량 파일을 그대로 두고 skip 판정이 **존재만** 봤다 —
            # 재실행하면 `생성 0 · skip N`이 되어 실패가 디스크에 고착됐다(자기 안내 "재실행하면
            # 실패분만 다시 위임"과 정면 모순). 사람이 산출물을 손으로 지우는 우회를 강요하던 자리다.
            # → 불량본을 `.rejected.png`로 개명해 **증거는 남기고 skip은 피한다**(마찰24 실물 증거 관례).
            rejected = _reject_artifact(out) if out.exists() else None
            report["failed"].append(s["out_name"])
            tail = f" - 불량본 보존: {rejected.name} (재실행하면 이 장을 다시 위임한다)" if rejected else ""
            notify(f"[FAIL] 장 {s['n']} 실측={dims} (기대 {exp['w']}x{exp['h']}){tail}")
    return report


def is_rejected(path: Path) -> bool:
    """`<이름>.rejected[N].png` = px 불일치로 반려된 증거본(마찰34). 생산물 집계에서 제외한다."""
    return path.stem.rsplit(".", 1)[-1].startswith("rejected") and "." in path.stem


def _reject_artifact(path: Path) -> "Path | None":
    """px 불일치 산출물을 `<이름>.rejected.png`로 개명(마찰34). 같은 이름이 있으면 번호를 붙인다.

    지우지 않는 이유: 무엇이 왜 틀렸는지 눈으로 대조할 실물이 있어야 프롬프트·모델 판단을 고칠 수
    있다(마찰24에서 파손 실물을 증거로 남긴 것과 같은 규율). 개명만으로 skip 검사(존재 기반)를
    벗어나므로 판정 로직 자체는 건드리지 않는다.
    """
    base = path.with_suffix(f".rejected{path.suffix}")
    target, i = base, 2
    while target.exists():
        target = path.with_suffix(f".rejected{i}{path.suffix}")
        i += 1
    try:
        path.rename(target)
    except OSError:
        return None   # 개명 실패(잠금 등)는 막지 않는다 — 종전처럼 FAIL만 보고한다.
    return target


# ---------------------------------------------------------------------------
# 수동 생산 루트 (W32) — codex/agy CLI 없는 사용자용 (감지 → 가이드 → adopt 수거)
#
# 파이프라인의 이미지 계약은 "정확 px의 PNG가 out_name으로 slides/에 존재"뿐이다 —
# codex는 그 파일을 만드는 드라이버 하나일 뿐이므로, CLI가 없으면 사람이 프롬프트를
# 이미지 생성 서비스(ChatGPT·Gemini 등)에 복붙해 같은 계약을 손으로 이행할 수 있다.
# 이 섹션은 그 절차를 ①감지(detect_producers) ②안내(write_manual_guide — 여정 09 폴더)
# ③수거 헬퍼(adopt — 다운로드 폴더의 이미지를 PNG 변환·정확 px 리사이즈·개명·배치)로
# 기계화한다. adopt 이후는 기존 collect/compose가 그대로 검증·조합한다(경로 무분기).
# ---------------------------------------------------------------------------

MANUAL_GUIDE_NAME = "이미지_수동생산_가이드.md"
# journey_folders.FOLDERS["09"]와 같은 값의 리터럴 — 이 모듈이 journey 상수를 문자열로
# 참조하는 기존 관례(PROMPTS_LOCAL_DIR 주석 참고)를 그대로 따른다(순환 import 회피).
JOURNEY_09_DIRNAME = "journey/09_이미지생산"
_ADOPT_SRC_EXTS = (".png", ".jpg", ".jpeg", ".webp")
# 다운로드 파일명 → 장 매칭: "장 번호로 시작"만 인정(05.png·5.jpg·05A.png·05_아무이름.png).
# 브라우저 기본명("image (3).png" 등)은 순서 추측으로 오배치될 수 있어 의도적으로 거부한다 —
# 결정론(파일명이 명시한 장에만 붙는다)이 편의보다 우선.
_ADOPT_NAME_RE = re.compile(r"^0*(\d+)\s*([A-Za-z]?)(?:[._\-\s(]|$)")


def detect_producers() -> dict:
    """이미지 생산 CLI 감지. codex=배선된 생산 러너(produce가 쓴다) / agy=참고 감지만
    (이미지 생산 러너 미배선 — 존재해도 produce는 codex 전용이다, 안내 문구용)."""
    return {"codex": shutil.which("codex"), "agy": shutil.which("agy")}


def _manual_guide_path(run: Path) -> Path:
    return Path(run) / JOURNEY_09_DIRNAME / MANUAL_GUIDE_NAME


def _slide_asset_status(slides_dir: Path, s: dict) -> str:
    """가이드 표의 상태 열 — 존재·px를 실측해 정직하게 표기(자기보고 없음)."""
    p = slides_dir / s["out_name"]
    if not p.exists():
        return "없음"
    dims = png_dims(p)
    if dims is None:
        return "있음(PNG 아님 - 변환 필요)"
    exp = s.get("expected_px") or {}
    if dims == (int(exp.get("w") or 0), int(exp.get("h") or 0)):
        return "OK"
    return f"있음(px 불일치 {dims[0]}x{dims[1]})"


def write_manual_guide(run: Path) -> Path:
    """수동 생산 가이드를 여정 09 폴더에 생성(매 호출 재생성 — 상태 열이 실측이라서).

    manifest가 정본(먼저 --bundle). 절차 5단계 + 장별 표(프롬프트 파일·최종 파일명·기대 px·
    현재 상태)를 자기완결로 담는다 — 이 파일만 보고 codex 없이 09~10 단계를 지날 수 있어야 한다.
    """
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(manifest_path)
    slides_dir = run / SLIDES_DIR
    image_slides = [s for s in manifest.get("slides", []) if s.get("render") != "html"]
    html_n = len(manifest.get("slides", [])) - len(image_slides)

    lines = [
        "<!-- 자동 생성 — 편집 금지. --produce(codex 미감지)/--manual-guide/--adopt 때마다 다시 만들어진다. -->",
        "> ⚠️ 자동 생성 — 상태 열은 생성 시점 실측이다. 이미지를 넣은 뒤 다시 만들면(--adopt가 자동 갱신) 최신이 된다.",
        "",
        "# 이미지 수동 생산 가이드 (codex CLI 없이)",
        "",
        "codex/agy CLI가 없어도 제안서 이미지를 만들 수 있다 — 파이프라인의 계약은",
        "\"정확한 px의 PNG가 정해진 파일명으로 `imagedeck/slides/`에 존재\"뿐이다.",
        "",
        "## 절차 (5단계)",
        "",
        "1. **프롬프트 열기** — `imagedeck_prompts/NN.md`(아래 표의 프롬프트 열). "
        "장별 추가 지시가 있으면 `imagedeck_prompts_local/NN.md`(오버레이)에 적고 재번들.",
        "2. **복붙 생성** — 이미지 생성이 되는 LLM(ChatGPT·Gemini 등)에 프롬프트 전체를 붙여넣는다.",
        "   - 프롬프트의 **Reference roles에 나열된 레퍼런스 이미지 파일을 채팅에 직접 첨부**하라"
        "(경로 문자열만으로는 웹 LLM이 읽지 못한다).",
        "   - \"저장 경로에 저장하라\"류 문구는 CLI용이니 무시한다(다운로드로 대체).",
        "3. **다운로드·모으기** — 생성 이미지를 한 폴더에 모은다. **파일명은 장 번호로 시작**하게"
        "(예: `05.png`, A/B 장은 `05A.png`) — 확장자는 png/jpg/webp 무엇이든 된다.",
        "4. **수거 헬퍼(adopt)** — 아래 한 줄이 PNG 변환·기대 px 리사이즈·최종 파일명 개명·배치를 다 한다"
        "(Pillow 필요: `pip install Pillow`):",
        "",
        f"   `python proposal_system/scripts/proposal_pipeline.py imagedeck --adopt <모은폴더> --run {run.name}`",
        "",
        f"5. **검증·조합** — `imagedeck --collect --run {run.name}` (px 실측·커버리지) → 통과 후 "
        f"`imagedeck --compose --run {run.name}`. 특정 장만 다시 만들려면 새 이미지를 같은 장 번호로 "
        "폴더에 넣고 --adopt를 다시 실행하면 교체된다.",
        "",
        "## 장별 목록",
        "",
        "| 장 | 프롬프트 | 최종 파일명 | 기대 px | 상태 |",
        "|---|---|---|---|---|",
    ]
    for s in image_slides:
        exp = s.get("expected_px") or {}
        lines.append(
            f"| {s['n']}{s.get('variant') or ''} | `{s.get('prompt_file') or '-'}` | "
            f"`{s['out_name']}` | {exp.get('w')}x{exp.get('h')} | {_slide_asset_status(slides_dir, s)} |"
        )
    if html_n:
        lines.append("")
        lines.append(f"> HTML 전용 장 {html_n}개(표지·목차·간지)는 프롬프트가 없다 — compose가 렌더한다(생성 불필요).")
    lines.append("")
    out = _manual_guide_path(run)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text("\n".join(lines), encoding="utf-8")
    return out


def _adopt_targets(manifest: dict) -> dict[str, dict]:
    """매칭 키("<n><VARIANT>" 대문자) → manifest 이미지 장. html 장은 대상 아님."""
    targets: dict[str, dict] = {}
    for s in manifest.get("slides", []):
        if s.get("render") == "html":
            continue
        key = f"{s['n']}{(s.get('variant') or '').upper()}"
        targets[str(key)] = s
    return targets


def adopt(run: Path, src: Path, only: "set | None" = None) -> dict:
    """(b′) 수동 루트 수거 헬퍼: src 폴더의 이미지를 장별로 매칭해 slides/에 계약대로 배치.

    - 매칭 = 파일명이 장 번호로 시작(_ADOPT_NAME_RE) 또는 out_name과 정확 일치. 그 외는
      unmatched로 보고만 한다(순서 추측 배치 금지 — 오배치가 침묵 통과하는 것보다 낫다).
    - 변환 = PNG 강제 + 기대 px로 리사이즈(비율 유지 cover-crop, 중앙 기준). Pillow 필수 —
      없으면 즉시 ImagedeckError(우아 강등 아님: 이 헬퍼의 본체가 리사이즈라서).
    - 이미 있는 장은 **교체**한다(replaced) — adopt는 사람이 파일을 폴더에 넣은 명시 행동이라
      produce의 skip-존재(재실행 안전) 문법과 다르다.
    - 끝에 가이드(write_manual_guide)를 재생성해 상태 열을 최신화한다(있을 때만 강제 아님 —
      manifest가 있으면 항상 만들 수 있으므로 그냥 만든다).
    반환: {"adopted": [...], "replaced": [...], "unmatched": [...], "failed": [...],
           "missing": [...]}  (missing = 여전히 파일 없는 장 out_name).
    """
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    src = Path(src)
    if not src.is_dir():
        raise ImagedeckError(f"adopt 소스 폴더가 없다: {src}")
    try:
        from PIL import Image, ImageOps
    except ImportError as e:
        raise ImagedeckError(
            "adopt는 Pillow가 필요하다(PNG 변환·정확 px 리사이즈가 본체) - pip install Pillow"
        ) from e

    manifest = _load_json(manifest_path)
    targets = _adopt_targets(manifest)
    by_out_name = {s["out_name"]: s for s in targets.values()}
    only_s = {str(x) for x in only} if only else None
    slides_dir = run / SLIDES_DIR
    slides_dir.mkdir(parents=True, exist_ok=True)

    report: dict[str, list] = {"adopted": [], "replaced": [], "unmatched": [], "failed": []}
    for f in sorted(src.iterdir()):
        if not f.is_file() or f.suffix.lower() not in _ADOPT_SRC_EXTS:
            continue
        s = by_out_name.get(f.name)
        if s is None:
            m = _ADOPT_NAME_RE.match(f.stem)
            if m:
                s = targets.get(f"{int(m.group(1))}{m.group(2).upper()}")
        if s is None:
            report["unmatched"].append(f.name)
            continue
        if only_s and str(s["n"]) not in only_s:
            continue
        exp = s.get("expected_px") or {}
        w, h = int(exp.get("w") or 0), int(exp.get("h") or 0)
        if not (w and h):
            report["failed"].append(f"{f.name}: manifest에 expected_px 없음")
            continue
        out = slides_dir / s["out_name"]
        existed = out.exists()
        try:
            with Image.open(f) as img:
                if img.mode not in ("RGB", "RGBA"):
                    img = img.convert("RGB")
                fitted = ImageOps.fit(img, (w, h), method=Image.LANCZOS)
                fitted.save(out, format="PNG")
        except (OSError, ValueError) as e:
            report["failed"].append(f"{f.name}: {e}")
            continue
        report["replaced" if existed else "adopted"].append(s["out_name"])
    report["missing"] = [
        s["out_name"] for s in targets.values() if not (slides_dir / s["out_name"]).exists()
    ]
    write_manual_guide(run)  # 상태 열 최신화(가이드 = 수동 루트의 진행판)
    return report


# ---------------------------------------------------------------------------
# collect — 수거 검증 (c)
# ---------------------------------------------------------------------------

def png_dims(path: Path) -> tuple[int, int] | None:
    """stdlib PNG 헤더 파싱(신규 의존성 금지). (width, height) 또는 None."""
    with open(path, "rb") as f:
        head = f.read(24)
    if len(head) < 24 or head[:8] != b"\x89PNG\r\n\x1a\n" or head[12:16] != b"IHDR":
        return None
    w, h = struct.unpack(">II", head[16:24])
    return int(w), int(h)


# ---------------------------------------------------------------------------
# W31 γ패킷(리허설 마찰24) — collect 결정론 픽셀 휴리스틱(PIL, 선택 의존성)
#
# design_checks.py가 HTML 마크업에 하는 일(문자수/줄길이 결정론 파싱)의 **이미지판**. 실제 텍스트
# 인식(OCR)은 하지 않는다(내용 안전 계약 — 창작·오독 방지) — 순수 "잉크(비배경색) 밀도"만 잰다.
# px 검사(존재·해상도)와 같은 등급의 collect 신호이되 등급은 **warn**(fail 아님 — 오탐 여지
# 감안, design_checks.py와 동일한 "이상치 탐지기이지 확정 판정기가 아니다" 원리). PIL이 없으면
# `available=False`로 우아하게 강등된다(가짜 pass 금지 — "안 봤음"과 "결함 없음"을 구분).
# ---------------------------------------------------------------------------
EDGE_BAND_PX = 12              # 캔버스 최외곽 띠 두께(px) — 잘림 의심 구간(카드 텍스트가 프레임
                               # 밖으로 넘치면 이 좁은 띠에 잉크가 몰린다 — 장 3 실증 사례 근거).
EDGE_INK_RATIO_BAD = 0.35      # 이 띠 안의 잉크 밀도 임계(35% 이상 = 가장자리 잘림 의심).
MARGIN_X_PX_DEFAULT = 64       # 예약 여백 존 기본값(skins/_neutral.json·inkline.json body_margin.x).
MARGIN_Y_PX_DEFAULT = 24       # 예약 여백 존 기본값(...body_margin.top) — 장별 실측값이 있으면 그걸 쓴다.
MARGIN_INK_RATIO_BAD = 0.30    # 예약 여백 존 잉크 밀도 임계(30% 이상 = 여백 침범 의심).
BG_SAMPLE_TOLERANCE = 24       # 배경 기준색(코너 픽셀 평균) 대비 허용 오차(RGB 채널별, 0~255).

# W28 마찰 L4(2026-07-24 강의 덱 run 실측): 한글 폰트 미적용(tofu — 글리프 없는 네모/물음표가
# 반복 출력) 탐지. OCR 없이 "글자상자로 보이는 폭의 잉크 블록이 줄 안에서 비트맵까지 동일하게
# 반복되는지"만 본다(창작·오독 방지 불변 유지 — 문자 인식이 아니라 상자 모양 중복 카운트).
# 임계값 근거: 실물 깨진 표본(<개발 원본 전용 경로>).png, 최대 중복비율
# 0.94) vs 같은 run의 정상 39장 전량 실측(최대 0.49) — 그 사이 0.6에 여유를 두고 임계 설정.
FONT_BREAKAGE_GLYPH_W_PX = (6, 90)   # 글자상자로 볼 잉크 블록 폭 범위(1792px 캔버스 실측 10~78px).
FONT_BREAKAGE_GLYPH_H_MIN_PX = 8     # 이보다 낮은 잉크 밴드는 줄로 보지 않는다(장식선 등 배제).
FONT_BREAKAGE_MIN_GLYPHS = 6         # 표본 부족(오탐 방지) - 이 개수 미만인 줄은 판정 보류.
FONT_BREAKAGE_DUP_RATIO_BAD = 0.6    # 동일 비트맵 중복 비율 임계(실측 0.94 vs 0.49 사이).


def _pil_available() -> bool:
    """PIL은 선택 의존성 — 없으면 픽셀 휴리스틱 전체를 우아하게 건너뛴다(강등, 크래시 아님)."""
    try:
        import PIL  # noqa: F401
        return True
    except Exception:
        return False


def _corner_bg_rgb(img) -> tuple[int, int, int]:
    """배경 기준색 — 네 모서리 픽셀의 평균(장식 배경이 코너까지 덮는 사례는 드물다는 전제)."""
    w, h = img.size
    pts = [(0, 0), (w - 1, 0), (0, h - 1), (w - 1, h - 1)]
    px = [img.getpixel(p) for p in pts]
    px = [(p if isinstance(p, tuple) else (p, p, p))[:3] for p in px]
    return tuple(sum(c[i] for c in px) // len(px) for i in range(3))


def _ink_ratio(img, box: tuple[int, int, int, int], bg: tuple[int, int, int],
              tol: int = BG_SAMPLE_TOLERANCE) -> float:
    """box(l,t,r,b) 영역에서 배경 기준색과 다른('잉크') 픽셀의 비율."""
    l, t, r, b = box
    l, t = max(0, int(l)), max(0, int(t))
    r, b = min(img.width, int(r)), min(img.height, int(b))
    if r <= l or b <= t:
        return 0.0
    data = img.crop((l, t, r, b)).getdata()
    total = len(data)
    if not total:
        return 0.0
    ink = sum(
        1 for p in data
        if abs(p[0] - bg[0]) > tol or abs(p[1] - bg[1]) > tol or abs(p[2] - bg[2]) > tol
    )
    return round(ink / total, 3)


def _font_breakage_heuristic(img, bg: tuple[int, int, int], tol: int = BG_SAMPLE_TOLERANCE) -> dict:
    """W28 마찰 L4(2026-07-24 강의 덱 run 실측): 한글 폰트 미적용(tofu) 의심 탐지.

    OCR 없이 결정론으로만 판단한다: 잉크가 있는 가로 밴드(줄)를 찾고, 그 안에서 글자상자로
    볼 만한 폭(FONT_BREAKAGE_GLYPH_W_PX)의 잉크 블록(글리프 후보)을 분리한 뒤, **비트맵이
    완전히 동일한 블록끼리 중복 카운트**한다. tofu(네모/물음표)는 같은 문자가 그대로 반복
    렌더링되므로 비트맵이 픽셀 단위로 동일하지만, 서로 다른 한글 음절은 각자 모양이 달라
    중복률이 낮다 — 이 차이만 잰다(글자를 읽지 않는다, 창작·오독 방지 불변 유지).

    실측: 깨진 27번 옛 PNG(<개발 원본 전용 경로>).png) 최대 중복비율
    0.94, 같은 run 정상 39장 전량 실측 최대 0.49 — 임계 0.6은 그 사이.
    """
    w, h = img.size
    r0, g0, b0 = bg
    data = img.getdata()
    ink = [1 if (abs(p[0] - r0) > tol or abs(p[1] - g0) > tol or abs(p[2] - b0) > tol) else 0
           for p in data]
    row_counts = [sum(ink[y * w:(y + 1) * w]) for y in range(h)]
    band_rows = [y for y, c in enumerate(row_counts) if c > 0]
    bands: list[tuple[int, int]] = []
    if band_rows:
        start = prev = band_rows[0]
        for y in band_rows[1:]:
            if y - prev > 2:
                bands.append((start, prev + 1))
                start = y
            prev = y
        bands.append((start, prev + 1))

    lo, hi = FONT_BREAKAGE_GLYPH_W_PX
    max_ratio = 0.0
    checked_bands = 0
    for (y0, y1) in bands:
        if y1 - y0 < FONT_BREAKAGE_GLYPH_H_MIN_PX:
            continue
        col_has = []
        for x in range(w):
            has = False
            for y in range(y0, y1):
                if ink[y * w + x]:
                    has = True
                    break
            col_has.append(has)
        segs = []
        x = 0
        while x < w:
            if col_has[x]:
                x0 = x
                while x < w and col_has[x]:
                    x += 1
                segs.append((x0, x))
            else:
                x += 1
        glyphs = [(a, b) for a, b in segs if lo <= (b - a) <= hi]
        if len(glyphs) < FONT_BREAKAGE_MIN_GLYPHS:
            continue
        checked_bands += 1
        counts: dict = {}
        for gx0, gx1 in glyphs:
            bm = tuple(tuple(ink[y * w + gx0:y * w + gx1]) for y in range(y0, y1))
            counts[bm] = counts.get(bm, 0) + 1
        dup = sum(v for v in counts.values() if v > 1)
        ratio = dup / len(glyphs)
        if ratio > max_ratio:
            max_ratio = ratio
    return {
        "checked_bands": checked_bands,
        "max_dup_ratio": round(max_ratio, 3),
        "dup_ratio_bad": FONT_BREAKAGE_DUP_RATIO_BAD,
        "suspect": max_ratio >= FONT_BREAKAGE_DUP_RATIO_BAD,
    }


def pixel_heuristics(path: Path, *, margin_x: int = MARGIN_X_PX_DEFAULT,
                     margin_y: int = MARGIN_Y_PX_DEFAULT, edge_px: int = EDGE_BAND_PX) -> dict:
    """W31 γ패킷(마찰24): PNG 1장의 결정론 픽셀 휴리스틱. PIL 없으면 {"available": False}."""
    if not _pil_available():
        return {"available": False}
    try:
        from PIL import Image
        with Image.open(path) as im:
            img = im.convert("RGB")
            w, h = img.size
            bg = _corner_bg_rgb(img)
            edge = {
                "top": _ink_ratio(img, (0, 0, w, edge_px), bg),
                "bottom": _ink_ratio(img, (0, h - edge_px, w, h), bg),
                "left": _ink_ratio(img, (0, 0, edge_px, h), bg),
                "right": _ink_ratio(img, (w - edge_px, 0, w, h), bg),
            }
            margin = {
                "top": _ink_ratio(img, (0, 0, w, margin_y), bg),
                "bottom": _ink_ratio(img, (0, h - margin_y, w, h), bg),
                "left": _ink_ratio(img, (0, 0, margin_x, h), bg),
                "right": _ink_ratio(img, (w - margin_x, 0, w, h), bg),
            }
            # W28 마찰 L4: px/여백 판정과 별개(warn 등급, fail에 관여 안 함) — 이상치 탐지기.
            font_breakage = _font_breakage_heuristic(img, bg)
    except Exception as exc:  # pragma: no cover - 손상 파일·디코딩 사고
        return {"available": False, "error": f"{type(exc).__name__}: {exc}"}
    flags: list[str] = []
    if any(v >= EDGE_INK_RATIO_BAD for v in edge.values()):
        flags.append("edge_ink_high")
    if any(v >= MARGIN_INK_RATIO_BAD for v in margin.values()):
        flags.append("margin_ink_high")
    if font_breakage["suspect"]:
        flags.append("font_breakage_suspect")
    return {
        "available": True,
        "bg_sample_rgb": list(bg),
        "edge_band_px": edge_px, "edge_ink_ratio_bad": EDGE_INK_RATIO_BAD, "edge_ink_ratio": edge,
        "margin_zone_px": {"x": margin_x, "y": margin_y}, "margin_ink_ratio_bad": MARGIN_INK_RATIO_BAD,
        "margin_ink_ratio": margin,
        "font_breakage": font_breakage,
        "flags": flags,
    }


def collect(run: Path) -> dict:
    """(c) 결정론 검증: 존재·해상도 px·커버리지·파일명 규약. 사람 정독 전 기계 관문."""
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(manifest_path)
    slides_dir = run / SLIDES_DIR
    exp = manifest.get("gen_canvas") or {}
    exp_w, exp_h = int(exp.get("w", 0)), int(exp.get("h", 0))

    results = []
    ok_count = 0
    image_total = 0
    for s in manifest.get("slides", []):
        if s.get("render") == "html":
            # HTML 전용 장(cover/toc/divider) - 이미지 검증 대상 아님(compose가 렌더).
            results.append({"n": s["n"], "variant": None, "deck_class": s.get("deck_class"),
                            "status": "html", "flags": s.get("flags", []), "missing_binds": []})
            continue
        image_total += 1
        out_name = s["out_name"]
        p = slides_dir / out_name
        # W29: 장별 기대 px(클래스·크롬 오버라이드 반영). 구 manifest는 전역 gen_canvas로 폴백.
        sexp = s.get("expected_px") or {}
        s_w, s_h = int(sexp.get("w") or exp_w), int(sexp.get("h") or exp_h)
        entry = {"n": s["n"], "variant": s.get("variant"), "out_name": out_name,
                 "expected_px": {"w": s_w, "h": s_h}, "deck_class": s.get("deck_class"),
                 "flags": s.get("flags", []),
                 "missing_binds": s.get("missing_binds", [])}
        if not p.exists():
            entry.update(status="missing", reason=f"파일 없음: {SLIDES_DIR}/{out_name}")
            results.append(entry)
            continue
        dims = png_dims(p)
        if dims is None:
            entry.update(status="not_png", reason="PNG 헤더 아님")
            results.append(entry)
            continue
        w, h = dims
        entry["actual_px"] = {"w": w, "h": h}
        # W31 γ패킷(마찰24): 결정론 픽셀 휴리스틱 — px 판정과 별개로 warn 등급 신호(오탐 여지
        # 감안, fail 아님). 파일을 읽을 수 있으면(px 일치 여부와 무관) 항상 시도한다.
        bm = s.get("body_margin") or {}
        px_heur = pixel_heuristics(
            p,
            margin_x=int(bm.get("x") or MARGIN_X_PX_DEFAULT),
            margin_y=int(bm.get("top") or MARGIN_Y_PX_DEFAULT),
        )
        entry["pixel_heuristics"] = px_heur
        if (w, h) != (s_w, s_h):
            entry.update(status="px_mismatch",
                         reason=f"기대 {s_w}x{s_h} != 실측 {w}x{h}",
                         regen_instruction=(f"장 {s['n']}{s.get('variant') or ''} 재생성: "
                                            f"정확히 {s_w}x{s_h} px, 텍스트만 유지."))
            results.append(entry)
            continue
        entry.update(status="ok")
        ok_count += 1
        results.append(entry)

    total = image_total
    pil_available = _pil_available()
    pixel_warn = [
        r for r in results
        if (r.get("pixel_heuristics") or {}).get("available") and (r["pixel_heuristics"].get("flags"))
    ]
    report = {
        "schema_version": SCHEMA_VERSION,
        "collected_at": _now(),
        "bundle_hash": manifest.get("bundle_hash"),
        "gen_canvas": exp,
        "coverage": {"ok": ok_count, "total": total},
        "flagged_slides": [r for r in results if r.get("flags") or r.get("missing_binds")],
        # W31 γ패킷(마찰24): px 검사와 같은 등급의 warn 신호 — fail이 아니라 pass/coverage에는
        # 관여하지 않는다(오탐 여지 감안 — design_checks.py와 동일 원칙, "이상치 탐지기").
        "pixel_heuristics_available": pil_available,
        "pixel_warnings": len(pixel_warn),
        "slides": results,
        "pass": ok_count == total and total > 0,
    }
    if not pil_available:
        report["pixel_heuristics_note"] = (
            "PIL(Pillow) 미설치 - 픽셀 휴리스틱 우아 강등(미측정, 'ok' 아님). "
            "설치하려면: pip install Pillow."
        )
    (run / COLLECT_NAME).write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (run / COLLECT_MD).parent.mkdir(parents=True, exist_ok=True)
    (run / COLLECT_MD).write_text(_collect_md(report), encoding="utf-8")
    # W31 γ패킷(마찰24): collect가 imagedeck_review.md scaffold를 자동 생성(기존 --review-scaffold
    # 로직 재사용 — 그 플래그도 계속 동작한다, 그냥 다시 만들 뿐이라 무해하다). warn 장은
    # review_scaffold() 안에서 표시된다.
    try:
        review_scaffold(run)
    except ImagedeckError:
        pass  # collect는 review scaffold 실패로 죽지 않는다(부가 산출물 — px 검증이 본체).
    return report


def _collect_md(report: dict) -> str:
    cov = report["coverage"]
    lines = [f"# imagedeck 수거 검증 (px 실측)",
             "",
             f"- 커버리지: {cov['ok']}/{cov['total']}",
             f"- 생성 캔버스: {report['gen_canvas'].get('w')} x {report['gen_canvas'].get('h')}",
             f"- 판정: {'PASS' if report['pass'] else 'FAIL'}",
             ]
    # W31 γ패킷(마찰24): 픽셀 휴리스틱 요약 — warn 등급(오탐 여지 있음, px 판정과는 별개).
    if report.get("pixel_heuristics_available"):
        lines.append(f"- 픽셀 휴리스틱 경고: {report.get('pixel_warnings', 0)}건(warn 등급 - fail 아님, 눈으로 확인 요망)")
    else:
        lines.append(f"- 픽셀 휴리스틱: {report.get('pixel_heuristics_note', '미측정')}")
    lines.append("")
    for r in report["slides"]:
        tag = r["status"].upper()
        if r["status"] == "html":
            lines.append(f"- [HTML] 장 {r['n']} ({r.get('deck_class')}) - compose가 렌더(이미지 검증 제외)")
            continue
        line = f"- [{tag}] 장 {r['n']}{r.get('variant') or ''} `{r['out_name']}`"
        if r.get("reason"):
            line += f" - {r['reason']}"
        lines.append(line)
        if r.get("flags"):
            lines.append(f"    flag 보존 확인 필요: {', '.join(r['flags'])}")
        if r.get("missing_binds"):
            lines.append(f"    누락 binds([검토요망]): {', '.join(r['missing_binds'])}")
        px_flags = ((r.get("pixel_heuristics") or {}).get("flags")) or []
        if px_flags:
            lines.append(f"    [WARN] 픽셀 휴리스틱: {', '.join(px_flags)}(오탐 여지 있음 - 직접 확인)")
    return "\n".join(lines) + "\n"


# ---------------------------------------------------------------------------
# compose — HTML 크롬 조합 (f)
# ---------------------------------------------------------------------------

def _brand(run: Path) -> dict:
    brief_path = run / "design_brief.json"
    if not brief_path.exists():
        return {}
    try:
        return (_load_json(brief_path) or {}).get("brand") or {}
    except (OSError, json.JSONDecodeError):
        return {}


def _logo_src(logo: str | None, run: Path) -> str | None:
    if not logo:
        return None
    p = Path(logo)
    if not p.is_absolute():
        p = run / logo
    if not p.exists():
        return None
    return p.resolve().as_uri()


def _select_final_slides(slides: list[dict]) -> list[dict]:
    """A/B variant 중 최종본 선택 — compose(HTML)·_compose_hybrid·compose_pptx가 공유하는 단일 규칙.

    두 갈래 다른 의미의 A/B가 있다: (1) 수동 wireframe 실험(ab_slides, A=off/B=on) — 최종본은
    하나만 고른다(A 우선, 종전 동작 그대로). (2) W31 γ패킷(마찰23) 자동 콘텐츠 사전 분할
    (content_split=True) — A/B는 서로 다른 절반의 실제 콘텐츠라 **둘 다** 최종 페이지에
    순서대로 들어가야 한다(선택이 아니라 이어붙이기). content_split 플래그로 두 갈래를 구분한다.
    """
    by_slide: dict = {}
    for s in slides:
        by_slide.setdefault(s["n"], []).append(s)
    chosen: list[dict] = []
    for n in sorted(by_slide, key=lambda x: (str(x).zfill(3))):
        cands = by_slide[n]
        if any(c.get("content_split") for c in cands):
            chosen.extend(sorted(cands, key=lambda c: (c.get("variant") or "")))
            continue
        pick = next((c for c in cands if c.get("variant") is None), None) \
            or next((c for c in cands if c.get("variant") == "A"), None) \
            or cands[0]
        chosen.append(pick)
    return chosen


def _resolve_chrome_skin(run: Path, manifest: dict) -> dict:
    """조립(compose/pptx)이 읽을 크롬 계약(W31 R5: chrome_contract만 소비).

    manifest.design_contract가 참이면 run/design_contract.json의 chrome_contract를 쓴다
    (image_contract는 프롬프트 전용이라 조립엔 관여하지 않는다). 계약이 없는 run(폴백)은
    종전처럼 manifest.skin_path 파일을 그대로 읽는다 — design_contract 도입 전 run과
    바이트 동일하게 동작해야 한다(회귀 방지)."""
    if manifest.get("design_contract"):
        contract = design_contract.load(run)
        if contract is not None:
            return contract.get("chrome_contract") or {}
    sp = manifest.get("skin_path")
    if sp and Path(sp).is_file():
        try:
            return _load_json(Path(sp))
        except (OSError, json.JSONDecodeError):
            return {}
    return {}


def compose(run: Path) -> dict:
    """(f) HTML 크롬 조합 -> deck.images.html. 크롬=제목/로고(D11), 본문=이미지."""
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(manifest_path)
    slides_dir = run / SLIDES_DIR
    exp = manifest.get("export") or {}
    ch = manifest.get("chrome") or {}
    header_h = int(ch.get("header_h") or 0)
    footer_h = int(ch.get("footer_h") or 0)
    W, H = int(exp.get("w", 0)), int(exp.get("h", 0))
    project = manifest.get("project") or run.name

    # W29 하이브리드 크롬: slide_classes 선언 스킨은 전용 컴포저로(장 클래스·variants·frame).
    # W31 R5: design_contract가 있으면 chrome_contract만 읽는다(_resolve_chrome_skin 참고).
    skin = _resolve_chrome_skin(run, manifest)
    if manifest.get("classes_enabled") and skin.get("slide_classes"):
        return _compose_hybrid(run, manifest, skin)

    brand = _brand(run)
    client_logo = _logo_src(brand.get("client_logo"), run)
    proposer_logo = _logo_src(brand.get("proposer_logo"), run)
    placement = brand.get("placement") or {}
    client_on = placement.get("client", "cover")   # cover=표지만, all=전장
    proposer_on = placement.get("proposer", "all")

    # A/B 최종본 선택(수동 실험=하나만, 마찰23 콘텐츠 분할=둘 다 — _select_final_slides 참고).
    chosen = _select_final_slides(manifest.get("slides", []))

    total = len(chosen)
    used, missing = 0, []
    slide_html = []
    for idx, s in enumerate(chosen, start=1):
        img = slides_dir / s["out_name"]
        section = ""  # storyline section은 이미지에 있으므로 크롬은 제목·페이지만.
        is_cover = idx == 1
        show_client = client_logo and (client_on == "all" or (client_on == "cover" and is_cover))
        show_proposer = proposer_logo and (proposer_on == "all" or (proposer_on == "cover" and is_cover))
        header = ""
        if header_h > 0:
            header = (f'<header class="chrome-h" style="height:{header_h}px">'
                      f'<span class="proj">{_esc(project)}</span>'
                      f'<span class="page">{idx} / {total}</span></header>')
        footer = ""
        if footer_h > 0:
            logos = ""
            if show_client:
                logos += f'<img class="logo client" src="{client_logo}" alt="client">'
            if show_proposer:
                logos += f'<img class="logo proposer" src="{proposer_logo}" alt="proposer">'
            footer = (f'<footer class="chrome-f" style="height:{footer_h}px">{logos}</footer>')
        if img.exists():
            used += 1
            body = f'<img class="content" src="{img.resolve().as_uri()}" alt="slide {s["n"]}">'
        else:
            missing.append(s["out_name"])
            body = (f'<div class="content missing">이미지 없음: {SLIDES_DIR}/{_esc(s["out_name"])}'
                    f'<br>(imagedeck --collect 로 확인)</div>')
        slide_html.append(
            f'<section class="slide" style="width:{W}px;height:{H}px">'
            f'{header}<div class="body" style="height:{H-header_h-footer_h}px">{body}</div>{footer}'
            f'</section>')

    _so, _sc, _scss = _stage_wrap(W, H)   # W32: 화면 맞춤 + 화면 단위 페이징 + 두 창 동기화
    html = _COMPOSE_SHELL.format(
        project=_esc(project), W=W, H=H,
        slides="\n".join(slide_html), stage_open=_so, stage_close=_sc, stage_css=_scss)
    out = run / COMPOSE_HTML
    out.write_text(html, encoding="utf-8")
    return {"out": str(out), "slides": total, "images_used": used, "missing": missing,
            "chrome": {"header_h": header_h, "footer_h": footer_h}}


def _esc(text: Any) -> str:
    return (str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


# ---------------------------------------------------------------------------
# compose (W29) — 하이브리드 크롬: 장 클래스·variants·frame
# ---------------------------------------------------------------------------

# 스킨 colors 키 -> CSS 변수 (variants 오버라이드 대상).
_VAR_MAP = [("ink", "--ink", "16324C"), ("ink_deep", "--ink-deep", "0F2438"),
            ("accent", "--accent", "E8590C"), ("bg", "--bg", "F8F7F4"),
            ("gray_text", "--gray", "5C6470"), ("line", "--line", "DDE1E6"),
            ("line_strong", "--line-strong", "C9D1DA")]


def _hexcolor(value: Any, default: str) -> str:
    c = str(value or default)
    return c if c.startswith("#") else "#" + c


def _variant_style(skin: dict, slide: dict) -> str:
    """storyline slide.style_variant -> skin.variants[이름].colors를 CSS 변수 인라인 오버라이드로.

    구간별 스타일 전환 계약(같은 덱 안에서 톤이 바뀌는 실전 사례 대응). 미지정/미등록 = 빈 문자열.
    """
    name = slide.get("style_variant")
    if not name:
        return ""
    var = (skin.get("variants") or {}).get(str(name))
    if not isinstance(var, dict):
        return ""
    colors = var.get("colors") or {}
    decls = []
    for key, cssvar, _d in _VAR_MAP:
        if key in colors:
            decls.append(f"{cssvar}:{_hexcolor(colors[key], _d)}")
    return ";".join(decls)


def _flag_pills(flags: list[str]) -> str:
    """flag를 HTML 딱지로(하이브리드에서는 이미지에 배너를 굽지 않는다 - 확정 시 재생성 없이 제거)."""
    pills = []
    for f in flags:
        label = str(f).strip("[]")
        cls = "flag example" if label == "예시" else "flag"
        pills.append(f'<span class="{cls}">{_esc(label)}</span>')
    return f'<div class="flags">{"".join(pills)}</div>' if pills else ""


def _bg_layer(slide: dict, run: Path) -> str:
    """HTML 장 전체 배경이미지 레이어(z-바닥) + 라이트 스크림(어두운 텍스트 가독성 확보).

    계약: storyline slide.fields.background_image (run 상대 또는 절대 경로). cover는
    _cover_html이 자체 다크 스크림으로 처리하므로 이 레이어를 쓰지 않는다.
    """
    src = _logo_src((slide.get("fields") or {}).get("background_image"), run)
    if not src:
        return ""
    return (f'<div class="bg" style="background-image:url({src})"></div>'
            '<div class="scrim light"></div>')


# ---------------------------------------------------------------------------
# DF2 — 자산 슬롯 계약: chrome_contract.decor_slots(코너 장식 등 배경 외 자산).
# 스펙: CONTEXT/DECK_FIRST_DESIGN.md §2-② · §3 DF2 행. HTML·PPTX가 같은 해석(_resolve_decor_slots)
# 을 공유해 두 렌더러가 어긋나지 않는다(L10류 이중 주입 재발 방지 — 이 재설계 전체의 취지).
# z-순서: 배경 위 · 크롬/본문 아래(§2-④③) — 두 렌더러 모두 배경 레이어 직후·크롬/본문 이전에
# 삽입한다. 겹침 판정은 하지 않는다(계약 작성자 책임 — 엔진은 위치를 그대로 그린다).
# ---------------------------------------------------------------------------
_DECOR_ANCHORS = {
    "top-left": ("top", "left"),
    "top-right": ("top", "right"),
    "bottom-left": ("bottom", "left"),
    "bottom-right": ("bottom", "right"),
}


def _resolve_decor_slots(skin: dict, cls: str, run: Path, warnings: list[str]) -> list[dict]:
    """decor_slots 계약을 검증·해석(HTML·PPTX 공용).

    미선언(빈/부재) = [] 즉시 반환(하위호환 — 조립 결과 바이트 동일). opt-out =
    slide_classes.<cls>.decor=false(기존 frame=false와 같은 문법 계열, cover 등에서 끔).
    항목이 무효(anchor 미지원·image/width 누락)이거나 자산 파일이 없으면 그 슬롯만 건너뛰되
    조용히 넘어가지 않고 warnings에 표면화한다(계약 — 자산 부재를 숨기지 않는다).
    """
    slots = skin.get("decor_slots")
    if not isinstance(slots, list) or not slots:
        return []
    cls_spec = (skin.get("slide_classes") or {}).get(cls) or {}
    if cls_spec.get("decor") is False:
        return []
    resolved: list[dict] = []
    for slot in slots:
        if not isinstance(slot, dict):
            continue
        slot_id = str(slot.get("id") or "decor")
        anchor = str(slot.get("anchor") or "")
        if anchor not in _DECOR_ANCHORS:
            warnings.append(f"decor_slots '{slot_id}': anchor 미지원({anchor!r}) - 건너뜀")
            continue
        image = slot.get("image")
        if not image:
            warnings.append(f"decor_slots '{slot_id}': image 경로 없음 - 건너뜀")
            continue
        width = slot.get("width")
        if not width:
            warnings.append(f"decor_slots '{slot_id}': width 없음 - 건너뜀")
            continue
        p = Path(str(image))
        if not p.is_absolute():
            p = run / image
        if not p.is_file():
            warnings.append(f"decor_slots '{slot_id}': 이미지 없음({image}) - 건너뜀")
            continue
        opacity = slot.get("opacity")
        try:
            op = float(opacity) if opacity is not None else 1.0
        except (TypeError, ValueError):
            op = 1.0
        op = min(1.0, max(0.0, op))
        resolved.append({
            "id": slot_id, "path": p, "anchor": anchor,
            "offset_x": int(slot.get("offset_x") or 0),
            "offset_y": int(slot.get("offset_y") or 0),
            "width": int(width), "opacity": op,
        })
    return resolved


def _decor_layer_html(skin: dict, cls: str, run: Path, warnings: list[str]) -> str:
    """decor_slots -> HTML 절대배치 레이어(.inner 기준 코너 앵커, 높이는 이미지 비율 자동)."""
    resolved = _resolve_decor_slots(skin, cls, run, warnings)
    if not resolved:
        return ""
    parts = []
    for d in resolved:
        v_side, h_side = _DECOR_ANCHORS[d["anchor"]]
        style = (f'position:absolute;{v_side}:{d["offset_y"]}px;{h_side}:{d["offset_x"]}px;'
                 f'width:{d["width"]}px;opacity:{d["opacity"]};pointer-events:none')
        parts.append(f'<img class="decor" src="{d["path"].resolve().as_uri()}" '
                     f'alt="{_esc(d["id"])}" style="{style}">')
    return "".join(parts)


def _section_order(slides: list) -> list:
    seen = []
    for s in slides:
        sec = s.get("section")
        if sec and sec not in seen:
            seen.append(sec)
    return seen


def _chrome_header(slide: dict, hh: int, proposer_logo: str | None,
                   proposer_name: str, flags: list[str]) -> str:
    if hh <= 0:
        return ""
    badge = f'<span class="badge">{_esc(slide.get("section"))}</span>' if slide.get("section") else ""
    title = f'<div class="title">{_esc(slide.get("title") or "")}</div>'
    subtitle = f'<div class="subtitle">{_esc(slide.get("message"))}</div>' if slide.get("message") else ""
    if proposer_logo:
        logo = f'<div class="logo"><img src="{proposer_logo}" alt="proposer"></div>'
    elif proposer_name:
        logo = f'<div class="proposer-name">{_esc(proposer_name)}</div>'
    else:
        logo = ""
    return (f'<header class="ck-h" style="height:{hh}px">'
            f'{badge}{title}{subtitle}{logo}{_flag_pills(flags)}</header>')


def _chrome_footer(fh: int, client_logo: str | None, project: str,
                   proposer_name: str, page: str) -> str:
    if fh <= 0:
        return ""
    logo = f'<img class="logo" src="{client_logo}" alt="client">' if client_logo else ""
    co = f'<span class="sep">|</span><span>{_esc(proposer_name)}</span>' if proposer_name else ""
    return (f'<footer class="ck-f" style="height:{fh}px">{logo}'
            f'<span class="proj">{_esc(project)}</span>{co}'
            f'<span class="page">{page}</span></footer>')


def _cover_html(slide: dict, meta: dict, run: Path, project: str, proposer_name: str,
                client_logo: str | None, proposer_logo: str | None,
                decor_html: str = "") -> str:
    fields = slide.get("fields") or {}
    bid_no = fields.get("bid_no") or meta.get("bid_no") or meta.get("bid") or ""
    date = fields.get("date") or meta.get("date") or ""
    title = slide.get("title") or fields.get("project_title") or project
    bg_src = _logo_src(fields.get("background_image"), run)
    bg = (f'<div class="bg" style="background-image:url({bg_src})"></div><div class="scrim"></div>'
          if bg_src else "")
    label = "PROPOSAL" + (f' · <span class="no">{_esc(bid_no)}</span>' if bid_no else "")
    meta_line = " · ".join(_esc(x) for x in (date, proposer_name) if x)
    logos = "".join(f'<img src="{src}" alt="logo">'
                    for src in (client_logo, proposer_logo) if src)
    logos_html = f'<div class="cv-logos">{logos}</div>' if logos else ""
    # decor_html은 배경 다음·표지 텍스트(cv-inner, z-index:2) 이전 - 항상 크롬/본문 아래(DF2).
    return (f'{bg}{decor_html}<div class="cv-inner"><div class="cv-label">{label}</div>'
            f'<div class="cv-title">{_esc(title)}</div><div class="cv-rule"></div>'
            + (f'<div class="cv-meta">{meta_line}</div>' if meta_line else "")
            + f'</div>{logos_html}')


def _toc_items(slide: dict, all_slides: list) -> list:
    """목차 항목 도출(HTML·pptx 공용): fields.items > storyline section 순서 자동 도출."""
    fields = slide.get("fields") or {}
    items = fields.get("items")
    rows = []
    if isinstance(items, list) and items:
        for i, it in enumerate(items, start=1):
            if isinstance(it, dict):
                rows.append((str(it.get("no") or i).zfill(2),
                             str(it.get("title") or ""),
                             str(it.get("sub") or it.get("subtitle") or "")))
            else:
                rows.append((str(i).zfill(2), str(it), ""))
    else:
        for i, sec in enumerate(_section_order(all_slides), start=1):
            rows.append((str(i).zfill(2), str(sec), ""))
    return rows


def _toc_html(slide: dict, all_slides: list) -> str:
    rows = _toc_items(slide, all_slides)
    cells = []
    for no, tt, ts in rows:
        sub = f'<div class="ts">{_esc(ts)}</div>' if ts else ""
        cells.append(f'<div class="toc-item"><div class="num">{_esc(no)}</div>'
                     f'<div><div class="tt">{_esc(tt)}</div>{sub}</div></div>')
    return f'<div class="toc-grid">{"".join(cells)}</div>'


def _divider_html(slide: dict, all_slides: list) -> str:
    fields = slide.get("fields") or {}
    sec = slide.get("section")
    order = _section_order(all_slides)
    no = fields.get("no") or ((order.index(sec) + 1) if sec in order else "")
    no_txt = str(no).zfill(2) if str(no).isdigit() else str(no)
    label = fields.get("label") or (f"SECTION {no_txt}" if no_txt else (sec or ""))
    sub = slide.get("message") or fields.get("sub") or ""
    ghost = f'<div class="dv-ghost">{_esc(no_txt)}</div>' if no_txt else ""
    return (f'{ghost}<div class="dv-inner"><div class="dv-label">{_esc(label)}</div>'
            f'<div class="dv-title">{_esc(slide.get("title") or sec or "")}</div>'
            + (f'<div class="dv-sub">{_esc(sub)}</div>' if sub else "")
            + '<div class="dv-rule"></div></div>')


def _compose_hybrid(run: Path, manifest: dict, skin: dict) -> dict:
    """(f-W29) 하이브리드 크롬 조합: 장 클래스별 HTML 셸 + 본문 이미지 -> deck.images.html."""
    slides_dir = run / SLIDES_DIR
    exp = manifest.get("export") or {}
    W, H = int(exp.get("w", 0)), int(exp.get("h", 0))
    project = manifest.get("project") or run.name

    try:
        storyline = _load_json(run / "storyline.json")
    except (OSError, json.JSONDecodeError):
        storyline = {}
    meta = storyline.get("meta") or {}
    all_slides = storyline.get("slides") or []
    # DF6: storyline 로드 직후 - bundle과 같은 헬퍼(단일 적용 지점). style_variant/
    # fields.background_image(안전 분류 ⓐ)는 이 재조립만으로 즉시 반영된다(재번들 불필요) -
    # chrome_override/deck_class(ⓑ)는 여기서도 얹히지만 크롬 높이·클래스·px는 manifest(bundle
    # 시점 값)에서 오므로, 실제 반영에는 재번들이 필요하다(위 bundle()의 동일 헬퍼 호출).
    all_slides, deck_overrides = _apply_deck_overrides(all_slides, run)
    by_n = {str(s.get("n")): s for s in all_slides}

    brand = _brand(run)
    client_logo = _logo_src(brand.get("client_logo"), run)
    proposer_logo = _logo_src(brand.get("proposer_logo"), run)
    proposer_name = brand.get("proposer_name") or ""
    placement = brand.get("placement") or {}
    client_on = placement.get("client", "all")
    proposer_on = placement.get("proposer", "all")

    # A/B 중 최종본 선택 — _select_final_slides(레거시와 동일 규칙 + 마찰23 콘텐츠 분할 이어붙이기).
    chosen = _select_final_slides(manifest.get("slides", []))

    # 바깥 프레임 띠: 슬라이드 padding = 띠, 안쪽은 .inner 래퍼(이미지 px는 bundle이 역산).
    frame_spec = (skin.get("chrome") or {}).get("frame") or {}
    frame_color = _hexcolor((skin.get("colors") or {}).get(frame_spec.get("color") or "", "")
                            or frame_spec.get("color"), "E8590C")
    # 바깥 배경이미지(덱 공통): 띠·본문 여백으로 자연스럽게 드러난다. 장별 = fields.background_image.
    frame_image = _logo_src(frame_spec.get("image"), run)

    total = len(chosen)
    used, missing = 0, []
    html_slides = 0
    sections = []
    warnings: list[str] = []
    for idx, s in enumerate(chosen, start=1):
        cls = s.get("deck_class") or "content"
        slide = by_n.get(str(s["n"])) or {}
        page = f"{idx:02d} / {total:02d}"
        vstyle = _variant_style(skin, slide)
        fw = int(s.get("frame_w") or 0)  # 바깥 프레임 띠 = padding(안쪽 전체가 줄어듦)
        style = (f"width:{W}px;height:{H}px;padding:{fw}px"
                 + (f";{vstyle}" if vstyle else ""))
        ch = s.get("chrome") or {}
        hh, fh = int(ch.get("header_h") or 0), int(ch.get("footer_h") or 0)
        cls_spec = (skin.get("slide_classes") or {}).get(cls) or {}
        bg_layer = _bg_layer(slide, run)  # HTML 장 전체 배경이미지(z-바닥) + 라이트 스크림
        # DF2: decor_slots(배경 외 장식) - 배경 위·크롬/본문 아래(항상 header/body보다 먼저 삽입).
        decor_html = _decor_layer_html(skin, cls, run, warnings)

        if s.get("render") == "html":
            html_slides += 1
            if cls == "cover":
                sc = client_on in ("all", "cover") and client_logo or None
                sp_ = proposer_on in ("all", "cover") and proposer_logo or None
                inner = _cover_html(slide, meta, run, project, proposer_name, sc, sp_,
                                    decor_html)
                body_cls = "slide cover"
            elif cls == "toc":
                inner = (bg_layer + decor_html
                         + _chrome_header(slide, hh, proposer_logo if proposer_on == "all" else None,
                                          proposer_name, s.get("flags") or [])
                         + _toc_html(slide, all_slides)
                         + _chrome_footer(fh, client_logo if client_on == "all" else None,
                                          project, proposer_name, page))
                body_cls = "slide toc"
            else:  # divider(기타 html 클래스 포함)
                inner = (bg_layer + decor_html + _divider_html(slide, all_slides)
                         + _chrome_footer(fh or int((skin.get("chrome") or {}).get("footer_h") or 0),
                                          None, project, proposer_name, page))
                body_cls = "slide divider"
            sections.append(f'<section class="{body_cls}" style="{style}">'
                            f'<div class="inner">{inner}</div></section>')
            continue

        img = slides_dir / s["out_name"]
        if img.exists():
            used += 1
            img_html = f'<img class="content" src="{img.resolve().as_uri()}" alt="slide {s["n"]}">'
        else:
            missing.append(s["out_name"])
            img_html = (f'<div class="content missing">이미지 없음: {SLIDES_DIR}/{_esc(s["out_name"])}'
                        f'<br>(imagedeck --collect 로 확인)</div>')

        if s.get("image_scope") == "full":
            overlay = ""
            if cls_spec.get("footer") == "overlay":
                overlay = (f'<div class="overlay-f"><span>{_esc(project)}</span>'
                           f'<span class="page">{page}</span></div>')
            sections.append(f'<section class="slide full" style="{style}"><div class="inner">'
                            f'{decor_html}<div class="body">{img_html}</div>{overlay}</div></section>')
        else:
            header = _chrome_header(slide, hh, proposer_logo if proposer_on == "all" else None,
                                    proposer_name, s.get("flags") or [])
            footer = _chrome_footer(fh, client_logo if client_on == "all" else None,
                                    project, proposer_name, page)
            bm = s.get("body_margin") or {}
            body_style = (f'padding:{int(bm.get("top") or 0)}px {int(bm.get("x") or 0)}px '
                          f'{int(bm.get("bottom") or 0)}px')
            # 바깥 배경(장별 > 덱 공통): .inner를 투명하게 - 띠·본문 여백으로 배경이 비친다.
            outer_bg = _logo_src((slide.get("fields") or {}).get("background_image"), run) \
                or frame_image
            sec_style, inner_open = style, '<div class="inner">'
            if outer_bg:
                sec_style += (f";background-image:url({outer_bg});"
                              "background-size:cover;background-position:center")
                inner_open = '<div class="inner" style="background:transparent">'
            sections.append(f'<section class="slide" style="{sec_style}">{inner_open}{decor_html}'
                            f'{header}<div class="body" style="{body_style}">{img_html}</div>'
                            f'{footer}</div></section>')

    typo = skin.get("typography") or {}
    family = typo.get("family") or "Pretendard"
    colors = skin.get("colors") or {}
    root_vars = ";".join(f"{cssvar}:{_hexcolor(colors.get(key), d)}" for key, cssvar, d in _VAR_MAP)
    _so, _sc, _scss = _stage_wrap(W, H)   # W32: 화면 맞춤 + 화면 단위 페이징 + 두 창 동기화
    html = _HYBRID_SHELL.format(project=_esc(project), root_vars=root_vars, family=family,
                                frame_color=frame_color,
                                slides="\n".join(sections), stage_open=_so, stage_close=_sc, stage_css=_scss)
    out = run / COMPOSE_HTML
    out.write_text(html, encoding="utf-8")
    ch0 = manifest.get("chrome") or {}
    # decor_slots는 덱 전체 선언(장별 아님) - 같은 슬롯이 같은 이유로 매 장 경고를 반복하면
    # 노이즈만 커진다. 순서 보존 중복 제거(서로 다른 슬롯/사유는 각각 그대로 남는다).
    warnings = list(dict.fromkeys(warnings))
    return {"out": str(out), "slides": total, "images_used": used, "missing": missing,
            "html_slides": html_slides, "warnings": warnings,
            "chrome": {"header_h": int(ch0.get("header_h") or 0),
                       "footer_h": int(ch0.get("footer_h") or 0)},
            # DF6: 빈 리스트도 항상 기록("0건"과 "미측정" 구분 - manifest와 같은 관례).
            "deck_overrides_applied": [d["n"] for d in deck_overrides],
            "deck_overrides": deck_overrides}


_HYBRID_SHELL = """<!doctype html>
<html lang="ko"><head><meta charset="utf-8"><title>{project} - 이미지 덱</title>
<style>
  :root {{ {root_vars}; }}
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ background:#222; font-family:'{family}','Noto Sans KR','Malgun Gothic',sans-serif; }}
  /* 레이어 구조: .slide(padding=바깥 프레임 띠, 배경=띠 색) > .inner(실면) > [.bg 배경층] > 크롬/본문.
     띠는 오버레이가 아니라 inset - 안쪽 전체(크롬+이미지 px)가 bundle 역산으로 함께 줄어든다. */
  .slide {{ position:relative; margin:0 auto 12px; background:{frame_color}; overflow:hidden; }}
  .inner {{ position:relative; width:100%; height:100%; background:var(--bg); overflow:hidden;
    display:flex; flex-direction:column; }}
  .bg {{ position:absolute; inset:0; background-size:cover; background-position:center; }}
  .scrim.light {{ position:absolute; inset:0; background:rgba(248,247,244,.82); }}
  /* 크롬 헤더: 섹션배지 - 제목 - 부제, 로고 우상단(고정틀), flag 딱지.
     하단 구분선 2px + 숨통(헤더 184px 기준 부제 아래 ~25px) - 2026-07-20 검토 반영 */
  .ck-h {{ flex:none; position:relative; padding:24px 64px 0; background:var(--bg);
    border-bottom:2px solid var(--line-strong); }}
  .ck-h .badge {{ display:inline-block; background:var(--ink); color:#fff; font-size:22px;
    font-weight:700; padding:8px 26px; border-radius:6px; letter-spacing:.06em; }}
  .ck-h .title {{ font-size:44px; font-weight:800; color:var(--ink-deep); margin-top:14px;
    line-height:1.15; }}
  .ck-h .subtitle {{ font-size:21px; color:var(--gray); margin-top:8px; font-weight:500; }}
  .ck-h .logo {{ position:absolute; top:30px; right:64px; height:52px; }}
  .ck-h .logo img {{ height:52px; width:auto; }}
  .ck-h .proposer-name {{ position:absolute; top:38px; right:64px; font-size:20px;
    font-weight:700; color:var(--ink); }}
  .flags {{ position:absolute; right:64px; top:96px; display:flex; gap:8px; }}
  .flag {{ font-size:16px; font-weight:700; padding:4px 14px; border-radius:999px;
    background:#FFF1E6; color:#C4441C; border:1px solid #F5C6A5; }}
  .flag.example {{ background:#EEF3F9; color:#3A5A80; border-color:#C6D4E4; }}
  /* 본문 = 생성 이미지 */
  .body {{ flex:1; display:flex; align-items:center; justify-content:center; min-height:0; }}
  .content {{ display:block; width:100%; height:100%; object-fit:contain; }}
  .content.missing {{ color:#a33; font-size:20px; font-weight:500; text-align:center;
    display:flex; flex-direction:column; justify-content:center; }}
  /* 크롬 푸터: 발주처 로고 - 프로젝트명 - 제안사명 | 페이지 */
  .ck-f {{ flex:none; position:relative; display:flex; align-items:center; padding:0 64px;
    background:var(--bg); border-top:2px solid var(--line-strong);
    font-size:18px; color:var(--gray); }}
  .ck-f .logo {{ height:34px; width:auto; margin-right:18px; }}
  .ck-f .proj {{ font-weight:700; color:var(--ink); }}
  .ck-f .sep {{ margin:0 14px; color:var(--line); }}
  .ck-f .page {{ margin-left:auto; font-weight:700; color:var(--ink); letter-spacing:.08em; }}
  /* full_image: 이미지 전체 + 오버레이 푸터 */
  .slide.full .inner {{ display:block; }}
  .slide.full .body {{ position:absolute; inset:0; }}
  .overlay-f {{ position:absolute; left:0; right:0; bottom:0; height:64px; z-index:5;
    background:linear-gradient(transparent, rgba(15,36,56,.55)); display:flex;
    align-items:center; padding:0 64px; color:#fff; font-size:18px; }}
  .overlay-f .page {{ margin-left:auto; font-weight:700; letter-spacing:.08em; }}
  /* cover */
  .slide.cover .inner {{ display:block; color:#fff;
    background:linear-gradient(135deg, var(--ink-deep) 0%, var(--ink) 70%); }}
  .cover .scrim {{ position:absolute; inset:0; background:rgba(15,36,56,.55); }}
  .cv-inner {{ position:absolute; left:110px; bottom:96px; z-index:2; }}
  .cv-label {{ font-size:22px; letter-spacing:.28em; font-weight:700; color:rgba(255,255,255,.75); }}
  .cv-label .no {{ color:var(--accent); }}
  .cv-title {{ font-size:76px; font-weight:800; line-height:1.18; margin-top:26px; }}
  .cv-rule {{ width:120px; height:4px; background:var(--accent); margin-top:38px; }}
  .cv-meta {{ margin-top:44px; font-size:22px; color:rgba(255,255,255,.85); }}
  .cv-logos {{ position:absolute; right:96px; bottom:88px; display:flex; gap:20px; z-index:2; }}
  .cv-logos img {{ height:52px; width:auto; background:rgba(255,255,255,.92);
    padding:6px 10px; border-radius:6px; }}
  /* toc */
  .toc-grid {{ flex:1; position:relative; display:grid; grid-template-columns:1fr 1fr;
    gap:34px 56px; padding:56px 110px; align-content:center; min-height:0; }}
  .toc-item {{ display:flex; gap:24px; align-items:flex-start;
    border-bottom:1px solid var(--line); padding-bottom:26px; }}
  .toc-item .num {{ width:56px; height:56px; border-radius:12px; background:var(--ink);
    color:#fff; display:flex; align-items:center; justify-content:center;
    font-size:26px; font-weight:800; flex:none; }}
  .toc-item .tt {{ font-size:28px; font-weight:700; color:var(--ink-deep); }}
  .toc-item .ts {{ font-size:19px; color:var(--gray); margin-top:6px; }}
  /* divider(간지): 대형 숫자 워터마크 + 좌하단 섹션, 푸터만 */
  .slide.divider .inner {{ display:block; }}
  .dv-ghost {{ position:absolute; right:-30px; top:-60px; font-size:560px; font-weight:800;
    color:rgba(22,50,76,.06); line-height:1; }}
  .dv-inner {{ position:absolute; left:110px; bottom:190px; }}
  .dv-label {{ font-size:24px; letter-spacing:.26em; color:var(--accent); font-weight:800; }}
  .dv-title {{ font-size:64px; font-weight:800; color:var(--ink-deep); margin-top:22px; }}
  .dv-sub {{ font-size:23px; color:var(--gray); margin-top:16px; }}
  .dv-rule {{ width:96px; height:4px; background:var(--ink); margin-top:34px; }}
  .slide.divider .ck-f {{ position:absolute; left:0; right:0; bottom:0; }}
{stage_css}
</style></head><body>
{stage_open}
{slides}
{stage_close}
</body></html>
"""


# ---------------------------------------------------------------------------
# DF4(2026-07-24, CONTEXT/DECK_FIRST_DESIGN.md §2-③·§3 DF4 행) — 덱 프리뷰 렌더.
#
# 계약 동결 후 "틀+배경(본문 비움)" 완성 슬라이드를 장 클래스별로 1장씩 실제로 굽는다
# (_compose_hybrid과 같은 조립 부품 재사용 — 크롬 헤더/푸터·decor_slots·바깥 배경·_HYBRID_SHELL).
# 본문 이미지가 있는 클래스(image=body/full)만 대상 — cover/toc/divider(image=none)는 프롬프트
# 자체가 없어 프리뷰가 무의미하므로 건너뛴다. HTML→PNG는 playwright(rasterize.html_to_png)를
# 정식 채택(§4-1, 2026-07-24) — 미설치면 우아 강등이 아니라 명확한 오류로 중단한다(설치 안내 포함).
# ---------------------------------------------------------------------------


def _viewer_mod():
    """app/render/viewer 지연 로드(`_rasterize()`와 같은 패턴). 부재 시 None — 우아 강등."""
    try:
        from render import viewer  # type: ignore
    except ImportError:
        import sys as _s
        _s.path.insert(0, str(Path(__file__).resolve().parents[2] / "app"))
        try:
            from render import viewer  # type: ignore
        except ImportError:  # pragma: no cover - 배포 변형 방어
            return None
    return viewer


def _stage_wrap(width: int, height: int, *, enabled: bool = True) -> "tuple[str, str, str]":
    """(여는 태그, 닫는 태그+스크립트, 셸 CSS) — W32 마찰35의 deck-stage 채택분.

    `enabled=False`면 전부 빈 문자열이다. 래스터라이즈(PNG 생성) 경로가 그 경우인데,
    deck-stage는 활성 장만 보이게 하고 화면에 맞춰 축소하므로 **실물 px가 어긋난다**.
    """
    if not enabled:
        return ("", "", "")
    viewer = _viewer_mod()
    if viewer is None:
        return ("", "", "")
    return (viewer.stage_open(width, height), viewer.stage_close(), viewer.STAGE_CSS)




def _rasterize():
    """app/render/rasterize 지연 로드(playwright는 이 기능에서만 필요 — `_prim()`과 같은 패턴)."""
    try:
        from render import rasterize as rz  # type: ignore
    except ImportError:
        import sys as _s
        _s.path.insert(0, str(Path(__file__).resolve().parents[2] / "app"))
        from render import rasterize as rz  # type: ignore
    return rz


def _preview_dummy_slide(cls: str) -> dict:
    """프리뷰용 표본 슬라이드 — 실 콘텐츠 없음(크롬이 어떻게 그려지는지 보여주는 자리만)."""
    return {"n": "P", "title": "덱 프리뷰 - 본문 예시 제목", "message": "부제 예시 텍스트",
            "section": "SAMPLE", "deck_class": cls, "fields": {}, "flags": []}


def _preview_placeholder_html(cls: str) -> str:
    """본문 자리 = 빈 플레이스홀더(점선 테두리 + 클래스명 라벨) — 실제 생성 이미지가 여기 낀다."""
    return (f'<div class="content placeholder" style="display:flex;align-items:center;'
            f'justify-content:center;width:100%;height:100%;'
            f'border:3px dashed rgba(120,120,120,.5);background:rgba(120,120,120,.06);'
            f'font-size:26px;font-weight:700;color:rgba(90,90,90,.55)">'
            f'[본문 이미지 자리 - {_esc(cls)}]</div>')


def _preview_section(cls: str, chrome_skin: dict, dims: dict, run: Path,
                     frame_image: str | None, warnings: list[str]) -> "tuple[str, dict] | None":
    """장 클래스 1개 -> 완성 슬라이드 섹션 HTML(본문=플레이스홀더). image=none 클래스는 None.

    `_compose_hybrid`의 content/full_image 분기와 같은 조립 부품(크롬·decor·바깥 배경)을 그대로
    쓴다 — 실제 조립과 프리뷰가 어긋나지 않게(L10류 이중 주입 재발 방지, DF2 docstring과 같은 취지).
    """
    dummy = _preview_dummy_slide(cls)
    layout = resolve_slide_layout(chrome_skin, dims, dummy)
    if layout["image"] == "none":
        return None
    exp = dims["export"]
    W, H = exp["w"], exp["h"]
    fw = layout["frame_w"]
    style = f"width:{W}px;height:{H}px;padding:{fw}px"
    hh, fh = layout["chrome"]["header_h"], layout["chrome"]["footer_h"]
    decor_html = _decor_layer_html(chrome_skin, cls, run, warnings)
    placeholder = _preview_placeholder_html(cls)
    if layout["image"] == "full":
        overlay = ""
        cls_spec = (chrome_skin.get("slide_classes") or {}).get(cls) or {}
        if cls_spec.get("footer") == "overlay":
            overlay = ('<div class="overlay-f"><span>DECK PREVIEW</span>'
                       '<span class="page">-- / --</span></div>')
        section = (f'<section class="slide full" style="{style}"><div class="inner">'
                  f'{decor_html}<div class="body">{placeholder}</div>{overlay}</div></section>')
    else:
        header = _chrome_header(dummy, hh, None, "", [])
        footer = _chrome_footer(fh, None, "DECK PREVIEW", "", "-- / --")
        bm = layout["body_margin"]
        body_style = f'padding:{bm["top"]}px {bm["x"]}px {bm["bottom"]}px'
        sec_style, inner_open = style, '<div class="inner">'
        if frame_image:
            sec_style += (f";background-image:url({frame_image});"
                          "background-size:cover;background-position:center")
            inner_open = '<div class="inner" style="background:transparent">'
        section = (f'<section class="slide" style="{sec_style}">{inner_open}{decor_html}'
                  f'{header}<div class="body" style="{body_style}">{placeholder}</div>'
                  f'{footer}</div></section>')
    return section, layout


def render_deck_preview(run: Path) -> dict:
    """DF4: 계약 동결 후 "틀+배경(본문 비움)" 프리뷰 PNG를 장 클래스별로 렌더 ->
    `imagedeck_refs/deck_preview/<class>.png`. `resolve_slide_refs`의 4계층 조회(slide>global>
    deck_preview>seed)가 이 산출을 이후 bundle 실행마다 자동 소비한다.

    전제(둘 다 필수 — 하나라도 없으면 사람 말 오류로 중단, 우아 강등 아님):
    ①design_contract.json 동결(chrome_contract.slide_classes 선언), ②playwright 설치.
    """
    run = Path(run)
    contract = design_contract.load(run)
    if contract is None:
        raise ImagedeckError(
            "design_contract.json 없음 - 계약이 동결되지 않았다(DF4 전제조건). "
            "먼저 go로 디자인 계약을 동결하라(CONTEXT/DECK_FIRST_DESIGN.md §3 DF4 행)."
        )
    chrome_skin = contract.get("chrome_contract") or {}
    dims = canvas_dims(chrome_skin)
    classes = chrome_skin.get("slide_classes") or {}
    if not classes:
        raise ImagedeckError(
            "chrome_contract.slide_classes 없음 - 장 클래스 선언이 없는 계약이라 덱 프리뷰를 "
            "만들 클래스가 없다(DF4는 slide_classes 선언 스킨 전용, W29 하이브리드 크롬)."
        )
    rz = _rasterize()
    if not rz.available():
        raise ImagedeckError(
            "playwright 미설치 - 덱 프리뷰 렌더에는 HTML->PNG 수단(playwright)이 필요하다. "
            "설치: pip install playwright && playwright install chromium "
            "(2026-07-24 정식 채택 결정, CONTEXT/DECK_FIRST_DESIGN.md §4-1)."
        )

    out_dir = run / "imagedeck_refs" / DECK_PREVIEW_DIRNAME
    out_dir.mkdir(parents=True, exist_ok=True)
    for old in _scan_ref_images(out_dir):
        Path(old).unlink()  # 스킨/계약 변경 후 잔재 방지(재렌더 때마다 클린 상태에서 시작)

    frame_spec = (chrome_skin.get("chrome") or {}).get("frame") or {}
    frame_image = _logo_src(frame_spec.get("image"), run)
    typo = chrome_skin.get("typography") or {}
    family = typo.get("family") or "Pretendard"
    colors = chrome_skin.get("colors") or {}
    root_vars = ";".join(f"{cssvar}:{_hexcolor(colors.get(key), d)}" for key, cssvar, d in _VAR_MAP)
    frame_color = _hexcolor((colors.get(frame_spec.get("color") or "")) or frame_spec.get("color"),
                            "E8590C")

    warnings: list[str] = []
    rendered: list[dict] = []
    skipped: list[str] = []
    exp = dims["export"]
    for cls in classes:
        result = _preview_section(cls, chrome_skin, dims, run, frame_image, warnings)
        if result is None:
            skipped.append(cls)
            continue
        section_html, layout = result
        html = _HYBRID_SHELL.format(project=_esc("Deck Preview"), root_vars=root_vars,
                                    family=family, frame_color=frame_color, slides=section_html,
                                    stage_open="", stage_close="", stage_css="")   # PNG 래스터라이즈 - deck-stage 비활성
        out_png = out_dir / f"{cls}.png"
        rz.html_to_png(html, out_png, width_px=exp["w"], height_px=exp["h"])
        rendered.append({"class": cls, "image_scope": layout["image"], "out": str(out_png)})

    warnings = list(dict.fromkeys(warnings))
    return {"out_dir": str(out_dir), "rendered": rendered, "skipped_classes": skipped,
            "warnings": warnings}


# ---------------------------------------------------------------------------
# compose_pptx (W30) — 하이브리드 pptx: 크롬·HTML 장 = 네이티브(수정 가능), 본문 = 이미지
# ---------------------------------------------------------------------------

COMPOSE_PPTX = "deck.images.pptx"


def _prim():
    """app/render/pptx_primitives 지연 로드(python-pptx는 이 기능에서만 필요)."""
    try:
        from render import pptx_primitives as prim  # type: ignore
    except ImportError:
        import sys
        sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "app"))
        from render import pptx_primitives as prim  # type: ignore
    return prim


def _slide_colors(skin: dict, slide: dict) -> dict:
    """스킨 colors + slide.style_variant 오버라이드(compose HTML의 variants와 동일 규칙)."""
    colors = {k: str(v).lstrip("#") for k, v in (skin.get("colors") or {}).items()
              if isinstance(v, str)}
    name = slide.get("style_variant")
    var = (skin.get("variants") or {}).get(str(name)) if name else None
    if isinstance(var, dict):
        for k, v in (var.get("colors") or {}).items():
            colors[k] = str(v).lstrip("#")
    return colors


def _fill_alpha(shape, alpha: float) -> None:
    """solid fill에 투명도(0=투명, 1=불투명). python-pptx 미노출 - XML 직접, 실패 시 무시."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        clr = shape.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        etree.SubElement(clr, qn("a:alpha")).set("val", str(int(alpha * 100000)))
    except Exception:
        pass


def _picture_alpha(picture, alpha: float) -> None:
    """그림(Picture) 투명도(0=투명, 1=불투명) - blipFill에 alphaModFix. 최선 노력, 실패 시 무시
    (decor_slots opacity는 선택 필드 - 실패해도 그림 자체는 정상 배치되므로 조립을 막지 않는다)."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        blip = picture._element.blipFill.find(qn("a:blip"))
        etree.SubElement(blip, qn("a:alphaModFix")).set("amt", str(int(alpha * 100000)))
    except Exception:
        pass


def _decor_pptx(sl, px, W: int, H: int, fw: int, resolved: list[dict]) -> None:
    """decor_slots -> python-pptx 그림 삽입(anchor -> EMU 좌표 환산). 호출 시점 = 배경 레이어
    직후·크롬/본문 삽입 이전(z-순서: shape 추가 순서가 곧 렌더 순서 - 나중에 추가한 도형이 위).
    """
    from pptx.util import Emu
    for d in resolved:
        width_px = d["width"]
        dims = png_dims(d["path"])
        if dims and dims[0]:
            height_px = width_px * dims[1] / dims[0]
        else:
            height_px = width_px  # PNG 헤더 파싱 실패(비-PNG 등) - 종횡비 추정 불가, 정사각형 가정
        anchor = d["anchor"]
        if anchor == "top-left":
            x, y = fw + d["offset_x"], fw + d["offset_y"]
        elif anchor == "top-right":
            x, y = W - fw - d["offset_x"] - width_px, fw + d["offset_y"]
        elif anchor == "bottom-left":
            x, y = fw + d["offset_x"], H - fw - d["offset_y"] - height_px
        else:  # bottom-right
            x, y = W - fw - d["offset_x"] - width_px, H - fw - d["offset_y"] - height_px
        pic = sl.shapes.add_picture(str(d["path"]), Emu(int(px(x) * 914400)),
                                    Emu(int(px(y) * 914400)),
                                    width=Emu(int(px(width_px) * 914400)),
                                    height=Emu(int(px(height_px) * 914400)))
        if d["opacity"] < 1.0:
            _picture_alpha(pic, d["opacity"])


def compose_pptx(run: Path) -> dict:
    """(f-pptx) 하이브리드 pptx 조합 -> deck.images.pptx.

    deck.images.html과 같은 정본(manifest·storyline·스킨·brand)에서 결정론 조립하되,
    크롬(헤더·푸터)과 HTML 전용 장(cover/toc/divider)은 **네이티브 텍스트박스·도형**으로
    그려 PowerPoint에서 수정 가능하고, 본문 콘텐츠만 이미지로 삽입한다.
    slide_classes 미선언 스킨(레거시)은 지원하지 않는다 - ship --pptx-mode image를 쓴다.
    """
    run = Path(run)
    manifest_path = run / MANIFEST_NAME
    if not manifest_path.exists():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(manifest_path)
    if not manifest.get("classes_enabled"):
        raise ImagedeckError("compose_pptx는 slide_classes 스킨(하이브리드) 전용 - "
                             "레거시 스킨은 ship --pptx-mode image를 쓰라.")
    # W31 R5: design_contract가 있으면 chrome_contract만 읽는다(_resolve_chrome_skin 참고).
    skin = _resolve_chrome_skin(run, manifest)

    prim = _prim()
    from pptx import Presentation
    from pptx.util import Emu

    exp = manifest.get("export") or {}
    W, H = int(exp.get("w", 1920)), int(exp.get("h", 1080))
    S = prim.EMU_W_IN / W  # px -> inch (16:9 균등)

    def px(v: float) -> float:
        return v * S

    def pt(v: float) -> float:
        return round(v * 0.75, 1)  # px(96dpi) -> pt

    try:
        storyline = _load_json(run / "storyline.json")
    except (OSError, json.JSONDecodeError):
        storyline = {}
    meta = storyline.get("meta") or {}
    all_slides = storyline.get("slides") or []
    # DF6: compose(HTML)와 동일한 공용 헬퍼·적용 지점(_compose_hybrid 참고 - 두 렌더러가 어긋나지
    # 않게 storyline 로드 직후 바로 통과시킨다).
    all_slides, deck_overrides = _apply_deck_overrides(all_slides, run)
    by_n = {str(s.get("n")): s for s in all_slides}
    project = manifest.get("project") or run.name

    brand = _brand(run)
    client_logo = brand.get("client_logo")
    proposer_logo = brand.get("proposer_logo")
    client_logo_p = Path(client_logo) if client_logo else None
    if client_logo_p and not client_logo_p.is_absolute():
        client_logo_p = run / client_logo
    proposer_logo_p = Path(proposer_logo) if proposer_logo else None
    if proposer_logo_p and not proposer_logo_p.is_absolute():
        proposer_logo_p = run / proposer_logo
    client_logo_p = client_logo_p if (client_logo_p and client_logo_p.is_file()) else None
    proposer_logo_p = proposer_logo_p if (proposer_logo_p and proposer_logo_p.is_file()) else None
    proposer_name = brand.get("proposer_name") or ""
    family = (skin.get("typography") or {}).get("family") or "Pretendard"

    frame_spec = (skin.get("chrome") or {}).get("frame") or {}
    frame_color_name = frame_spec.get("color") or "accent"

    # A/B 최종본 선택(HTML compose와 동일 규칙 — _select_final_slides)
    chosen = _select_final_slides(manifest.get("slides", []))

    prs = Presentation()
    prs.slide_width = Emu(int(prim.EMU_W_IN * 914400))
    prs.slide_height = Emu(int(prim.EMU_H_IN * 914400))
    slides_dir = run / SLIDES_DIR
    total = len(chosen)
    used, missing = 0, []
    html_native = 0
    warnings: list[str] = []

    def _chrome_header_pptx(sl, colors, fw, hh, slide_data, flags):
        c = colors.get
        if hh <= 0:
            return
        sec = slide_data.get("section")
        if sec:
            bw = max(118, 52 + int(len(str(sec)) * 22 * 1.05))
            prim.add_rounded(sl, px(fw + 64), px(fw + 26), px(bw), px(44),
                             fill=c("section_badge", c("ink", "16324C")))
            prim.add_text(sl, px(fw + 64), px(fw + 26), px(bw), px(44), str(sec),
                          size=pt(22), color="FFFFFF", family=family, bold=True,
                          align="center", anchor="middle")
        # 제목 폭: 우측 예약(로고/제안사명)이 있을 때만 좁힌다 - HTML 크롬과 동일 전폭 원칙.
        right_reserve = 250 if proposer_logo_p else (400 if proposer_name else 64)
        title = slide_data.get("title") or ""
        tw = W - 2 * fw - 64 - right_reserve
        # 한 줄 안전망: 전각 기준 수용량 초과 시 폰트 단계 축소(부제와의 겹침 방지).
        tsize = 44
        for cand in (44, 36, 30):
            tsize = cand
            if len(title) * cand <= tw:
                break
        prim.add_text(sl, px(fw + 64), px(fw + 78), px(tw), px(58), title,
                      size=pt(tsize), color=c("ink_deep", "0F2438"),
                      family=family, bold=True)
        if slide_data.get("message"):
            # 폰트 폴백 시 줄높이 팽창을 감안해 HTML(140)보다 6px 아래 - 제목 하단과 겹침 방지.
            prim.add_text(sl, px(fw + 64), px(fw + 146), px(W - 2 * fw - 128), px(30),
                          slide_data.get("message"), size=pt(21),
                          color=c("gray_text", "5C6470"), family=family)
        if proposer_logo_p:
            sl.shapes.add_picture(str(proposer_logo_p), Emu(int(px(W - fw - 64 - 170) * 914400)),
                                  Emu(int(px(fw + 30) * 914400)),
                                  height=Emu(int(px(52) * 914400)))
        elif proposer_name:
            prim.add_text(sl, px(W - fw - 64 - 320), px(fw + 34), px(320), px(30),
                          proposer_name, size=pt(20), color=c("ink", "16324C"),
                          family=family, bold=True, align="right")
        # flag 딱지(HTML과 동일: 이미지에 굽지 않는다)
        fx = W - fw - 64
        for f in reversed(flags or []):
            label = str(f).strip("[]")
            pw = 30 + int(len(label) * 16 * 1.05)
            fx -= pw + 8
            fill, txt = ("EEF3F9", "3A5A80") if label == "예시" else ("FFF1E6", "C4441C")
            prim.add_rounded(sl, px(fx), px(fw + 96), px(pw), px(30), fill=fill)
            prim.add_text(sl, px(fx), px(fw + 96), px(pw), px(30), label, size=pt(16),
                          color=txt, family=family, bold=True, align="center", anchor="middle")
        prim.add_line(sl, px(fw), px(fw + hh), px(W - fw), px(fw + hh),
                      color=c("line_strong", "C9D1DA"), width_pt=1.5)

    def _chrome_footer_pptx(sl, colors, fw, fh, page, with_client=True):
        c = colors.get
        if fh <= 0:
            return
        fy = H - fw - fh
        prim.add_line(sl, px(fw), px(fy), px(W - fw), px(fy),
                      color=c("line_strong", "C9D1DA"), width_pt=1.5)
        tx = fw + 64
        if with_client and client_logo_p:
            sl.shapes.add_picture(str(client_logo_p), Emu(int(px(tx) * 914400)),
                                  Emu(int(px(fy + (fh - 34) / 2) * 914400)),
                                  height=Emu(int(px(34) * 914400)))
            tx += 128
        label = project + (f"   |   {proposer_name}" if proposer_name else "")
        prim.add_text(sl, px(tx), px(fy), px(W - fw - 300 - tx), px(fh), label,
                      size=pt(18), color=c("ink", "16324C"), family=family,
                      bold=True, anchor="middle")
        prim.add_text(sl, px(W - fw - 64 - 220), px(fy), px(220), px(fh), page,
                      size=pt(18), color=c("ink", "16324C"), family=family,
                      bold=True, align="right", anchor="middle")

    for idx, s in enumerate(chosen, start=1):
        cls = s.get("deck_class") or "content"
        slide_data = by_n.get(str(s["n"])) or {}
        colors = _slide_colors(skin, slide_data)
        c = colors.get
        page = f"{idx:02d} / {total:02d}"
        ch = s.get("chrome") or {}
        hh, fh = int(ch.get("header_h") or 0), int(ch.get("footer_h") or 0)
        fw = int(s.get("frame_w") or 0)
        sl = prim.blank_slide(prs)

        # 레이어 1: 바깥 띠/배경 + 실면
        fields = slide_data.get("fields") or {}
        outer_bg = fields.get("background_image") or frame_spec.get("image")
        outer_bg_p = None
        if outer_bg and cls not in ("cover",):  # cover는 자체 배경 처리
            p = Path(outer_bg)
            outer_bg_p = (p if p.is_absolute() else run / outer_bg)
            outer_bg_p = outer_bg_p if outer_bg_p.is_file() else None
        if outer_bg_p and cls == "content":
            sl.shapes.add_picture(str(outer_bg_p), 0, 0,
                                  width=prs.slide_width, height=prs.slide_height)
            prim.add_rect(sl, px(fw), px(fw), px(W - 2 * fw), px(hh),
                          fill=c("bg", "F8F7F4"))
            if fh > 0:
                prim.add_rect(sl, px(fw), px(H - fw - fh), px(W - 2 * fw), px(fh),
                              fill=c("bg", "F8F7F4"))
        else:
            if fw > 0:
                prim.add_rect(sl, 0, 0, prim.EMU_W_IN, prim.EMU_H_IN,
                              fill=c(frame_color_name, c("accent", "E8590C")))
            prim.add_rect(sl, px(fw), px(fw), px(W - 2 * fw), px(H - 2 * fw),
                          fill=c("bg", "F8F7F4"))

        # DF2: decor_slots(배경 외 장식) - 배경 위·크롬/본문 아래. cover는 자체 배경(그림/그라디언트+
        # 스크림)이 아래에서 더 그려지므로 그 뒤(제목/로고보다 앞)에 따로 삽입한다(HTML과 동일 순서).
        decor_slots = _resolve_decor_slots(skin, cls, run, warnings)
        if cls != "cover":
            _decor_pptx(sl, px, W, H, fw, decor_slots)

        if cls == "cover":
            html_native += 1
            bg_p = None
            if fields.get("background_image"):
                p = Path(fields["background_image"])
                bg_p = p if p.is_absolute() else run / fields["background_image"]
                bg_p = bg_p if bg_p.is_file() else None
            if bg_p:
                sl.shapes.add_picture(str(bg_p), 0, 0,
                                      width=prs.slide_width, height=prs.slide_height)
                scrim = prim.add_rect(sl, 0, 0, prim.EMU_W_IN, prim.EMU_H_IN,
                                      fill=c("ink_deep", "0F2438"))
                _fill_alpha(scrim, 0.55)
            else:
                prim.add_gradient(sl, 0, 0, prim.EMU_W_IN, prim.EMU_H_IN,
                                  [(0.0, c("ink_deep", "0F2438")), (1.0, c("ink", "16324C"))],
                                  angle=135.0)
            _decor_pptx(sl, px, W, H, fw, decor_slots)
            bid_no = fields.get("bid_no") or meta.get("bid_no") or meta.get("bid") or ""
            label = "PROPOSAL" + (f" · {bid_no}" if bid_no else "")
            prim.add_text(sl, px(110), px(H - 380), px(W - 400), px(34), label,
                          size=pt(22), color=c("accent", "E8590C"), family=family, bold=True)
            tb = prim.add_textbox(sl, px(110), px(H - 336), px(W - 400), px(120))
            for i, line in enumerate(str(slide_data.get("title") or project).splitlines()):
                prim.write_para(tb.text_frame, line, size=pt(76), color="FFFFFF",
                                family=family, bold=True, first=(i == 0))
            prim.add_rect(sl, px(110), px(H - 196), px(120), px(4),
                          fill=c("accent", "E8590C"))
            date = fields.get("date") or meta.get("date") or ""
            meta_line = " · ".join(x for x in (date, proposer_name) if x)
            if meta_line:
                prim.add_text(sl, px(110), px(H - 160), px(W - 400), px(30), meta_line,
                              size=pt(22), color="E7EBF0", family=family)
            lx = W - 96
            for lp in (proposer_logo_p, client_logo_p):
                if lp:
                    lx -= 190
                    sl.shapes.add_picture(str(lp), Emu(int(px(lx) * 914400)),
                                          Emu(int(px(H - 140) * 914400)),
                                          height=Emu(int(px(52) * 914400)))
            continue

        if cls == "toc":
            html_native += 1
            _chrome_header_pptx(sl, colors, fw, hh, slide_data, s.get("flags") or [])
            rows = _toc_items(slide_data, all_slides)
            x0, col_gap = fw + 110, 56
            col_w = (W - 2 * fw - 220 - col_gap) / 2
            y0, row_h = fw + hh + 90, 118
            for i, (no, tt, ts) in enumerate(rows):
                cx = x0 + (i % 2) * (col_w + col_gap)
                cy = y0 + (i // 2) * row_h
                prim.add_rounded(sl, px(cx), px(cy), px(56), px(56), fill=c("ink", "16324C"))
                prim.add_text(sl, px(cx), px(cy), px(56), px(56), no, size=pt(26),
                              color="FFFFFF", family=family, bold=True,
                              align="center", anchor="middle")
                prim.add_text(sl, px(cx + 80), px(cy), px(col_w - 80), px(34), tt,
                              size=pt(28), color=c("ink_deep", "0F2438"),
                              family=family, bold=True)
                if ts:
                    prim.add_text(sl, px(cx + 80), px(cy + 40), px(col_w - 80), px(26), ts,
                                  size=pt(19), color=c("gray_text", "5C6470"), family=family)
                prim.add_line(sl, px(cx), px(cy + row_h - 26), px(cx + col_w),
                              px(cy + row_h - 26), color=c("line", "DDE1E6"), width_pt=0.75)
            _chrome_footer_pptx(sl, colors, fw, fh, page)
            continue

        if s.get("render") == "html":  # divider(기타 html 클래스 포함)
            html_native += 1
            order = _section_order(all_slides)
            sec = slide_data.get("section")
            no = fields.get("no") or ((order.index(sec) + 1) if sec in order else "")
            no_txt = str(no).zfill(2) if str(no).isdigit() else str(no)
            if no_txt:
                prim.add_text(sl, px(W - 900), px(-60), px(880), px(760), no_txt,
                              size=pt(560), color="E9EAEB", family=family,
                              bold=True, align="right")
            label = fields.get("label") or (f"SECTION {no_txt}" if no_txt else (sec or ""))
            prim.add_text(sl, px(110), px(H - 390), px(W - 400), px(34), label,
                          size=pt(24), color=c("accent", "E8590C"), family=family, bold=True)
            prim.add_text(sl, px(110), px(H - 346), px(W - 400), px(80),
                          slide_data.get("title") or sec or "", size=pt(64),
                          color=c("ink_deep", "0F2438"), family=family, bold=True)
            sub = slide_data.get("message") or fields.get("sub") or ""
            if sub:
                prim.add_text(sl, px(110), px(H - 250), px(W - 400), px(32), sub,
                              size=pt(23), color=c("gray_text", "5C6470"), family=family)
            prim.add_rect(sl, px(110), px(H - 190), px(96), px(4), fill=c("ink", "16324C"))
            _chrome_footer_pptx(sl, colors, fw, fh or
                                int((skin.get("chrome") or {}).get("footer_h") or 0),
                                page, with_client=False)
            continue

        # 이미지 장 (content / full_image)
        img = slides_dir / s["out_name"]
        epx = s.get("expected_px") or {}
        if s.get("image_scope") == "full":
            if img.exists():
                used += 1
                sl.shapes.add_picture(str(img), Emu(int(px(fw) * 914400)),
                                      Emu(int(px(fw) * 914400)),
                                      width=Emu(int(px(W - 2 * fw) * 914400)),
                                      height=Emu(int(px(H - 2 * fw) * 914400)))
            else:
                missing.append(s["out_name"])
            band = prim.add_rect(sl, px(fw), px(H - fw - 64), px(W - 2 * fw), px(64),
                                 fill=c("ink_deep", "0F2438"))
            _fill_alpha(band, 0.55)
            prim.add_text(sl, px(fw + 64), px(H - fw - 64), px(W - 2 * fw - 400), px(64),
                          project, size=pt(18), color="FFFFFF", family=family,
                          anchor="middle")
            prim.add_text(sl, px(W - fw - 64 - 220), px(H - fw - 64), px(220), px(64),
                          page, size=pt(18), color="FFFFFF", family=family,
                          bold=True, align="right", anchor="middle")
            continue

        bm = s.get("body_margin") or {}
        bx, btop = int(bm.get("x") or 0), int(bm.get("top") or 0)
        _chrome_header_pptx(sl, colors, fw, hh, slide_data, s.get("flags") or [])
        if img.exists():
            used += 1
            sl.shapes.add_picture(str(img), Emu(int(px(fw + bx) * 914400)),
                                  Emu(int(px(fw + hh + btop) * 914400)),
                                  width=Emu(int(px(int(epx.get("w") or 0)) * 914400)),
                                  height=Emu(int(px(int(epx.get("h") or 0)) * 914400)))
        else:
            missing.append(s["out_name"])
            prim.add_rect(sl, px(fw + bx), px(fw + hh + btop),
                          px(int(epx.get("w") or 100)), px(int(epx.get("h") or 100)),
                          fill="EDEFF2")
            prim.add_text(sl, px(fw + bx), px(fw + hh + btop + 40),
                          px(int(epx.get("w") or 100)), px(40),
                          f"이미지 없음: {SLIDES_DIR}/{s['out_name']}", size=pt(20),
                          color="AA3333", family=family, align="center")
        _chrome_footer_pptx(sl, colors, fw, fh, page)

    out = run / COMPOSE_PPTX
    prs.save(str(out))
    warnings = list(dict.fromkeys(warnings))  # HTML compose와 동일 규칙(순서 보존 중복 제거)
    return {"out": str(out), "slides": total, "images_used": used, "missing": missing,
            "html_native": html_native, "editable_chrome": True, "warnings": warnings,
            # DF6: 빈 리스트도 항상 기록("0건"과 "미측정" 구분 - manifest와 같은 관례).
            "deck_overrides_applied": [d["n"] for d in deck_overrides],
            "deck_overrides": deck_overrides}


# ---------------------------------------------------------------------------
# review scaffold — Claude 검수 계약 (d, Q2)
# ---------------------------------------------------------------------------

REVIEW_MD = "imagedeck_review.md"


def _storyline_by_n(run: Path) -> dict:
    try:
        sl = _load_json(run / "storyline.json")
    except (OSError, json.JSONDecodeError):
        return {}
    return {str(s.get("n")): s for s in (sl.get("slides") or [])}


_KNOWLEDGE_RECORD_HEADER = "## 지식 사용 기록 (원장 — 다음 go가 이 섹션을 수거한다)"


def _knowledge_check_lines(run: Path) -> list[str]:
    """KC 패킷 ③(2026-07-24 확정) → ε패킷(2026-07-23)에서 config 표 소비로 일반화.

    관문 프로파일 express는 생략(빈 리스트), standard는 기본 대조, full은 장별 샘플링 심화
    안내가 하나 더 붙는다. vault를 파일시스템으로 읽지 않는다 — pull 요구 문구만 심는다
    (세션이 obsidian_search로 능동 조회, KC 패킷 제약 그대로 승계).

    ε패킷 이전에는 폴더 경로(`ref/디자인지식/테마/`·`ref/경험설계지식/`)가 이 함수에 직접
    하드코딩돼 있었다(vault 재편 전 경로 — 지금은 `ref/기획지식/경험설계/`로 이동). 지금은
    `pipeline.config.json`의 `knowledge_stages.imagedeck_review.pull` 표를 읽는다.
    """
    profile = gates.load_config(run)["profile"]
    if profile == "express":
        return []
    lines = [
        "## 지식 대조 (사람 전속 확인 — Claude 대리 금지, 검토_체크.md와 동일 등급)",
        "완성 렌더를 아래 지식 계열 기준으로 역검사하라(pull 프로토콜 — obsidian_search로 "
        "능동 조회, 자동 주입 아님):",
    ]
    folders = knowledge_ledger.pull_folders("imagedeck_review")
    if folders:
        for f in folders:
            lines.append(f"- [ ] `ref/{f}/` 기준 역검사(덱 전체 흐름·강조·위계가 카드와 정합하는지)")
    else:
        lines.append("- [ ] (config knowledge_stages.imagedeck_review에 pull 폴더가 설정되지 않음)")
    if profile == "full":
        lines.append(
            "- [ ] (full 프로파일) 장별 샘플링 심화 — 위 대조를 장마다 개별 수행하고, "
            "이상 있는 장은 해당 장 '지적' 칸에 근거 카드 슬러그를 남겨라(표본 생략 금지)."
        )
    lines.append("")
    lines.append(_KNOWLEDGE_RECORD_HEADER)
    lines.append("- cards: (사용한 카드 슬러그를 쉼표로 나열 — 없으면 '없음')")
    lines.append("- web:")
    lines.append("  - (형식: `URL — 용도 한 줄` 한 줄에 하나씩 — 없으면 '없음' 한 줄)")
    lines.append("")
    return lines


_CARDS_LINE_RE = re.compile(r"^\s*-\s*cards:\s*(.+)$", re.MULTILINE)
_WEB_ITEM_RE = re.compile(r"^\s*-\s*(https?://\S+)\s*(?:[-—]+\s*(.*))?$", re.MULTILINE)


def collect_review_knowledge(run: Path) -> dict:
    """ε패킷: imagedeck_review.md의 '지식 사용 기록' 섹션을 파싱해 원장에 기록한다.

    이 단계(imagedeck_review, KC③)는 검수 자체가 **선택**이다(MANUAL §9.9 — "바로 정독·채택"
    도 유효한 경로) — 그래서 다른 5단계와 달리 knowledge_used 누락을 차단하지 않는다(soft).
    파일이 없거나 섹션이 비어 있으면 조용히 {"found": False}를 반환한다.
    """
    run = Path(run)
    path = run / REVIEW_MD
    if not path.is_file():
        return {"found": False, "recorded": False}
    try:
        text = path.read_text(encoding="utf-8")
    except OSError:
        return {"found": False, "recorded": False}
    if _KNOWLEDGE_RECORD_HEADER not in text:
        return {"found": False, "recorded": False}
    block = text.split(_KNOWLEDGE_RECORD_HEADER, 1)[1]
    next_header = block.find("\n## ")
    if next_header != -1:
        block = block[:next_header]
    cards: list[str] = []
    m = _CARDS_LINE_RE.search(block)
    if m:
        raw = m.group(1).strip()
        if raw and raw not in ("없음", "(없음)"):
            cards = [c.strip() for c in raw.split(",") if c.strip()]
    web: list[dict] = []
    for wm in _WEB_ITEM_RE.finditer(block):
        web.append({"url": wm.group(1).strip(), "purpose": (wm.group(2) or "").strip()})
    if not cards and not web:
        return {"found": True, "recorded": False, "errors": []}
    ku = {"knowledge_used": {"cards": cards, "web": web}}
    errors, warnings = knowledge_ledger.validate_knowledge_used(ku, "imagedeck_review")
    if errors:  # 안전장치②의 web_search=false 위반 등 — 검수 자체는 선택이라도 이건 차단한다.
        return {"found": True, "recorded": False, "errors": errors, "warnings": warnings}
    knowledge_ledger.record(run, "imagedeck_review", {"cards": cards, "web": web}, source_file=str(path))
    return {"found": True, "recorded": True, "cards": cards, "web": web, "errors": [], "warnings": warnings}


def review_scaffold(run: Path) -> dict:
    """(d) Claude 검수 계약 scaffold — 세션이 이미지를 Read하고 채울 장별 대조표.

    deck_review.md의 자매 형식: 장별 [정본 문구·수치 대조 / 누락 binds / 금지 스타일 / verdict].
    자동 OCR을 하지 않는다(내용 안전) — 사람/세션의 정독을 돕는 체크리스트를 결정론 조립할 뿐이다.
    """
    run = Path(run)
    mp = run / MANIFEST_NAME
    if not mp.is_file():
        raise ImagedeckError(f"{MANIFEST_NAME} 없음 - 먼저 imagedeck --bundle")
    manifest = _load_json(mp)
    by_n = _storyline_by_n(run)
    # collect 결과가 있으면 px/커버리지 상태를 머리말에 반영.
    collect = None
    cp = run / COLLECT_NAME
    if cp.is_file():
        try:
            collect = _load_json(cp)
        except (OSError, json.JSONDecodeError):
            collect = None

    lines = ["# 이미지 장표 검수 (Claude 세션 대조 · 사람 채택)", "",
             "> 각 이미지를 Read로 열어 아래 '정본'과 대조하라. 자동 OCR 금지 - 눈으로 확인.",
             "> verdict: ok(정본 일치) / fix(오탈자·수치 오류·누락) / style(금지 스타일 위반).",
             f"> 금지 스타일: {', '.join(FORBIDDEN_STYLE)}", ""]
    collect_by_key: dict = {}
    if collect:
        cov = collect.get("coverage") or {}
        lines.append(f"- 수거 검증: {cov.get('ok')}/{cov.get('total')} "
                     f"· 판정 {'PASS' if collect.get('pass') else 'FAIL'}")
        if collect.get("pixel_heuristics_available"):
            lines.append(f"- 픽셀 휴리스틱 경고(warn 등급 - fail 아님): {collect.get('pixel_warnings', 0)}건")
        lines.append("")
        for r in collect.get("slides", []):
            collect_by_key[(str(r.get("n")), r.get("variant") or "")] = r
    lines.extend(_knowledge_check_lines(run))
    for s in manifest.get("slides", []):
        n = s.get("n")
        variant = s.get("variant") or ""
        slide = by_n.get(str(n)) or {}
        if s.get("render") == "html":
            lines.append(f"## 장 {n} - {_esc(slide.get('title', ''))} (HTML 전용 장 - {s.get('deck_class')})")
            lines.append("- 이미지 검수 대상 아님 - compose 산출(deck.images.html)에서 크롬·문구를 확인.")
            lines.append("")
            continue
        lines.append(f"## 장 {n}{variant} - {_esc(slide.get('title', ''))}")
        lines.append(f"- 이미지: `{SLIDES_DIR}/{s.get('out_name')}`")
        if s.get("content_split"):
            lines.append(f"- ⚠️ 사전 분할(마찰23): {s.get('split_reason') or ''} — 나머지 절반은 다른 variant 참고.")
        px_flags = ((collect_by_key.get((str(n), variant)) or {}).get("pixel_heuristics") or {}).get("flags") or []
        if px_flags:
            lines.append(f"- ⚠️ 픽셀 휴리스틱 경고: {', '.join(px_flags)}(오탐 여지 있음 - 직접 확인 요망).")
        ch = s.get("chrome") or {}
        if s.get("image_scope") == "body" and (int(ch.get("header_h") or 0) > 0
                                               or int(ch.get("footer_h") or 0) > 0):
            lines.append("- 크롬 분리(W29): 섹션배지·제목·부제·페이지·flag 배너가 이미지 안에 "
                         "**없어야** 한다(HTML 크롬이 그림 - 중복이면 fix). 본문 텍스트의 [예시] 접두는 보존.")
        title = slide.get("title")
        message = slide.get("message")
        if title:
            lines.append(f"- 정본 제목: {title}")
        if message:
            lines.append(f"- 정본 메시지: {message}")
        fields = slide.get("fields") or {}
        if fields:
            for k, v in fields.items():
                lines.append(f"- 정본 field `{k}`: {json.dumps(v, ensure_ascii=False)[:120]}")
        if s.get("flags"):
            lines.append(f"- 보존 필수 flag(이미지에 그대로 있어야 함): {', '.join(s['flags'])}")
        if s.get("missing_binds"):
            lines.append(f"- [검토요망] 누락 binds(빈칸 아닌 표시여야): {', '.join(s['missing_binds'])}")
        lines.append("- verdict: __ (ok/fix/style)")
        lines.append("- 지적(있으면 정확한 위치·올바른 문구): ")
        lines.append("")
    (run / REVIEW_MD).write_text("\n".join(lines) + "\n", encoding="utf-8")
    return {"out": str(run / REVIEW_MD), "slides": len(manifest.get("slides", []))}


def _needs_copy(src: Path, dst: Path) -> bool:
    """dst가 없거나, 있어도 크기·mtime이 src와 다르면(구본) 복사가 필요하다고 본다."""
    if not dst.exists():
        return True
    ss, ds = src.stat(), dst.stat()
    return ss.st_size != ds.st_size or int(ss.st_mtime) > int(ds.st_mtime)


def export_outputs(run: Path, dest: Path) -> dict:
    """(g) W28 마찰 L3(2026-07-24 강의 덱 run 실측): run/imagedeck 산출물을 워크스페이스 밖
    지정 폴더로 내보내는 공식 통로.

    엔진은 `workspace/runs/`로만 쓰는데, 덱을 이식하려고 사람이 로컬 사본을 따로 두면
    재생성분이 사본에 반영되지 않아 **옛 이미지가 조용히 남는** 사고가 난다(강의 덱 27번 장
    사례 — 사람 육안으로만 발견됐다). 이 함수가 그 "엔진 -> 사본" 동기화를 대신한다.

    - 대상: `imagedeck/slides/*.png`, `imagedeck_manifest.json`, `deck.images.html`(있으면),
      `deck.images.pptx`(있으면).
    - 신규/갱신분만 복사(크기+mtime 비교) - dest에 이미 최신본이 있으면 skip.
    - 단방향: run(원본)은 절대 쓰지 않는다. dest는 없으면 생성한다.
    """
    import shutil
    run = Path(run)
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)

    candidates: list[tuple[Path, Path]] = []
    slides_src = run / SLIDES_DIR
    if slides_src.is_dir():
        for p in sorted(slides_src.glob("*.png")):
            if is_rejected(p):
                continue   # 마찰34 반려 증거본은 run에만 남기고 밖으로 내보내지 않는다
            candidates.append((p, Path(SLIDES_DIR) / p.name))
    for name in (MANIFEST_NAME, COMPOSE_HTML, COMPOSE_PPTX):
        p = run / name
        if p.is_file():
            candidates.append((p, Path(name)))

    copied: list[str] = []
    skipped: list[str] = []
    for src, rel in candidates:
        dst = dest / rel
        if _needs_copy(src, dst):
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src, dst)
            copied.append(str(rel))
        else:
            skipped.append(str(rel))

    return {
        "dest": str(dest),
        "copied": copied,
        "skipped": skipped,
        "total": len(candidates),
    }


_COMPOSE_SHELL = """<!doctype html>
<html lang="ko"><head><meta charset="utf-8"><title>{project} - 이미지 덱</title>
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ background:#222; }}
  .slide {{ position:relative; margin:0 auto 12px; background:#fff; overflow:hidden; }}
  .chrome-h {{ display:flex; align-items:center; justify-content:space-between;
    padding:0 32px; background:#1F3864; color:#fff; font:600 22px/1 'Malgun Gothic',sans-serif; }}
  .chrome-h .page {{ opacity:.75; font-weight:400; font-size:18px; }}
  .body {{ display:flex; align-items:center; justify-content:center; }}
  .content {{ display:block; width:100%; height:100%; object-fit:contain; }}
  .content.missing {{ color:#a33; font:500 20px 'Malgun Gothic',sans-serif; text-align:center;
    display:flex; flex-direction:column; justify-content:center; }}
  .chrome-f {{ display:flex; align-items:center; justify-content:flex-end; gap:24px;
    padding:0 32px; background:#f4f4f4; }}
  .chrome-f .logo {{ max-height:70%; width:auto; }}
{stage_css}
</style></head><body>
{stage_open}
{slides}
{stage_close}
</body></html>
"""
