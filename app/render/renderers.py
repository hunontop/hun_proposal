# -*- coding: utf-8 -*-
"""PPTX generic 렌더러 (선택 백엔드) — 노하우 0, 팩 토큰 주입.

HTML(htmlgen.py)이 기본 타깃이고 PPTX는 '편집본 필요시' 백엔드.
per-template 정교화 대신 cover + generic content + fallback로 정본 deck를 안전히 .pptx화.
색은 팩 tokens에서 시맨틱 해석(primary/accent/ink/muted). 하드코딩 금지.
시그니처: fn(prs, tokens, slide, f) → 슬라이드 1장 추가.
"""
from __future__ import annotations

from typing import Any, Callable

from . import pptx_primitives as P

EMU_W_IN, EMU_H_IN = 13.333, 7.5


def _colors(tokens: dict) -> dict:
    return tokens.get("colors") or {}


def _pick(tokens: dict, *keys, default="404040") -> str:
    c = _colors(tokens)
    for k in keys:
        if k in c:
            return c[k].lstrip("#")
    return default


def _sem(tokens: dict) -> dict:
    return {
        "primary": _pick(tokens, "navy", "dark_navy", "deep_navy", default="1F3864"),
        "accent": _pick(tokens, "orange", "bright_blue", "orange_bright", default="2E5496"),
        "ink": _pick(tokens, "black", "text_dark", default="1A1A1A"),
        "muted": _pick(tokens, "gray_text", "footer_gray", default="808080"),
        "flag": _pick(tokens, "red", "section_concern", "status_red", default="C00000"),
        "cover_bg": _pick(tokens, "section_cover", "deep_navy", "navy", default="1F3864"),
    }


def _fsize(tokens: dict, key: str, default: int) -> int:
    return ((tokens.get("fonts") or {}).get("sizes") or {}).get(key, default)


def _family(tokens: dict) -> str:
    return (tokens.get("fonts") or {}).get("family", "Arial")


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, l, t, w, h, text, *, size, color, family, bold=False, align=None):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = family
    p.font.color.rgb = RGBColor.from_string(color)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    return tb


def _review(slide, tokens, sem, fam, slide_data):
    rn = (slide_data.get("review_needed") or []) + [f"(미결) {q}" for q in (slide_data.get("open_questions") or [])]
    if not rn:
        return
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    for i, r in enumerate(rn):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "🔴 " + str(r)
        p.font.size = Pt(_fsize(tokens, "small", 12)); p.font.name = fam
        p.font.color.rgb = RGBColor.from_string(sem["flag"])


def _present(value: Any) -> bool:
    return value is not None and value != "" and value != [] and value != {} and value != ()


def _sequence(value: Any) -> list:
    if isinstance(value, (list, tuple)):
        return list(value)
    return []


def _as_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, dict):
        for key in ("label", "name", "title", "text", "value", "description", "detail"):
            if _present(value.get(key)):
                return _as_text(value.get(key))
        return " / ".join(f"{k}: {_as_text(v)}" for k, v in value.items() if _present(v))
    if isinstance(value, (list, tuple)):
        return " / ".join(_as_text(v) for v in value if _present(v))
    return str(value)


def _first(mapping: Any, *keys: str) -> Any:
    if not isinstance(mapping, dict):
        return None
    for key in keys:
        value = mapping.get(key)
        if _present(value):
            return value
    return None


def _number(value: Any, default: float = 0.0) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        cleaned = value.replace(",", "").replace("%", "").strip()
        try:
            return float(cleaned)
        except ValueError:
            return default
    return default


def _body_items(slide: dict) -> list:
    body = slide.get("body")
    if isinstance(body, (list, tuple)):
        return list(body)
    return [body] if _present(body) else []


def _title(s, slide: dict, skin: P.Skin):
    P.add_title(s, slide.get("title", ""), skin)
    if slide.get("key_message"):
        P.add_text(
            s, 0.6, 1.22, 12.0, 0.45, slide.get("key_message", ""),
            size=skin.fsize("body", 15), color=skin.accent, family=skin.family(),
            bold=True,
        )


def _bullet_list(s, skin: P.Skin, left, top, width, height, items, *, size=None, color=None):
    clean = [_as_text(item) for item in items if _present(item)]
    if not clean:
        return None
    tb = P.add_textbox(s, left, top, width, height)
    tf = tb.text_frame
    for i, item in enumerate(clean[:8]):
        P.write_para(
            tf, item, size=size or skin.fsize("body", 13), color=color or skin.ink,
            family=skin.family(), first=(i == 0), bullet=True,
        )
    return tb


def _chart_payload(f: dict) -> tuple[str, list[str], list[tuple[str, list[float]]]]:
    chart_type = str(f.get("chart_type") or f.get("type") or "column").lower()
    if chart_type not in {"bar", "column", "line", "pie"}:
        chart_type = "column"

    categories = [_as_text(c) for c in _sequence(f.get("categories")) if _present(c)]
    series_data = f.get("series")
    series: list[tuple[str, list[float]]] = []
    if isinstance(series_data, list):
        for idx, item in enumerate(series_data):
            if isinstance(item, dict):
                vals = [_number(v) for v in _sequence(item.get("values"))]
                name = _as_text(_first(item, "name", "label", "title")) or f"Series {idx + 1}"
            else:
                vals = [_number(item)]
                name = f"Series {idx + 1}"
            if vals:
                series.append((name, vals[:len(categories)] if categories else vals))

    comparison = _sequence(f.get("comparison"))
    if comparison and not series:
        categories = []
        values = []
        for idx, item in enumerate(comparison):
            if isinstance(item, dict):
                categories.append(_as_text(_first(item, "label", "name", "title")) or str(idx + 1))
                values.append(_number(_first(item, "value", "score", "amount")))
            else:
                categories.append(str(idx + 1))
                values.append(_number(item))
        series.append((_as_text(f.get("metric")) or "Value", values))

    values = _sequence(f.get("values"))
    if categories and values and not series:
        series.append((_as_text(f.get("metric")) or "Value", [_number(v) for v in values]))

    if series and not categories:
        categories = [str(i + 1) for i in range(max(len(vals) for _, vals in series))]

    if categories and series:
        width = len(categories)
        series = [(name, (vals + [0.0] * width)[:width]) for name, vals in series]

    return chart_type, categories, series


def _score_text(value: Any) -> str:
    if isinstance(value, str):
        return value
    try:
        score = int(value)
    except Exception:
        return _as_text(value)
    return "●" * max(0, min(4, score)) + "○" * max(0, 4 - max(0, min(4, score)))


def _table_rows_from_comparison(f: dict) -> tuple[list[str], list[list[str]], int | None, str]:
    raw_options = f.get("options")
    raw_criteria = f.get("criteria")
    options = (
        [{"name": name, "values": values} for name, values in raw_options.items()]
        if isinstance(raw_options, dict) else _sequence(raw_options)
    )
    criteria = (
        [{"name": name, "values": values} for name, values in raw_criteria.items()]
        if isinstance(raw_criteria, dict) else _sequence(raw_criteria)
    )
    recommendation = _as_text(f.get("recommendation"))
    try:
        recommended_index = int(f.get("recommended_index")) if f.get("recommended_index") is not None else None
    except (TypeError, ValueError):
        recommended_index = None

    option_labels = [_as_text(_first(o, "name", "label", "title", "option") if isinstance(o, dict) else o) for o in options]
    option_labels = [label for label in option_labels if label]
    rows: list[list[str]] = []

    for criterion in criteria:
        if isinstance(criterion, dict):
            label = _as_text(_first(criterion, "name", "label", "title", "criterion"))
            scores = _sequence(_first(criterion, "scores", "values", "options"))
            notes = _sequence(criterion.get("notes"))
            values = []
            for idx, option in enumerate(options):
                val = scores[idx] if idx < len(scores) else None
                if not _present(val) and isinstance(option, dict):
                    opt_values = _first(option, "criteria", "values", "scores", "details")
                    if isinstance(opt_values, dict):
                        val = opt_values.get(label)
                    elif isinstance(opt_values, (list, tuple)) and len(rows) < len(opt_values):
                        val = opt_values[len(rows)]
                note = notes[idx] if idx < len(notes) else None
                values.append((_score_text(val) + (f" - {_as_text(note)}" if _present(note) else "")).strip())
        else:
            label = _as_text(criterion)
            values = ["" for _ in option_labels]
        if label and any(_present(v) for v in values):
            rows.append([label] + values[:len(option_labels)])

    if not rows and option_labels and recommendation:
        rows.append(["Recommendation"] + [recommendation if i == recommended_index or recommended_index is None else "" for i in range(len(option_labels))])

    return option_labels, rows, recommended_index, recommendation


def _band(value: Any, *, reverse_numeric: bool = False) -> str:
    if isinstance(value, (int, float)):
        idx = int(value)
        if reverse_numeric:
            return ["high", "mid", "low"][max(0, min(2, idx))]
        return ["low", "mid", "high"][max(0, min(2, idx))]
    text = str(value or "").lower()
    aliases = {
        "0": "low", "1": "mid", "2": "high",
        "low": "low", "l": "low",
        "medium": "mid", "med": "mid", "mid": "mid", "m": "mid",
        "high": "high", "h": "high",
    }
    return aliases.get(text, text)


def _aligned(source: Any, index: int, key: Any, count: int) -> Any:
    if isinstance(source, dict):
        candidates = (key, str(key), index, str(index), index + 1, str(index + 1))
        for candidate in candidates:
            if candidate in source and _present(source[candidate]):
                return source[candidate]
        return None
    values = _sequence(source)
    if index < len(values) and _present(values[index]):
        return values[index]
    if count == 1 and _present(source) and not values:
        return source
    return None


# --- 렌더러 ----------------------------------------------------------------

def cover(prs, tokens, slide, f):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    s = _blank(prs); sem = _sem(tokens); fam = _family(tokens)
    bg = s.shapes.add_shape(1, 0, 0, Inches(EMU_W_IN), Inches(EMU_H_IN))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor.from_string(sem["cover_bg"]); bg.line.fill.background()
    title = slide.get("title") or f.get("project_title", "")
    msg = slide.get("key_message") or f.get("concept_message", "")
    _box(s, 1, 2.6, 11.3, 2, title, size=_fsize(tokens, "title", 40), color="FFFFFF", family=fam, bold=True, align="center")
    if msg:
        _box(s, 1, 4.7, 11.3, 1, msg, size=_fsize(tokens, "body", 20), color="D6DCE5", family=fam, align="center")
    return s


def content(prs, tokens, slide, f):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    s = _blank(prs); sem = _sem(tokens); fam = _family(tokens)
    # 좌측 액센트 바
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.22), Inches(EMU_H_IN))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(sem["accent"]); bar.line.fill.background()
    _box(s, 0.55, 0.4, 12.2, 1.0, slide.get("title", ""), size=_fsize(tokens, "section", 28), color=sem["ink"], family=fam, bold=True)
    if slide.get("key_message"):
        _box(s, 0.55, 1.5, 12.2, 0.7, "▶ " + slide["key_message"], size=_fsize(tokens, "body", 16), color=sem["accent"], family=fam, bold=True)
    # 본문: fields 안의 리스트형들 + slide.body
    items: list[str] = list(slide.get("body") or [])
    for v in f.values():
        if isinstance(v, list):
            items += [str(x.get("name") if isinstance(x, dict) else x) for x in v]
        elif isinstance(v, str) and v not in (slide.get("title"), slide.get("key_message")):
            items.append(v)
    if items:
        tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12.0), Inches(3.7))
        tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(items[:10]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + str(b)
            p.font.size = Pt(_fsize(tokens, "body", 16)); p.font.name = fam
            p.font.color.rgb = RGBColor.from_string(sem["ink"])
    _review(s, tokens, sem, fam, slide)
    return s


def _use_content(prs, tokens, slide, f, reason="필드 shape 불일치"):
    """네이티브 렌더러가 generic content로 폴백할 때 사유를 slide에 표시한 뒤 렌더.

    dispatch가 렌더 후 `slide['_fallback_reason']`을 읽어 사용자에게 경고로 알린다
    (조용한 폴백 금지 — 어떤 슬라이드가 적합 템플릿을 못 써서 밋밋해졌는지 보고).
    """
    slide["_fallback_reason"] = reason
    return content(prs, tokens, slide, f)


def fallback(prs, tokens, slide, f):
    return _use_content(prs, tokens, slide, f, "적합 렌더러 없음(명시적 fallback)")


def data_chart(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    chart_type, categories, series = _chart_payload(f)
    interpretations = _sequence(f.get("interpretation")) or _body_items(slide)
    metric = _as_text(f.get("metric"))

    if categories and series:
        if metric:
            P.add_text(s, 0.75, 1.9, 7.6, 0.35, metric, size=skin.fsize("small", 12),
                       color=skin.muted, family=skin.family(), bold=True)
        P.add_chart(s, chart_type, 0.75, 2.25, 7.7, 3.55,
                    categories=categories, series=series, skin=skin)
        _bullet_list(s, skin, 8.8, 2.05, 3.8, 3.9, interpretations, size=skin.fsize("body", 13))
    else:
        P.add_rounded(s, 0.8, 2.0, 4.3, 2.0, fill=skin.panel, line=skin.rule)
        P.add_text(s, 1.05, 2.28, 3.8, 0.45, "Metric", size=skin.fsize("small", 11),
                   color=skin.muted, family=skin.family(), bold=True)
        P.add_text(s, 1.05, 2.78, 3.8, 0.9, metric or slide.get("key_message", ""),
                   size=skin.fsize("section", 24), color=skin.primary, family=skin.family(),
                   bold=True)
        _bullet_list(s, skin, 5.6, 2.0, 6.8, 3.2, interpretations, size=skin.fsize("body", 14))
    P.add_review_note(s, skin, slide)
    return s


def kpi_dashboard(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    kpis = [k for k in _sequence(f.get("kpis")) if isinstance(k, dict)]
    if not kpis:
        return _use_content(prs, tokens, slide, f)

    cols = 4 if len(kpis) > 4 else max(1, len(kpis))
    rows = 2 if len(kpis) > 4 else 1
    gap = 0.28
    left, top = 0.75, 2.05
    width = (11.9 - gap * (cols - 1)) / cols
    height = (3.85 - gap * (rows - 1)) / rows
    for idx, item in enumerate(kpis[:8]):
        col = idx % cols
        row = idx // cols
        x = left + col * (width + gap)
        y = top + row * (height + gap)
        P.add_rounded(s, x, y, width, height, fill=skin.panel, line=skin.rule)
        P.add_text(s, x + 0.2, y + 0.18, width - 0.4, 0.28, _as_text(item.get("label")),
                   size=skin.fsize("small", 10), color=skin.muted, family=skin.family(), bold=True)
        P.add_text(s, x + 0.2, y + 0.56, width - 0.4, 0.52, _as_text(item.get("value")),
                   size=skin.fsize("section", 24), color=skin.primary, family=skin.family(), bold=True)
        delta = _as_text(item.get("delta"))
        if delta:
            direction = str(item.get("delta_dir") or "flat").lower()
            delta_color = skin.accent if direction == "up" else skin.flag if direction == "down" else skin.muted
            marker = "UP" if direction == "up" else "DOWN" if direction == "down" else "FLAT"
            P.add_text(s, x + 0.2, y + height - 0.52, width - 0.4, 0.24, f"{marker} {delta}",
                       size=skin.fsize("small", 10), color=delta_color, family=skin.family(), bold=True)
        context = _as_text(item.get("context"))
        if context:
            P.add_text(s, x + 0.2, y + height - 0.28, width - 0.4, 0.22, context,
                       size=skin.fsize("footer", 9), color=skin.muted, family=skin.family())
    P.add_review_note(s, skin, slide)
    return s


def comparison_table(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    option_labels, rows, recommended_index, recommendation = _table_rows_from_comparison(f)
    if not option_labels or not rows:
        return _use_content(prs, tokens, slide, f)

    total_rows = min(len(rows), 7) + 1
    total_cols = min(len(option_labels), 4) + 1
    table = P.add_table(s, 0.65, 1.95, 12.0, min(3.9, 0.46 * total_rows + 0.4),
                        total_rows, total_cols, skin=skin, header=True)
    P.fill_cell(table.cell(0, 0), "Criteria", skin=skin, size=10, color=skin.white,
                bold=True, fill=skin.primary, align="center")
    for col, label in enumerate(option_labels[:total_cols - 1], 1):
        fill = skin.accent if recommended_index == col - 1 else skin.primary
        P.fill_cell(table.cell(0, col), label, skin=skin, size=10, color=skin.white,
                    bold=True, fill=fill, align="center")
    for row_idx, row in enumerate(rows[:total_rows - 1], 1):
        P.fill_cell(table.cell(row_idx, 0), row[0], skin=skin, size=10, color=skin.primary,
                    bold=True, fill=skin.panel)
        for col in range(1, total_cols):
            value = row[col] if col < len(row) else ""
            fill = skin.panel if recommended_index == col - 1 else None
            P.fill_cell(table.cell(row_idx, col), value, skin=skin, size=9,
                        color=skin.ink, fill=fill, align="center")
    if recommendation:
        P.add_text(s, 0.85, 6.05, 11.5, 0.45, f"Recommendation: {recommendation}",
                   size=skin.fsize("body", 13), color=skin.primary, family=skin.family(), bold=True)
    P.add_review_note(s, skin, slide)
    return s


def strategy_pillars(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    raw = _sequence(f.get("pillars")) or _body_items(slide)
    lines = _sequence(f.get("one_line_per_pillar"))
    pillars = []
    for idx, item in enumerate(raw[:5]):
        if isinstance(item, dict):
            title = _as_text(_first(item, "title", "name", "label", "pillar"))
            detail = _as_text(_first(item, "description", "detail", "body", "text"))
        else:
            title = _as_text(item)
            detail = _as_text(lines[idx]) if idx < len(lines) else ""
        if title:
            pillars.append((title, detail))
    if not pillars:
        return _use_content(prs, tokens, slide, f)

    gap = 0.25
    card_w = (11.9 - gap * (len(pillars) - 1)) / len(pillars)
    for idx, (title, detail) in enumerate(pillars):
        x = 0.75 + idx * (card_w + gap)
        P.add_rounded(s, x, 2.05, card_w, 3.6, fill=skin.panel, line=skin.rule)
        P.add_rect(s, x, 2.05, card_w, 0.12, fill=skin.accent, line=None)
        P.add_text(s, x + 0.22, 2.38, card_w - 0.44, 0.75, title,
                   size=skin.fsize("body", 15), color=skin.primary, family=skin.family(), bold=True)
        if detail:
            P.add_text(s, x + 0.22, 3.28, card_w - 0.44, 1.6, detail,
                       size=skin.fsize("small", 11), color=skin.ink, family=skin.family())
        P.add_text(s, x + 0.22, 5.12, card_w - 0.44, 0.26, f"{idx + 1:02d}",
                   size=skin.fsize("small", 10), color=skin.accent, family=skin.family(), bold=True)
    P.add_review_note(s, skin, slide)
    return s


def matrix_priority(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    items = [item for item in _sequence(f.get("items")) if isinstance(item, dict)]
    if not items:
        return _use_content(prs, tokens, slide, f)

    left, top, width, height = 1.25, 1.85, 10.6, 4.25
    cell_w, cell_h = width / 3.0, height / 3.0
    bands = ["low", "mid", "high"]
    for row, y_band in enumerate(["high", "mid", "low"]):
        for col, x_band in enumerate(bands):
            x = left + col * cell_w
            y = top + row * cell_h
            hot = row == 0 and col == 0
            P.add_rect(s, x, y, cell_w, cell_h, fill=skin.panel if hot else None,
                       line=skin.accent if hot else skin.rule, line_width=1.2 if hot else 0.5)
            cell_items = [
                _as_text(_first(item, "name", "label", "title"))
                for item in items
                if _band(item.get("x_band")) == x_band and _band(item.get("y_band"), reverse_numeric=True) == y_band
            ]
            for i, label in enumerate([v for v in cell_items if v][:4]):
                P.add_rounded(s, x + 0.12, y + 0.14 + i * 0.3, min(cell_w - 0.24, 2.9), 0.23,
                              fill=skin.primary, line=None)
                P.add_text(s, x + 0.22, y + 0.17 + i * 0.3, min(cell_w - 0.44, 2.7), 0.16,
                           label, size=skin.fsize("footer", 8), color=skin.white,
                           family=skin.family(), bold=True)
    P.add_text(s, left, top + height + 0.18, width, 0.35,
               f"{_as_text(f.get('x_axis')) or 'X axis'}: low -> high",
               size=skin.fsize("small", 10), color=skin.muted, family=skin.family(), align="center")
    P.add_text(s, 0.35, top + 1.55, 0.75, 0.6,
               f"{_as_text(f.get('y_axis')) or 'Y'} ↑",
               size=skin.fsize("small", 10), color=skin.accent, family=skin.family(), bold=True)
    P.add_review_note(s, skin, slide)
    return s


def roadmap_gantt(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    units = _sequence(f.get("time_units")) or _sequence(f.get("weeks"))
    streams = _sequence(f.get("workstreams"))
    if not units or not streams:
        return _use_content(prs, tokens, slide, f)

    units = units[:8]
    left, top, width, row_h = 0.75, 1.9, 11.9, 0.58
    name_w = 2.2
    col_w = (width - name_w) / len(units)
    P.add_rect(s, left, top, width, row_h, fill=skin.primary, line=skin.primary)
    P.add_text(s, left + 0.12, top + 0.18, name_w - 0.24, 0.2, "Workstream",
               size=skin.fsize("small", 9), color=skin.white, family=skin.family(), bold=True)
    for idx, unit in enumerate(units):
        P.add_text(s, left + name_w + idx * col_w, top + 0.18, col_w, 0.2, _as_text(unit),
                   size=skin.fsize("small", 9), color=skin.white, family=skin.family(),
                   bold=True, align="center")

    unit_keys = [str(u) for u in units]
    for row_idx, stream in enumerate(streams[:6], 1):
        y = top + row_idx * row_h
        P.add_rect(s, left, y, width, row_h, fill=None, line=skin.rule)
        name = _as_text(_first(stream, "name", "label", "title", "workstream") if isinstance(stream, dict) else stream)
        P.add_text(s, left + 0.12, y + 0.18, name_w - 0.24, 0.2, name,
                   size=skin.fsize("small", 9), color=skin.primary, family=skin.family(), bold=True)
        cells = _first(stream, "cells", "schedule", "periods", "time_units") if isinstance(stream, dict) else None
        start = _first(stream, "start_week", "start", "start_unit", "from") if isinstance(stream, dict) else None
        end = _first(stream, "end_week", "end", "end_unit", "to") if isinstance(stream, dict) else None
        active_cols = []
        if _present(start) and _present(end):
            start_idx = unit_keys.index(str(start)) if str(start) in unit_keys else 0
            end_idx = unit_keys.index(str(end)) if str(end) in unit_keys else start_idx
            active_cols = list(range(start_idx, end_idx + 1))
        else:
            for idx, unit in enumerate(units):
                if _present(_aligned(cells, idx, unit, len(units))):
                    active_cols.append(idx)
        for idx in range(len(units)):
            x = left + name_w + idx * col_w
            P.add_line(s, x, y, x, y + row_h, color=skin.rule, width_pt=0.4)
        if active_cols:
            first, last = min(active_cols), max(active_cols)
            P.add_rounded(s, left + name_w + first * col_w + 0.08, y + 0.17,
                          (last - first + 1) * col_w - 0.16, 0.22,
                          fill=skin.accent if row_idx % 2 else skin.primary, line=None)

    milestones = _sequence(f.get("milestones"))
    if milestones:
        base_y = top + (min(len(streams), 6) + 1) * row_h + 0.2
        P.add_text(s, left, base_y, name_w - 0.2, 0.25, "Milestones",
                   size=skin.fsize("small", 10), color=skin.primary, family=skin.family(), bold=True)
        for ms in milestones[:len(units)]:
            period = _as_text(_first(ms, "time_unit", "period", "unit", "week", "phase") if isinstance(ms, dict) else None)
            label = _as_text(_first(ms, "label", "name", "title", "milestone") if isinstance(ms, dict) else ms)
            idx = unit_keys.index(period) if period in unit_keys else milestones.index(ms) % len(units)
            x = left + name_w + idx * col_w + col_w / 2 - 0.1
            P.add_oval(s, x, base_y + 0.05, 0.2, 0.2, fill=skin.accent, line=None)
            P.add_text(s, x - col_w / 2 + 0.05, base_y + 0.32, col_w - 0.1, 0.32, label,
                       size=skin.fsize("footer", 8), color=skin.ink, family=skin.family(), align="center")
    P.add_review_note(s, skin, slide)
    return s


# 렌더러 이름/역할 → 함수. cover류만 cover, 나머지는 generic content.
def problem_intro(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    core = _as_text(_first(f, "core_question", "question", "main_question")) or slide.get("key_message", "")
    raw_questions = _first(f, "sub_questions", "questions", "question_list")
    questions = _sequence(raw_questions) or _body_items(slide)

    if not core and not questions:
        return _use_content(prs, tokens, slide, f)

    P.add_rounded(s, 0.8, 1.85, 11.75, 1.35, fill=skin.primary, line=None)
    P.add_text(s, 1.1, 2.08, 0.6, 0.55, "Q", size=skin.fsize("section", 24),
               color=skin.white, family=skin.family(), bold=True, align="center")
    P.add_text(s, 1.85, 2.02, 10.25, 0.72, core,
               size=skin.fsize("body", 17), color=skin.white,
               family=skin.family(), bold=True)
    if slide.get("key_message") and slide.get("key_message") != core:
        P.add_text(s, 1.85, 2.7, 10.25, 0.28, slide.get("key_message", ""),
                   size=skin.fsize("small", 10), color=skin.white, family=skin.family())

    cards = []
    for item in questions[:4]:
        if isinstance(item, dict):
            label = _as_text(_first(item, "label", "name", "title", "question"))
            detail = _as_text(_first(item, "description", "detail", "body", "answer"))
        else:
            label = _as_text(item)
            detail = ""
        if label:
            cards.append((label, detail))

    if cards:
        gap = 0.22
        card_w = (11.75 - gap * (len(cards) - 1)) / len(cards)
        for idx, (label, detail) in enumerate(cards):
            x = 0.8 + idx * (card_w + gap)
            P.add_rounded(s, x, 3.65, card_w, 1.85, fill=skin.panel, line=skin.rule)
            P.add_oval(s, x + 0.22, 3.92, 0.38, 0.38, fill=skin.accent, line=None)
            P.add_text(s, x + 0.22, 3.99, 0.38, 0.14, str(idx + 1),
                       size=skin.fsize("footer", 8), color=skin.white,
                       family=skin.family(), bold=True, align="center")
            P.add_text(s, x + 0.72, 3.88, card_w - 0.96, 0.58, label,
                       size=skin.fsize("body", 13), color=skin.primary,
                       family=skin.family(), bold=True)
            if detail:
                P.add_text(s, x + 0.25, 4.66, card_w - 0.5, 0.52, detail,
                           size=skin.fsize("small", 10), color=skin.ink,
                           family=skin.family())

    P.add_review_note(s, skin, slide)
    return s


def process_steps(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    raw_steps = f.get("steps")
    outputs = f.get("outputs")
    if isinstance(raw_steps, dict):
        steps = [{"label": label, "description": description} for label, description in raw_steps.items()]
    else:
        steps = _sequence(raw_steps)
    if not steps:
        return _use_content(prs, tokens, slide, f)

    count = min(len(steps), 5)
    gap = 0.18
    left, top, width, height = 0.75, 2.0, 11.85, 3.45
    card_w = (width - gap * (count - 1)) / count
    for idx, step in enumerate(steps[:count]):
        if isinstance(step, dict):
            label = _as_text(_first(step, "label", "name", "title", "step"))
            detail = _as_text(_first(step, "description", "body", "detail"))
            embedded_output = _first(step, "output", "deliverable")
        else:
            label, detail, embedded_output = _as_text(step), "", None
        output = _aligned(outputs, idx, label, len(steps))
        if not _present(output):
            output = embedded_output
        x = left + idx * (card_w + gap)
        P.add_rounded(s, x, top, card_w, height, fill=skin.panel, line=skin.rule)
        P.add_rect(s, x, top, card_w, 0.16, fill=skin.accent if idx == 0 else skin.primary, line=None)
        P.add_text(s, x + 0.2, top + 0.38, card_w - 0.4, 0.26, f"STEP {idx + 1}",
                   size=skin.fsize("small", 9), color=skin.accent,
                   family=skin.family(), bold=True)
        P.add_text(s, x + 0.2, top + 0.82, card_w - 0.4, 0.68, label,
                   size=skin.fsize("body", 13), color=skin.primary,
                   family=skin.family(), bold=True)
        if detail:
            P.add_text(s, x + 0.2, top + 1.64, card_w - 0.4, 0.78, detail,
                       size=skin.fsize("small", 10), color=skin.ink, family=skin.family())
        if _present(output):
            P.add_line(s, x + 0.2, top + 2.65, x + card_w - 0.2, top + 2.65,
                       color=skin.rule, width_pt=0.6)
            P.add_text(s, x + 0.2, top + 2.82, card_w - 0.4, 0.38, _as_text(output),
                       size=skin.fsize("footer", 9), color=skin.muted, family=skin.family())
        if idx < count - 1:
            P.add_line(s, x + card_w + 0.03, top + height / 2,
                       x + card_w + gap - 0.03, top + height / 2,
                       color=skin.accent, width_pt=1.2)

    P.add_review_note(s, skin, slide)
    return s


def risk_dashboard(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    raw_risks = f.get("risks")
    severities = f.get("severity")
    mitigations = f.get("mitigations")
    if isinstance(raw_risks, dict):
        risks = [{"name": name, "detail": detail} for name, detail in raw_risks.items()]
    else:
        risks = _sequence(raw_risks)
    if not risks:
        return _use_content(prs, tokens, slide, f)

    rows = []
    for idx, risk in enumerate(risks[:6]):
        if isinstance(risk, dict):
            label = _as_text(_first(risk, "name", "label", "title", "risk"))
            detail = _as_text(_first(risk, "description", "detail", "body"))
            embedded_severity = _first(risk, "severity")
            embedded_mitigation = _first(risk, "mitigation", "response")
        else:
            label, detail, embedded_severity, embedded_mitigation = _as_text(risk), "", None, None
        severity = _aligned(severities, idx, label, len(risks))
        mitigation = _aligned(mitigations, idx, label, len(risks))
        if not _present(severity):
            severity = embedded_severity
        if not _present(mitigation):
            mitigation = embedded_mitigation
        if label:
            rows.append((label, detail, _as_text(severity), _as_text(mitigation)))
    if not rows:
        return _use_content(prs, tokens, slide, f)

    cols = 3 if len(rows) > 2 else max(1, len(rows))
    gap = 0.26
    left, top = 0.75, 2.0
    card_w = (11.85 - gap * (cols - 1)) / cols
    card_h = 1.55
    for idx, (label, detail, severity, mitigation) in enumerate(rows):
        col = idx % cols
        row = idx // cols
        x = left + col * (card_w + gap)
        y = top + row * (card_h + 0.28)
        sev_lower = severity.lower()
        sev_fill = skin.flag if "high" in sev_lower or "critical" in sev_lower else skin.accent if "mid" in sev_lower or "medium" in sev_lower else skin.panel
        sev_color = skin.white if sev_fill in {skin.flag, skin.accent} else skin.primary
        P.add_rounded(s, x, y, card_w, card_h, fill=skin.panel, line=skin.rule)
        P.add_rect(s, x, y, card_w, 0.12, fill=sev_fill, line=None)
        if severity:
            P.add_rounded(s, x + card_w - 1.05, y + 0.25, 0.82, 0.26,
                          fill=sev_fill, line=None)
            P.add_text(s, x + card_w - 1.0, y + 0.32, 0.72, 0.1, severity,
                       size=skin.fsize("footer", 7), color=sev_color,
                       family=skin.family(), bold=True, align="center")
        P.add_text(s, x + 0.2, y + 0.28, card_w - 1.35, 0.44, label,
                   size=skin.fsize("body", 12), color=skin.primary,
                   family=skin.family(), bold=True)
        if detail:
            P.add_text(s, x + 0.2, y + 0.78, card_w - 0.4, 0.24, detail,
                       size=skin.fsize("footer", 8), color=skin.muted,
                       family=skin.family())
        if mitigation:
            P.add_line(s, x + 0.2, y + card_h - 0.44, x + card_w - 0.2, y + card_h - 0.44,
                       color=skin.rule, width_pt=0.5)
            P.add_text(s, x + 0.2, y + card_h - 0.28, card_w - 0.4, 0.18,
                       mitigation, size=skin.fsize("footer", 8), color=skin.ink,
                       family=skin.family())

    P.add_review_note(s, skin, slide)
    return s


def table_block(prs, tokens, slide, f):
    if slide.get("role") == "risk" or all(_present(f.get(key)) for key in ("risks", "severity", "mitigations")):
        return risk_dashboard(prs, tokens, slide, f)
    if all(_present(f.get(key)) for key in ("options", "criteria")):
        return comparison_table(prs, tokens, slide, f)
    return _use_content(prs, tokens, slide, f)


def closing_matrix(prs, tokens, slide, f):
    s = P.blank_slide(prs)
    skin = P.Skin(tokens)
    _title(s, slide, skin)
    commitments = f.get("commitments")
    proof_points = f.get("proof_points")
    if isinstance(commitments, dict):
        items = [{"commitment": label, "proof": proof} for label, proof in commitments.items()]
    else:
        items = _sequence(commitments)
    if not items and _body_items(slide):
        items = _body_items(slide)
    if not items:
        return _use_content(prs, tokens, slide, f)

    rows = []
    for idx, item in enumerate(items[:5]):
        if isinstance(item, dict):
            commitment = _as_text(_first(item, "commitment", "label", "name", "title", "text"))
            embedded_proof = _first(item, "proof", "proof_point", "evidence", "detail")
        else:
            commitment, embedded_proof = _as_text(item), None
        proof = _aligned(proof_points, idx, commitment, len(items))
        if not _present(proof):
            proof = embedded_proof
        if commitment:
            rows.append((commitment, _as_text(proof)))
    if not rows:
        return _use_content(prs, tokens, slide, f)

    P.add_rect(s, 0.78, 1.86, 11.78, 0.34, fill=skin.primary, line=None)
    P.add_text(s, 1.0, 1.94, 5.1, 0.16, "Commitment",
               size=skin.fsize("small", 9), color=skin.white, family=skin.family(), bold=True)
    P.add_text(s, 6.7, 1.94, 5.1, 0.16, "Proof point",
               size=skin.fsize("small", 9), color=skin.white, family=skin.family(), bold=True)

    row_h = min(0.8, 3.9 / len(rows))
    for idx, (commitment, proof) in enumerate(rows):
        y = 2.35 + idx * (row_h + 0.12)
        P.add_rounded(s, 0.8, y, 5.55, row_h, fill=skin.panel, line=skin.rule)
        P.add_rounded(s, 6.55, y, 6.0, row_h, fill=None, line=skin.rule)
        P.add_oval(s, 1.02, y + 0.18, 0.28, 0.28, fill=skin.accent, line=None)
        P.add_text(s, 1.02, y + 0.23, 0.28, 0.1, str(idx + 1),
                   size=skin.fsize("footer", 7), color=skin.white, family=skin.family(),
                   bold=True, align="center")
        P.add_text(s, 1.42, y + 0.18, 4.65, row_h - 0.24, commitment,
                   size=skin.fsize("body", 12), color=skin.primary,
                   family=skin.family(), bold=True)
        if proof:
            P.add_text(s, 6.82, y + 0.18, 5.36, row_h - 0.24, proof,
                       size=skin.fsize("small", 10), color=skin.ink, family=skin.family())

    P.add_review_note(s, skin, slide)
    return s


_COVER = {"cover", "cover_cinematic", "cover_slide"}
REGISTRY: dict[str, Callable] = {name: cover for name in _COVER}
REGISTRY["content"] = content
REGISTRY["fallback"] = fallback
REGISTRY["data_chart"] = data_chart
REGISTRY["data_interpretation"] = data_chart
REGISTRY["data"] = data_chart
REGISTRY["kpi_dashboard"] = kpi_dashboard
REGISTRY["comparison_table"] = comparison_table
REGISTRY["table_block"] = table_block
REGISTRY["comparison"] = comparison_table
REGISTRY["strategy_pillars"] = strategy_pillars
REGISTRY["card_grid"] = strategy_pillars
REGISTRY["strategy"] = strategy_pillars
REGISTRY["matrix_priority"] = matrix_priority
REGISTRY["decision"] = matrix_priority
REGISTRY["roadmap_gantt"] = roadmap_gantt
REGISTRY["timeline_matrix"] = roadmap_gantt
REGISTRY["roadmap"] = roadmap_gantt
REGISTRY["problem_intro"] = problem_intro
REGISTRY["problem_questions"] = problem_intro
REGISTRY["process_steps"] = process_steps
REGISTRY["process"] = process_steps
REGISTRY["risk_dashboard"] = risk_dashboard
REGISTRY["closing_matrix"] = closing_matrix
REGISTRY["closing"] = closing_matrix

# --- 카탈로그 role 분류 동의어 정규화 (팩 대칭 커버리지) --------------------
# 통합 카탈로그(layout_templates.json)의 role은 팩마다 같은 시각 archetype을
# 다른 이름으로 부른다(house_a=data/roadmap/decision, house_b=chart/timeline/matrix).
# dispatch가 `renderer` 미매칭 시 `role`로 폴백하므로, house_b role을 동일 archetype
# 렌더러로 라우팅한다. 필드 shape가 다른 템플릿은 렌더러 내부에서 content로 우아하게 폴백.
REGISTRY["chart"] = data_chart          # house_b: column_*/line_chart (categories+values/series) 호환
REGISTRY["timeline"] = roadmap_gantt    # house_b: gantt_timeline (weeks+workstreams) 호환
REGISTRY["matrix"] = matrix_priority    # house_b: prioritization_matrix (items{x_band,y_band}) 호환
REGISTRY["kpi"] = kpi_dashboard         # house_b: stat_hero (필드 다름 → 보강 대상, 폴백 무해)
REGISTRY["table"] = table_block         # house_b: assessment_table (필드 다름 → 보강 대상)
REGISTRY["assessment"] = table_block
REGISTRY["trends"] = strategy_pillars   # house_b: three_trends_* (필드 다름 → 보강 대상)
REGISTRY["areas"] = strategy_pillars    # house_b: five_key_areas/overview_areas
REGISTRY["risk"] = risk_dashboard       # 카탈로그 role='risk' → risk_dashboard(렌더러명만 있던 누락 보정)
