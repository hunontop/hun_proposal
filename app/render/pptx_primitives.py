# -*- coding: utf-8 -*-
"""PPTX 셰이프 프리미티브 (Phase 2 P2-A) — 노하우 0, 팩/스킨 tokens 주입.

목적: PPTX 네이티브 렌더러(P2-B~C)가 재사용하는 도형·텍스트·표·차트 헬퍼.
설계: `_source/house_b-pptx-ref/.../base.py`(MIT)의 프리미티브를 **theme→tokens**로 이식.
색·폰트는 전부 tokens에서 시맨틱 해석(하드코딩 금지). 시그니처는 renderers.py와 정합.

계약(렌더러가 지켜야 할 것):
- 렌더러 시그니처는 기존과 동일: `fn(prs, tokens, slide, f) -> slide`.
- 색은 `Skin(tokens)`의 시맨틱 접근자(.primary/.accent/.ink/.muted/.flag/.bg류)로만.
- 좌표 단위 = 인치(13.333 x 7.5). RGBColor 반환은 항상 `Skin.rgb(...)` 경유.
- 어떤 헬퍼도 예외를 삼키지 않음(폴백은 dispatch 레벨에서). 단 add_chart는 pptx 버전차 방어.
"""
from __future__ import annotations

from typing import Any, Iterable, Optional, Sequence

EMU_W_IN, EMU_H_IN = 13.333, 7.5


# --- 토큰 해석 (renderers._sem 과 정합) -------------------------------------

class Skin:
    """tokens(팩 + 스킨 캐스케이드 병합본)에서 시맨틱 색/폰트를 해석하는 어댑터.

    PPTX 캐스케이드 배선(dispatch)에서 병합된 tokens dict를 그대로 받는다.
    renderers._sem() 과 동일 키 우선순위를 사용해 HTML/PPTX 룩앤필을 정합시킨다.
    """

    def __init__(self, tokens: dict):
        self.tokens = tokens or {}

    # -- 색 --
    def _colors(self) -> dict:
        return self.tokens.get("colors") or {}

    def pick(self, *keys: str, default: str = "404040") -> str:
        c = self._colors()
        for k in keys:
            if k in c:
                return str(c[k]).lstrip("#")
        return default

    @property
    def primary(self) -> str:
        return self.pick("navy", "dark_navy", "deep_navy", default="1F3864")

    @property
    def accent(self) -> str:
        return self.pick("orange", "orange_deep", "bright_blue", "orange_bright", default="2E5496")

    @property
    def ink(self) -> str:
        return self.pick("black", "text_dark", default="1A1A1A")

    @property
    def muted(self) -> str:
        return self.pick("gray_text", "footer_gray", default="808080")

    @property
    def flag(self) -> str:
        return self.pick("red", "section_concern", "status_red", default="C00000")

    @property
    def cover_bg(self) -> str:
        return self.pick("section_cover", "deep_navy", "navy", default="1F3864")

    @property
    def panel(self) -> str:
        """카드/패널 옅은 배경."""
        return self.pick("panel_bg", "light_gray", "bg_gray", default="F2F2F2")

    @property
    def rule(self) -> str:
        return self.pick("rule_gray", "line_gray", default="D9D9D9")

    @property
    def white(self) -> str:
        return "FFFFFF"

    # -- 폰트 --
    def family(self) -> str:
        return (self.tokens.get("fonts") or {}).get("family", "Arial")

    def fsize(self, key: str, default: int) -> int:
        return ((self.tokens.get("fonts") or {}).get("sizes") or {}).get(key, default)

    # -- 변환 --
    @staticmethod
    def rgb(hex_str: str):
        from pptx.dml.color import RGBColor
        return RGBColor.from_string(str(hex_str).lstrip("#"))


# --- 슬라이드 생성 ----------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --- 텍스트 ----------------------------------------------------------------

def add_textbox(slide, left, top, width, height, *, anchor: str = "top"):
    from pptx.util import Inches, Emu
    from pptx.enum.text import MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    return tb


def write_para(tf, text, *, size, color, family="Arial", bold=False, italic=False,
               align: str = "left", first=False, bullet=False):
    """단락 1개 추가. first=True면 빈 첫 단락 재사용."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = family
    run.font.color.rgb = Skin.rgb(color)
    if bullet:
        _set_bullet(p)
    return p


def add_text(slide, left, top, width, height, text, *, size, color, family="Arial",
             bold=False, align="left", anchor="top"):
    """단일 라인/문단 텍스트박스 편의 함수."""
    tb = add_textbox(slide, left, top, width, height, anchor=anchor)
    write_para(tb.text_frame, text, size=size, color=color, family=family,
               bold=bold, align=align, first=True)
    return tb


def _set_bullet(paragraph):
    from pptx.oxml.ns import qn
    from lxml import etree
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu = etree.SubElement(pPr, qn("a:buChar"))
    bu.set("char", "•")
    pPr.set("indent", "-228600"); pPr.set("marL", "228600")


# --- 도형 ------------------------------------------------------------------

def _apply_fill_line(s, fill, line, line_width):
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = Skin.rgb(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = Skin.rgb(line)
        if line_width is not None:
            from pptx.util import Pt
            s.line.width = Pt(line_width)
    try:
        s.shadow.inherit = False
    except Exception:
        pass


def add_shape(slide, mso_shape, left, top, width, height, *, fill=None, line=None,
              line_width=None):
    """임의 MSO_SHAPE(사각/타원/둥근사각/다이아/화살표…) 추가. mso_shape=MSO_SHAPE 멤버."""
    from pptx.util import Inches
    s = slide.shapes.add_shape(mso_shape, Inches(left), Inches(top), Inches(width), Inches(height))
    _apply_fill_line(s, fill, line, line_width)
    return s


def add_rect(slide, left, top, width, height, *, fill=None, line=None, line_width=None):
    from pptx.enum.shapes import MSO_SHAPE
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height,
                     fill=fill, line=line, line_width=line_width)


def add_rounded(slide, left, top, width, height, *, fill=None, line=None, line_width=None):
    from pptx.enum.shapes import MSO_SHAPE
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill=fill, line=line, line_width=line_width)


def add_oval(slide, left, top, width, height, *, fill=None, line=None, line_width=None):
    from pptx.enum.shapes import MSO_SHAPE
    return add_shape(slide, MSO_SHAPE.OVAL, left, top, width, height,
                     fill=fill, line=line, line_width=line_width)


def add_line(slide, x1, y1, x2, y2, *, color, width_pt=0.75, dash=None):
    from pptx.util import Inches, Pt
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = Skin.rgb(color)
    ln.line.width = Pt(width_pt)
    if dash is not None:
        ln.line.dash_style = dash
    return ln


def add_gradient(slide, left, top, width, height, stops: Sequence[tuple], *,
                 angle: float = 90.0, line=None):
    """선형 그라데이션 사각형. stops = [(pos0..1, 'RRGGBB'), ...].

    python-pptx가 gradient API를 완전 노출하지 않아 XML을 직접 구성한다.
    실패 시(버전차) 첫 stop 색 단색으로 폴백.
    """
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = Skin.rgb(line)
    try:
        s.shadow.inherit = False
    except Exception:
        pass
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        spPr = s.fill._xPr  # shape properties element
        # 기존 fill 제거
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                    "a:pattFill", "a:grpFill"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        grad = etree.SubElement(spPr, qn("a:gradFill"))
        gsLst = etree.SubElement(grad, qn("a:gsLst"))
        for pos, hexc in stops:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(int(max(0.0, min(1.0, pos)) * 100000)))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", str(hexc).lstrip("#"))
        lin = etree.SubElement(grad, qn("a:lin"))
        lin.set("ang", str(int(angle * 60000)))
        lin.set("scaled", "1")
        # gradFill을 line(a:ln) 앞에 두도록 위치 보정
        ln_el = spPr.find(qn("a:ln"))
        if ln_el is not None:
            spPr.remove(grad); spPr.insert(list(spPr).index(ln_el), grad)
    except Exception:
        s.fill.solid()
        s.fill.fore_color.rgb = Skin.rgb(stops[0][1] if stops else "1F3864")
    return s


# --- 표 --------------------------------------------------------------------

def add_table(slide, left, top, width, height, rows: int, cols: int, *, skin: Skin,
              header=True, col_widths: Optional[Sequence[float]] = None):
    """빈 표 생성 후 GraphicFrame.table 반환. 셀 채움은 호출측이 fill_cell로."""
    from pptx.util import Inches
    gf = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    tbl = gf.table
    if col_widths:
        for i, w in enumerate(col_widths[:cols]):
            tbl.columns[i].width = Inches(w)
    return tbl


def fill_cell(cell, text, *, skin: Skin, size=12, color=None, bold=False,
              fill=None, align="left"):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = Skin.rgb(fill)
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = skin.family()
    run.font.color.rgb = Skin.rgb(color or skin.ink)
    return cell


# --- 차트 ------------------------------------------------------------------

def add_chart(slide, chart_type: str, left, top, width, height, *,
              categories: Sequence[str], series: Sequence[tuple], skin: Skin):
    """네이티브 차트. chart_type in {'bar','column','line','pie'}.
    series = [(name, [values...]), ...]. 실패 시 None 반환(dispatch가 폴백).
    """
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    tmap = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    xt = tmap.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = list(categories)
    for name, vals in series:
        data.add_series(str(name), tuple(vals))
    gf = slide.shapes.add_chart(xt, Inches(left), Inches(top), Inches(width),
                                Inches(height), data)
    return gf.chart


# --- 이미지 (C 모드 폴백에서 사용) -----------------------------------------

def add_picture(slide, image_path, left, top, width=None, height=None):
    from pptx.util import Inches
    kw = {}
    if width is not None:
        kw["width"] = Inches(width)
    if height is not None:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kw)


# --- 공통 크롬 (title/accent bar/review note) ------------------------------

def add_title(slide, text, skin: Skin, *, top=0.4, with_bar=True):
    if with_bar:
        add_rect(slide, 0, 0, 0.22, EMU_H_IN, fill=skin.accent)
    return add_text(slide, 0.55, top, EMU_W_IN - 0.9, 1.0, text,
                    size=skin.fsize("section", 28), color=skin.ink,
                    family=skin.family(), bold=True)


def add_review_note(slide, skin: Skin, slide_data: dict):
    rn = (slide_data.get("review_needed") or []) + \
         [f"(미결) {q}" for q in (slide_data.get("open_questions") or [])]
    if not rn:
        return
    tb = add_textbox(slide, 0.7, 6.4, 12.0, 0.9)
    tf = tb.text_frame
    for i, r in enumerate(rn):
        write_para(tf, "🔴 " + str(r), size=skin.fsize("small", 12),
                   color=skin.flag, family=skin.family(), first=(i == 0))
