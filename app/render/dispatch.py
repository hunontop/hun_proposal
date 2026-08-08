# -*- coding: utf-8 -*-
"""렌더 디스패치 — 정본 SlideModel → 노하우 팩 → PPTX.

인터페이스(plan.md "★ M1 정본 계약 (3)"): `add(template_id, **fields)` + 적응형 추론.
(house_b PresentationBuilder.add / add_specs 패턴 채택 — 출처 seulee26/house_b-pptx, MIT.)

엔진은 노하우 0: 색·폰트는 **팩 tokens에서 주입**, 템플릿→렌더러는 **팩 templates.json**이 정한다.
미지원 template_id → fallback 레이아웃 + review_needed 노트.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Callable

ROOT = Path(__file__).resolve().parents[2]
PACKS = ROOT / "packs"
PACKS_EXCLUDED = ROOT / "packs_excluded"  # 배제 하우스 자산 — --pack 명시 로드만(결정 11·12). W31 E3: 실물은 <개발 원본 전용 경로> 격리, 이 경로는 상시 부재(is_dir() 폴백만 수행)


def _pack_dir(name: str) -> Path:
    base = PACKS / name
    return base if base.is_dir() else PACKS_EXCLUDED / name


class Pack:
    """노하우 팩 로더(②→ 주입). tokens + templates."""

    def __init__(self, name: str):
        base = _pack_dir(name)
        self.name = name
        self.tokens: dict[str, Any] = _load(base / "tokens.json")
        templates = _load(base / "templates.json")
        # templates.json은 배열 또는 {templates:[...]} 둘 다 허용
        items = templates.get("templates", templates) if isinstance(templates, dict) else templates
        self.templates: dict[str, dict] = {t["id"]: t for t in items}

    def template(self, template_id: str) -> dict | None:
        return self.templates.get(template_id)


def _load(path: Path) -> Any:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _cascade_tokens(base: dict, skins: "list | None") -> dict:
    """PPTX 토큰 캐스케이드 — pack.tokens(base) 위에 skins를 순서대로 병합.

    HTML(htmlgen._cascade_skins)과 동일 규칙(뒤가 앞을 덮음, 부분 스킨 허용)으로
    룩앤필을 정합시킨다. 단 HTML은 skins=[pack]부터 시작하는 반면, PPTX는
    항상 pack.tokens를 base로 깔고 skins를 오버라이드로 얹어 하위호환을 보장한다.
    """
    from copy import deepcopy
    from .htmlgen import _deep_merge, _resolve_skin

    merged = deepcopy(base or {})
    for src in (skins or []):
        _deep_merge(merged, _resolve_skin(src))
    return merged


class Deck:
    """렌더 빌더. add(template_id, **fields)로 슬라이드를 쌓고 save()."""

    def __init__(self, pack: str | Pack, skins: "list | None" = None,
                 mode: str = "native"):
        self.pack = pack if isinstance(pack, Pack) else Pack(pack)
        self.skins = skins
        self.mode = mode or "native"
        # ③스킨 축: pack.tokens(base) + skins 캐스케이드 → 렌더러가 소비할 tokens.
        self.tokens = _cascade_tokens(self.pack.tokens, skins)
        self._prs = None  # python-pptx Presentation (지연 생성)
        self.warnings: list[str] = []

    def _ensure_prs(self):
        if self._prs is None:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            layout = self.tokens.get("layout", {})
            prs.slide_width = Inches(layout.get("slide_w_in", 13.333))
            prs.slide_height = Inches(layout.get("slide_h_in", 7.5))
            self._prs = prs
        return self._prs

    def add_slide(self, slide: dict):
        """정본 SlideModel slide 1개 → 렌더러 디스패치(renderer→role→content)."""
        from . import renderers

        prs = self._ensure_prs()
        tid = slide.get("template_id")
        tdef = self.pack.template(tid or "")
        rname = (tdef or {}).get("renderer")
        role = slide.get("role", "")
        matched = renderers.REGISTRY.get(rname) or renderers.REGISTRY.get(role)
        fn: Callable = matched or renderers.REGISTRY["content"]
        sid = slide.get("slide_id")
        if tdef:
            missing = [k for k in tdef.get("required_fields", []) if k not in (slide.get("fields") or {})]
            if missing:
                self.warnings.append(f"slide {sid} ({tid}): required_fields 누락 {missing}")
        # 디스패치 폴백: renderer명/role에 맞는 네이티브 렌더러가 없어 generic 사용.
        is_cover = matched is renderers.REGISTRY.get("cover")
        if matched is None and not is_cover:
            self.warnings.append(
                f"slide {sid} ({tid}): 적합 네이티브 렌더러 없음"
                f"(renderer={rname!r}/role={role!r}) → generic 렌더 (밋밋)"
            )
        fn(prs, self.tokens, slide, dict(slide.get("fields") or {}))
        # 렌더러 내부 폴백(네이티브 렌더러가 필드 shape 불일치로 content로 떨어짐).
        reason = slide.pop("_fallback_reason", None)
        if reason and matched is not None:
            self.warnings.append(
                f"slide {sid} ({tid}): 네이티브 렌더러가 {reason} → generic 폴백 (밋밋)"
            )
        return self

    def add(self, template_id: str, **fields):
        """하위호환: template_id + fields → 합성 slide로 add_slide."""
        slide = {"template_id": template_id, "role": "", "fields": fields}
        for k in ("title", "key_message", "body"):
            if k in fields:
                slide[k] = fields[k]
        return self.add_slide(slide)

    def render_deck(self, deck: dict):
        for s in deck.get("slides", []):
            self.add_slide(s)
        return self

    def save(self, out_path: str | Path):
        prs = self._ensure_prs()
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return out_path


def add_specs(deck: dict, pack: str, out_path: str | Path,
              skins: "list | None" = None, mode: str = "native") -> dict:
    """정본 deck → PPTX 한 방. 경고 리포트 반환.

    - skins: ③스킨 축 캐스케이드 소스 리스트(dict/경로/스킨이름/팩이름). pack.tokens 위에 오버라이드.
    - mode: 'native'(기본, 네이티브 셰이프) | 'image'(전 슬라이드 HTML→PNG) | 'hybrid'(네이티브 우선·미지원만 이미지).
            image/hybrid는 HTML→PNG 래스터라이저가 있어야 활성 — 없으면 경고 후 native로 폴백.
    """
    mode = (mode or "native").lower()
    d = Deck(pack, skins=skins, mode=mode)
    if mode in ("image", "hybrid"):
        from . import rasterize
        if not rasterize.available():
            d.warnings.append(
                f"pptx-mode={mode} 요청됐으나 HTML→PNG 래스터라이저(playwright) 미설치 "
                f"→ native로 폴백. 활성화하려면 'pip install playwright && playwright install chromium'."
            )
            mode = "native"; d.mode = "native"
    d.render_deck(deck).save(out_path)
    return {"out": str(out_path), "warnings": d.warnings, "mode": mode}


def images_to_pptx(image_paths: "list", pack: str, out_path: str | Path,
                    skins: "list | None" = None) -> dict:
    """슬라이드 PNG 목록(순서대로) → 슬라이드당 1이미지 풀블리드 PPTX.

    W4(S6-2 본체): `add_specs`와 달리 deck를 렌더러에 태우지 않는다 — 이미 승인된
    화면(rasterize.html_to_slide_pngs의 캡처)을 그대로 담는다. 그래서 override·
    image_slots·manual_layer가 전부 "찍혀서" 실린다(텍스트가 아니라 픽셀로).
    대가: 이 PPTX는 편집 불가·접근성 없음·native보다 용량이 크다.
    """
    from pptx import Presentation
    from pptx.util import Inches
    from . import pptx_primitives as prim

    p = Pack(pack)
    tokens = _cascade_tokens(p.tokens, skins)
    layout = tokens.get("layout", {})
    w_in = layout.get("slide_w_in", 13.333)
    h_in = layout.get("slide_h_in", 7.5)

    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    for img in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        prim.add_picture(slide, img, 0, 0, width=w_in, height=h_in)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {
        "out": str(out), "mode": "image", "slides": len(image_paths),
        "warnings": ["이미지 PPTX: 텍스트 편집 불가·스크린리더 접근성 없음·native 대비 파일 용량 증가"],
    }
