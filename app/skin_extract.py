# -*- coding: utf-8 -*-
"""스킨 추출기 (Phase 1-B B2) — 레퍼런스 pptx → skin tokens.json (결정론적, LLM 미개입).

1차 예시스튜디오 수동추출(`CONTEXT/research/PHASE1B_token_extraction.md`)의 로직을 스크립트화.
python-pptx로 전 슬라이드 shape/run을 순회하며 fill color·font color·font size·shape
geometry 빈도를 집계한 뒤, 결정론적 규칙으로 렌더러가 읽는 tokens 스키마
(`colors`/`fonts.sizes`/`fonts.family`)로 정규화한다.

출력 tokens는 render()의 skins= 인자에 바로 넣을 수 있는 **부분(partial) 스킨**이다
(캐스케이드 병합 대상). 색 이름은 렌더러 어휘(navy/orange/gray_*/…)에 맞춰 자동 분류.

사용: python app/skin_extract.py <ref.pptx> [out_tokens.json] [--report report.json]
"""
from __future__ import annotations

import json
import sys
from collections import Counter
from pathlib import Path
from typing import Any, Iterable


# --- pptx 순회 -------------------------------------------------------------

def _iter_shapes(shapes) -> Iterable:
    """그룹을 재귀적으로 펼쳐 모든 leaf shape를 산출."""
    for sh in shapes:
        yield sh
        if getattr(sh, "shape_type", None) is not None and sh.shape_type == 6:  # GROUP
            try:
                yield from _iter_shapes(sh.shapes)
            except Exception:
                continue


def _rgb_hex(color) -> str | None:
    """python-pptx ColorFormat → 6자리 대문자 hex(RGB일 때만; theme/None은 건너뜀)."""
    try:
        if color is None or color.type is None:
            return None
        rgb = color.rgb  # RGBColor; theme color면 예외
        return str(rgb).upper()
    except Exception:
        return None


def _collect(pptx_path: Path) -> dict[str, Counter]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    fill_colors: Counter = Counter()
    font_colors: Counter = Counter()
    font_sizes: Counter = Counter()
    families: Counter = Counter()
    geometry: Counter = Counter()

    for slide in prs.slides:
        for sh in _iter_shapes(slide.shapes):
            # fill color
            try:
                if getattr(sh, "fill", None) is not None:
                    hx = _rgb_hex(sh.fill.fore_color)
                    if hx:
                        fill_colors[hx] += 1
            except Exception:
                pass
            # shape geometry preset
            try:
                geom = getattr(sh, "adjustments", None)
                st = getattr(sh, "auto_shape_type", None)
                if st is not None:
                    geometry[str(st)] += 1
            except Exception:
                pass
            # text runs
            tf = getattr(sh, "text_frame", None)
            if tf is None and not getattr(sh, "has_text_frame", False):
                continue
            try:
                tf = sh.text_frame
            except Exception:
                continue
            for para in tf.paragraphs:
                for run in para.runs:
                    f = run.font
                    hx = _rgb_hex(f.color)
                    if hx:
                        font_colors[hx] += 1
                    if f.size is not None:
                        font_sizes[round(f.size.pt, 1)] += 1
                    if f.name:
                        families[f.name] += 1
    return {
        "fill_colors": fill_colors,
        "font_colors": font_colors,
        "font_sizes": font_sizes,
        "families": families,
        "geometry": geometry,
    }


# --- 색 분류(결정론적) -----------------------------------------------------

def _hex_to_rgb(hx: str) -> tuple[int, int, int]:
    hx = hx.lstrip("#")
    return int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)


def _classify(hx: str) -> str:
    """hex → 렌더러 어휘 기본 이름(변형 접미사는 _build_colors가 부여)."""
    r, g, b = _hex_to_rgb(hx)
    mx, mn = max(r, g, b), min(r, g, b)
    light = mx / 255.0
    sat = 0.0 if mx == 0 else (mx - mn) / mx
    # 무채색
    if sat < 0.12:
        if light > 0.85:
            return "white" if light > 0.96 else "gray_bg"
        if light > 0.55:
            return "gray_line"
        if light > 0.20:
            return "gray_text"
        return "black"
    # 유채색: hue 계산
    mxf, mnf = mx / 255.0, mn / 255.0
    rf, gf, bf = r / 255.0, g / 255.0, b / 255.0
    d = mxf - mnf
    if mxf == rf:
        h = ((gf - bf) / d) % 6
    elif mxf == gf:
        h = (bf - rf) / d + 2
    else:
        h = (rf - gf) / d + 4
    h *= 60
    if h < 15 or h >= 345:
        return "red"
    if h < 45:
        return "orange"
    if h < 70:
        return "yellow"
    if h < 170:
        return "green"
    if h < 200:
        return "blue2"
    if h < 255:
        return "navy" if light < 0.55 else "blue2"
    if h < 290:
        return "purple"
    return "red"  # magenta계열은 red로


# 저빈도 색 잡음 컷 + 기본 이름당 변형 상한(잡음 억제, 결정론적).
_COLOR_MIN_COUNT = 8
_MAX_VARIANTS = 2  # base + 최대 2개(_deep/_bright)


def _build_colors(fill: Counter, font: Counter) -> dict[str, str]:
    """fill+font 빈도 합산 → 이름별 변형 접미사(_deep 어둡게 / _bright 밝게) 부여.

    저빈도 색은 컷(_COLOR_MIN_COUNT), 이름당 변형은 상한(_MAX_VARIANTS)으로 잡음 억제.
    """
    total: Counter = Counter()
    total.update(fill)
    total.update(font)
    # 기본 이름별로 후보 색을 모아 빈도순 정렬(저빈도 컷)
    buckets: dict[str, list[tuple[str, int]]] = {}
    for hx, cnt in total.items():
        if cnt < _COLOR_MIN_COUNT:
            continue
        buckets.setdefault(_classify(hx), []).append((hx, cnt))
    colors: dict[str, str] = {}
    for base, cands in buckets.items():
        cands.sort(key=lambda x: (-x[1], x[0]))  # 빈도 desc, hex asc(안정정렬)
        base_hex = cands[0][0]
        colors[base] = base_hex
        base_light = max(_hex_to_rgb(base_hex)) / 255.0
        used_suffix: set[str] = set()
        for hx, _cnt in cands[1:]:
            if hx == base_hex or len(used_suffix) >= _MAX_VARIANTS:
                break
            lt = max(_hex_to_rgb(hx)) / 255.0
            suffix = "deep" if lt < base_light else "bright"
            if suffix in used_suffix:
                continue
            used_suffix.add(suffix)
            colors[f"{base}_{suffix}"] = hx
    return colors


# --- 타입스케일(결정론적 밴드) --------------------------------------------

# (슬롯 이름, 하한 pt) — 큰 것부터. 각 밴드의 최빈 사이즈를 대표값으로.
_SIZE_BANDS = [
    ("title", 40), ("section", 28), ("body", 15),
    ("label", 13), ("small", 11), ("caption", 9), ("footer", 0),
]


def _build_sizes(sizes: Counter) -> dict[str, float | int]:
    out: dict[str, float | int] = {}
    for name, low in _SIZE_BANDS:
        # 이 밴드 [low, 다음밴드상한) 안의 사이즈들
        upper = next((l for n, l in _SIZE_BANDS if l > low), 10_000)
        band = {sz: c for sz, c in sizes.items() if low <= sz < upper}
        if not band:
            continue
        rep = max(band.items(), key=lambda x: (x[1], -x[0]))[0]  # 최빈, 동률이면 작은 값
        out[name] = int(rep) if float(rep).is_integer() else rep
    return out


def _clean_family(families: Counter) -> str | None:
    """최빈 폰트명에서 굵기 접미사를 떼 패밀리 루트 추정."""
    if not families:
        return None
    name = families.most_common(1)[0][0]
    for suf in (" Thin", " ExtraLight", " Light", " Medium", " SemiBold",
                " ExtraBold", " Black", " Bold", " Regular"):
        if name.endswith(suf):
            name = name[: -len(suf)]
    return name.strip() or None


# --- 조립 -----------------------------------------------------------------

def extract_skin(pptx_path: str | Path) -> dict[str, Any]:
    """레퍼런스 pptx → {'tokens': <skin>, 'report': <raw 빈도>}."""
    stats = _collect(Path(pptx_path))
    colors = _build_colors(stats["fill_colors"], stats["font_colors"])
    sizes = _build_sizes(stats["font_sizes"])
    fonts: dict[str, Any] = {}
    fam = _clean_family(stats["families"])
    if fam:
        fonts["family"] = fam
    if sizes:
        fonts["sizes"] = sizes
    tokens: dict[str, Any] = {}
    if colors:
        tokens["colors"] = colors
    if fonts:
        tokens["fonts"] = fonts
    report = {
        "source": str(pptx_path),
        "fill_colors": stats["fill_colors"].most_common(15),
        "font_colors": stats["font_colors"].most_common(15),
        "font_sizes": sorted(stats["font_sizes"].items(), key=lambda x: -x[1])[:15],
        "families": stats["families"].most_common(10),
        "geometry": stats["geometry"].most_common(10),
    }
    return {"tokens": tokens, "report": report}


def main(argv: list[str]) -> int:
    args = [a for a in argv if not a.startswith("--")]
    report_path = None
    for a in argv:
        if a.startswith("--report="):
            report_path = a.split("=", 1)[1]
    if not args:
        print("사용: python app/skin_extract.py <ref.pptx> [out_tokens.json] [--report=report.json]")
        return 1
    src = args[0]
    result = extract_skin(src)
    out = args[1] if len(args) > 1 else None
    if out:
        Path(out).parent.mkdir(parents=True, exist_ok=True)
        Path(out).write_text(json.dumps(result["tokens"], ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"skin tokens → {out}")
    else:
        print(json.dumps(result["tokens"], ensure_ascii=False, indent=2))
    if report_path:
        Path(report_path).write_text(json.dumps(result["report"], ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"report → {report_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
