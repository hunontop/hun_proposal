# -*- coding: utf-8 -*-
"""Network-free dashboard smoke tests.

Run from the repository root:
    python dashboard/test_smoke.py
"""
from __future__ import annotations

import io
import json
import os
import re
import shutil
import sqlite3
import tempfile
import unittest
import uuid
from pathlib import Path
from types import SimpleNamespace
from unittest import mock

import build_data
import server

import sys as _sys
_sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "app"))
_sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "app" / "render"))
import image_slots  # noqa: E402  (stage9 B-9)

_sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "proposal_system" / "scripts"))
import deck_review  # noqa: E402  (W3c 승인 전 덱 평가)
import design_brief  # noqa: E402  (W3a 브리핑 / W3b 슬롯 동기화)
import review_resolve  # noqa: E402  (W5 검토요망 해소)
import pipeline_state  # noqa: E402  (W7-C2 단계 순서)
import proposal_pipeline  # noqa: E402  (W7-C3 브리핑 정지)
import skeleton  # noqa: E402  (W10 표준 시나리오 스켈레톤 역제안)
import message_map  # noqa: E402  (W15 메시지맵 결정 게이트 산출물)
import bind as _bind  # noqa: E402  (app/ — W5 재투영용 바인더 / W7-C1 자동배정 폴백)
import adapt_storyline as _adapt  # noqa: E402  (W7-C1 자동배정 provenance)
import ingest as _ingest  # noqa: E402  (W7-C1 템플릿 카탈로그)
import enrich as _enrich  # noqa: E402  (W7-C1 폴백 슬라이드 상호작용)
import htmlgen  # noqa: E402  (W8-D1 r_data 방어)
import docgen  # noqa: E402  (W31 리허설 마찰3 — deck.doc.html 문서형 파생 뷰)
import storyline_prompt  # noqa: E402  (W8-D1 schema 예시)
import knowledge_gaps  # noqa: E402  (가공 어휘 갭 로그 — find/update_status)
import curate  # noqa: E402  (큐레이션 생애주기 — 라이브러리·register·참고자료 반입)
import imagedeck  # noqa: E402  (W28 이미지 렌더 트랙 — bundle/collect/compose·라우트 분기)
import journey_check  # noqa: E402  (W31 리허설 마찰4 — journey 폴더 검토_체크.md 채널)
import journey_folders  # noqa: E402  (W31 R7 단계 폴더 여정 — journey/ 파생 뷰·매뉴얼)
import design_contract  # noqa: E402  (W31 R-B(R2·R5) run별 디자인 계약)
import gates  # noqa: E402  (W31 리허설 마찰2 — 관문 다이얼 + 조건부 승격)
import company  # noqa: E402  (W31 리허설 마찰6 — 제안사 자사 프로필 창고)
import archive  # noqa: E402  (W31 리허설 마찰9 — run 보관소 왕복)
import design_knowledge_cards  # noqa: E402  (δ패킷 — A6 선택 디자인지식 카드 결정론 운반)
import knowledge_ledger  # noqa: E402  (ε패킷 — 지식 소비 체계: config 1점화·원장·지시 오버레이·안전장치)


REPO_ROOT = Path(__file__).resolve().parents[1]

# 공개 배포판: CLAUDE.md 등 개발 원본 전용 파일은 재배포에서 제외됐다.
# 이를 읽는 검증(Claude 대리 금지 문구 등)은 원본 파일이 있을 때만 돈다.
_HAS_CLAUDE_MD = (REPO_ROOT / "CLAUDE.md").exists()
_requires_claude_md = unittest.skipUnless(
    _HAS_CLAUDE_MD, "CLAUDE.md는 개발 원본 전용 — 배포판에서는 검증 생략"
)

# W31 E3: 하우스 팩(house_a/house_b)은 <개발 원본 전용 경로> 격리됐다.
# --pack 명시 시 하우스 스타일이 살아있는지를 검증하는 테스트는 팩이 로컬에 있을 때만 돈다
# (없으면 skip — core 경로는 다른 테스트들이 그대로 검증). 재사용하려면 위 경로의
# house_a/를 packs_excluded/house_a/로 되돌려라(README 참고).
_HAS_HOUSE_A = (
    (REPO_ROOT / "packs" / "house_a" / "pack.json").exists()
    or (REPO_ROOT / "packs_excluded" / "house_a" / "pack.json").exists()
)
_requires_house_a = unittest.skipUnless(
    _HAS_HOUSE_A, "house_a 하우스 팩이 로컬에 없음(Reuse 격리) — 코어 경로만 검증"
)
REQUIRED_BID_KEYS = {
    "dday", "close_dt", "bid_name", "inst_name", "budget_label",
    "budget_num", "bid_no", "detail_url", "id", "digest_date",
}


def _create_bids_db(path: Path, rows: list[tuple]) -> None:
    con = sqlite3.connect(str(path))
    try:
        con.execute(
            """
            CREATE TABLE bids (
                bid_no TEXT PRIMARY KEY,
                bid_ord TEXT,
                bid_name TEXT,
                inst_name TEXT,
                demand_inst TEXT,
                notice_dt TEXT,
                close_dt TEXT,
                budget TEXT,
                detail_url TEXT,
                collected_at TEXT,
                raw TEXT
            )
            """
        )
        con.executemany("INSERT INTO bids VALUES (?,?,?,?,?,?,?,?,?,?,?)", rows)
        con.commit()
    finally:
        con.close()


def _post_json(path: str, payload: dict) -> tuple[int, dict]:
    """Invoke a JSON POST Handler branch without opening a socket."""
    encoded = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    handler = server.Handler.__new__(server.Handler)
    handler.path = path
    handler.headers = {"Content-Length": str(len(encoded))}
    handler.rfile = io.BytesIO(encoded)
    captured: dict = {}

    def capture(code, obj):
        captured["code"] = code
        captured["body"] = obj

    handler._json = capture
    handler.do_POST()
    return captured["code"], captured["body"]


def _feedback_post(payload: dict) -> tuple[int, dict]:
    return _post_json("/api/feedback", payload)


class LoadBidsFallbackSmoke(unittest.TestCase):
    def test_empty_newest_digest_falls_back_to_prior_nonempty(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            digests = root / "digest"
            digests.mkdir()
            (digests / "2026-06-29.md").write_text(
                "# 입찰공고 다이제스트\n\n> 신규 공고 없음", encoding="utf-8"
            )
            (digests / "2026-06-28.md").write_text(
                "\n".join([
                    "# 입찰공고 다이제스트",
                    "| D-day | 마감일시 | 사업명 | 공고기관 | 사업금액 |",
                    "|---|---|---|---|---|",
                    "| D-3 | 2026-07-02 10:00 | 안전한 사업 | 테스트기관 | 1000원 |",
                ]),
                encoding="utf-8",
            )
            db = root / "bids.db"
            _create_bids_db(db, [(
                "SAFE-001", "00", "안전한 사업", "테스트기관", "수요기관",
                "2026-06-28 09:00", "2026-07-02 10:00", "1000",
                "https://example.invalid/bid", "2026-06-28T09:00:00", "{}",
            )])
            with (
                mock.patch.object(build_data, "DIGEST_DIR", digests),
                mock.patch.object(build_data, "DB", db),
            ):
                bids = build_data.load_bids()
            self.assertEqual(1, len(bids))
            self.assertEqual("SAFE-001", bids[0]["bid_no"])
            self.assertEqual("2026-06-28", bids[0]["digest_date"])
            self.assertTrue(REQUIRED_BID_KEYS.issubset(bids[0]))

    def test_db_fallback_excludes_replacement_character_rows(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            digests = root / "digest"
            digests.mkdir()
            (digests / "2026-06-29.md").write_text(
                "# 빈 digest\n\n> 신규 공고 없음", encoding="utf-8"
            )
            db = root / "bids.db"
            _create_bids_db(db, [
                (
                    "SAFE-002", "00", "정상 한글 사업", "정상기관", "수요기관",
                    "2026-06-28 08:00", "2026-07-03 10:00", "2000",
                    "https://example.invalid/safe", "2026-06-28T08:00:00", "{}",
                ),
                (
                    "BAD-001", "00", "손상\ufffd사업", "기관", "수요기관",
                    "2026-06-28 07:00", "2026-07-04 10:00", "3000",
                    "https://example.invalid/bad", "2026-06-28T07:00:00", "{}",
                ),
            ])
            with (
                mock.patch.object(build_data, "DIGEST_DIR", digests),
                mock.patch.object(build_data, "DB", db),
            ):
                bids = build_data.load_bids()
            self.assertEqual(["SAFE-002"], [row["bid_no"] for row in bids])
            self.assertTrue(REQUIRED_BID_KEYS.issubset(bids[0]))
            self.assertFalse(build_data._safe_text("\ud800"))


class PersistenceSmoke(unittest.TestCase):
    def test_save_card_feedback_and_reviewed_merge_roundtrip(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            analysis = root / "analysis"
            feedback = root / "feedback.json"
            analyzer = SimpleNamespace(os=os, ANALYSIS_DIR=str(analysis))
            with mock.patch.object(server, "_load_analyzer", return_value=analyzer):
                saved = server.save_card("CARD-001", "# 분석카드\n정상 내용")
            self.assertTrue(Path(saved["path"]).is_file())
            self.assertIn("정상 내용", Path(saved["path"]).read_text(encoding="utf-8"))

            feedback.write_text(
                json.dumps(
                    {"CARD-001": {"decision": "go", "memo": "유지할 메모"}},
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            with mock.patch.object(server, "FEEDBACK", feedback):
                legacy = server.load_feedback()
                self.assertFalse(legacy["CARD-001"]["reviewed"])

                code, response = _feedback_post({"bid_no": "CARD-001", "reviewed": True})
                self.assertEqual(200, code)
                self.assertTrue(response["entry"]["reviewed"])
                self.assertEqual("go", response["entry"]["decision"])
                self.assertEqual("유지할 메모", response["entry"]["memo"])

                code, response = _feedback_post({"bid_no": "CARD-001", "memo": "새 메모"})
                self.assertEqual(200, code)
                self.assertTrue(response["entry"]["reviewed"])
                self.assertEqual("go", response["entry"]["decision"])
                self.assertEqual("새 메모", response["entry"]["memo"])
                persisted = server.load_feedback()["CARD-001"]
                self.assertTrue(persisted["reviewed"])
                self.assertEqual("go", persisted["decision"])

    def test_run_analyze_uses_stubbed_attachment_boundary(self):
        with tempfile.TemporaryDirectory() as td:
            combined_text = mock.Mock(return_value=("첨부 본문", ["local.txt"]))
            analyzer = SimpleNamespace(
                find_bid=lambda _bid: (
                    "STUB-001", json.dumps({"bidNtceNm": "스텁 공고"}, ensure_ascii=False)
                ),
                combined_text=combined_text,
                build_prompt=lambda raw, body, manifest: "로컬 프롬프트",
                deterministic_facts=lambda raw: [("공고명", raw["bidNtceNm"])],
                os=os,
                ANALYSIS_DIR=td,
            )
            with mock.patch.object(server, "_load_analyzer", return_value=analyzer):
                result = server.run_analyze("STUB-001")
            self.assertEqual("STUB-001", result["bid_no"])
            self.assertEqual(["local.txt"], result["manifest"])
            self.assertEqual("로컬 프롬프트", result["prompt"])
            combined_text.assert_called_once()


class ImageSlotsSmoke(unittest.TestCase):
    """stage9 B-9: 이미지 슬롯 렌더(degrade)+생성 게이트. W27 D6: evidence도 생성 허용(딱지 표시)·tier<2 degrade."""

    def _ov(self):
        return {"version": 1, "slides": {"6": {"image_slots": [
            {"id": "m1", "role": "mood", "prompt": "무드"},
            {"id": "c1", "role": "conceptual", "prompt": "개념"},
            {"id": "ev", "role": "evidence", "prompt": "실적"},
        ]}}}

    def test_placeholder_and_embed_gate(self):
        for role in ("mood", "conceptual", "evidence"):
            h = image_slots.render_slot_html({"id": "s", "role": role})
            self.assertIn(f"dov-slot--{role}", h)
        # S5a: evidence 도 **실제 자산**은 임베드 허용(생성 금지와 별개 게이트 — fill 쪽은 test_fill_gating).
        h = image_slots.render_slot_html({"id": "e", "role": "evidence"}, asset_svg="<svg><rect/></svg>")
        self.assertIn("<svg", h)
        # 래스터 자산 → 상대경로 <img>(base64 아님).
        h = image_slots.render_slot_html({"id": "e", "role": "evidence"},
                                         asset_src="stage9_design/slots/slide6_e.png")
        self.assertIn('<img src="stage9_design/slots/slide6_e.png"', h)
        # mood/conceptual은 <svg>만 추출해 인라인(주변 텍스트 제거).
        h = image_slots.render_slot_html({"id": "c", "role": "conceptual"},
                                         asset_svg="pre<svg id='g'><rect/></svg>post")
        self.assertIn("<svg id='g'>", h)
        self.assertNotIn("pre", h)

    def test_fill_gating(self):
        run = tempfile.mkdtemp()
        try:
            calls = []

            def fake(prompt, meta):
                calls.append(meta["slot"])
                return "x <svg viewBox='0 0 8 6'><circle/></svg> y"

            # tier 0 → 전부 degrade, 러너 미호출.
            r0 = image_slots.fill_images(json.loads(json.dumps(self._ov())), run, tier=0, runner=fake)
            self.assertEqual([], calls)
            self.assertEqual(0, len(r0["skipped_evidence"]))   # W27 D6: 더 이상 자동 skip 없음
            self.assertEqual(3, len(r0["degraded"]))
            # tier 2 → mood/conceptual/evidence 전부 생성(W27 D6: evidence도 개방).
            r2 = image_slots.fill_images(json.loads(json.dumps(self._ov())), run, tier=2, runner=fake)
            self.assertEqual(3, len(r2["generated"]))
            self.assertEqual(0, len(r2["skipped_evidence"]))
            self.assertIn("slide6:ev", calls)
            # 재실행 → 자산 캐시(재렌더 결정론).
            r2b = image_slots.fill_images(json.loads(json.dumps(self._ov())), run, tier=2, runner=fake)
            self.assertEqual(3, len(r2b["cached"]))
            self.assertEqual(0, len(r2b["generated"]))
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_allow_generate_overrides_tier_for_all_roles(self):
        """W3b/N3-4: 단발 위임(--fill-images)은 tier 다이얼과 무관 — W27 D6로 evidence도 포함."""
        run = tempfile.mkdtemp()
        try:
            calls = []

            def fake(prompt, meta):
                calls.append(meta["slot"])
                return "<svg viewBox='0 0 8 6'><circle/></svg>"

            # tier 0 이어도 allow_generate=True 면 mood/conceptual/evidence 전부 생성.
            r = image_slots.fill_images(json.loads(json.dumps(self._ov())), run, tier=0,
                                        runner=fake, allow_generate=True)
            self.assertEqual(3, len(r["generated"]))
            self.assertEqual(0, len(r["skipped_evidence"]))
            self.assertIn("slide6:ev", calls)
            # 반대로 tier 2 여도 allow_generate=False 면 생성 안 함(명시가 tier를 이긴다).
            calls.clear()
            run2 = tempfile.mkdtemp()
            try:
                r2 = image_slots.fill_images(json.loads(json.dumps(self._ov())), run2, tier=2,
                                             runner=fake, allow_generate=False)
                self.assertEqual([], calls)
                self.assertEqual(3, len(r2["degraded"]))
                self.assertEqual(0, len(r2["skipped_evidence"]))
            finally:
                shutil.rmtree(run2, ignore_errors=True)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_empty_prompt_slot_skips_generation_not_generic(self):
        """결정 2026-07-15: 빈 프롬프트 mood/conceptual 슬롯은 제네릭 생성 대신 skip+경고.

        재현(run gen_R26BK01631547): 빈 프롬프트 bg 슬롯이 다크 '플렉서스' 제네릭을 뽑았다.
        이제 프롬프트가 비면 러너에 도달조차 안 하고 skipped_no_prompt로 표면화(placeholder degrade).
        """
        run = tempfile.mkdtemp()
        try:
            calls = []

            def fake(prompt, meta):
                calls.append(meta["slot"])
                return "<svg viewBox='0 0 8 6'><circle/></svg>"

            ov = {"version": 1, "slides": {"1": {"image_slots": [
                {"id": "bg", "role": "mood", "layer": "background", "format": "svg"},  # 빈 프롬프트
                {"id": "c", "role": "conceptual", "prompt": "실주제"},                  # 프롬프트 있음
            ]}}}
            rep = image_slots.fill_images(ov, run, tier=2, runner=fake, allow_generate=True)
            self.assertEqual(["slide1:bg"], rep["skipped_no_prompt"])  # 빈 프롬프트 → 생성 안 함
            self.assertEqual(["slide1:c"], rep["generated"])           # 프롬프트 있는 건 생성
            self.assertNotIn("slide1:bg", calls)                       # 러너에 도달조차 안 함
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_generated_evidence_persists_flag_and_renders_tag(self):
        """③ W27 D6: evidence 생성 시 slot dict에 generated:true 영속 +
        render_slot_html에 "AI 생성 예시" 딱지. generated_resolved 있으면 딱지 생략."""
        run = tempfile.mkdtemp()
        try:
            def fake(prompt, meta):
                return "<svg viewBox='0 0 8 6'><circle/></svg>"

            ov = {"version": 1, "slides": {"1": {"image_slots": [
                {"id": "ev", "role": "evidence", "prompt": "실적 그래프"},
            ]}}}
            rep = image_slots.fill_images(ov, run, tier=2, runner=fake, allow_generate=True)
            self.assertEqual(["slide1:ev"], rep["generated"])
            slot = ov["slides"]["1"]["image_slots"][0]
            self.assertTrue(slot.get("generated"))  # override dict에 영속(재렌더 결정론)

            asset_path = image_slots.slot_asset_path(run, "1", "ev", fmt="svg")
            asset_svg = asset_path.read_text(encoding="utf-8")
            h = image_slots.render_slot_html(slot, asset_svg=asset_svg)
            self.assertIn("AI 생성 예시", h)
            self.assertIn("dov-slot__tag--gen", h)

            # 사람이 generated_resolved를 명시하면 딱지 생략(코드가 임의로 안 지운다).
            slot["generated_resolved"] = "실제 실적 데이터로 확인, 표기 유지 불필요"
            h2 = image_slots.render_slot_html(slot, asset_svg=asset_svg)
            self.assertNotIn("AI 생성 예시", h2)

            # mood/conceptual 생성물은 generated:true라도 딱지 없음(장식이라 기만 소지 없음).
            mood_slot = {"id": "m", "role": "mood", "generated": True}
            h3 = image_slots.render_slot_html(mood_slot, asset_svg="<svg><rect/></svg>")
            self.assertNotIn("AI 생성 예시", h3)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_compute_image_provenance_counts(self):
        """④ image_provenance 카운트 정확 — 생성 evidence 미해소 N 포함."""
        run = Path(tempfile.mkdtemp())
        try:
            real_asset_path = run / "design_refs" / "real.jpg"
            real_asset_path.parent.mkdir(parents=True, exist_ok=True)
            real_asset_path.write_bytes(b"fake-jpg")

            ov = {"version": 1, "slides": {
                "1": {"image_slots": [
                    {"id": "ev1", "role": "evidence", "generated": True},          # 미해소
                    {"id": "ev2", "role": "evidence", "generated": True,
                     "generated_resolved": "해소 사유"},                            # 해소됨
                    {"id": "m1", "role": "mood", "generated": True},               # 생성(비-evidence)
                ]},
                "2": {"image_slots": [
                    {"id": "web1", "role": "conceptual", "source_url": "https://example.com/a.jpg"},
                    {"id": "web2", "role": "conceptual", "source_route": "web_sample"},  # 출처 미기록
                    {"id": "asset1", "role": "evidence", "path": "design_refs/real.jpg"},
                ]},
            }}
            prov = image_slots.compute_image_provenance(ov, run_dir=run)
            self.assertEqual(6, prov["total"])
            self.assertEqual(3, prov["generated"])
            self.assertEqual(2, prov["generated_evidence"])
            self.assertEqual(1, prov["generated_evidence_unresolved"])
            self.assertEqual(2, prov["web_sample"])
            self.assertEqual(1, prov["web_sample_sourced"])
            self.assertEqual(1, prov["real_asset"])
            self.assertEqual(2, prov["placeholder"])  # web1·web2: 자산 없음, generated 아님
        finally:
            shutil.rmtree(run, ignore_errors=True)


class ImageProvenanceSurfaceSmoke(unittest.TestCase):
    """W27 P2 D6·D7: image_provenance → ship 경고 문자열 + web_sample 임베드 경고(비차단)."""

    def test_ship_warnings_surface_generated_tag_and_web_sample(self):
        """⑤ ship 경고 문자열 표면화(생성 딱지 잔존·web_sample) — 둘 다 0이면 경고 없음."""
        gating_empty = {"image_provenance": {
            "generated_evidence_unresolved": 0, "web_sample": 0, "web_sample_sourced": 0,
        }}
        self.assertEqual([], proposal_pipeline._image_provenance_ship_warnings(gating_empty))

        gating = {"image_provenance": {
            "generated_evidence_unresolved": 2, "web_sample": 3, "web_sample_sourced": 1,
        }}
        warnings = proposal_pipeline._image_provenance_ship_warnings(gating)
        self.assertEqual(2, len(warnings))
        self.assertIn("AI 생성 evidence 딱지 잔존 2건", warnings[0])
        self.assertIn("웹 수급 자산 3건", warnings[1])
        self.assertIn("출처 기록 1건", warnings[1])

        # 게이트 리포트에 image_provenance 자체가 없으면 조용히 빈 리스트(부가 관측·차단자 아님).
        self.assertEqual([], proposal_pipeline._image_provenance_ship_warnings({}))

    def test_embed_time_web_sample_missing_source_url_warns_not_blocks(self):
        """web_sample 슬롯인데 source_url 없으면 경고 1줄 — 검증 오류(errors)는 아니다."""
        overrides = self._import_overrides()
        ov = {"version": 1, "slides": {"1": {"image_slots": [
            {"id": "w1", "role": "conceptual", "source_route": "web_sample"},          # 출처 없음
            {"id": "w2", "role": "conceptual", "source_route": "web_sample",
             "source_url": "https://example.com/b.jpg"},                              # 출처 있음
        ]}}}
        warns = overrides.image_slot_warnings(ov)
        self.assertEqual(1, len(warns))
        self.assertIn("w1", warns[0])
        self.assertIn("source_url", warns[0])

    @staticmethod
    def _import_overrides():
        import overrides  # noqa: (app/render 경로는 모듈 상단에서 이미 sys.path에 실려있다)
        return overrides


class DesignBriefSlotSyncSmoke(unittest.TestCase):
    """W3b: design_brief.image_slots_plan → design_overrides.json(정본) 추가 동기화."""

    def _brief(self):
        return {"image_slots_plan": {"slots": [
            {"slide_id": 10, "id": "conceptual_1", "role": "conceptual", "format": "svg",
             "treatment": "밝게", "prompt": "", "source": "기본값"},
            {"slide_id": 3, "id": "mood_1", "role": "mood", "format": "svg",
             "treatment": "어둡게", "prompt": "표지 무드"},
        ]}}

    def test_adds_only_missing_and_strips_plan_fields(self):
        ov = {"version": 1, "slides": {}}
        added = design_brief.sync_slots_into_overrides(self._brief(), ov)
        self.assertEqual(["slide10:conceptual_1", "slide3:mood_1"], added)
        slot = ov["slides"]["10"]["image_slots"][0]
        self.assertEqual("conceptual", slot["role"])
        self.assertNotIn("slide_id", slot)     # 계획 전용 필드는 override에 새지 않는다
        self.assertNotIn("source", slot)
        self.assertNotIn("prompt", slot)       # 빈 문자열은 안 싣는다(fill 쪽이 빈 프롬프트를 skip)
        self.assertEqual("표지 무드", ov["slides"]["3"]["image_slots"][0]["prompt"])

    def test_idempotent_and_preserves_human_edits(self):
        ov = {"version": 1, "slides": {"10": {"image_slots": [
            {"id": "conceptual_1", "role": "mood", "prompt": "사람이 고침"}]}}}
        added = design_brief.sync_slots_into_overrides(self._brief(), ov)
        self.assertEqual(["slide3:mood_1"], added)                       # 10번은 이미 있음 → 추가 안 함
        self.assertEqual(1, len(ov["slides"]["10"]["image_slots"]))
        self.assertEqual("mood", ov["slides"]["10"]["image_slots"][0]["role"])   # 편집본 보존
        self.assertEqual([], design_brief.sync_slots_into_overrides(self._brief(), ov))  # 멱등


class DesignChecksSmoke(unittest.TestCase):
    """W3b/N3-5: 결정론 디자인 게이트 — deck.html 마크업 실측(0토큰·차단 없음)."""

    def _html(self, body: str) -> str:
        return f'<html><body>{body}</body></html>'

    def _slide(self, n: int, inner: str) -> str:
        return f'<section id="slide-{n}" class="slide">{inner}</section>'

    def test_density_bounds_and_overflow_risk(self):
        import design_checks
        thin = self._slide(1, "<h2>짧다</h2>")
        fat = self._slide(2, "<div>" + ("가" * 950) + "</div>")
        long_line = self._slide(3, "<p>" + ("나" * 130) + "</p>")
        bullets = self._slide(4, "<ul>" + "<li>항목입니다 조금 길게</li>" * 9 + "</ul>")
        dc = design_checks.compute_design_checks(self._html(thin + fat + long_line + bullets))
        by = {s["slide_id"]: s for s in dc["slides"]}
        self.assertIn("density_under", by["1"]["flags"])
        self.assertIn("density_over", by["2"]["flags"])
        self.assertIn("overflow_risk", by["3"]["flags"])
        self.assertIn("bullets_over", by["4"]["flags"])
        self.assertEqual(9, by["4"]["bullets"])
        self.assertEqual(1, dc["summary"]["overflow_risk"])
        self.assertEqual("warn", dc["status"])
        self.assertNotIn("fail", dc["status"])  # 게이트는 차단하지 않는다

    def test_image_slot_fulfillment_measured_from_markup(self):
        import design_checks
        ph = ('<div class="dov-slot dov-slot--ph dov-slot--mood" data-role="mood">'
              '<span class="dov-slot__tag">MOOD</span><span class="dov-slot__cap">무드</span></div>')
        filled = ('<div class="dov-slot dov-slot--conceptual" data-role="conceptual">'
                  '<span class="dov-slot__tag">CONCEPT</span><svg><rect/></svg></div>')
        html = self._html(self._slide(1, "<p>본문이 충분히 길게 들어간 슬라이드 본문 텍스트입니다.</p>" + ph)
                          + self._slide(2, "<p>본문이 충분히 길게 들어간 슬라이드 본문 텍스트입니다.</p>" + filled))
        dc = design_checks.compute_design_checks(html)
        im = dc["summary"]["image_slots"]
        self.assertEqual({"total": 2, "filled": 1, "placeholder": 1, "fulfillment": 0.5}, im)
        # 슬롯(장식)은 본문 텍스트 밀도에 섞이지 않는다 — 두 슬라이드 본문 길이가 같아야 한다.
        self.assertEqual(dc["slides"][0]["text_chars"], dc["slides"][1]["text_chars"])

    def test_no_slots_means_null_fulfillment_not_zero(self):
        import design_checks
        dc = design_checks.compute_design_checks(self._html(self._slide(1, "<p>" + "다" * 100 + "</p>")))
        self.assertIsNone(dc["summary"]["image_slots"]["fulfillment"])  # 지어내지 않는다
        self.assertEqual("pass", dc["status"])

    def test_unparsed_html_is_unmeasured_not_pass(self):
        """조용한 0 슬라이드 = 가짜 pass. 측정 실패는 실패라고 말해야 한다."""
        import design_checks
        dc = design_checks.compute_design_checks("<html><body><div>슬라이드 아님</div></body></html>")
        self.assertEqual("unmeasured", dc["status"])
        self.assertEqual(0, dc["summary"]["slides"])

    def test_legacy_sections_without_id_are_measured(self):
        """구 렌더본은 `<section class="slide">`(id 없음) — 인덱스로 번호를 매겨 측정한다."""
        import design_checks
        html = ('<html><body><section class="slide cover"><p>' + "가" * 100 + '</p></section>'
                '<section class="slide "><p>' + "나" * 100 + '</p></section></body></html>')
        dc = design_checks.compute_design_checks(html)
        self.assertEqual(2, dc["summary"]["slides"])
        self.assertEqual(["1", "2"], [s["slide_id"] for s in dc["slides"]])


class LayoutProbeSmoke(unittest.TestCase):
    """W6-A1: 브라우저 실측 계층 — 2026-07-08 정주행 마찰(슬롯이 차트를 덮는데 overflow_risk=0)의 재현.

    각 픽스처는 **먼저 정적 계층이 못 본다는 것**을 단언하고, 그 다음 브라우저가 본다는 것을 단언한다.
    (정적 계층이 언젠가 이걸 보게 되면 이 테스트가 먼저 깨진다 — 그건 좋은 실패다.)
    """

    SLIDE_CSS = (
        "<style>"
        "body{margin:0} section.slide{position:relative;width:1280px;height:720px;"
        "overflow:hidden;box-sizing:border-box;padding:40px;font-family:sans-serif}"
        ".chart{position:relative;width:600px;height:400px;border:1px solid #333}"
        ".dov-slot{position:absolute;left:100px;top:300px;width:300px;height:150px;"
        "background:#eee;z-index:5}"
        ".wash{position:absolute;left:0;right:0;bottom:0;height:40%;background:#f5f5f5;z-index:1}"
        ".fg{position:relative;z-index:9}"
        "</style>"
    )

    def _probe(self, body: str) -> dict:
        import layout_probe
        run = Path(tempfile.mkdtemp())
        html = f"<html><head>{self.SLIDE_CSS}</head><body>{body}</body></html>"
        p = run / "deck.html"
        p.write_text(html, encoding="utf-8")
        try:
            return layout_probe.probe_html(p)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def _skip_without_browser(self):
        import layout_probe
        if not layout_probe.available():
            self.skipTest(layout_probe.unavailable_reason() + " — 브라우저 계층 미측정(가짜 pass 금지)")

    # 재현 ① 슬롯이 차트 사각형을 침범(절대배치) — 정적 파싱은 원리적으로 못 본다.
    SLOT_OVER_CHART = (
        '<section id="slide-9" class="slide">'
        '<p>데이터 슬라이드 본문입니다. 평가 구조가 말하는 것 — 기술 80 대 가격 20.</p>'
        '<p>가격 경쟁이 아니라 기획의 설득력이 승부처입니다. 배점표는 과업지시서 확인.</p>'
        '<div class="chart"><span>80</span></div>'
        '<div class="dov-slot dov-slot--ph dov-slot--conceptual" data-role="conceptual">'
        '<span class="dov-slot__tag">CONCEPT</span></div>'
        '</section>'
    )

    def test_static_layer_misses_slot_over_chart(self):
        """정적 design_checks는 절대배치 슬롯이 차트를 침범해도 overflow_risk=0·플래그 없음."""
        import design_checks
        dc = design_checks.compute_design_checks(f"<html><body>{self.SLOT_OVER_CHART}</body></html>")
        self.assertEqual(0, dc["summary"]["overflow_risk"])
        self.assertEqual(["slot_placeholder"], dc["slides"][0]["flags"])  # 겹침은 안 보인다

    def test_browser_layer_catches_slot_over_chart(self):
        self._skip_without_browser()
        probe = self._probe(self.SLOT_OVER_CHART)
        self.assertEqual("warn", probe["status"])
        row = probe["slides"][0]
        self.assertIn("slot_overlaps_content", row["flags"])
        boxes = [o["box"] for o in row["content_overlaps"]]
        self.assertIn("div.chart", boxes)
        self.assertEqual(1, probe["summary"]["content_overlap"])

    # 재현 ② scrollHeight > clientHeight (마찰 로그가 쓴 프로브) + 경계 밖 자손 사각형.
    OVERFLOWING = (
        '<section id="slide-7" class="slide" style="overflow:visible">'
        '<p style="height:900px">본문이 슬라이드보다 길다 — 브라우저에서만 보이는 사실.</p>'
        '</section>'
    )

    def test_static_layer_misses_pixel_overflow(self):
        import design_checks
        dc = design_checks.compute_design_checks(f"<html><body>{self.OVERFLOWING}</body></html>")
        self.assertEqual(0, dc["summary"]["overflow_risk"])  # 짧은 문장 = long_lines 0

    def test_browser_layer_measures_overflow_px(self):
        self._skip_without_browser()
        probe = self._probe(self.OVERFLOWING)
        row = probe["slides"][0]
        self.assertIn("overflow_measured", row["flags"])
        self.assertGreater(row["scroll_overflow_px"], 100)
        self.assertGreater(row["content_overflow_px"], 100)
        self.assertEqual(1, probe["summary"]["overflow"])

    def test_background_wash_slot_is_not_a_false_positive(self):
        """커버의 전면 워시 슬롯(본문 아래·의도된 장식)은 신고하지 않는다 — 오탐이 게이트를 죽인다."""
        self._skip_without_browser()
        probe = self._probe(
            '<section id="slide-1" class="slide">'
            '<h1 class="fg">표지 제목입니다 — 충분한 길이의 본문 텍스트를 위해 늘립니다.</h1>'
            '<p class="fg">부제 문장. 여기에 본문이 조금 더 들어갑니다. 밀도 하한 통과용.</p>'
            '<div class="dov-slot wash" data-role="mood"></div>'
            '</section>'
        )
        self.assertEqual("pass", probe["status"])
        self.assertEqual([], probe["slides"][0]["flags"])

    def test_missing_playwright_is_unmeasured_not_pass(self):
        """의존성 부재는 '결함 없음'이 아니라 '안 봤음'이다."""
        import layout_probe
        run = Path(tempfile.mkdtemp())
        (run / "deck.html").write_text("<html><body></body></html>", encoding="utf-8")
        try:
            with mock.patch.object(layout_probe, "available", return_value=False):
                probe = layout_probe.probe_html(run / "deck.html")
        finally:
            shutil.rmtree(run, ignore_errors=True)
        self.assertEqual("unmeasured", probe["status"])
        self.assertIn("playwright", probe["reason"])
        self.assertNotEqual("pass", probe["status"])

    def test_attach_browser_layer_preserves_keys_and_escalates_status(self):
        import design_checks
        static = design_checks.compute_design_checks(
            f"<html><body>{self.SLOT_OVER_CHART}</body></html>")
        static["status"] = "pass"  # 정적 계층이 통과라고 말한 상황을 만든다
        merged = design_checks.attach_browser_layer(static, {
            "status": "warn", "method": "브라우저 실측",
            "summary": {"slides": 1, "overflow": 0, "occlusion": 0,
                        "content_overlap": 1, "void": 2},
            "slides": [],
        })
        self.assertEqual("warn", merged["status"])                    # 승격
        self.assertEqual(1, merged["summary"]["browser"]["content_overlap"])
        self.assertEqual(2, merged["summary"]["browser"]["void"])
        self.assertIn("overflow_risk", merged["summary"])             # 기존 소비 키 보존
        self.assertIn("image_slots", merged["summary"])
        self.assertIn("브라우저 실측", merged["method"])

        unmeasured = design_checks.attach_browser_layer(
            {"status": "pass", "method": "정적", "summary": {}},
            {"status": "unmeasured", "reason": "playwright 미설치", "summary": {}})
        self.assertEqual("pass", unmeasured["status"])  # 미측정은 status를 흔들지 않는다
        self.assertEqual("unmeasured", unmeasured["summary"]["browser"]["status"])


class LayoutProbeVoidSmoke(unittest.TestCase):
    """W27 P1b: 브라우저 원시 행의 영역 공허 임계·집계·수리 배선을 결정론 검증한다."""

    @staticmethod
    def _row(sid, *, void_candidates=None, scroll=0, content=0, occlusions=None):
        return {
            "slide_id": sid,
            "scroll_overflow_px": scroll,
            "content_overflow_px": content,
            "content_overflow_x_px": 0,
            "occlusions": occlusions or [],
            "content_overlaps": [],
            "decoration_occlusions": [],
            "readability": [],
            "void_candidates": void_candidates or [],
        }

    @staticmethod
    def _candidate(selector="div.card", *, height=400, ratio=0.75,
                   has_text=True, background=False):
        return {"selector": selector, "height_px": height, "void_ratio": ratio,
                "has_visible_text": has_text, "is_background": background}

    def _browser(self, *rows):
        import layout_probe
        slides, summary = layout_probe._summarize_rows(list(rows))
        return {"status": "warn" if any(r["flags"] for r in slides) else "pass",
                "viewport": {"width": 1280, "height": 720}, "method": "fake rows",
                "summary": summary, "slides": slides}

    def test_tall_high_void_text_block_is_flagged(self):
        browser = self._browser(self._row("3", void_candidates=[self._candidate()]))
        row = browser["slides"][0]
        self.assertIn("void_measured", row["flags"])
        self.assertEqual([{"selector": "div.card", "height_px": 400, "void_ratio": 0.75}],
                         row["void_blocks"])
        self.assertEqual(1, browser["summary"]["void"])

    def test_below_threshold_and_decorative_blocks_are_not_flagged(self):
        candidates = [
            self._candidate("div.short", height=239, ratio=0.9),
            self._candidate("div.filled", ratio=0.49),
            self._candidate("div.deco", ratio=0.9, has_text=False),
            self._candidate("div.background", ratio=0.9, background=True),
        ]
        browser = self._browser(self._row("4", void_candidates=candidates))
        self.assertNotIn("void_measured", browser["slides"][0]["flags"])
        self.assertEqual([], browser["slides"][0]["void_blocks"])
        self.assertEqual(0, browser["summary"]["void"])

    def test_void_is_a_typed_repair_target_and_review_quotes_ratio(self):
        import layout_probe
        browser = self._browser(self._row("8", void_candidates=[self._candidate(ratio=0.6)]))
        targets = layout_probe.repair_targets(browser)
        self.assertEqual("void", targets[0]["kind"])
        self.assertEqual(0.6, targets[0]["void_blocks"][0]["void_ratio"])
        measured = "\n".join(deck_review._render_browser_layer(browser))
        self.assertIn("void_ratio=0.6", measured)

    def test_existing_overflow_and_occlusion_aggregation_is_unchanged(self):
        browser = self._browser(
            self._row("7", scroll=9),
            self._row("9", occlusions=[{"target": "p", "ratio": 0.2}]),
        )
        self.assertEqual(1, browser["summary"]["overflow"])
        self.assertEqual(1, browser["summary"]["occlusion"])
        self.assertEqual(0, browser["summary"]["content_overlap"])
        self.assertIn("overflow_measured", browser["slides"][0]["flags"])
        self.assertIn("slot_occlusion", browser["slides"][1]["flags"])


class Stage9ScreenshotSmoke(unittest.TestCase):
    """W6-A2: stage9 번들이 "PNG를 보라"고 하면서 스크린샷을 안 붙이던 공백(MANUAL §10 불일치)."""

    def _pp(self):
        import proposal_pipeline
        return proposal_pipeline

    def test_block_says_no_screenshots_honestly(self):
        pp = self._pp()
        run = Path(tempfile.mkdtemp())
        try:
            block = pp._stage9_screenshot_block(run, [], "playwright 미설치", ["7"])
        finally:
            shutil.rmtree(run, ignore_errors=True)
        self.assertIn("스크린샷 없음", block)
        self.assertIn("텍스트", block)
        self.assertIn("playwright 미설치", block)

    def test_block_attaches_paths_and_marks_targets(self):
        pp = self._pp()
        run = Path(tempfile.mkdtemp())
        shots = [run / "assets" / "slides" / f"slide-{i:02d}.png" for i in (7, 9)]
        try:
            block = pp._stage9_screenshot_block(run, shots, None, ["9"])
        finally:
            shutil.rmtree(run, ignore_errors=True)
        self.assertIn("slide-07.png", block)
        self.assertIn("slide-09.png ← 대상", block)
        self.assertNotIn("slide-07.png ← 대상", block)
        self.assertIn("direct", block)

    def test_screenshots_are_actually_rasterized_into_the_run(self):
        pp = self._pp()
        import rasterize
        if not rasterize.available():
            self.skipTest(rasterize.unavailable_reason())
        run = Path(tempfile.mkdtemp())
        try:
            (run / "deck.html").write_text(
                '<html><body><section id="slide-1" class="slide" '
                'style="width:1280px;height:720px">A</section>'
                '<section id="slide-2" class="slide" '
                'style="width:1280px;height:720px">B</section></body></html>',
                encoding="utf-8")
            shots, reason = pp._stage9_screenshots(run)
            self.assertIsNone(reason)
            self.assertEqual(2, len(shots))
            for p in shots:
                self.assertTrue(p.is_file(), f"PNG 실존해야 한다: {p}")
                self.assertGreater(p.stat().st_size, 0)
            self.assertEqual(run / "assets" / "slides", shots[0].parent)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_missing_html_reports_reason_not_crash(self):
        pp = self._pp()
        run = Path(tempfile.mkdtemp())
        try:
            shots, reason = pp._stage9_screenshots(run)
        finally:
            shutil.rmtree(run, ignore_errors=True)
        self.assertEqual([], shots)
        self.assertIn("deck.html 없음", reason)


class DeckReviewSmoke(unittest.TestCase):
    """W3c/N6-1: 승인 전 덱 평가 — 산출물 계약 검증(결정론). LLM 호출 없음."""

    VALID = ("# 덱 평가 — r\n\n## 총평\n" + "가" * 320 +
             "\n\n## 슬라이드별\n- slide 1: 문제\n\n## 리스크\n- 없음\n\n"
             "## 승인 권고\n- verdict: approve\n- reasons: 이유\n")

    def test_valid_review_passes_and_verdict_extracted(self):
        self.assertEqual([], deck_review.validate(self.VALID))
        self.assertEqual("approve", deck_review.verdict_of(self.VALID))

    def test_missing_sections_and_bad_verdict_are_rejected(self):
        errs = deck_review.validate("# 덱 평가\n\n좋아 보인다.\n")
        self.assertEqual(4 + 1, len(errs))                       # 섹션 3 + verdict + 길이
        bad = self.VALID.replace("- verdict: approve", "- verdict: 좋음")
        self.assertTrue(any("verdict" in e for e in deck_review.validate(bad)))
        self.assertIsNone(deck_review.verdict_of(bad))           # 계약 밖 값은 verdict가 아니다

    def test_collect_reports_missing_file_without_pretending(self):
        run = tempfile.mkdtemp()
        try:
            rep = deck_review.collect(Path(run))
            self.assertFalse(rep["found"])
            self.assertIsNone(rep["verdict"])
            self.assertTrue(rep["errors"])
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_measurement_mismatch_is_surfaced_not_smoothed(self):
        """기록본과 지금 잰 값이 다르면 프롬프트가 '불일치'라고 말해야 한다(정합 신호)."""
        measured = {"measured_from": "m", "method": "p", "status": "pass",
                    "summary": {"slides": 2, "density_over": 0, "density_under": 0, "bullets_over": 0,
                                "overflow_risk": 0,
                                "image_slots": {"total": 0, "filled": 0, "placeholder": 0, "fulfillment": None}},
                    "slides": []}
        same = json.loads(json.dumps(measured))
        self.assertIn("**일치**", deck_review.render_measurements(measured, same))
        stale = json.loads(json.dumps(measured))
        stale["summary"]["density_under"] = 3
        self.assertIn("**불일치**", deck_review.render_measurements(measured, stale))
        self.assertIn("미측정", deck_review.render_measurements(measured, None))
        # 슬롯 0이면 충족률은 0.0이 아니라 "미정의"로 인용된다(지어내지 않는다).
        self.assertIn("미정의(슬롯 0)", deck_review.render_measurements(measured, same))


class ReviewResolveSmoke(unittest.TestCase):
    """W5 검토요망 해소 — 불변식(사람 기록 없으면 태그 불변) 중심."""

    TAG_FREE = "제도 인지도 수치는 근거 확인 필요"
    TAG_FIELD = "[필수입력 미확보] lead — 구조 데이터 필요(지어내지 않음)"

    def _deck(self):
        return {"slides": [
            {"slide_id": 3, "template_id": "problem_questions", "role": "문제",
             "title": "문제", "key_message": "핵심 질문",
             "body": ["기존 불릿"],
             "fields": {"core_question": "핵심 질문", "sub_questions": ["기존 불릿"]},
             "review_needed": [self.TAG_FREE, "다른 태그 — 결정 없음"]},
            {"slide_id": 8, "template_id": "org_roles", "role": "조직",
             "title": "조직", "key_message": "팀", "body": ["역할 A"],
             "fields": {"roles": ["역할 A"]},
             "review_needed": [self.TAG_FIELD]},
        ]}

    def _doc(self, items):
        return {"schema_version": 1, "items": items}

    def test_skeleton_covers_every_tag_and_infers_field_target(self):
        with tempfile.TemporaryDirectory() as td:
            doc, added = review_resolve.build_skeleton(Path(td) / "run", self._deck())
        self.assertEqual(3, len(doc["items"]))          # 태그 전수
        self.assertEqual(3, added)
        by_tag = {i["tag"]: i for i in doc["items"]}
        # 필드명은 태그 문자열에서 결정론 추론(bind가 박아 넣은 것) — 나머지는 body.
        self.assertEqual({"kind": "field", "name": "lead"}, by_tag[self.TAG_FIELD]["target"])
        self.assertEqual({"kind": "body"}, by_tag[self.TAG_FREE]["target"])
        self.assertTrue(all(i["decision"] == "" for i in doc["items"]))   # 코드는 결정하지 않는다

    def test_skeleton_preserves_existing_decisions_and_marks_stale(self):
        existing = self._doc([
            {"slide_id": 3, "tag": self.TAG_FREE, "decision": "deferred", "decided_at": "2026-07-08T00:00:00"},
            {"slide_id": 99, "tag": "사라진 태그", "decision": "no_basis_confirmed"},
        ])
        with tempfile.TemporaryDirectory() as td:
            doc, added = review_resolve.build_skeleton(Path(td) / "run", self._deck(), existing)
        by_key = {(i["slide_id"], i["tag"]): i for i in doc["items"]}
        self.assertEqual("deferred", by_key[(3, self.TAG_FREE)]["decision"])       # 사람 편집본 보존
        self.assertEqual("2026-07-08T00:00:00", by_key[(3, self.TAG_FREE)]["decided_at"])
        self.assertTrue(by_key[(99, "사라진 태그")]["stale"])                       # 기록은 지우지 않는다
        self.assertEqual(2, added)                                                 # 신규는 나머지 2건뿐

    def test_no_tag_removed_without_a_recorded_decision(self):
        """불변식: 결정이 없으면 코드는 태그를 절대 지우지 않는다(창작금지의 대칭)."""
        deck = self._deck()
        for doc in (None, self._doc([]),
                    self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": ""}]),
                    self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": "deferred"}]),
                    # fact 없는 fact_supplied = 무효 → 태그 유지
                    self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": "fact_supplied", "fact": ""}]),
                    # 미상 decision → 무효
                    self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": "resolved"}]),
                    # 태그 문자열 불일치 → 근사 매칭 없음
                    self._doc([{"slide_id": 3, "tag": "제도 인지도 수치는", "decision": "no_basis_confirmed"}]),
                    # stale 기록은 무시
                    self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": "no_basis_confirmed", "stale": True}])):
            with self.subTest(doc=doc):
                d = json.loads(json.dumps(deck))
                rep = review_resolve.apply(d, doc, binders=_bind.BINDERS)
                self.assertEqual(0, rep["tags_removed"])
                self.assertEqual([self.TAG_FREE, "다른 태그 — 결정 없음"], d["slides"][0]["review_needed"])
                self.assertEqual([self.TAG_FIELD], d["slides"][1]["review_needed"])

    def test_fact_supplied_writes_body_and_reprojects_bound_field(self):
        """body 기입분이 bind 파생 필드에 재투영돼야 렌더러(fields 소비)에 도달한다."""
        deck = self._deck()
        fact = "2025년 인지도 41.2%(발주처 자료)"
        rep = review_resolve.apply(deck, self._doc([
            {"slide_id": 3, "tag": self.TAG_FREE, "decision": "fact_supplied", "fact": fact},
        ]), binders=_bind.BINDERS)
        s3 = deck["slides"][0]
        self.assertEqual(1, rep["tags_removed"])
        self.assertEqual(1, rep["facts_applied"])
        self.assertIn(fact, s3["body"])
        self.assertIn(fact, s3["fields"]["sub_questions"])          # ← 재투영(없으면 html에 안 나온다)
        self.assertEqual(["다른 태그 — 결정 없음"], s3["review_needed"])  # 결정 없는 태그는 보존
        self.assertEqual([], rep["notes"])

    def test_fact_supplied_field_target_writes_field_only(self):
        deck = self._deck()
        rep = review_resolve.apply(deck, self._doc([
            {"slide_id": 8, "tag": self.TAG_FIELD, "decision": "fact_supplied",
             "fact": "홍길동 팀장", "target": {"kind": "field", "name": "lead"}},
        ]), binders=_bind.BINDERS)
        s8 = deck["slides"][1]
        self.assertEqual("홍길동 팀장", s8["fields"]["lead"])
        self.assertEqual(["역할 A"], s8["body"])                     # body는 안 건드린다
        self.assertEqual([], s8["review_needed"])
        self.assertEqual("fields.lead", rep["applied"][0]["wrote"])

    def test_no_basis_confirmed_removes_tag_but_keeps_content(self):
        deck = self._deck()
        before = json.loads(json.dumps(deck["slides"][1]))
        rep = review_resolve.apply(deck, self._doc([
            {"slide_id": 8, "tag": self.TAG_FIELD, "decision": "no_basis_confirmed"},
        ]), binders=_bind.BINDERS)
        s8 = deck["slides"][1]
        self.assertEqual(1, rep["tags_removed"])
        self.assertEqual(0, rep["facts_applied"])                    # 아무것도 기입하지 않는다
        for key in ("body", "fields", "key_message", "title"):
            self.assertEqual(before[key], s8[key])
        self.assertEqual([], s8["review_needed"])

    def test_apply_is_idempotent_and_report_is_measured_not_selfreported(self):
        deck = self._deck()
        doc = self._doc([{"slide_id": 3, "tag": self.TAG_FREE, "decision": "no_basis_confirmed"}])
        first = review_resolve.apply(deck, doc, binders=_bind.BINDERS)
        second = review_resolve.apply(deck, doc, binders=_bind.BINDERS)   # 태그가 이미 없다
        self.assertEqual(1, first["tags_removed"])
        self.assertEqual(0, second["tags_removed"])                        # 결정 수(1)를 되풀이 보고하지 않는다
        self.assertEqual(1, second["counts"]["no_basis_confirmed"])
        self.assertEqual(1, len(second["unmatched"]))

    def test_stamp_decisions_records_time_only_for_decided_items(self):
        doc = self._doc([
            {"slide_id": 3, "tag": self.TAG_FREE, "decision": "deferred", "decided_at": None},
            {"slide_id": 8, "tag": self.TAG_FIELD, "decision": "", "decided_at": None},
        ])
        self.assertEqual(1, review_resolve._stamp_decisions(doc))
        self.assertTrue(doc["items"][0]["decided_at"])
        self.assertIsNone(doc["items"][1]["decided_at"])

    def test_summarize_counts_match_resolution_records(self):
        deck = self._deck()
        doc = self._doc([
            {"slide_id": 3, "tag": self.TAG_FREE, "decision": "no_basis_confirmed"},
            {"slide_id": 8, "tag": self.TAG_FIELD, "decision": "deferred"},
        ])
        summary = review_resolve.summarize(review_resolve.apply(deck, doc, binders=_bind.BINDERS), "p.json")
        self.assertEqual(1, summary["resolved"])
        self.assertEqual(1, summary["deferred"])
        self.assertEqual(1, summary["tags_removed"])
        # "숫자가 좋아지면 의심하라": 제거 수는 해소 기록 수를 넘을 수 없다.
        self.assertLessEqual(summary["tags_removed"], summary["resolved"])
        # 3개 태그 중 해소 1건만 빠진다 — deferred·결정없음은 남는다.
        self.assertEqual(2, sum(len(s["review_needed"]) for s in deck["slides"]))
        self.assertIn(self.TAG_FIELD, deck["slides"][1]["review_needed"])


class TemplateFallbackSmoke(unittest.TestCase):
    """W7-C1 자동배정 폴백 — 재현: 근거 없는 '일정' 슬라이드가 강제 gantt 배정으로 태그 3건을 낳는다.

    폴백은 **검사를 없애는 게 아니라 거짓 배정을 없앤다**: 태그가 줄어든 자리에는 반드시
    template_fallback warning이 1건 생긴다(조용한 폴백 금지).
    """

    PACK = "core"  # W31 E3: house_a는 Reuse 격리 — core도 동일 계약(roadmap_gantt 등)을 갖는다
    GANTT_REQUIRED = ["time_units", "workstreams", "milestones"]

    def _storyline(self):
        return {"meta": {"project": "재현"}, "slides": [
            # 1: 자동배정 cover — 필수필드 일부(project_title/concept_message)가 채워짐 → 템플릿 유지
            {"n": 1, "section": "표지", "title": "제안서", "message": "우리가 한다", "bullets": []},
            # 2: 자동배정 gantt — 불릿이 workstreams를 채움 → 부분 충족이므로 템플릿 유지
            {"n": 2, "section": "일정", "title": "추진 일정", "message": "3개월 안에 끝낸다",
             "bullets": ["1개월차 착수", "2개월차 개발", "3개월차 검수"]},
            # 3: 자동배정 gantt — 근거 0 → 필수필드 3개 전부 미확보 → generic 폴백 (재현 대상)
            {"n": 3, "section": "일정", "title": "세부 일정", "message": "단계별로 진행", "bullets": []},
            # 4: 스토리라인이 **명시 지정** — 근거가 없어도 폴백 대상이 아니다
            {"n": 4, "section": "일정", "title": "명시 지정 일정", "message": "명시", "bullets": [],
             "template_id": "roadmap_gantt"},
        ]}

    def _bound(self, **kw):
        deck = _adapt.adapt_storyline(self._storyline(), project="재현", pack=self.PACK)
        report = _bind.bind_deck(deck, _ingest._load_templates(self.PACK), **kw)
        return deck, report, {s["slide_id"]: s for s in deck["slides"]}

    def test_provenance_marks_only_inferred_templates(self):
        deck = _adapt.adapt_storyline(self._storyline(), project="재현", pack=self.PACK)
        # 명시 지정(4)은 provenance에 없다 → 폴백 면제의 유일한 근거.
        self.assertEqual([1, 2, 3], deck["meta"][_adapt.AUTO_TEMPLATE_META_KEY])

    def test_unfillable_auto_template_falls_back_to_generic_with_a_warning(self):
        _, report, by_id = self._bound()
        # 재현: 폴백 전이라면 slide 3은 필수필드 3개가 전부 비어 태그 3건이었다.
        self.assertEqual(3, len(self.GANTT_REQUIRED))
        self.assertIsNone(by_id[3]["template_id"])           # generic (스키마의 null=미정)
        self.assertEqual([], by_id[3]["review_needed"])      # 태그 3건 → 0건
        self.assertEqual(1, len(report["template_fallback"]))  # …그 자리에 warning 1건
        fb = report["template_fallback"][0]
        self.assertEqual((3, "roadmap_gantt"), (fb["slide_id"], fb["from"]))
        self.assertEqual(self.GANTT_REQUIRED, fb["missing"])
        self.assertIn("generic 폴백", fb["warning"])

    def test_explicit_template_id_is_never_downgraded(self):
        _, report, by_id = self._bound()
        # slide 4는 slide 3과 콘텐츠가 동일(근거 0)하지만 명시 지정이므로 템플릿·태그가 그대로다.
        self.assertEqual("roadmap_gantt", by_id[4]["template_id"])
        self.assertEqual(3, len(by_id[4]["review_needed"]))
        self.assertNotIn(4, [f["slide_id"] for f in report["template_fallback"]])

    def test_partially_bound_auto_template_is_kept(self):
        _, report, by_id = self._bound()
        # 하나라도 채워지면(workstreams) 템플릿을 유지한다 — 검사를 없애지 않는다.
        self.assertEqual("roadmap_gantt", by_id[2]["template_id"])
        self.assertEqual(2, len(by_id[2]["review_needed"]))   # time_units, milestones
        self.assertEqual("cover_cinematic", by_id[1]["template_id"])
        self.assertEqual(1, len(by_id[1]["review_needed"]))   # visual_subject
        self.assertEqual([3], [f["slide_id"] for f in report["template_fallback"]])

    def test_fallback_is_idempotent_and_deferrable_for_enrich(self):
        deck, report, _ = self._bound(allow_fallback=False)
        self.assertEqual([], report["template_fallback"])              # 미룸 = 폴백 없음
        self.assertEqual(3, len(deck["slides"][2]["review_needed"]))   # 태그는 아직 붙어 있다
        # enrich가 필드를 채울 기회를 준 뒤 호출부가 판정한다.
        first = _bind.apply_template_fallback(deck, _ingest._load_templates(self.PACK))
        self.assertEqual([3], [f["slide_id"] for f in first])
        self.assertEqual([], deck["slides"][2]["review_needed"])       # 그 템플릿의 필수필드 태그만 회수
        # 멱등: 두 번째 호출은 아무것도 하지 않는다(template_id가 이미 없다).
        self.assertEqual([], _bind.apply_template_fallback(deck, _ingest._load_templates(self.PACK)))

    def test_fallback_never_touches_unrelated_tags(self):
        """W5 불변식: 폴백은 **그 템플릿의 필수필드 태그**만 회수한다(태그 임의 제거 금지)."""
        deck, _, _ = self._bound(allow_fallback=False)
        deck["slides"][2]["review_needed"].append("사람이 붙인 다른 태그")
        _bind.apply_template_fallback(deck, _ingest._load_templates(self.PACK))
        self.assertEqual(["사람이 붙인 다른 태그"], deck["slides"][2]["review_needed"])

    def test_enrich_does_not_retag_a_generic_slide(self):
        """폴백이 만든 generic 슬라이드를 enrich가 다시 태그하면 폴백이 무의미해진다."""
        deck, _, _ = self._bound()
        _enrich.enrich_deck(deck, {}, "")
        by_id = {s["slide_id"]: s for s in deck["slides"]}
        self.assertEqual([], by_id[3]["review_needed"])
        self.assertIsNone(by_id[3]["template_id"])


class ConfirmOrderSmoke(unittest.TestCase):
    """W7-C2/C3 — 해소→브리핑 순서, 그리고 브리핑 생성 시 1회 정지."""

    def _deck(self):
        return {"meta": {"project": "순서"}, "slides": [
            {"slide_id": 1, "template_id": "org_roles", "role": "조직", "title": "조직",
             "key_message": "팀 구성", "body": ["역할 A"], "fields": {"roles": ["역할 A"]},
             "review_needed": ["[필수입력 미확보] lead — 구조 데이터 필요(지어내지 않음)"]},
            {"slide_id": 2, "template_id": "org_roles", "role": "조직", "title": "조직 2",
             "key_message": "팀 구성 2", "body": ["역할 B"], "fields": {"roles": ["역할 B"]},
             "review_needed": ["[필수입력 미확보] teams — 구조 데이터 필요(지어내지 않음)"]},
        ]}

    def test_resolution_runs_before_the_brief(self):
        order = pipeline_state.STAGE_ORDER
        self.assertLess(order.index("review_resolve"), order.index("design_brief"))

    def test_evidence_candidates_exclude_resolved_slides(self):
        """재현: 해소 **전** 덱으로 브리핑을 만들면 서명이 끝난 슬라이드가 후보로 잔존한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            before = design_brief.build_default(run, self._deck())
            stale = [e["slide_id"] for e in before["image_slots_plan"]["evidence_candidates"]]
            self.assertEqual([1, 2], stale)          # 현행(역순) — 둘 다 후보

            deck = self._deck()
            doc = {"schema_version": 1, "items": [
                {"slide_id": 1, "tag": deck["slides"][0]["review_needed"][0],
                 "decision": "no_basis_confirmed"},
            ]}
            review_resolve.apply(deck, doc, binders=_bind.BINDERS)
            after = design_brief.build_default(run, deck)
            fresh = [e["slide_id"] for e in after["image_slots_plan"]["evidence_candidates"]]
        self.assertEqual([2], fresh)                 # 교정(해소→브리핑) — 해소 슬라이드 제외

    def test_brief_creation_pauses_go_once_then_proceeds(self):
        """C3: 첫 confirm은 브리핑에서 정지(created=True), 두 번째 go는 진행(멱등)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            (run / "deck.json").write_text(json.dumps(self._deck(), ensure_ascii=False), encoding="utf-8")
            args = SimpleNamespace(design_guide=None)
            with mock.patch.object(proposal_pipeline, "_record_state"):
                first = proposal_pipeline._go_design_brief(run, args)
                second = proposal_pipeline._go_design_brief(run, args)
        self.assertTrue(first)      # 생성 → go 정지
        self.assertFalse(second)    # 기존 보존(사람 편집본) → 정지 없이 진행


class RDataStrDefenseSmoke(unittest.TestCase):
    """W8-D1: 재현 — LLM이 comparison에 리스트 대신 문자열 요약("70% vs 30%")을 넣으면
    r_data가 AttributeError로 죽는다. 방어 후에는 폴백 렌더 + warning이어야 한다(조용한 폴백 금지)."""

    def _deck(self, comparison):
        return {"meta": {"project": "P"}, "slides": [
            {"slide_id": 1, "role": "data", "template_id": "data_interpretation",
             "title": "지표", "message": "핵심", "bullets": [],
             "fields": {"metric": "성장률", "comparison": comparison, "interpretation": ["해석"]}},
        ]}

    def test_str_comparison_still_raises_from_r_data_itself(self):
        # r_data 단독 호출은 방어 전과 동일하게 AttributeError를 낸다 — 방어는 호출부(render())의 몫이다.
        slide, fields = {"title": "t"}, {"metric": "m", "comparison": "70% vs 30%", "interpretation": ["x"]}
        with self.assertRaises(AttributeError):
            htmlgen.r_data(slide, fields)

    def test_str_comparison_falls_back_with_visible_warning_through_full_render(self):
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html(self._deck("70% vs 30%"), "core", out)  # 방어 전엔 여기서 크래시
            self.assertTrue(out.is_file())
        self.assertEqual(1, len(rep["warnings"]))
        self.assertIn("렌더 실패", rep["warnings"][0])
        self.assertIn("AttributeError", rep["warnings"][0])

    def test_valid_comparison_still_renders_chart(self):
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html(
                self._deck([{"label": "A", "value": 70}, {"label": "B", "value": 30}]),
                "core", out)
        self.assertEqual([], rep["warnings"])

    def test_schema_block_documents_data_interpretation_shape(self):
        block = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn("data_interpretation", block)
        self.assertIn('"comparison"', block)
        self.assertIn('"label"', block)
        self.assertIn('"value"', block)


class Stage9SlotPlacementPromptSmoke(unittest.TestCase):
    """W8-D2: 재현 — 슬롯 선언에 배치 CSS가 없으면 여러 슬롯이 본문 위에 쌓인다(overlap).
    프롬프트가 배치 CSS를 명시적으로 요구하는지 검증(계약은 텍스트지만 회귀 감지는 가능)."""

    def test_prompt_requires_placement_css_for_declared_slots(self):
        from proposal_pipeline import PROMPTS
        text = (PROMPTS / "stage9_design_director.md").read_text(encoding="utf-8")
        self.assertIn("#slide-N .dov-slot", text)
        self.assertIn("배치", text)


class DeckReviewRecollectSmoke(unittest.TestCase):
    """W8-B1: 재현 — deck_review.md가 수거 이후 사람 손으로 갱신되면(revise→approve) 낡은
    verdict가 영원히 남는다(재수거 경로가 없었다). mtime 비교로 재수거하되 계약 검증은 그대로 통과해야 한다."""

    VALID = ("# 덱 평가 — r\n\n## 총평\n" + "가" * 320 +
             "\n\n## 슬라이드별\n- slide 1: 문제\n\n## 리스크\n- 없음\n\n"
             "## 승인 권고\n- verdict: {v}\n- reasons: 이유\n")

    def _run_with_recorded_verdict(self, td, verdict):
        run = Path(td)
        old_stamp = (
            __import__("datetime").datetime.now() - __import__("datetime").timedelta(minutes=10)
        ).isoformat(timespec="seconds")
        state = pipeline_state._blank(run)
        state["stages"]["deck_review"] = {"at": old_stamp, "source": "recorded",
                                           "verdict": verdict, "chars": 500}
        (run / pipeline_state.STATE_NAME).write_text(
            json.dumps(state, ensure_ascii=False, indent=2), encoding="utf-8")
        return run

    def test_verdict_recollected_when_file_updated_after_recording(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run_with_recorded_verdict(td, "revise")
            (run / "deck_review.md").write_text(self.VALID.format(v="approve"), encoding="utf-8")
            view = pipeline_state.resolve(run)
        self.assertEqual("approve", view["stages"]["deck_review"]["verdict"])
        self.assertIn("recollected_at", view["stages"]["deck_review"])

    def test_contract_violation_keeps_old_verdict_not_silently_updated(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run_with_recorded_verdict(td, "revise")
            (run / "deck_review.md").write_text("# 계약 위반\n너무 짧다.\n", encoding="utf-8")
            view = pipeline_state.resolve(run)
        self.assertEqual("revise", view["stages"]["deck_review"]["verdict"])
        self.assertNotIn("recollected_at", view["stages"]["deck_review"])

    def test_unchanged_file_does_not_recollect(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run_with_recorded_verdict(td, "revise")
            # 파일 없음 → 재수거 트리거 자체가 없어야 한다(recorded 그대로).
            view = pipeline_state.resolve(run)
        self.assertEqual("revise", view["stages"]["deck_review"]["verdict"])
        self.assertNotIn("recollected_at", view["stages"]["deck_review"])


class GatingRenderBytesRefreshSmoke(unittest.TestCase):
    """W8-B2: 재현 — stage9 --apply 후 gating_report.applied_axes는 갱신되지만
    render.html(bytes 등)은 render 시점 값에 낡은 채로 남는다(74867 vs 실제 86171B 실측)."""

    def _deck(self):
        return {"meta": {"project": "P"}, "slides": [
            {"slide_id": 1, "role": "data", "template_id": "data_interpretation",
             "title": "지표", "message": "핵심", "bullets": [],
             "fields": {"metric": "m", "comparison": [{"label": "A", "value": 1}],
                        "interpretation": ["x"]}},
        ]}

    def test_apply_stage9_refreshes_render_html_bytes_not_just_applied_axes(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "deck.json").write_text(json.dumps(self._deck(), ensure_ascii=False), encoding="utf-8")

            stale_rep = htmlgen.render_html(self._deck(), "core", run / "deck.html")
            stale_bytes = stale_rep["bytes"]
            gating = {
                "applied_axes": {
                    "html": {"pack": "core", "skins": [], "overrides": False,
                              "image_slots": 0, "manual_layer": False},
                    "pptx": None,
                },
                "render": {"html": stale_rep, "pptx": None},
            }
            (run / "gating_report.json").write_text(
                json.dumps(gating, ensure_ascii=False, indent=2), encoding="utf-8")

            overrides = {"version": 1, "slides": {"1": {
                "css": "#slide-1 .dov-cover{padding:2vw;border:1px solid red;display:block}"}}}
            (run / "design_overrides.json").write_text(
                json.dumps(overrides, ensure_ascii=False), encoding="utf-8")

            args = SimpleNamespace(overrides=None, pack="core")
            rep, ov_path, _skins = proposal_pipeline.apply_stage9(run, args)
            proposal_pipeline._update_applied_axes(
                run, overrides_path=ov_path, html_path=Path(rep["out"]), render_rep=rep)

            report = json.loads((run / "gating_report.json").read_text(encoding="utf-8"))
            # bytes 규약은 프로젝트 전역과 동일하게 len(doc.encode("utf-8"))이다(rep["bytes"]) —
            # OS 파일 크기(stat().st_size)는 Windows에서 개행 변환(\n→\r\n)으로 더 커져 다른 수치다.
            self.assertEqual(rep["bytes"], report["render"]["html"]["bytes"])
            self.assertNotEqual(stale_bytes, report["render"]["html"]["bytes"])


class DocViewSmoke(unittest.TestCase):
    """W31 리허설 마찰3(CONTEXT/REHEARSAL_FRICTIONS_W31.md 3행): deck.html(정본·덱 형태)은
    검토·정독 가독성이 나쁘다 — 같은 SlideModel을 문서형으로 재조판한 deck.doc.html을
    render_html과 별도로 생성한다(app/render/docgen.py). deck.html 바이트는 절대 불변."""

    def _deck(self, *, review=False, example=False, bad_comparison=None):
        fields = {"main_claim": "핵심 주장", "supporting_points": ["근거1", "근거2"]}
        slide = {"slide_id": 1, "role": "제안개요", "template_id": "executive_summary",
                  "title": "제목", "key_message": "메시지", "body": [], "fields": fields}
        if review:
            slide["review_needed"] = ["실적 교체 필요"]
        if example:
            slide["example"] = True
        slides = [slide]
        if bad_comparison is not None:
            slides.append({"slide_id": 2, "role": "데이터", "template_id": "data_interpretation",
                            "title": "데이터", "key_message": "m", "body": [],
                            "fields": {"metric": "m", "comparison": bad_comparison,
                                       "interpretation": ["해석"]}})
        return {"meta": {"project": "테스트 프로젝트"}, "slides": slides}

    def test_render_html_bytes_unchanged_when_docgen_also_called(self):
        """정본 불변 회귀: 같은 입력이면 docgen 호출 여부와 무관하게 render_html 바이트가 같다."""
        deck = self._deck(review=True)
        with tempfile.TemporaryDirectory() as td:
            out1 = Path(td) / "a" / "deck.html"
            rep1 = htmlgen.render_html(deck, "core", out1)
            bytes1 = out1.read_bytes()
        with tempfile.TemporaryDirectory() as td:
            out2 = Path(td) / "b" / "deck.html"
            rep2 = htmlgen.render_html(deck, "core", out2)
            docgen.render_doc(deck, Path(td) / "b" / "deck.doc.html")  # 같은 호출 시퀀스에서 병행
            bytes2 = out2.read_bytes()
        self.assertEqual(rep1["bytes"], rep2["bytes"])
        self.assertEqual(bytes1, bytes2)

    def test_doc_view_creates_badge_and_section_headers(self):
        deck = self._deck()
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.doc.html"
            rep = docgen.render_doc(deck, out)
            self.assertTrue(out.is_file())
            text = out.read_text(encoding="utf-8")
        self.assertEqual(1, rep["slides"])
        self.assertIn("검토·회의용 문서 뷰", text)
        self.assertIn("제출물 아님", text)
        self.assertIn("deck.html", text)
        self.assertIn('<h1 class="doc-section">제안개요</h1>', text)

    def test_doc_view_shows_flag_badges_for_review_and_example(self):
        deck = self._deck(review=True, example=True)
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.doc.html"
            docgen.render_doc(deck, out)
            text = out.read_text(encoding="utf-8")
        self.assertIn("doc-badge--review", text)
        self.assertIn("검토요망", text)
        self.assertIn("doc-badge--example", text)
        self.assertIn("예시 데이터", text)

    def test_malformed_field_is_gracefully_omitted_not_crashed(self):
        """조판 실패 필드(예: comparison이 문자열)는 크래시 없이 그 필드만 생략한다."""
        deck = self._deck(bad_comparison="70% vs 30%")
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.doc.html"
            rep = docgen.render_doc(deck, out)  # 크래시하면 여기서 예외
            self.assertTrue(out.is_file())
        self.assertEqual(2, rep["slides"])
        self.assertEqual([], rep["warnings"])  # 필드 단위 생략은 슬라이드 실패가 아니다(경고 없음)

    def test_render_run_generates_doc_view_alongside_deck_html(self):
        """본선 render 경로(go가 부르는 render_run) 훅 검증 — deck.html과 함께 deck.doc.html이
        같은 run 루트에 생긴다."""
        with tempfile.TemporaryDirectory() as td:
            runs = Path(td) / "workspace" / "runs"
            runs.mkdir(parents=True)
            run = runs / "run"
            run.mkdir()
            storyline = {"meta": {"project": "P"}, "slides": [
                {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": []},
                {"n": 2, "section": "제안개요", "title": "개요", "message": "핵심",
                 "template_id": "executive_summary", "bullets": [],
                 "fields": {"main_claim": "주장", "supporting_points": ["근거"]}},
            ]}
            (run / "storyline.json").write_text(json.dumps(storyline, ensure_ascii=False), encoding="utf-8")
            args = SimpleNamespace(
                run_dir=str(run), stage6=None, stage7=None, stage8=None, storyline=None,
                project=None, pack="core", pattern_sets=None, pptx=False, pptx_mode="native",
                skins=None, analysis=None, rfp=None, anonymize_config=None, json=False,
            )
            with mock.patch.object(proposal_pipeline, "RUNS", runs):
                rc = proposal_pipeline.render_run(args)
            self.assertEqual(0, rc)
            self.assertTrue((run / "deck.html").is_file())
            self.assertTrue((run / "deck.doc.html").is_file())
            doc_text = (run / "deck.doc.html").read_text(encoding="utf-8")
        self.assertIn("검토·회의용 문서 뷰", doc_text)


class AppliedAxesPptxRasterSmoke(unittest.TestCase):
    """W8-B3: 재현 — ship --pptx-mode image가 만든 이미지 PPTX는 pptx_raster에만 기록되고
    gating_report.applied_axes.pptx는 계속 null이었다("경로별 실제 적용 축" 정의 위반)."""

    def _deck(self):
        return {"meta": {"project": "P"}, "slides": [
            {"slide_id": 1, "role": "data", "template_id": "data_interpretation",
             "title": "지표", "message": "핵심", "bullets": [],
             "fields": {"metric": "m", "comparison": [{"label": "A", "value": 1}],
                        "interpretation": ["x"]}},
        ]}

    def test_image_mode_pptx_records_applied_axes_pptx_not_null(self):
        import rasterize
        if not rasterize.available():
            self.skipTest(rasterize.unavailable_reason())
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "deck.json").write_text(json.dumps(self._deck(), ensure_ascii=False), encoding="utf-8")
            overrides = {"version": 1, "slides": {"1": {
                "image_slots": [{"id": "s1", "role": "mood"}]}}}
            (run / "design_overrides.json").write_text(
                json.dumps(overrides, ensure_ascii=False), encoding="utf-8")
            args = SimpleNamespace(overrides=None, pack="core")
            rep, ov_path, _skins = proposal_pipeline.apply_stage9(run, args)
            gating = {"applied_axes": {
                "html": {"pack": "core", "skins": [], "overrides": True, "image_slots": 1,
                          "manual_layer": False},
                "pptx": None,
            }}
            (run / "gating_report.json").write_text(
                json.dumps(gating, ensure_ascii=False, indent=2), encoding="utf-8")

            view = {"stages": {"render": {"pack": "core", "skins": None}}}
            pptx_rep = proposal_pipeline._ship_pptx_image(run, view)
            report = json.loads((run / "gating_report.json").read_text(encoding="utf-8"))
        self.assertIsNotNone(report["applied_axes"]["pptx"])
        self.assertEqual("image", report["applied_axes"]["pptx"]["mode"])
        self.assertGreaterEqual(report["applied_axes"]["pptx"]["image_slots"], 1)
        self.assertTrue(report["applied_axes"]["pptx"]["overrides"])
        self.assertGreater(pptx_rep["slides"], 0)


class ShipNextExcludesExistingDerivativesSmoke(unittest.TestCase):
    """W8-B4: 재현 — 승인 완료 후 "다음"이 이미 만든 파생물(pptx)도 계속 제안했다.
    실제 파일 존재로 판정한다(approve 단계 자기보고는 ship 호출마다 덮여 신뢰할 수 없다).
    (W31 처분: cinematic 파생 폐기 — pptx만 남는다.)"""

    def _approved_run(self, td):
        run = Path(td)
        pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
        pipeline_state.clear_checkpoint(run, "decision")
        pipeline_state.record(run, "render")
        pipeline_state.record(run, "review_resolve")
        pipeline_state.record(run, "design_brief")
        (run / "design_overrides.json").write_text('{"version":1,"slides":{}}', encoding="utf-8")
        pipeline_state.record(run, "stage9_apply")
        pipeline_state.clear_checkpoint(run, "design")
        pipeline_state.record(run, "approve")
        return run

    def test_no_derivatives_suggests_pptx(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._approved_run(td)
            view = pipeline_state.resolve(run)
        self.assertEqual("done", view["next"]["kind"])
        self.assertIn("--pptx", view["next"]["command"])

    def test_pptx_already_made_reports_fully_done(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._approved_run(td)
            (run / "deck.pptx").write_text("x", encoding="utf-8")
            view = pipeline_state.resolve(run)
        self.assertNotIn("--pptx", view["next"]["command"])


class WireframeGateSmoke(unittest.TestCase):
    """③ go 자동 편입(2026-07-14): [3]와이어프레임은 decision 체크포인트(✋② 내용 동결)로
    게이트된다 — 동결 전 탐색 루프에선 next로 안 뜨고(탐색 보호), 청산 후에만 자동 진입.
    회귀 방지: STAGE_ORDER/라우팅을 만져 이 게이트가 조용히 깨지는 것을 잡는다."""

    def _at_wireframe_boundary(self, td, *, decision_cleared):
        run = Path(td)
        pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
        if decision_cleared:
            pipeline_state.clear_checkpoint(run, "decision")
        pipeline_state.record(run, "render")          # 콘텐츠 루프 통과("render" in stages)
        pipeline_state.record(run, "review_resolve")  # 해소 통과 → 다음은 [3] 경계
        return run

    def test_wireframe_gated_before_freeze(self):
        # 동결 전: 결정 체크포인트에서 멈춘다 — wireframe이 next가 아니다
        with tempfile.TemporaryDirectory() as td:
            view = pipeline_state.resolve(self._at_wireframe_boundary(td, decision_cleared=False))
        nxt = view["next"]
        self.assertEqual("checkpoint", nxt["kind"])
        self.assertEqual("decision", nxt.get("checkpoint"))
        self.assertNotIn(nxt.get("stage"), ("wireframe_bundle", "wireframe_apply"))

    def test_wireframe_is_next_after_freeze(self):
        # 청산 후: [3] 뼈대(wireframe_bundle)가 자동 next로 뜬다
        with tempfile.TemporaryDirectory() as td:
            view = pipeline_state.resolve(self._at_wireframe_boundary(td, decision_cleared=True))
        self.assertEqual("wireframe_bundle", view["next"].get("stage"))


class ImageRouteSmoke(unittest.TestCase):
    """W28(D8~D13): 렌더 2분기 — image_infographic 라우트가 stage9/refine/deck_review 대신
    imagedeck(bundle→collect→ack→compose)를 관통한다. 회귀 방지: 라우트 분기를 만져도
    (①기본 html_editable 무영향 ②image/off는 wireframe 스킵 ③사람 ack 관문 ④재무장 ⑤불합격 정지)가
    조용히 깨지지 않게 잡는다."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _png(self, path, w, h):
        import struct
        import zlib
        def ch(t, d):
            c = t + d
            return struct.pack(">I", len(d)) + c + struct.pack(">I", zlib.crc32(c) & 0xffffffff)
        path.parent.mkdir(parents=True, exist_ok=True)
        raw = b"".join(b"\x00" + b"\xff" * (w * 3) for _ in range(h))
        path.write_bytes(b"\x89PNG\r\n\x1a\n" + ch(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
                         + ch(b"IDAT", zlib.compress(raw)) + ch(b"IEND", b""))

    def _at_design_boundary(self, td, *, route=None, wireframe_mode="off"):
        """decision 청산·render/review_resolve/design_brief 기록, design 미청산 = 디자인 루프 입구."""
        run = Path(td)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")
        (run / "design_brief.json").write_text(json.dumps({"brand": {"placement": {}}}, ensure_ascii=False),
                                               encoding="utf-8")
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        pipeline_state.record(run, "render")
        pipeline_state.record(run, "review_resolve")
        pipeline_state.record(run, "design_brief")
        pipeline_state.clear_checkpoint(run, "decision")  # 파일 뒤 마이크로초 clearance(재무장 회피)
        if route:
            pipeline_state.set_render_route(run, route, wireframe_mode)
        else:
            # W29 승격 후 init은 image를 명시 기록한다 - route=None 케이스는 "레거시 run
            # (render_route.json 없음)" 하위호환 시나리오이므로 파일을 지워 시뮬레이션한다.
            (run / "render_route.json").unlink(missing_ok=True)
        return run

    def _produce_images(self, run):
        manifest = json.loads((run / "imagedeck_manifest.json").read_text(encoding="utf-8"))
        gen = manifest["gen_canvas"]
        for s in manifest["slides"]:
            self._png(run / "imagedeck" / "slides" / s["out_name"], gen["w"], gen["h"])

    def _bundle(self, run):
        skin = imagedeck.resolve_skin("quartz_infographic", self.SKINS_DIR)
        imagedeck.bundle(run, skin, wireframe_mode="off")
        pipeline_state.record(run, "imagedeck_bundle")

    def _ack(self, run):
        # 사람 관문 통과 시뮬레이션 — go_cmd가 대시보드 ack 파일을 읽고 clear_checkpoint 하는 것과 동치
        # (ack파일→clearance 경로는 CLI 통합에서 별도 검증). design 관문 테스트와 같은 관례.
        pipeline_state.clear_checkpoint(run, "imagedeck_ack")

    def _ack_prompt(self, run):
        # W30: 생산 전 프롬프트·레퍼런스 확인 관문 통과 시뮬레이션.
        pipeline_state.clear_checkpoint(run, "imagedeck_prompt_ack")

    def test_default_route_untouched_html_editable(self):
        # 레거시 run(render_route.json 없음) → 기존 흐름 불변: 디자인 루프 입구는 stage9, imagedeck 아님.
        with tempfile.TemporaryDirectory() as td:
            view = pipeline_state.resolve(self._at_design_boundary(td, route=None))
        self.assertEqual("html_editable", view["render_route"]["route"])
        self.assertNotIn("imagedeck", str(view["next"].get("stage") or ""))

    def test_init_defaults_new_run_to_image_main_route(self):
        # W29 메인 루트 승격(사용자 결정 2026-07-20): 신규 run은 init이 image_infographic을
        # 명시 기록한다(✋②에서 사람이 html로 바꾸는 자유 유지). 레거시 폴백은 위 테스트가 보증.
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            self.assertTrue((run / "render_route.json").is_file())
            route, wf_mode = pipeline_state.render_route(run)
        self.assertEqual("image_infographic", route)
        self.assertEqual("auto", wf_mode)

    def test_image_off_skips_wireframe_and_bundles(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            view = pipeline_state.resolve(run)
            self.assertEqual("image_infographic", view["render_route"]["route"])
            # 뼈대(wireframe_bundle/apply)는 건너뛴다 — 대신 W31 R2·R3(B1 테마 확정: design_contract
            # 동결 → theme_confirm)를 먼저 지나 imagedeck_bundle로 간다(뼈대 스킵과 무관한 새 게이트).
            self.assertEqual("design_contract", view["next"].get("stage"))
            design_contract.save(run, design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR))
            pipeline_state.record(run, "design_contract")
            pipeline_state.clear_checkpoint(run, "theme_confirm")
            self.assertEqual("imagedeck_bundle", pipeline_state.resolve(run)["next"].get("stage"))

    def test_prompt_ack_gates_production_and_rearms_on_rebundle(self):
        # W30: 번들 직후 생산 전 사람 관문(프롬프트·레퍼런스 확인). 재번들 = 재무장.
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            self._bundle(run)
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt["kind"])
            self.assertEqual("imagedeck_prompt_ack", nxt["checkpoint"])
            self.assertTrue(any("imagedeck_prompts" in p for p in nxt.get("review", [])))
            self._ack_prompt(run)
            self.assertEqual("llm", pipeline_state.resolve(run)["next"]["kind"])  # 이제 생산 정지
            self._bundle(run)  # 재번들 → manifest 갱신 → 재무장
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("imagedeck_prompt_ack", nxt.get("checkpoint"))
            # 건너뛸 수 없는 사람 전속 관문
            self.assertIn("imagedeck_prompt_ack", pipeline_state.HUMAN_CHECKPOINTS)
            self.assertNotIn("imagedeck_prompt_ack", server.SKIPPABLE_ACK_GATES)

    def test_full_walk_bundle_collect_ack_compose_to_design_gate(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            self._bundle(run)
            self._ack_prompt(run)  # W30: 생산 전 프롬프트 확인 관문 통과
            # 이미지 생산 전: Codex 핸드오프(llm)에서 정지
            self.assertEqual("llm", pipeline_state.resolve(run)["next"]["kind"])
            self._produce_images(run)
            rep = imagedeck.collect(run)
            self.assertTrue(rep["pass"])
            pipeline_state.record(run, "imagedeck_collect", passed=rep["pass"])
            # collect PASS 후: 사람 ack 관문
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt["kind"])
            self.assertEqual("imagedeck_ack", nxt["checkpoint"])
            self._ack(run)
            # ack 후: compose 단계가 next
            self.assertEqual("imagedeck_compose", pipeline_state.resolve(run)["next"].get("stage"))
            imagedeck.compose(run)
            pipeline_state.record(run, "imagedeck_compose")
            # compose 후: 디자인 게이트(완성 덱=deck.images.html 검토)
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("design", nxt.get("checkpoint"))
            self.assertTrue(any("deck.images.html" in p for p in nxt.get("review", [])))

    def test_ack_rearms_on_rebundle(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            self._bundle(run)
            self._ack_prompt(run)
            self._produce_images(run)
            rep = imagedeck.collect(run)
            pipeline_state.record(run, "imagedeck_collect", passed=rep["pass"])
            self._ack(run)
            # 재번들 → manifest 갱신 → imagedeck_ack 재무장(감시=manifest/collect)
            self._bundle(run)
            nxt = pipeline_state.resolve(run)["next"]
            # 재수거 또는 재무장된 ack 대기 — 어느 쪽이든 compose로 새지 않아야 한다.
            self.assertNotEqual("imagedeck_compose", nxt.get("stage"))

    def test_bad_px_stops_not_advances(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            self._bundle(run)
            self._ack_prompt(run)
            manifest = json.loads((run / "imagedeck_manifest.json").read_text(encoding="utf-8"))
            gen = manifest["gen_canvas"]
            # 첫 장 정합, 둘째 장 잘못된 px
            self._png(run / "imagedeck" / "slides" / manifest["slides"][0]["out_name"], gen["w"], gen["h"])
            self._png(run / "imagedeck" / "slides" / manifest["slides"][1]["out_name"], 640, 480)
            rep = imagedeck.collect(run)
            self.assertFalse(rep["pass"])
            pipeline_state.record(run, "imagedeck_collect", passed=rep["pass"])
            nxt = pipeline_state.resolve(run)["next"]
            # 불합격 → 정지(llm), ack/compose로 전진 안 함
            self.assertEqual("llm", nxt["kind"])
            self.assertNotIn("imagedeck_compose", str(nxt.get("stage") or ""))

    def test_review_scaffold_lists_source_of_truth(self):
        # Q2 Claude 검수 계약: 장별 정본(제목·field)·보존 flag·verdict 자리를 결정론 조립.
        with tempfile.TemporaryDirectory() as td:
            run = self._at_design_boundary(td, route="image_infographic", wireframe_mode="off")
            self._bundle(run)
            rep = imagedeck.review_scaffold(run)
            md = (run / "imagedeck_review.md").read_text(encoding="utf-8")
        self.assertEqual(2, rep["slides"])
        self.assertIn("정본 제목", md)
        self.assertIn("project_title", md)   # 정본 field 대조 대상
        self.assertIn("verdict", md)          # 채울 자리
        self.assertIn("자동 OCR 금지", md)     # 내용 안전 원칙

    def test_imagedeck_ack_is_mandatory_human_gate(self):
        # imagedeck_ack = 사람 전속·건너뛸 수 없음(필수 채택). 대시보드 skip 거부의 근거.
        self.assertIn("imagedeck_ack", pipeline_state.HUMAN_CHECKPOINTS)
        self.assertNotIn("imagedeck_ack", server.SKIPPABLE_ACK_GATES)


class HybridChromeSmoke(unittest.TestCase):
    """W29 하이브리드 크롬(inkline): 장 클래스(content/full_image/cover/toc/divider)·장별 기대 px·
    HTML 전용 장 제외·크롬 오버라이드·compose 크롬(헤더 제목/배지·푸터 페이지·variants)이 계약대로
    동작하고, slide_classes 미선언 스킨(quartz)은 현행 그대로임을 잡는다."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _png(self, path, w, h):
        import struct
        import zlib
        def ch(t, d):
            c = t + d
            return struct.pack(">I", len(d)) + c + struct.pack(">I", zlib.crc32(c) & 0xffffffff)
        path.parent.mkdir(parents=True, exist_ok=True)
        raw = b"".join(b"\x00" + b"\xff" * (w * 3) for _ in range(h))
        path.write_bytes(b"\x89PNG\r\n\x1a\n" + ch(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
                         + ch(b"IDAT", zlib.compress(raw)) + ch(b"IEND", b""))

    def _run(self, td):
        run = Path(td)
        (run / "storyline.json").write_text(json.dumps({
            "meta": {"project": "T프로젝트", "bid_no": "R26-001", "date": "2026. 07."},
            "slides": [
                {"n": 1, "title": "T프로젝트 제안서", "template_id": "cover_slide", "fields": {}},
                {"n": 2, "title": "목차", "template_id": "toc", "fields": {}},
                {"n": 3, "title": "운영 방안", "template_id": "divider", "section": "운영", "fields": {}},
                {"n": 4, "title": "검증된 수급", "section": "프로세스", "message": "[예시] 3단계 수급",
                 "template_id": "process_steps", "fields": {}, "example": True},
                {"n": 5, "title": "히어로", "deck_class": "full_image",
                 "template_id": "strategy_pillars", "fields": {}},
                {"n": 6, "title": "푸터확장", "template_id": "process_steps",
                 "chrome_override": {"footer_h": 160}, "fields": {}},
            ]}, ensure_ascii=False), encoding="utf-8")
        return run

    def _produce(self, run, manifest):
        for s in manifest["slides"]:
            if s.get("render") == "html":
                continue
            px = s["expected_px"]
            self._png(run / "imagedeck" / "slides" / s["out_name"], px["w"], px["h"])

    def test_inkline_bundle_slide_classes_and_per_slide_px(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            by_n = {s["n"]: s for s in manifest["slides"]}
            self.assertTrue(manifest["classes_enabled"])
            # cover/toc/divider = HTML 전용 장(프롬프트 없음)
            for n, cls in ((1, "cover"), (2, "toc"), (3, "divider")):
                self.assertEqual("html", by_n[n].get("render"))
                self.assertEqual(cls, by_n[n]["deck_class"])
                self.assertNotIn("prompt_file", by_n[n])
            # 장별 기대 px: content=본문만(크롬+본문여백 64/24/24 역산) / full=전체 / 푸터 가변
            self.assertEqual({"w": 1792, "h": 784}, by_n[4]["expected_px"])
            self.assertEqual({"w": 1920, "h": 1080}, by_n[5]["expected_px"])
            self.assertEqual({"w": 1792, "h": 688}, by_n[6]["expected_px"])
            # 프롬프트 계약: content=BODY ONLY + edge-to-edge(여백은 HTML 예약) / full=FULL canvas
            p4 = (run / by_n[4]["prompt_file"]).read_text(encoding="utf-8")
            self.assertIn("BODY ONLY", p4)
            self.assertIn("Do NOT draw", p4)
            self.assertIn("EDGE-TO-EDGE", p4)
            p5 = (run / by_n[5]["prompt_file"]).read_text(encoding="utf-8")
            self.assertIn("FULL canvas", p5)

    def test_quartz_unaffected_all_content(self):
        # 하위호환: slide_classes 미선언 스킨은 전 장 content·전역 gen_canvas(현행 동작 불변).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("quartz_infographic", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self.assertFalse(manifest["classes_enabled"])
            gen = manifest["gen_canvas"]
            for s in manifest["slides"]:
                self.assertNotEqual("html", s.get("render"))
                self.assertEqual(gen, s["expected_px"])  # chrome_override도 무시

    def test_collect_per_slide_px_and_html_exclusion(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.collect(run)
            self.assertTrue(rep["pass"])
            self.assertEqual(3, rep["coverage"]["total"])  # 이미지 장(4·5·6)만 검증 대상
            # 생산 완료 판정도 HTML 장 제외
            self.assertTrue(pipeline_state.imagedeck_images_present(run))

    def test_compose_hybrid_chrome_and_classes(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.compose(run)
            self.assertEqual(3, rep["html_slides"])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertIn('class="slide cover"', html)   # 표지 = 풀 HTML
            self.assertIn("R26-001", html)               # 공고번호 라벨
            self.assertIn('class="toc-grid"', html)      # 목차 그리드
            self.assertIn('class="dv-ghost"', html)      # 간지 대형 숫자 워터마크
            self.assertIn("검증된 수급", html)            # 장 제목 = HTML 크롬 헤더
            self.assertIn('class="badge"', html)         # 섹션 배지
            self.assertIn('class="flag example"', html)  # [예시] flag = HTML 딱지
            self.assertIn("overlay-f", html)             # full_image 오버레이 푸터
            self.assertIn("04 / 06", html)               # 페이지 번호(크롬 푸터)
            self.assertIn("padding:24px 64px 24px", html)  # 본문 여백 = HTML 고정 예약(일관성)

    def test_produce_delegates_missing_only_and_verifies_px(self):
        # W29 승격: produce = 미생산 장만 runner에 위임(재실행 안전), 종료 즉시 px 실측.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")

            def fake_runner(prompt, meta):  # 계약: runner가 meta["out"]에 파일을 쓴다
                px = meta["expected_px"]
                self._png(Path(meta["out"]), px["w"], px["h"])
                return ""

            rep = imagedeck.produce(run, fake_runner)
            self.assertEqual(3, len(rep["generated"]))  # 이미지 장(4·5·6)만, HTML 장 제외
            self.assertEqual([], rep["failed"])
            rep2 = imagedeck.produce(run, fake_runner)   # 재실행 = 전량 skip
            self.assertEqual(3, len(rep2["skipped"]))
            self.assertEqual([], rep2["generated"])
            self.assertTrue(imagedeck.collect(run)["pass"])

    def test_produce_reports_px_mismatch_as_failed(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            imagedeck.bundle(run, skin, wireframe_mode="off")

            def bad_runner(prompt, meta):
                self._png(Path(meta["out"]), 640, 480)  # 계약 위반 px
                return ""

            rep = imagedeck.produce(run, bad_runner, only={4})
            self.assertEqual(["04_검증된-수급.png"], rep["failed"])
            self.assertEqual([], rep["generated"])

    def test_compose_pptx_native_chrome_plus_image_body(self):
        # W30: 하이브리드 pptx - 크롬·표지·목차·간지=네이티브 텍스트(수정 가능), 본문=이미지 1장.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.compose_pptx(run)
            self.assertTrue(rep["editable_chrome"])
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(rep["out"])
            self.assertEqual(6, len(prs.slides))
            all_text = "\n".join(s.text_frame.text for sl in prs.slides
                                 for s in sl.shapes if s.has_text_frame)
            self.assertIn("PROPOSAL", all_text)          # 표지 라벨(네이티브)
            self.assertIn("검증된 수급", all_text)        # content 크롬 제목(네이티브 - 수정 가능)
            self.assertIn("04 / 06", all_text)           # 크롬 푸터 페이지
            pics = [sum(1 for s in sl.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
                    for sl in prs.slides]
            self.assertEqual(3, sum(pics))               # 본문 이미지 3장(4·5·6)만 픽셀

    def test_frame_band_insets_generation_px(self):
        # 바깥 프레임 띠 = inset 밴드(오버레이 금지): 안쪽 전체가 줄고 이미지 px가 함께 역산된다.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
            skin["chrome"]["frame"]["width"] = 14
            sp = run / "skin_frame.json"
            sp.write_text(json.dumps(skin, ensure_ascii=False), encoding="utf-8")
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            by_n = {s["n"]: s for s in manifest["slides"]}
            self.assertEqual({"w": 1764, "h": 756}, by_n[4]["expected_px"])   # 1792-28 x 784-28
            self.assertEqual({"w": 1892, "h": 1052}, by_n[5]["expected_px"])  # full도 띠만큼 축소
            self._produce(run, manifest)
            self.assertTrue(imagedeck.collect(run)["pass"])
            imagedeck.compose(run)
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertIn("padding:14px", html)  # 띠 = slide padding(inset), 콘텐츠 가림 없음

    def test_html_slide_background_image_layer(self):
        # HTML 장(간지 등) 전체 배경 = z-바닥 레이어+스크림. content 장 배경 = 바깥층(.slide)에
        # 깔리고 .inner가 투명해져 띠·본문 여백으로 비친다(일반 이미지 배경의 자연스러운 좌우 노출).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            self._png(run / "assets" / "bg.png", 8, 8)
            story = json.loads((run / "storyline.json").read_text(encoding="utf-8"))
            story["slides"][2]["fields"]["background_image"] = "assets/bg.png"   # 간지(n=3)
            story["slides"][3]["fields"]["background_image"] = "assets/bg.png"   # content(n=4)
            (run / "storyline.json").write_text(json.dumps(story, ensure_ascii=False), encoding="utf-8")
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            imagedeck.compose(run)
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertIn('class="bg"', html)                                  # 간지: 바닥 레이어
            self.assertIn("scrim light", html)
            self.assertIn("background-image:url(", html)                       # content: 바깥층 배경
            self.assertIn('class="inner" style="background:transparent"', html)

    def test_style_variant_overrides_css_vars(self):
        # 같은 덱 안 구간별 스타일 전환: skin.variants + slide.style_variant → CSS 변수 오버라이드.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
            skin["variants"]["green"] = {"colors": {"ink": "1B7A43", "accent": "FFD43B"}}
            sp = run / "skin_green.json"
            sp.write_text(json.dumps(skin, ensure_ascii=False), encoding="utf-8")
            story = json.loads((run / "storyline.json").read_text(encoding="utf-8"))
            story["slides"][3]["style_variant"] = "green"
            (run / "storyline.json").write_text(json.dumps(story, ensure_ascii=False), encoding="utf-8")
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest)
            imagedeck.compose(run)
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertIn("--ink:#1B7A43", html)
            self.assertIn("--accent:#FFD43B", html)

    # -----------------------------------------------------------------------
    # DF2(자산 슬롯 계약) — chrome_contract.decor_slots. 스펙 CONTEXT/DECK_FIRST_DESIGN.md
    # §2-②·§2-④·§3 DF2 행. 배경 레이어(기존 구현, 회귀 금지)와 별개인 배경 "외" 슬롯.
    # -----------------------------------------------------------------------

    def _skin_with_decor(self, run, slots, extra=None):
        skin = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
        skin["decor_slots"] = slots
        if extra:
            for path, value in extra:
                node = skin
                for key in path[:-1]:
                    node = node[key]
                node[path[-1]] = value
        sp = run / "skin_decor.json"
        sp.write_text(json.dumps(skin, ensure_ascii=False), encoding="utf-8")
        return sp

    def test_decor_slots_render_in_html_and_pptx(self):
        # 자산 슬롯 렌더 확인(완료 조건 ①): HTML에 배치 + PPTX에 그림 삽입, 양쪽 warning 없음.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            self._png(run / "assets" / "decor.png", 40, 20)
            sp = self._skin_with_decor(run, [
                {"id": "corner", "image": "assets/decor.png", "anchor": "top-right",
                 "offset_x": 20, "offset_y": 20, "width": 80, "opacity": 0.9},
            ])
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.compose(run)
            self.assertEqual([], rep["warnings"])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            # 6장(표지·목차·간지·본문x2·전체이미지) 모두 opt-out 없이 슬롯 1개씩 - 6회 등장.
            self.assertEqual(6, html.count('class="decor"'))
            self.assertIn("top:20px;right:20px;width:80px;opacity:0.9", html)

            rep_p = imagedeck.compose_pptx(run)
            self.assertEqual([], rep_p["warnings"])
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(rep_p["out"])
            pics = sum(1 for sl in prs.slides for s in sl.shapes
                      if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            self.assertEqual(9, pics)  # 본문 이미지 3(4·5·6) + 장식 6(전 장)

    def test_decor_slots_class_opt_out(self):
        # slide_classes.<cls>.decor=false = opt-out(기존 frame=false와 같은 문법 계열).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            self._png(run / "assets" / "decor.png", 40, 20)
            sp = self._skin_with_decor(
                run,
                [{"id": "corner", "image": "assets/decor.png", "anchor": "bottom-left",
                  "width": 60}],
                extra=[(["slide_classes", "cover", "decor"], False)])
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.compose(run)
            self.assertEqual([], rep["warnings"])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertEqual(5, html.count('class="decor"'))  # 표지만 opt-out - 나머지 5장

    def test_decor_slots_missing_asset_surfaces_warning(self):
        # 완료 조건 ③: 자산 부재는 조용히 skip하지 않고 warnings로 표면화(HTML·PPTX 둘 다).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            sp = self._skin_with_decor(run, [
                {"id": "ghost", "image": "assets/nope.png", "anchor": "top-left", "width": 50},
            ])
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest)
            rep = imagedeck.compose(run)
            self.assertEqual(1, len(rep["warnings"]))
            self.assertIn("ghost", rep["warnings"][0])
            self.assertIn("이미지 없음", rep["warnings"][0])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertNotIn('class="decor"', html)

            rep_p = imagedeck.compose_pptx(run)
            self.assertEqual(1, len(rep_p["warnings"]))
            self.assertIn("ghost", rep_p["warnings"][0])

    def test_decor_slots_undeclared_is_byte_identical(self):
        # 완료 조건 ②: decor_slots 미선언(기존 모든 스킨) - 조립 결과 바이트 동일. 빈 배열 선언과도
        # 동일해야 한다(진짜 신규 코드 경로는 비어있지 않은 decor_slots에서만 열린다).
        with tempfile.TemporaryDirectory() as td:
            # 같은 run(같은 절대경로) 안에서 비교 - 서로 다른 temp 디렉터리를 쓰면 이미지 src의
            # file:// URI 자체가 달라져 "바이트 동일" 비교가 무의미해진다.
            run = self._run(td)
            skin_plain = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin_plain, wireframe_mode="off")
            self._produce(run, manifest)
            rep_plain = imagedeck.compose(run)
            html_plain = (run / "deck.images.html").read_text(encoding="utf-8")

            sp = self._skin_with_decor(run, [])  # 명시적 빈 배열
            manifest2 = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest2)
            rep_empty = imagedeck.compose(run)
            html_empty = (run / "deck.images.html").read_text(encoding="utf-8")

        self.assertEqual(html_plain, html_empty)
        self.assertEqual([], rep_plain["warnings"])
        self.assertEqual([], rep_empty["warnings"])
        self.assertNotIn('class="decor"', html_plain)


class DeckOverridesSmoke(HybridChromeSmoke):
    """DF6(CONTEXT/DECK_FIRST_DESIGN.md §2-⑦ 경로 B): run/deck_overrides.json 마무리 오버라이드
    채널. 동결된 storyline을 재편집하지 않고 특정 장만 바꾸는 "오버라이드 수정 -> 재조립" 루프.
    안전 분류 2급 - ⓐ(style_variant/background_image, 재조립만) / ⓑ(chrome_override/deck_class,
    px 변형 가능 - 재번들 필요·collect가 불일치 검출)를 각각 실증한다. HybridChromeSmoke의
    _run/_png 픽스처(6장 storyline, inkline 스킨)를 그대로 재사용한다."""

    def _skin_with_variant(self, run, name="green", colors=None):
        skin = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
        skin["variants"][name] = {"colors": colors or {"ink": "1B7A43", "accent": "FFD43B"}}
        sp = run / "skin_green.json"
        sp.write_text(json.dumps(skin, ensure_ascii=False), encoding="utf-8")
        return sp

    def test_style_variant_override_recompose_only(self):
        # 완료 조건 ①: style_variant 오버라이드가 재조립(compose)만으로 색을 반영한다 - bundle을
        # 다시 돌릴 필요가 없다(ⓐ 재조립만으로 끝나는 변형).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            sp = self._skin_with_variant(run)
            manifest = imagedeck.bundle(run, sp, wireframe_mode="off")
            self._produce(run, manifest)
            rep0 = imagedeck.compose(run)
            self.assertEqual([], rep0.get("deck_overrides_applied"))
            html0 = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertNotIn("--ink:#1B7A43", html0)

            (run / "deck_overrides.json").write_text(
                json.dumps({"slides": {"4": {"style_variant": "green"}}}, ensure_ascii=False),
                encoding="utf-8")
            rep = imagedeck.compose(run)  # bundle 재실행 없음 - 재조립만.
            self.assertEqual([4], rep.get("deck_overrides_applied"))
            self.assertEqual("a", rep["deck_overrides"][0]["category"])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertIn("--ink:#1B7A43", html)
            self.assertIn("--accent:#FFD43B", html)

            # PPTX도 같은 헬퍼를 공유 - 색이 반영된다(재번들 없이).
            rep_p = imagedeck.compose_pptx(run)
            self.assertEqual([4], rep_p.get("deck_overrides_applied"))

    def test_chrome_override_needs_rebundle_and_collect_catches_stale_px(self):
        # 완료 조건 ②: chrome_override(ⓑ 구조 변형)는 compose 재실행만으로는 px가 바뀌지 않는다
        # (manifest가 bundle 시점 값을 들고 있으므로) - bundle을 다시 돌려야 새 px가 반영되고,
        # 그 시점부터 옛(재생성 전) 이미지는 px 불일치로 collect가 잡아낸다(기존 안전망 우회 없음).
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            self.assertTrue(imagedeck.collect(run)["pass"])
            by_n0 = {s["n"]: s for s in manifest["slides"]}
            old_px = by_n0[4]["expected_px"]

            (run / "deck_overrides.json").write_text(
                json.dumps({"slides": {"4": {"chrome_override": {"header_h": 300}}}},
                          ensure_ascii=False),
                encoding="utf-8")

            manifest2 = imagedeck.bundle(run, skin, wireframe_mode="off")
            self.assertEqual([4], manifest2.get("deck_overrides_applied"))
            self.assertEqual("b", manifest2["deck_overrides"][0]["category"])
            by_n2 = {s["n"]: s for s in manifest2["slides"]}
            new_px = by_n2[4]["expected_px"]
            self.assertNotEqual(old_px, new_px)
            self.assertEqual(300, by_n2[4]["chrome"]["header_h"])

            # 옛(재생성 전) 이미지가 아직 남아있는 상태 - collect가 px 불일치를 검출해야 한다.
            rep = imagedeck.collect(run)
            self.assertFalse(rep["pass"])
            statuses = {r["n"]: r["status"] for r in rep["slides"]}
            self.assertEqual("px_mismatch", statuses.get(4))

            # 재생성(신규 px로) 후에는 다시 PASS - 루프가 수렴한다.
            self._produce(run, manifest2)
            self.assertTrue(imagedeck.collect(run)["pass"])

    def test_unknown_key_rejected_with_clear_message(self):
        # 완료 조건 ③: 허용 외 키(콘텐츠 필드 등)는 명확한 오류 메시지로 거부된다.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            (run / "deck_overrides.json").write_text(
                json.dumps({"slides": {"4": {"title": "몰래 바꾼 제목"}}}, ensure_ascii=False),
                encoding="utf-8")
            with self.assertRaises(imagedeck.ImagedeckError) as cm:
                imagedeck.bundle(run, skin, wireframe_mode="off")
            msg = str(cm.exception)
            self.assertIn("title", msg)
            self.assertIn("deck_overrides.json", msg)

    def test_missing_or_empty_overrides_byte_identical(self):
        # 완료 조건 ④: 미존재/빈 파일 - 조립 결과(HTML/manifest 오버라이드 필드) 바이트·값 동일.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest)
            imagedeck.compose(run)
            html_missing = (run / "deck.images.html").read_text(encoding="utf-8")
            self.assertEqual([], manifest.get("deck_overrides_applied"))

            (run / "deck_overrides.json").write_text(
                json.dumps({"slides": {}}, ensure_ascii=False), encoding="utf-8")
            manifest2 = imagedeck.bundle(run, skin, wireframe_mode="off")
            self._produce(run, manifest2)
            imagedeck.compose(run)
            html_empty = (run / "deck.images.html").read_text(encoding="utf-8")

            self.assertEqual(html_missing, html_empty)
            self.assertEqual([], manifest2.get("deck_overrides_applied"))

    def test_unknown_slide_number_rejected(self):
        # 오타 방지: storyline에 없는 장 번호를 가리키면 조용히 무시하지 않고 에러.
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            (run / "deck_overrides.json").write_text(
                json.dumps({"slides": {"99": {"style_variant": "green"}}}, ensure_ascii=False),
                encoding="utf-8")
            with self.assertRaises(imagedeck.ImagedeckError) as cm:
                imagedeck.bundle(run, skin, wireframe_mode="off")
            self.assertIn("99", str(cm.exception))


class ExampleDataPolicySmoke(unittest.TestCase):
    """W9: 예시 데이터 정책(창작금지의 정교화) + 청중 계약.

    재현 = 구조 슬라이드를 '예시'로 명시하면 (①렌더 라벨 ②자동 태그 ③ship 경고)의 3중
    안전장치가 걸리고, W5 해소지 fact_supplied가 예시 마크까지 제거해야 한다. 예시 없는 덱은
    바이트 불변. 청중 계약 = '가격 배점 20%' 내부 전략이 슬라이드로 유출되지 않도록 프롬프트가 막는다.
    """

    LABEL = "예시 데이터"
    REPL = "실데이터로 교체 필요"
    PRICE_LEAK = "가격 배점 20%"

    def _example_storyline(self):
        return {"meta": {"project": "P"}, "slides": [
            {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": []},
            {"n": 2, "section": "데이터", "title": "시장", "message": "규모",
             "template_id": "data_interpretation", "example": True,
             "fields": {"metric": "m", "comparison": [{"label": "A", "value": 70},
                                                       {"label": "B", "value": 30}],
                        "interpretation": ["예시 해석"]}},
        ]}

    def _render_bytes(self, deck):
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            htmlgen.render_html(deck, "core", out)
            return out.read_bytes()

    # --- 계약 & 안전장치 ① 렌더 라벨 / ② 자동 태그 ---
    def test_example_flag_survives_adapt_and_only_when_true(self):
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        self.assertIs(True, deck["slides"][1].get("example"))
        self.assertNotIn("example", deck["slides"][0])   # false는 키 자체를 심지 않는다(바이트 보존)

    def test_bind_auto_tags_example_slide_idempotently(self):
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        templates = _ingest._load_templates("core")
        rep1 = _bind.bind_deck(deck, templates)
        self.assertEqual(1, rep1["example_tagged"])
        self.assertIn(_bind.EXAMPLE_REVIEW_TAG, deck["slides"][1]["review_needed"])
        self.assertIn(self.REPL, _bind.EXAMPLE_REVIEW_TAG)
        rep2 = _bind.bind_deck(deck, templates)
        self.assertEqual(0, rep2["example_tagged"])       # 이미 붙은 태그는 다시 안 붙인다

    def test_render_stamps_visual_label_bytes(self):
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        html = self._render_bytes(deck)
        self.assertIn(self.LABEL.encode("utf-8"), html)   # "예시 데이터" 라벨 바이트 존재
        self.assertIn(b"example-watermark", html)         # 워터마크 CSS도 이 덱에만 실린다

    # --- 안전장치 ③ 교체: W5 fact_supplied가 예시 마크 제거(불변식 유지) ---
    def test_fact_supplied_clears_example_mark_and_label(self):
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        doc = {"schema_version": 1, "items": [
            {"slide_id": 2, "tag": _bind.EXAMPLE_REVIEW_TAG, "decision": "fact_supplied",
             "fact": {"metric": "실측", "comparison": [{"label": "A", "value": 55}]},
             "target": {"kind": "field", "name": "comparison"}}]}
        rep = review_resolve.apply(deck, doc, binders=_bind.BINDERS)
        self.assertEqual(1, rep["tags_removed"])
        self.assertTrue(rep["applied"][0].get("example_cleared"))
        self.assertNotIn("example", deck["slides"][1])           # 마크 제거
        self.assertNotIn(_bind.EXAMPLE_REVIEW_TAG, deck["slides"][1]["review_needed"])
        self.assertNotIn(self.LABEL.encode("utf-8"), self._render_bytes(deck))  # 라벨도 사라진다

    def test_no_basis_confirmed_keeps_example_mark(self):
        """근거 없음 확정은 태그만 지운다 — 예시 마크(라벨/경고)는 유지(불변식 비대칭)."""
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        doc = {"schema_version": 1, "items": [
            {"slide_id": 2, "tag": _bind.EXAMPLE_REVIEW_TAG, "decision": "no_basis_confirmed"}]}
        rep = review_resolve.apply(deck, doc, binders=_bind.BINDERS)
        self.assertEqual(1, rep["tags_removed"])
        self.assertIsNone(rep["applied"][0].get("example_cleared"))
        self.assertIs(True, deck["slides"][1].get("example"))    # 예시는 유지된다

    def test_example_mark_not_cleared_without_recorded_decision(self):
        """불변식: 사람 결정 없으면 예시 마크도 태그도 그대로(창작금지의 대칭)."""
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        rep = review_resolve.apply(deck, {"schema_version": 1, "items": []}, binders=_bind.BINDERS)
        self.assertEqual(0, rep["tags_removed"])
        self.assertIs(True, deck["slides"][1].get("example"))

    # --- 안전장치 ④ 관측 + 바이트 불변 ---
    def test_plain_deck_has_no_example_markup(self):
        """예시 없는 덱: 라벨·배지 CSS 어느 것도 실리지 않는다(바이트 불변 보증)."""
        plain = {"meta": {"project": "P"}, "slides": [
            {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": ["a"]}]}
        deck = _adapt.adapt_storyline(plain, pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        html = self._render_bytes(deck)
        self.assertNotIn(self.LABEL.encode("utf-8"), html)
        self.assertNotIn(b"example-badge", html)

    # --- 청중 계약: 내부 배점표 유출 차단 ---
    def test_prompt_carries_audience_contract_against_score_leak(self):
        block = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn(self.PRICE_LEAK, block)         # 유출 사례를 반례로 명시
        self.assertIn("배점", block)
        self.assertIn("수신 기관", block)

    def test_prompt_carries_example_data_policy(self):
        block = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn('"example": true', block)
        self.assertIn(self.REPL, block)
        self.assertIn("예시임을 명시", block)          # 창작금지 개정 문구


class SkeletonReverseProposalSmoke(unittest.TestCase):
    """W10: 탐색 루프 시작 = 백지가 아니라 역제안(표준 시나리오 더미 즉시 렌더·0토큰)."""

    def test_generator_marks_every_slide_as_example_in_sequence(self):
        sc = skeleton.load_scenario(None)                 # 기본 시나리오(public_proposal)
        doc = skeleton.build_skeleton(sc, project="[예시] 테스트")
        self.assertTrue(doc["meta"]["skeleton"])
        self.assertGreaterEqual(len(doc["slides"]), 1)
        for i, s in enumerate(doc["slides"], 1):
            self.assertEqual(i, s["n"])                    # 연속 번호
            self.assertTrue(s["example"])                  # 전 장표 예시(W9 3중 안전장치 진입)
            self.assertEqual([skeleton.EXAMPLE_FLAG], s["flag"])

    def test_unknown_scenario_is_a_surfaced_error_not_silent(self):
        with self.assertRaises(skeleton.SkeletonError):
            skeleton.load_scenario("does_not_exist")

    def test_structure_block_injects_confirmed_shape_for_fill_handoff(self):
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc)
        block = skeleton.structure_block(doc)
        self.assertIn("확정된 스켈레톤 구조", block)          # 채움 핸드오프에 구조 주입
        self.assertIn("슬라이드 수·순서", block)
        for s in doc["slides"]:                            # 사용자가 뺀/바꾼 장표가 그대로 반영
            self.assertIn(f"슬라이드 {s['n']}", block)

    def test_go_skeleton_pauses_once_and_writes_editable_file(self):
        """역제안 정지: _go_skeleton은 True(1회 멈춤)를 돌려주고 skeleton.json을 쓴다(편집 UI)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="demo.md")
            args = SimpleNamespace(scenario=None, pack="house_a", skins=None)
            with mock.patch.object(proposal_pipeline, "render_run") as render:
                paused = proposal_pipeline._go_skeleton(run, args)
            self.assertTrue(paused)                        # 사용자가 모양을 볼 틈에 정지
            self.assertTrue(render.called)                 # 더미를 즉시 렌더(역제안)
            self.assertTrue(skeleton.skeleton_path(run).is_file())  # 파일이 편집 UI


class ImageFormatFreedomSmoke(unittest.TestCase):
    """결정 5: 포맷은 내용이 정한다 — jpg/jpeg를 1급으로(예전엔 svg로 강등). 이모지 금지·팔레트 전달."""

    def test_jpg_and_jpeg_are_first_class_raster_not_downgraded_to_svg(self):
        for f in ("jpg", "jpeg", "JPG", ".jpg", "png"):
            self.assertTrue(image_slots._is_raster(image_slots._slot_format({"format": f})))
        self.assertEqual("jpg", image_slots._slot_format({"format": "JPG"}))
        self.assertEqual("svg", image_slots._slot_format({"format": "webp"}))   # 미지원 → svg 폴백
        self.assertEqual("svg", image_slots._slot_format({}))

    def test_raster_runner_contract_writes_jpg_file_and_ignores_return(self):
        """래스터 계약(png/jpg/jpeg): 러너가 target_path에 파일을 쓴다. 반환값 아닌 파일로 검증."""
        run = tempfile.mkdtemp()
        try:
            ov = {"version": 1, "slides": {"1": {"image_slots": [
                {"id": "bg", "role": "mood", "layer": "background", "format": "jpg",
                 "prompt": "안동 한옥 배경", "treatment": "오버레이 55%"}]}}}

            def raster_runner(prompt, meta):
                # 실 Codex 러너처럼 target_path에 파일을 쓰고 자기보고는 무의미한 텍스트만 반환.
                Path(meta["target_path"]).write_bytes(b"\xff\xd8\xff\xe0fake-jpeg-bytes")
                self.assertEqual("jpg", meta["format"])
                self.assertIn("JPG", prompt)              # 포맷이 프롬프트에 명시
                self.assertIn("이모지", prompt)           # 이모지 금지 계약
                return "다른 출력"
            rep = image_slots.fill_images(ov, run, tier=2, runner=raster_runner,
                                          palette={"navy": "#1E4A8C"})
            self.assertEqual(["slide1:bg"], rep["generated"])
            asset = image_slots.slot_asset_path(run, "1", "bg", fmt="jpg")
            self.assertTrue(asset.exists() and asset.stat().st_size > 0)
            self.assertEqual(".jpg", asset.suffix)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_raster_degrades_when_runner_writes_nothing(self):
        run = tempfile.mkdtemp()
        try:
            ov = {"version": 1, "slides": {"1": {"image_slots": [
                {"id": "bg", "role": "mood", "format": "jpg", "prompt": "산맥 실루엣"}]}}}
            rep = image_slots.fill_images(ov, run, tier=2, runner=lambda p, m: "경로만",
                                          allow_generate=True)
            self.assertEqual(["slide1:bg"], rep["degraded"])   # 파일 없음 → degrade(자기보고 불신)
        finally:
            shutil.rmtree(run, ignore_errors=True)

    def test_palette_from_skin_tokens_feeds_prompt(self):
        # W31 E3: house_a(navy 하우스 색)는 Reuse 격리 — core(무채 팔레트, gray_text)로 동일 배선 검증.
        pal = proposal_pipeline._skin_palette("core")
        self.assertIn("gray_text", pal)
        prompt = image_slots.build_fill_prompt(
            {"id": "i", "role": "conceptual", "format": "svg"}, palette=pal)
        self.assertIn(pal["gray_text"], prompt)             # 스킨 토큰이 프롬프트 팔레트에
        self.assertIn("이모지", prompt)                    # 이모지 금지


class LayeredBackgroundSmoke(unittest.TestCase):
    """결정 6: layer=background 선언 계약 + 가독성 실측 대체(선언만으로 pass 금지)."""

    def test_background_slot_renders_layer_marker_and_scrim(self):
        h = image_slots.render_slot_html(
            {"id": "bg", "role": "mood", "layer": "background", "treatment": "오버레이"},
            asset_src="stage9_design/slots/slide1_bg.jpg")
        self.assertIn('data-layer="background"', h)        # probe가 읽는 계약 표면
        self.assertIn("dov-slot--bg", h)
        self.assertIn("dov-slot__scrim", h)                # treatment의 실체(DOM 오버레이)

    def test_foreground_default_has_no_background_markers(self):
        h = image_slots.render_slot_html({"id": "c", "role": "conceptual"},
                                         asset_src="x.png")
        self.assertNotIn("data-layer", h)
        self.assertNotIn("dov-slot__scrim", h)

    def test_validator_requires_treatment_for_declared_background(self):
        import overrides as ov_mod
        deck = {"slides": [{"slide_id": 1, "title": "표지"}]}
        bad = {"version": 1, "slides": {"1": {"image_slots": [
            {"id": "x", "role": "mood", "layer": "background"}]}}}     # treatment 없음
        errs = ov_mod.validate_overrides(bad, deck)
        self.assertTrue(any("treatment 필수" in e for e in errs))
        ok = {"version": 1, "slides": {"1": {"image_slots": [
            {"id": "x", "role": "mood", "layer": "background", "treatment": "오버레이 55%"}]}}}
        self.assertEqual([], ov_mod.validate_overrides(ok, deck))

    def test_validator_rejects_bad_format_and_layer(self):
        import overrides as ov_mod
        deck = {"slides": [{"slide_id": 1, "title": "표지"}]}
        errs = ov_mod.validate_overrides({"version": 1, "slides": {"1": {"image_slots": [
            {"id": "x", "role": "mood", "format": "webp", "layer": "middle"}]}}}, deck)
        self.assertTrue(any("format" in e for e in errs))
        self.assertTrue(any("layer" in e for e in errs))


class LayoutProbeW11Smoke(unittest.TestCase):
    """결정 6: probe 확장 — 디렉터 장식의 텍스트 가림 + 선언 배경의 가독성 실측."""

    _BASE = ("section.slide{position:relative;width:1280px;height:720px;overflow:hidden;"
             "background:#fff}.msg{position:absolute;top:340px;left:200px;width:500px;"
             "font-size:40px;color:#fff}")
    _IMG = ("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='10' "
            "height='10'%3E%3Crect width='10' height='10' fill='%23333'/%3E%3C/svg%3E")

    def _probe(self, body_css, section):
        import layout_probe
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "deck.html"
            p.write_text(f"<!doctype html><html><head><style>{self._BASE}{body_css}</style>"
                         f"</head><body>{section}</body></html>", encoding="utf-8")
            return layout_probe.probe_html(p)

    def test_director_decoration_over_text_is_detected(self):
        import layout_probe
        if not layout_probe.available():
            self.skipTest(layout_probe.unavailable_reason() + " — 브라우저 실측 불가")
        css = ".dov-deco{position:absolute;top:330px;left:180px;width:520px;height:90px;background:#F37321;z-index:5}"
        sec = ('<section class="slide" id="slide-5"><div class="msg" style="color:#111">'
               '전환 이제 비전으로</div><div class="dov-deco"></div></section>')
        rep = self._probe(css, sec)
        self.assertEqual("warn", rep["status"])
        self.assertIn("decoration_occlusion", rep["slides"][0]["flags"])

    def test_declared_background_without_treatment_is_flagged(self):
        import layout_probe
        if not layout_probe.available():
            self.skipTest(layout_probe.unavailable_reason() + " — 브라우저 실측 불가")
        css = (".dov-slot--bg{position:absolute;inset:0;z-index:0}"
               ".dov-slot--bg img{width:100%;height:100%;object-fit:cover}"
               ".dov-slot__scrim{position:absolute;inset:0;background:none}")   # 스크림 죽임
        sec = (f'<section class="slide" id="slide-1">'
               f'<div class="dov-slot dov-slot--bg" data-layer="background" data-treatment="선언만">'
               f'<img src="{self._IMG}"/><div class="dov-slot__scrim" aria-hidden="true"></div></div>'
               f'<div class="msg">배경 위 한 문장</div></section>')
        rep = self._probe(css, sec)
        self.assertIn("background_no_treatment", rep["slides"][0]["flags"])
        self.assertEqual([], rep["slides"][0]["content_overlaps"])   # 배경은 겹침 결함에서 제외

    def test_declared_background_with_real_scrim_passes(self):
        import layout_probe
        if not layout_probe.available():
            self.skipTest(layout_probe.unavailable_reason() + " — 브라우저 실측 불가")
        css = (".dov-slot--bg{position:absolute;inset:0;z-index:0}"
               ".dov-slot--bg img{width:100%;height:100%;object-fit:cover}"
               ".dov-slot__scrim{position:absolute;inset:0;"
               "background:linear-gradient(180deg,rgba(5,5,5,.2),rgba(5,5,5,.6))}")
        sec = (f'<section class="slide" id="slide-1">'
               f'<div class="dov-slot dov-slot--bg" data-layer="background" data-treatment="오버레이 55%">'
               f'<img src="{self._IMG}"/><div class="dov-slot__scrim" aria-hidden="true"></div></div>'
               f'<div class="msg">배경 위 한 문장</div></section>')
        rep = self._probe(css, sec)
        self.assertNotIn("background_no_treatment", rep["slides"][0]["flags"])


class HeroRhythmSmoke(unittest.TestCase):
    """결정 1: hero 리듬 = 스토리 피크 기반(상한 2~3). '밋밋→hero' 폐기(리듬 사망 방지)."""

    def _deck(self):
        titles = ["표지", "목차", "사업 배경", "현황 문제", "추진 전략", "차별화 방안",
                  "핵심 약속", "추진 일정", "기대효과", "조직", "마무리"]
        return {"slides": [{"slide_id": i, "role": t, "title": t, "body": ["a"],
                            "key_message": "짧", "fields": {}} for i, t in enumerate(titles, 1)]}

    def test_thin_slides_do_not_all_become_hero(self):
        with tempfile.TemporaryDirectory() as td:
            brief = design_brief.build_default(Path(td), self._deck())
        heroes = [r for r in brief["page_rhythm"]["slides"] if r["density"] == "hero"]
        self.assertLessEqual(len(heroes), design_brief._HERO_CAP)   # 상한 준수(리듬 생존)
        self.assertGreaterEqual(len(heroes), 1)
        ids = {r["slide_id"] for r in heroes}
        self.assertIn(1, ids)                                       # 표지 = 스토리 피크
        self.assertIn(7, ids)                                       # 핵심 약속 = 스토리 피크

    def test_hero_background_is_candidate_not_auto_live_slot(self):
        """결정 2026-07-15: hero 배경은 **자동 라이브 슬롯이 아니라 제안(candidate)** 으로만 표면화.
        빈 프롬프트 자동 슬롯이 제네릭을 뽑거나 디렉터 선언 배경을 조용히 덮던 문제를 원천 차단."""
        with tempfile.TemporaryDirectory() as td:
            brief = design_brief.build_default(Path(td), self._deck())
        plan = brief["image_slots_plan"]
        self.assertEqual([], plan["slots"])                        # 자동 라이브 슬롯 0(주입 없음)
        cands = plan["background_candidates"]
        self.assertTrue(cands)                                     # 대신 제안으로 표면화
        for c in cands:                                            # 권장 연출 = 배경 이미지(결정 1·5·6)
            self.assertEqual("background", c["suggested_layer"])
            self.assertEqual("jpg", c["suggested_format"])
            self.assertTrue(c["suggested_treatment"])
            self.assertNotIn("id", c)                              # 라이브 슬롯 id 없음(sync 대상 아님)

    def test_no_peak_vocab_means_no_forced_hero(self):
        """피크 어휘가 없으면 hero 0 — 얇다고 hero로 승격하지 않는다(폐기된 규칙 회귀 감지)."""
        deck = {"slides": [{"slide_id": 1, "role": "일정", "title": "추진 일정",
                            "body": ["a"], "key_message": "짧", "fields": {}}]}
        with tempfile.TemporaryDirectory() as td:
            brief = design_brief.build_default(Path(td), deck)
        self.assertEqual([], [r for r in brief["page_rhythm"]["slides"]
                              if r["density"] == "hero"])


class W12GatePromotionSmoke(unittest.TestCase):
    """W12: 실측 게이트 승격 — 실결함 계열 flag가 조용한 warn을 지나가지 못하게 1급 표시.

    마찰 로그 핵심 교훈: 실측 게이트가 실결함을 잡았으나(warn) 비차단이라 조용히 통과 —
    육안 검증 없었으면 놓칠 뻔. 여기서 그 '조용한 통과'가 소멸했음을 회귀로 고정한다.
    """

    def _row(self, sid, flags):
        return {"slide_id": sid, "flags": flags, "overflow_px": 0,
                "content_overlaps": [], "occlusions": []}

    def _browser(self, *rows):
        return {"status": "warn", "viewport": {"width": 1280, "height": 720},
                "summary": {"slides": len(rows)}, "slides": list(rows)}

    def test_repair_targets_keeps_only_real_defect_flags(self):
        import layout_probe
        browser = self._browser(
            self._row("1", ["background_no_treatment"]),
            self._row("5", ["decoration_occlusion", "slot_occlusion"]),
            self._row("7", ["overflow_measured"]),
            # slot_overlaps_content는 더 보수적 신호 — 수리 대상에서 제외한다.
            self._row("9", ["slot_overlaps_content"]),
            self._row("3", []),
        )
        targets = layout_probe.repair_targets(browser)
        ids = [t["slide_id"] for t in targets]
        self.assertEqual(["1", "5", "7"], ids)                 # 9(overlaps_content)·3(무결) 제외
        self.assertIn("slot_occlusion", targets[1]["flags"])   # 5는 decoration+slot 둘 다 유효

    def test_unmeasured_browser_yields_no_repair_targets(self):
        import layout_probe
        self.assertEqual([], layout_probe.repair_targets(None))
        self.assertEqual([], layout_probe.repair_targets(
            {"status": "unmeasured", "reason": "playwright 미설치", "slides": []}))

    def _gate_with_defects(self, run):
        import design_checks
        browser = self._browser(self._row("1", ["background_no_treatment"]),
                                 self._row("5", ["decoration_occlusion"]))
        # 정적 계층은 실제 계산본을 써서 _print_design_checks가 읽는 요약 키를 모두 채운다.
        checks = design_checks.compute_design_checks(
            "<html><body><section id='slide-1' class='slide'><p>x</p></section></body></html>")
        checks = design_checks.attach_browser_layer(checks, browser)
        (run / "gating_report.json").write_text(
            json.dumps({"design_checks": checks}, ensure_ascii=False, indent=2), encoding="utf-8")

    def test_design_defect_warning_is_surfaced_in_status(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._gate_with_defects(run)
            warns = pipeline_state._design_defect_warnings(run)
        self.assertTrue(warns)
        self.assertIn("수리 대상 2건", warns[0])
        self.assertIn("slide 1", warns[0])
        self.assertIn("slide 5", warns[0])

    def test_no_gating_report_means_no_defect_warning(self):
        with tempfile.TemporaryDirectory() as td:
            self.assertEqual([], pipeline_state._design_defect_warnings(Path(td)))

    def test_resolve_carries_the_defect_warning(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._gate_with_defects(run)
            view = pipeline_state.resolve(run)
        self.assertTrue(any("수리 대상 2건" in w for w in view["warnings"]))

    def test_stage9_apply_prints_repair_targets_first_class(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._gate_with_defects(run)
            buf = io.StringIO()
            with mock.patch("sys.stdout", buf):
                proposal_pipeline._print_design_checks(run)
            out = buf.getvalue()
        self.assertIn("수리 대상 2건", out)                     # 1급 표시(조용한 warn 아님)

    def test_deck_review_bundle_foregrounds_repair_targets(self):
        browser = self._browser(self._row("1", ["background_no_treatment"]))
        block = "\n".join(deck_review._render_browser_layer(browser))
        self.assertIn("수리 대상 1건", block)


class W13ManualLayerDiffSmoke(unittest.TestCase):
    """W13(결정 7): Claude Design 자유편집 가드를 차단에서 표면화로 — freeze는 항상 진행,
    변경 명세만 diff로 남긴다. '차단이 자유를 되조인다'는 기각(선택지 b)의 반대편 회귀."""

    def _deck(self):
        return {"slides": [{
            "slide_id": "1",
            "title": "[예시] 표지 이미지 주제",
            "key_message": "핵심 메시지 문장입니다 충분히 길게",
            "review_needed": ["[예시] 표지 이미지 주제"],
        }]}

    def test_fact_deletion_is_recorded_but_does_not_block_freeze(self):
        import overrides as ov_mod
        deck = self._deck()
        baseline = "<html><body><h1>[예시] 표지 이미지 주제</h1><p>핵심 메시지 문장입니다 충분히 길게</p></body></html>"
        edited = "<html><body><p>핵심 메시지 문장입니다 충분히 길게</p></body></html>"  # 제목 삭제
        diff = ov_mod.manual_layer_diff(edited, deck, baseline)
        self.assertEqual(1, len(diff["removed"]))
        self.assertTrue(diff["removed"][0]["review_flag"])       # review_needed [예시]와 겹침
        self.assertEqual(0, len(diff["added"]))
        md = ov_mod.render_manual_layer_diff_md(diff)
        self.assertIn("삭제 1건", md)
        self.assertIn("예시 태그 겹침", md)

    def test_no_change_reports_zero_explicitly(self):
        import overrides as ov_mod
        deck = self._deck()
        baseline = "<html><body><h1>[예시] 표지 이미지 주제</h1></body></html>"
        diff = ov_mod.manual_layer_diff(baseline, deck, baseline)
        self.assertEqual([], diff["removed"])
        self.assertEqual([], diff["added"])
        md = ov_mod.render_manual_layer_diff_md(diff)
        self.assertIn("변경 0건", md)                             # 빈 파일 금지 — 명시적으로 0건 기록

    def test_added_text_is_detected_by_reverse_comparison(self):
        import overrides as ov_mod
        deck = self._deck()
        baseline = "<html><body><h1>[예시] 표지 이미지 주제</h1></body></html>"
        edited = ("<html><body><h1>[예시] 표지 이미지 주제</h1>"
                  "<p>출처 없이 새로 끼워넣은 문장입니다 테스트용</p></body></html>")
        diff = ov_mod.manual_layer_diff(edited, deck, baseline)
        self.assertEqual([], diff["removed"])
        self.assertEqual(1, len(diff["added"]))
        self.assertIn("새로 끼워넣은", diff["added"][0]["text"])
        self.assertFalse(diff["added"][0]["review_flag"])         # review_needed와 안 겹침 → 확인 필요 태그

    def test_manual_layer_diff_warning_surfaces_only_when_nonzero(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self.assertEqual([], pipeline_state._manual_layer_diff_warning(run))  # 파일 없음
            (run / "manual_layer_diff.md").write_text(
                "# ...\n\n변경 0건 — 편집본이 베이스라인 대비 사실 텍스트를 삭제·추가하지 않았다.\n",
                encoding="utf-8")
            self.assertEqual([], pipeline_state._manual_layer_diff_warning(run))  # 0건은 경고 아님
            (run / "manual_layer_diff.md").write_text(
                "# ...\n\n삭제 1건 · 추가 0건\n", encoding="utf-8")
            warns = pipeline_state._manual_layer_diff_warning(run)
            self.assertTrue(warns)
            self.assertIn("결정 7", warns[0])


class W14SelectionProvenanceSmoke(unittest.TestCase):
    """W14(결정 8): 공고 선택 출처(selected_by) 기록·표면화 — 무엇에 입찰하나는 사람 전속 판단.
    차단 없이 기록만 하고, agent/unspecified만 1급 경고로 띄운다. 기대값은 라이브 검증
    run 3종(workspace/runs/w14_test_agent·dashboard·unspecified)의 selection 블록과 동일."""

    def _start(self, td, bid, *, selected_by=None, feedback=None):
        """실제 start_cmd를 임시 REPO_ROOT/RUNS로 관통 — 자동 실측(파일 대조)까지 태운다."""
        root = Path(td)
        runs = root / "workspace" / "runs"
        runs.mkdir(parents=True)
        if feedback is not None:
            (root / "dashboard").mkdir()
            (root / "dashboard" / "feedback.json").write_text(
                json.dumps(feedback, ensure_ascii=False), encoding="utf-8")
        args = SimpleNamespace(bid=bid, brief=None, mode="direct",
                               run_name=None, selected_by=selected_by)
        buf = io.StringIO()
        with mock.patch.object(proposal_pipeline, "REPO_ROOT", root), \
             mock.patch.object(proposal_pipeline, "RUNS", runs), \
             mock.patch("sys.stdout", buf):
            proposal_pipeline.start_cmd(args)
        run = next(runs.iterdir())
        state = json.loads((run / pipeline_state.STATE_NAME).read_text(encoding="utf-8"))
        return run, state, buf.getvalue()

    def test_agent_selection_is_recorded_and_warned_first_class(self):
        with tempfile.TemporaryDirectory() as td:
            run, state, out = self._start(td, "TEST-W14-AGENT-001", selected_by="agent")
            self.assertEqual({"selected_by": "agent", "bid": "TEST-W14-AGENT-001",
                              "feedback_match": False}, state["selection"])
            view = pipeline_state.resolve(run)
        self.assertTrue(any("공고 선택 출처=agent" in w for w in view["warnings"]),
                        msg=str(view["warnings"]))       # status/go가 읽는 경고에 1급 노출
        self.assertIn("공고 선택 출처=agent", out)        # start 직후 status 출력에도 즉시 표시

    def test_feedback_go_bid_auto_measures_dashboard_and_stays_silent(self):
        feedback = {"R26BK01608833-000": {"decision": "go"},
                    "R26BK-OTHER-000": {"decision": "no_go"}}
        with tempfile.TemporaryDirectory() as td:
            run, state, out = self._start(td, "R26BK01608833-000", feedback=feedback)
            self.assertEqual("dashboard", state["selection"]["selected_by"])
            self.assertIs(True, state["selection"]["feedback_match"])  # 자기보고 아님 — 파일 대조 실측
            view = pipeline_state.resolve(run)
        self.assertFalse(any("공고 선택 출처" in w for w in view["warnings"]),
                         msg=str(view["warnings"]))       # 사람 Go 기록 있는 bid는 경고 침묵
        self.assertNotIn("[!] 누가 골랐는지", out)

    def test_unrecorded_bid_defaults_to_unspecified_with_notice(self):
        with tempfile.TemporaryDirectory() as td:
            run, state, out = self._start(td, "TEST-W14-NOFEEDBACK-001")  # feedback.json 자체가 없음
            self.assertEqual("unspecified", state["selection"]["selected_by"])
            self.assertIs(False, state["selection"]["feedback_match"])
            view = pipeline_state.resolve(run)
        self.assertIn("누가 골랐는지 기록되지 않음", out)   # start 직후 안내
        self.assertTrue(any("--selected-by user" in w for w in view["warnings"]),
                        msg=str(view["warnings"]))        # status에도 명시 유도 경고


class W15MessageMapSmoke(unittest.TestCase):
    """W15(결정 9①③④): message_map을 1급 산출물로 배선 — 메시지 우선 공정.

    검증 철학(결정 7~8): 구조 위반(governing 0/2+)만 차단, 나머지는 표면화(경고 1급).
    스토리라인 핸드오프에 message_map을 메시지 계약으로 주입한다."""

    def _map(self, *, governing="발주처의 신뢰 회복을 위한 통합 소통체계", axes=2, empty=False):
        ax = []
        for i in range(1, axes + 1):
            ax.append({
                "id": f"axis{i}", "message": f"하위 메시지 {i}",
                "evidence_slots": [{
                    "type": "데이터", "desc": f"근거 {i}",
                    "status": "empty" if empty else "filled",
                    "source": None if empty else "출처",
                }],
            })
        return {"governing_message": governing, "strategy_axes": ax,
                "audience_note": "심사위원은 실현가능성을 본다",
                "knowledge_used": {"cards": [], "web": []}}  # ε패킷 안전장치① — 생략하면 수거 차단

    def _write_map(self, run, doc):
        run.mkdir(parents=True, exist_ok=True)
        message_map.map_path(run).write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

    def test_two_governing_messages_is_a_hard_block(self):
        """① governing 2개 → 오류(구조 위반이라 이것만 차단)."""
        doc = self._map(governing=["주장 A", "주장 B"])          # 2개
        errors, _ = message_map.validate(doc)
        self.assertTrue(any("정확히 1개" in e for e in errors), msg=str(errors))
        self.assertFalse(message_map.gating_block(doc)["governing_ok"])
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            self._write_map(run, doc)
            with self.assertRaises(proposal_pipeline.PipelineInputError):
                proposal_pipeline._go_message_map_collect(run, SimpleNamespace())

    def test_five_axes_is_surfaced_as_warning_not_blocked(self):
        """② 축 5개 → 경고 1급 표면화(차단 없음)."""
        doc = self._map(axes=5)
        errors, warns = message_map.validate(doc)
        self.assertEqual([], errors)                              # 차단 아님
        self.assertTrue(any("전략 축 5개" in w for w in warns), msg=str(warns))
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            self._write_map(run, doc)
            view = pipeline_state.resolve(run)
        self.assertTrue(any("전략 축 5개" in w for w in view["warnings"]), msg=str(view["warnings"]))

    def test_empty_slot_is_tagged_review_needed(self):
        """③ empty 슬롯 → review_needed 계열로 표면화(창작금지 대칭·차단 없음)."""
        doc = self._map(empty=True)
        self.assertEqual(2, message_map.gating_block(doc)["slots"]["empty"])
        self.assertEqual(2, len(message_map.empty_slots(doc)))
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            self._write_map(run, doc)
            view = pipeline_state.resolve(run)
            buf = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_record_state"), \
                 mock.patch("sys.stdout", buf):
                proposal_pipeline._go_message_map_collect(run, SimpleNamespace())
        self.assertTrue(any("empty" in w and "검토요망" in w for w in view["warnings"]),
                        msg=str(view["warnings"]))
        self.assertIn("[검토요망]", buf.getvalue())               # collect도 검토요망으로 노출

    def test_storyline_handoff_injects_axis_contract(self):
        """④ 스토리라인 핸드오프에 message_map(축 계약)이 주입된다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir(parents=True)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="demo.md")
            (run / "brief.md").write_text("발주처 요구사항 문서", encoding="utf-8")
            self._write_map(run, self._map(axes=3))
            out = proposal_pipeline.bundle_storyline_from_brief(
                run, SimpleNamespace(pack="core"))
            text = out.read_text(encoding="utf-8")
        self.assertIn("메시지 계약", text)                        # 계약 헤더
        self.assertIn("발주처의 신뢰 회복", text)                 # governing 주입
        self.assertIn("axis3", text)                              # 축 id 주입
        # W32 마찰29: 축 추적은 supports_axis 필드 전담 — message 본문 표기 지시는 제거됐다
        # (내부 표기가 청중 장표에 그대로 조판되던 결함). 계약 자체는 유지: 축 지지 의무 + 필드 명시.
        self.assertIn("전략 축 중 하나를 지지해야 한다", text)
        self.assertIn("message 본문에 축 id를 적지 마라", text)
        self.assertIn("supports_axis 필드로 명시하라", text)
        self.assertNotIn("축의 id를 함께 명시", text)


class W16SkeletonSubordinationSmoke(unittest.TestCase):
    """W16(결정 9①④⑤): 스켈레톤·장표 도출을 message_map에 종속화 + 대목차 의무 슬롯 + 분량 리듬.

    ① map 있으면 축별 조립+supports_axis 운반 ② map 없으면 기존 경로 바이트 동일
    ③ 기본 시나리오에 closing_matrix 부재+G6 존재 ④ 분량 밴드 위반→경고 표면화
    ⑤ Level 1 골격 순서 고정."""

    def _map(self, axes=3):
        ax = []
        for i in range(1, axes + 1):
            ax.append({
                "id": f"axis{i}", "message": f"하위 메시지 {i}",
                "evidence_slots": [{"type": "데이터", "desc": f"근거 {i}",
                                    "status": "example", "source": None}],
            })
        return {"governing_message": "발주처를 위한 통합 소통체계", "strategy_axes": ax}

    def test_axis_assembly_carries_supports_axis_into_deck(self):
        """① map 있으면 축별 조립 — supports_axis가 스켈레톤→deck.json까지 운반된다."""
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map(axes=3))
        self.assertTrue(doc["meta"]["message_driven"])
        # 축별 전략+근거 장표에 supports_axis가 실린다.
        axis_slides = [s for s in doc["slides"] if s.get("supports_axis")]
        self.assertTrue(any(s["supports_axis"] == "axis3" for s in axis_slides))
        # 근거 슬롯 type=데이터 → data_interpretation로 매핑.
        self.assertTrue(any(s["template_id"] == "data_interpretation" for s in axis_slides))
        # deck.json까지 운반: adapt_storyline 통과 후에도 supports_axis 보존.
        deck = _adapt.adapt_storyline(doc, pack="house_a")
        carried = [s for s in deck["slides"] if s.get("supports_axis")]
        self.assertTrue(carried)                                   # 최소 1장 운반
        self.assertTrue(any(s["supports_axis"] == "axis1" for s in carried))

    def test_map_absent_path_is_byte_identical_to_legacy_logic(self):
        """② map 없으면 기존 시나리오 경로 바이트 동일(축 로직 미주입·supports_axis 없음)."""
        # 밴드 선언 없는 고정 픽스처(레거시 시나리오와 동형) — 코드 경로가 불변임을 격리 검증.
        scenario = {"id": "fx", "label": "픽스처", "slides": [
            {"section": "표지", "title": "T", "message": "M", "template_id": "cover_cinematic",
             "bullets": ["b"], "note": "n"},
            {"section": "전략", "title": "S", "message": "MM", "template_id": "strategy_pillars",
             "bullets": []},
        ]}
        expected = {
            "meta": {"project": "픽스처", "scenario": "fx", "skeleton": True},
            "slides": [
                {"n": 1, "section": "표지", "title": "T", "message": "M", "bullets": ["b"],
                 "template_id": "cover_cinematic", "example": True,
                 "flag": [skeleton.EXAMPLE_FLAG], "note": "n"},
                {"n": 2, "section": "전략", "title": "S", "message": "MM", "bullets": [],
                 "template_id": "strategy_pillars", "example": True,
                 "flag": [skeleton.EXAMPLE_FLAG]},
            ],
        }
        got = skeleton.build_skeleton(scenario, message_map_doc=None)
        self.assertEqual(expected, got)                            # 바이트(구조) 동일
        # 밴드/축 키가 새지 않는다(map 없음 → 순수 시나리오 경로).
        for s in got["slides"]:
            self.assertNotIn("supports_axis", s)
            self.assertNotIn("length_band", s)

    def test_default_scenario_drops_closing_matrix_for_g6(self):
        """③ public_proposal(레거시·opt-in 시나리오)에서 closing_matrix 부재 + G6(성과관리·효과조사)
        존재(W28: 기본 시나리오는 public_proposal_core로 전환돼 G6 표현이 없다 — 이 레거시 검증은
        --scenario public_proposal 명시로 고정)."""
        sc = skeleton.load_scenario("public_proposal")
        tids = [s.get("template_id") for s in sc["slides"]]
        self.assertNotIn("closing_matrix", tids)                   # 결정 9④ — 홍보용역 기본 제외
        titles = " ".join(str(s.get("title", "")) + str(s.get("note", "")) for s in sc["slides"])
        self.assertIn("성과관리", titles)                          # G6 클로징으로 교체
        # 축 조립 경로에서도 마무리가 G6(성과관리)이고 closing_matrix가 아니다.
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map())
        assembled = [s.get("template_id") for s in doc["slides"]]
        self.assertNotIn("closing_matrix", assembled)

    def test_band_violation_surfaces_as_warning_not_block(self):
        """④ 분량 밴드 위반 → gating_report.length_rhythm 실측 + 경고 1급 표면화(차단 없음)."""
        # 얇은 밴드를 선언한 슬라이드에 과한 분량 → over 위반. 동적 범위도 함께 실측.
        deck = {"meta": {"project": "P"}, "slides": [
            {"slide_id": 1, "role": "opening", "template_id": "problem_questions",
             "key_message": "짧다", "body": [], "fields": {}, "length_band": [20, 70]},
            {"slide_id": 2, "role": "evidence", "template_id": "data_interpretation",
             "key_message": " ".join(["어절"] * 120), "body": [], "fields": {},
             "length_band": [20, 70]},   # 120 > 70 → over 위반
        ]}
        rhythm = proposal_pipeline._compute_length_rhythm(deck)
        self.assertIsNotNone(rhythm)
        self.assertTrue(any(v["kind"] == "over" and v["slide_id"] == 2
                            for v in rhythm["band_violations"]))
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "gating_report.json").write_text(
                json.dumps({"length_rhythm": rhythm}, ensure_ascii=False), encoding="utf-8")
            warns = pipeline_state._length_rhythm_warnings(run)
        self.assertTrue(any("분량 밴드 위반" in w for w in warns), msg=str(warns))

    def test_level1_backbone_order_is_fixed(self):
        """⑤ Level 1 대목차 골격 순서 고정 — 고민→제안개요(복창의무)→제안업체→사업관리→(사업내용)→마무리.

        W28: public_proposal(레거시·자체 G1 declare)로 명시 고정 — 기본 시나리오는
        public_proposal_core로 전환돼 "고민"이 기본 경로에 없다(별도 W28 스모크가 검증)."""
        sc = skeleton.load_scenario("public_proposal")
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map(axes=2))
        sections = [s["section"] for s in doc["slides"]]
        # Level 1 고정 슬롯이 선언 순서대로 등장(사업내용 축 그룹은 사업관리 뒤·마무리 앞).
        self.assertEqual("고민", sections[0])
        self.assertEqual("제안개요", sections[1])
        self.assertEqual("제안업체", sections[2])
        self.assertEqual("사업관리", sections[3])
        self.assertEqual("마무리", sections[-1])                   # G6 클로징이 맨 끝
        self.assertLess(sections.index("사업관리"), sections.index("마무리"))
        self.assertIn("전략", sections)                           # 사업내용 = 축별 전략 장표
        # 제안개요는 복창 의무 슬롯.
        overview = next(s for s in doc["slides"] if s["section"] == "제안개요")
        self.assertTrue(overview.get("mandatory"))


class W27GseriesOptInSmoke(unittest.TestCase):
    """W27(목표조정 2): 수주덱 채굴 지식(G계열)은 기본 골격에서 제외 — house_knowledge="gseries"
    명시 시에만 opt-in 복원된다. public_proposal.json은 자체 level1을 선언해 무영향(경계 확인은
    별도 픽스처로 DEFAULT_LEVEL1 폴백 경로를 직접 실측)."""

    def _map(self, axes=2):
        ax = []
        for i in range(1, axes + 1):
            ax.append({
                "id": f"axis{i}", "message": f"하위 메시지 {i}",
                "evidence_slots": [{"type": "데이터", "desc": f"근거 {i}",
                                    "status": "example", "source": None}],
            })
        return {"governing_message": "발주처를 위한 통합 소통체계", "strategy_axes": ax}

    def _no_level1_scenario(self):
        return {"id": "fx_no_level1", "label": "픽스처(level1 미선언)", "slides": []}

    def test_default_skeleton_has_no_gseries_content(self):
        """① 기본 skeleton 산출(house_knowledge 미지정)에 '고민' 섹션·G1/G8 문자열이 없다."""
        doc = skeleton.build_skeleton(self._no_level1_scenario(), message_map_doc=self._map())
        sections = [s.get("section") for s in doc["slides"]]
        self.assertNotIn("고민", sections)
        blob = " ".join(str(s.get("title", "")) + str(s.get("note", "")) for s in doc["slides"])
        self.assertNotIn("G1", blob)
        self.assertNotIn("G8", blob)

    def test_house_knowledge_gseries_restores_overlay(self):
        """② house_knowledge="gseries" 지정 시 고민 슬롯(G1)이 제안개요 앞에 등장한다."""
        doc = skeleton.build_skeleton(
            self._no_level1_scenario(), message_map_doc=self._map(), house_knowledge="gseries",
        )
        sections = [s.get("section") for s in doc["slides"]]
        self.assertIn("고민", sections)
        self.assertLess(sections.index("고민"), sections.index("제안개요"))
        blob = " ".join(str(s.get("title", "")) + str(s.get("note", "")) for s in doc["slides"])
        self.assertIn("G1", blob)
        self.assertIn("G8", blob)                                  # 제안업체 노트 복원
        self.assertIn("G6", blob)                                  # 마무리 노트 복원


class W28NeutralDefaultScenarioSmoke(unittest.TestCase):
    """W28(목표조정 2 완결): 기본 시나리오 자체가 수주덱 채굴 원산이었다(W27은 skeleton.py의
    DEFAULT_LEVEL1 폴백만 중립화했지만, public_proposal.json이 자체 level1을 선언해 폴백을
    가려버렸다) — 기본을 public_proposal_core로 전환하고, 레거시는 --scenario 명시로 opt-in."""

    def _map(self, axes=2):
        ax = []
        for i in range(1, axes + 1):
            ax.append({
                "id": f"axis{i}", "message": f"하위 메시지 {i}",
                "evidence_slots": [{"type": "데이터", "desc": f"근거 {i}",
                                    "status": "example", "source": None}],
            })
        return {"governing_message": "발주처를 위한 통합 소통체계", "strategy_axes": ax}

    def test_default_path_has_no_gomin_or_g_series_content(self):
        """① 기본 경로(스켈레톤 기본 시나리오) 산출에 "고민"·"G1"·"수주작" 부재, 의무 골격 순서 유지."""
        sc = skeleton.load_scenario(None)
        self.assertEqual("public_proposal_core", sc.get("id"))
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map())
        sections = [s.get("section") for s in doc["slides"]]
        self.assertNotIn("고민", sections)
        blob = " ".join(str(s.get("title", "")) + str(s.get("note", "")) + str(s.get("message", ""))
                         for s in doc["slides"])
        self.assertNotIn("G1", blob)
        self.assertNotIn("수주작", blob)
        # 문서 프레임 의무 슬롯(D23 개정): 표지·목차는 맨 앞, 끝인사는 맨 끝.
        self.assertEqual("표지", sections[0])
        self.assertEqual("목차", sections[1])
        self.assertEqual("끝인사", sections[-1])
        # RFP 의무 골격(대목차 순서)은 프레임 안쪽에서 그대로 — 제안개요→제안업체→사업관리→(사업내용)→마무리.
        self.assertLess(sections.index("제안개요"), sections.index("제안업체"))
        self.assertLess(sections.index("제안업체"), sections.index("사업관리"))
        self.assertLess(sections.index("사업관리"), sections.index("마무리"))
        overview = next(s for s in doc["slides"] if s["section"] == "제안개요")
        self.assertTrue(overview.get("mandatory"))

    def test_scenario_public_proposal_explicit_keeps_legacy_gomin(self):
        """② --scenario public_proposal 명시 시 기존 G 골격(고민 등) 그대로(레거시 재현성)."""
        sc = skeleton.load_scenario("public_proposal")
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map())
        sections = [s.get("section") for s in doc["slides"]]
        self.assertIn("고민", sections)
        self.assertEqual("고민", sections[0])

    def test_house_knowledge_gseries_restores_overlay_on_declared_core_level1(self):
        """③ house_knowledge=gseries + 기본(core) 시나리오 조합 — core가 level1을 선언해도
        오버레이가 그 위에 얹혀 G1이 삽입된다(선언 시나리오 위 오버레이 — W28 수정 대상)."""
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map(), house_knowledge="gseries")
        sections = [s.get("section") for s in doc["slides"]]
        self.assertIn("고민", sections)
        self.assertLess(sections.index("고민"), sections.index("제안개요"))
        blob = " ".join(str(s.get("title", "")) + str(s.get("note", "")) for s in doc["slides"])
        self.assertIn("G1", blob)

    def test_legacy_scenario_gseries_overlay_does_not_duplicate_gomin(self):
        """레거시 public_proposal(자체 G1 declare)에 house_knowledge=gseries를 얹어도 "고민"이
        중복 삽입되지 않는다(멱등 가드)."""
        sc = skeleton.load_scenario("public_proposal")
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map(), house_knowledge="gseries")
        sections = [s.get("section") for s in doc["slides"]]
        self.assertEqual(1, sections.count("고민"))


class D23DocumentFrameSlotsSmoke(unittest.TestCase):
    """D23 개정(2026-07-14): 문서 프레임 의무 슬롯 3종(표지·목차·끝인사)을 기본 시나리오
    (public_proposal_core) level1 골격에 배선. 표지=role cover(브랜드 마크 자리)·목차=대목차
    자동 조립·끝인사=정적 문구(마무리 G6 뒤 별도). 렌더 기계는 재사용(새 기계 아님)."""

    def _map(self, axes=2):
        ax = []
        for i in range(1, axes + 1):
            ax.append({
                "id": f"axis{i}", "message": f"하위 메시지 {i}",
                "evidence_slots": [{"type": "데이터", "desc": f"근거 {i}",
                                    "status": "example", "source": None}],
            })
        return {"governing_message": "발주처를 위한 통합 소통체계", "strategy_axes": ax}

    def test_frame_slots_present_in_default_skeleton(self):
        """① 기본 시나리오 축 조립 골격에 표지(cover)·목차(agenda)·끝인사(endcard)가 나온다."""
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map())
        by_section = {s["section"]: s for s in doc["slides"]}
        # 표지: 맨 앞 · template cover_slide(텍스트 우선) · G6 마무리보다 앞.
        self.assertEqual("표지", doc["slides"][0]["section"])
        self.assertEqual("cover_slide", by_section["표지"]["template_id"])
        # 목차: 표지 다음 · template agenda · items = 대목차 섹션 자동 조립(지어내기 0).
        self.assertEqual("목차", doc["slides"][1]["section"])
        self.assertEqual("agenda", by_section["목차"]["template_id"])
        items = by_section["목차"]["fields"]["items"]
        self.assertIn("제안개요", items)
        self.assertIn("사업내용", items)
        self.assertIn("마무리", items)
        self.assertNotIn("표지", items)          # 프레임은 목차에 자기 자신을 열거하지 않는다
        self.assertNotIn("목차", items)
        self.assertNotIn("끝인사", items)
        # 끝인사: 맨 끝 · template closing_thanks · 마무리(G6)는 그 앞에 유지(둘 다 존재).
        self.assertEqual("끝인사", doc["slides"][-1]["section"])
        self.assertEqual("closing_thanks", by_section["끝인사"]["template_id"])
        sections = [s["section"] for s in doc["slides"]]
        self.assertLess(sections.index("마무리"), sections.index("끝인사"))

    def test_frame_slots_render_with_client_brand_on_cover(self):
        """② 기본 시나리오 골격을 core 팩+브랜드 스킨으로 렌더 — 표지/목차/끝인사가 폴백 경고
        없이 그려지고, 클라이언트 브랜드 마크가 표지 슬라이드에만 뜬다."""
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map())
        deck = _adapt.adapt_storyline(doc, pack="core")
        roles = {s["role"] for s in deck["slides"]}
        self.assertIn("cover", roles)                    # 표지 section → role cover
        skin = {"brand": {"client_name": "테스트기관",
                          "placement": {"client": "cover", "proposer": "all"}}}
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render(deck, catalog="core", skins=[skin], out_path=out)
            html = out.read_text(encoding="utf-8")
        # 프레임 슬롯 3종이 미지원 폴백(🔴)으로 떨어지지 않는다.
        self.assertNotIn("미지원 템플릿", html)
        # 클라이언트 마크가 표지(cover)에 뜬다(placement.client=cover).
        self.assertIn("slide__brand--client", html)
        self.assertIn("테스트기관", html)
        # 목차·끝인사 렌더 마크업이 실제로 존재(조각 재사용 확인).
        self.assertIn("core-agenda", html)
        self.assertIn("core-endcard", html)

    def test_gseries_overlay_keeps_frame_slots_ahead_of_gomin(self):
        """③ house_knowledge=gseries 오버레이 시 고민(G1)은 프레임(표지·목차) 뒤·제안개요 앞."""
        sc = skeleton.load_scenario(None)
        doc = skeleton.build_skeleton(sc, message_map_doc=self._map(), house_knowledge="gseries")
        sections = [s["section"] for s in doc["slides"]]
        self.assertEqual("표지", sections[0])
        self.assertEqual("목차", sections[1])
        self.assertLess(sections.index("목차"), sections.index("고민"))
        self.assertLess(sections.index("고민"), sections.index("제안개요"))
        self.assertEqual("끝인사", sections[-1])


class W12EvalBundleFreshnessSmoke(unittest.TestCase):
    """W12/마찰 D3: stage9 --apply 이후 기존 deck_review 번들이 stale이면 재번들 안내(mtime 계열)."""

    def _run_with_apply(self, td, *, apply_recent=True):
        run = Path(td)
        import datetime as _dt
        recent = _dt.datetime.now().isoformat(timespec="seconds")
        state = pipeline_state._blank(run)
        state["mode"] = "direct"
        state["input"] = {"kind": "brief", "ref": "b"}
        for cp in ("start", "decision"):
            state["checkpoints"][cp] = {"cleared_at": recent}
        state["stages"]["render"] = {"at": recent, "source": "recorded"}
        state["stages"]["stage9_apply"] = {"at": recent, "source": "recorded"}
        (run / pipeline_state.STATE_NAME).write_text(
            json.dumps(state, ensure_ascii=False, indent=2), encoding="utf-8")
        (run / "deck.json").write_text("{}", encoding="utf-8")
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        return run

    def _write_bundle(self, run, *, minutes_old):
        prompt = run / "deck_review" / "deck_review_prompt.md"
        prompt.parent.mkdir(parents=True, exist_ok=True)
        prompt.write_text("bundle", encoding="utf-8")
        old = prompt.stat().st_mtime - minutes_old * 60
        os.utime(prompt, (old, old))
        return prompt

    def test_bundle_older_than_apply_warns_to_rebundle(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run_with_apply(td)
            self._write_bundle(run, minutes_old=10)
            view = pipeline_state.resolve(run)
        self.assertTrue(any("재번들" in w for w in view["warnings"]),
                        msg=str(view["warnings"]))

    def test_fresh_bundle_does_not_warn(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run_with_apply(td)
            self._write_bundle(run, minutes_old=-10)   # 번들이 apply보다 새것
            view = pipeline_state.resolve(run)
        self.assertFalse(any("재번들" in w for w in view["warnings"]),
                         msg=str(view["warnings"]))


class W12StrategyPillarsSmoke(unittest.TestCase):
    """W12/마찰 C: strategy_pillars가 card_grid(다른 필드)로 매핑돼 계약 필드
    (pillars/one_line_per_pillar)를 채워도 본문이 빈 채 warnings=0 통과하던 조용한 폴백.
    """

    def _render(self):
        deck = {"meta": {"project": "P"}, "slides": [{
            "slide_id": "s1", "role": "strategy", "template_id": "strategy_pillars",
            "title": "핵심 추진 전략", "key_message": "실행을 3개 축으로",
            "fields": {"pillars": ["전략가", "전략나", "전략다"],
                       "one_line_per_pillar": ["설명가", "설명나", "설명다"]}}]}
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render(deck, catalog="core", skins=["core"], out_path=out)
            return rep, out.read_text(encoding="utf-8")

    def test_contract_fields_actually_render(self):
        rep, html = self._render()
        self.assertEqual([], rep["warnings"])
        self.assertEqual(3, html.count('class="card"'))        # 회귀 전엔 0(빈 본문)
        for tok in ("전략가", "전략나", "전략다", "설명가", "설명나", "설명다"):
            self.assertIn(tok, html, msg=f"계약 필드 '{tok}'가 렌더 안 됨(조용한 폴백 회귀)")


class W12HeroScrimSmoke(unittest.TestCase):
    """W12/마찰 D2: 배경 스크림이 하단 편중이라 상단(메타·eyebrow) 텍스트 보호가 약했다."""

    def test_scrim_protects_top_information_zone(self):
        css = image_slots.SLOT_CSS
        self.assertIn(".dov-slot__scrim", css)
        self.assertIn("rgba(5,5,5,.5) 0%", css)                # 상단(0%) 강한 보호 stop
        self.assertNotIn("rgba(5,5,5,.15) 0%", css)            # 종전 약한 상단 stop 소멸
        self.assertIn("linear-gradient", css)                  # 여전히 실체 있는 treatment

    def test_background_prompt_rule_mentions_top_zone(self):
        slot = {"role": "conceptual", "layer": "background", "prompt": "도시 야경",
                "format": "jpg", "treatment": "저대비"}
        prompt = image_slots.build_fill_prompt(slot)
        self.assertIn("상단", prompt)                          # 규칙층 한 줄: 상단 정보영역 보호


class W18WatermarkBleedSmoke(unittest.TestCase):
    """W18: 예시 워터마크(안전장치 ①)의 바운딩박스가 슬라이드 폭 밖으로 수평 블리드하던 결함.

    재현 근원 = `.example-watermark`가 `inset:0`(슬라이드 전체 크기 박스)인 채로 그 박스 자체를
    회전(rotate)해서, 회전 후 축정렬 바운딩박스가 원래 슬라이드보다 커졌다(회전한 각도만큼
    필연). 수리 = 텍스트 콘텐츠 크기만큼만 박스를 잡고(`translate(-50%,-50%) rotate(...)`)
    그것만 회전 — 봉쇄가 아니라 애초에 안 새게. 삭제로 통과하는 우회를 막기 위해 워터마크
    DOM 존재·가시성도 같이 잰다(정직성 장치 보존 확인).
    """

    def _example_storyline(self):
        return {"meta": {"project": "P"}, "slides": [
            {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": []},
            {"n": 2, "section": "데이터", "title": "시장", "message": "규모",
             "template_id": "data_interpretation", "example": True,
             "fields": {"metric": "m", "comparison": [{"label": "A", "value": 70},
                                                       {"label": "B", "value": 30}],
                        "interpretation": ["예시 해석"]}},
        ]}

    def _render(self, td):
        deck = _adapt.adapt_storyline(self._example_storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        out = Path(td) / "deck.html"
        htmlgen.render_html(deck, "core", out)
        return out

    def test_example_watermark_probe_overflow_is_zero(self):
        """① 워터마크가 있는 example 슬라이드를 실제로 렌더해 layout_probe로 재면 overflow 0."""
        import layout_probe
        if not layout_probe.available():
            self.skipTest(layout_probe.unavailable_reason() + " — 브라우저 실측 불가")
        with tempfile.TemporaryDirectory() as td:
            out = self._render(td)
            probe = layout_probe.probe_html(out)
        watermark_slide = next(s for s in probe["slides"] if s["slide_id"] == "2")
        self.assertNotIn("overflow_measured", watermark_slide["flags"])
        self.assertNotIn("overflow_measured_x", watermark_slide["flags"])
        self.assertEqual(0, watermark_slide["content_overflow_x_px"])

    def test_watermark_dom_exists_and_visible_not_deleted_to_pass(self):
        """② overflow 0을 워터마크 삭제로 달성하는 우회 방지 — DOM이 실제로 존재·가시해야 한다."""
        with tempfile.TemporaryDirectory() as td:
            out = self._render(td)
            html = out.read_text(encoding="utf-8")
        self.assertIn("예시 데이터".encode("utf-8").decode("utf-8"), html)
        self.assertIn('class="example-watermark"', html)
        self.assertNotIn('.example-watermark { display: none', html)
        self.assertNotIn('.example-watermark{display:none', html)
        self.assertNotIn("visibility: hidden", html.split(".example-watermark")[1][:200])


class NeutralCorePackSmoke(unittest.TestCase):
    """W20(NORTHSTAR 결정 11): 하우스 스타일 제로 기본 — 중립 코어 팩.

    기본 파이프라인은 packs/core(무채·형태 어휘 카탈로그)만으로 완주하고,
    하우스 색(오렌지)은 기본 렌더 CSS에 없다. house_a/house_b는 --pack 명시 시에만
    쓰는 플러그인 예제였다 — W31 E3에서 <개발 원본 전용 경로> 격리(로컬에
    없으면 하우스 동작 재현 테스트만 skip, `_requires_house_a` 참고).
    """

    @staticmethod
    def _render(pack: str, td: str):
        import re
        sc = skeleton.load_scenario(None)          # public_proposal(중립 어휘 template_id)
        doc = skeleton.build_skeleton(sc)
        out = Path(td) / "deck.html"
        rep = htmlgen.render_html(doc, pack, out)
        html = out.read_text(encoding="utf-8")
        css = re.search(r"<style>(.*?)</style>", html, re.S).group(1)
        return doc, rep, html, css

    def test_default_pack_core_renders_skeleton_to_completion(self):
        """① 기본 팩=중립 코어로 표준 시나리오 스켈레톤이 폴백 경고 없이 완주한다."""
        with tempfile.TemporaryDirectory() as td:
            doc, rep, html, _css = self._render("core", td)
        self.assertGreaterEqual(len(doc["slides"]), 1)
        # 폴백 가시화: 적합 렌더러 없어 generic으로 떨어진 슬라이드가 없어야 한다(밋밋 경고 0).
        self.assertEqual([], rep.get("warnings") or [],
                         f"중립 코어 렌더에 폴백 경고: {rep.get('warnings')}")
        # 정직성 장치가 무채에서도 마크업으로 존재(엔진 하드코딩 대비 — [예시] 워터마크).
        self.assertIn('class="example-watermark"', html)

    def test_house_color_tokens_absent_from_default_render_css(self):
        """② 하우스 색 토큰(오렌지 등)이 기본 렌더 CSS에 없다 — 무채 팔레트."""
        with tempfile.TemporaryDirectory() as td:
            _doc, _rep, html, css = self._render("core", td)
        # house_a 하우스 색 hex가 어디에도 나타나지 않는다(대소문자 무관).
        low = html.lower()
        for house_hex in ("f37321", "ff7900", "ea5628", "1e4a8c"):   # orange 3종 + house_a navy
            self.assertNotIn(house_hex, low, f"하우스 색 {house_hex}가 기본 렌더에 누출")
        # 하우스 색 토큰 '정의'(--c-orange:)가 CSS에 없다(엔진 폴백 '참조'는 무해 — 정의만 검사).
        self.assertNotIn("--c-orange:", css)
        self.assertNotIn("--c-orange-deep:", css)
        self.assertNotIn("--c-navy:", css)
        # 중립 primary 토큰은 있다(무채 코어가 실제로 주입됐다는 증거).
        self.assertIn("--c-ink:", css)

    @_requires_house_a
    def test_pack_house_a_still_house_styled_when_explicit(self):
        """③ --pack house_a 명시 시 기존 하우스 동작 유지 — 플러그인 예제 존치·오렌지 복귀."""
        with tempfile.TemporaryDirectory() as td:
            _doc, rep, html, css = self._render("house_a", td)
        self.assertEqual([], rep.get("warnings") or [])
        # 명시 시엔 하우스 색이 돌아온다(중립화가 house_a를 훼손하지 않았다).
        self.assertIn("--c-orange:", css)
        self.assertIn("f37321", html.lower())


class ComposeEngineSmoke(unittest.TestCase):
    """W21-0 2a(결정 12): 골격(frame)×조각(piece) 조합 렌더러.

    frame 선언 슬라이드는 template_id 대신 compose로 렌더된다.
    계약 = packs/core/frames.json·pieces.json. R2(빈 슬롯 배치 금지)·
    R4(출처요망 딱지 — 강등 아님)·R6(공통 베이스라인 구조 강제)·격리(하우스 모듈 미로드).
    """

    @staticmethod
    def _deck(slides):
        return {"meta": {"project": "compose smoke"}, "slides": slides}

    def _render(self, slides, td):
        out = Path(td) / "deck.html"
        rep = htmlgen.render_html(self._deck(slides), "core", out)
        return rep, out.read_text(encoding="utf-8")

    def test_frame_slide_renders_via_compose_with_r2_collapse(self):
        """① frame 슬라이드가 조합으로 완주하고, 빈 슬롯은 배치되지 않는다(R2 — cols 축소)."""
        slides = [{"slide_id": "1", "title": "t", "frame": "row_n", "layout_group": "g1", "slots": [
            {"piece": "stat_card", "size": "third", "data": {"label": "인원", "value": "1.2만", "source_note": "결산"}},
            {"piece": "stat_card", "size": "third", "data": {"label": "시간", "value": "6개월", "source_note": "계획"}},
            {"piece": "stat_card", "size": "third", "data": {}},   # 빈 슬롯 — 배치 금지 대상
        ]}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertIn("compose--row_n cols-2", html)               # 3슬롯 중 2개만 유효 → cols-2
        self.assertIn('data-layout-group="g1"', html)              # R9 감시용 어트리뷰트
        self.assertIn(".compose--row_n", html)                     # 조합 CSS 주입

    def test_missing_source_gets_badge_not_demotion(self):
        """② R4: 출처 없는 수치는 '출처요망' 딱지 — 감춤·강등 아님(수치는 그대로 렌더)."""
        slides = [{"slide_id": "1", "title": "t", "frame": "full", "slots": [
            {"piece": "big_number", "size": "full", "data": {"value": "7.8", "unit": "배", "label": "도달"}},
        ]}]
        with tempfile.TemporaryDirectory() as td:
            _rep, html = self._render(slides, td)
        self.assertIn("출처요망", html)
        self.assertIn("7.8", html)                                 # 수치 자체는 온전히 표시(강등 없음)

    def test_unknown_piece_surfaces_review_not_silent(self):
        """③ 미구현 piece는 조용한 폴백이 아니라 검토요망+경고로 표면화(catalog_gap 런타임 짝)."""
        slides = [{"slide_id": "1", "title": "t", "frame": "full", "slots": [
            {"piece": "no_such_piece", "size": "full", "data": {"x": 1}},
        ]}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertIn("검토요망", html)
        self.assertTrue(any("no_such_piece" in w for w in rep.get("warnings") or []),
                        f"미구현 piece 경고 없음: {rep.get('warnings')}")

    def test_core_deck_does_not_absorb_house_layout_css(self):
        """④ 격리: core 덱 렌더는 하우스 레이아웃 **모듈**(layouts_house_a/house_b)을 흡수하지 않는다.

        주의: 엔진 _base_css에 body.pack-house_a 스코프 CSS가 박혀 있는 것은 별개
        잔재(비활성 — 백로그: 플러그인 이관). 여기선 모듈 흡수 차단만 고정한다.
        """
        slides = [{"slide_id": "1", "title": "t", "frame": "full", "slots": [
            {"piece": "text_block", "size": "full", "data": {"body": ["a"]}},
        ]}]
        with tempfile.TemporaryDirectory() as td:
            _rep, html = self._render(slides, td)
        # layouts_house_a.py CSS 고유 클래스(엔진 base엔 없음) — 미주입 확인.
        self.assertNotIn(".house_a-table", html)
        self.assertNotIn(".house_a-step", html)
        # layouts_core는 로드된다(core 렌더러 구현 소유).
        self.assertIn(".core-table", html)

    def test_rendition_swaps_visual_metaphor_not_logic(self):
        """⑤ rendition = 시각 은유 계층 — frame(논리)·슬롯 불변, 은유 클래스만 바뀐다.

        [3]는 기본 boxed로 내고 [4]가 T2 조정권으로 갈아끼운다(사용자 우려
        "CSS가 레이아웃을 고착" 처방 — 같은 순서 논리를 spine 은유로도 그린다).
        """
        base = {"slide_id": "1", "title": "t", "frame": "flow_seq", "slots": [
            {"piece": "text_block", "size": "half", "data": {"body": ["a"]}},
            {"piece": "text_block", "size": "half", "data": {"body": ["b"]}},
        ]}
        with tempfile.TemporaryDirectory() as td:
            _rep, html_default = self._render([dict(base)], td)
        with tempfile.TemporaryDirectory() as td:
            _rep, html_spine = self._render([dict(base, rendition="spine")], td)
        self.assertIn("rend-boxed", html_default)                  # 기본 은유 = boxed
        self.assertIn("compose--flow_seq cols-2 rend-spine", html_spine)
        # 논리 구조(슬롯 내용·순서)는 두 은유에서 동일하게 존재.
        for html in (html_default, html_spine):
            self.assertIn(">a<", html.replace("<li>a</li>", ">a<"))
            self.assertIn("compose--flow_seq", html)

    # 계약(pieces.json) 25종 전수 렌더 픽스처 — 최소 유효 데이터(requires 충족).
    PIECE_FIXTURES = {
        "text_block": {"body": ["a", "b"]},
        "big_number": {"value": "7.8", "unit": "배", "label": "도달", "source_note": "리포트"},
        "stat_card": {"label": "인원", "value": "1.2만", "source_note": "결산"},
        "calc_arrow": {"narrative": "지연 1일당", "result_value": "4.2배", "formula_footnote": "식"},
        "contrast_pair": {"reject": "A", "adopt": "B"},
        "compare_table": {"axes": [{"label": "투명성", "voice": "돈이 잘 가나?"}], "columns": ["자사", "A"],
                          "cells": [["공개", "부분"]], "selection_criteria": ["", "1위"], "source_note": "조사"},
        "before_after": {"metrics": ["도달"], "before": ["12만"], "after": ["94만"]},
        "loop_pair": {"vicious_loop": ["불신", "회피"], "virtuous_loop": ["신뢰", "참여"]},
        "part_of_whole": {"whole": "전체 예산", "part": "1.5억", "part_label": "모니터링", "ratio": 15, "source_note": "RFP"},
        "flow_arrow": {"stages": [{"label": "전", "description": "d"}, {"label": "후"}]},
        "matrix_2x2": {"axis_x": "비용", "axis_y": "효과", "quadrants": ["q1", "q2", "q3", "q4"], "current": 3, "target": 2},
        "connect_diagram": {"items": ["a", "b", "c"], "relation": "그룹핑",
                            "groups": [{"name": "G1", "items": ["a", "b"]}]},
        "group_naming": {"collected": ["x", "y"], "groups": [{"name": "이름", "items": ["x"]}], "names": ["이름"]},
        "analogy_hero": {"analogy": "이건 마치 OO이다", "proof": "숫자 근거"},
        "journey_flow": {"nodes": [{"label": "문제", "content": "c"}, {"label": "사용 후"}]},
        "match_pairs": {"complaints": ["불만"], "proposals": ["제안"],
                        "pairs": [{"complaint": "불만", "proposal": "제안"}]},
        "claim_proof_split": {"claim": "총평", "components": ["요소1", "요소2"], "proofs": ["95%"]},
        "funnel_3layer": {"tam": {"value": "5조", "source_note": "통계"}, "sam": {"value": "8000억"}, "som": "300억"},
        "chart": {"chart_type": "bar", "series": [{"label": "전", "value": 45}, {"label": "후", "value": 80}], "source_note": "조사"},
        "quote": {"text": "가상 반응", "attribution": "예상 헤드라인", "is_example": True},
        "image_evidence": {"asset_path": "assets/x.png", "caption": "실측 화면"},
        "timeline_gantt": {"tasks": [{"label": "과업A", "start": 1, "end": 2}], "period_labels": ["3월", "4월", "5월"], "period": ["3월", "4월", "5월"]},
        "org_table": {"lead": "PM", "roles": [{"team": "기획", "role": "전략"}]},
        "case_card": {"cases": [{"client": "C사", "description": "d", "metric": "97%", "is_example": True}]},
        "agenda": {"items": [{"title": "왜 하나", "relief": "원인"}, {"title": "성과는"}]},
        "pillar_card": {"claim": "3대 실행 전략", "pillars": [
            {"label": "진단", "line": "현행 격차를 데이터로 규명"},
            {"label": "설계", "line": "표준 프로세스로 재정렬"},
            {"label": "정착", "line": "운영 체계로 이관·지속"}]},
    }

    def test_every_contract_piece_has_renderer_and_renders(self):
        """⑥ 계약-구현 완전성: pieces.json의 전 조각이 렌더러를 갖고, 픽스처로 실제 렌더된다.

        계약(어휘)과 구현(PIECES 레지스트리)의 드리프트를 구조적으로 차단 — 조각을
        계약에 추가하면 이 테스트가 구현·픽스처를 강제한다.
        """
        import compose  # app/render는 htmlgen과 같은 경로로 이미 sys.path에 있음
        contract = json.loads((Path(server.ROOT) / "packs" / "core" / "pieces.json").read_text(encoding="utf-8"))
        contract_ids = [p["id"] for p in contract["pieces"]]
        missing_impl = [pid for pid in contract_ids if pid not in compose.PIECES]
        self.assertEqual([], missing_impl, f"계약에 있으나 렌더러 없음: {missing_impl}")
        missing_fixture = [pid for pid in contract_ids if pid not in self.PIECE_FIXTURES]
        self.assertEqual([], missing_fixture, f"픽스처 없음: {missing_fixture}")
        slides = [{"slide_id": str(i + 1), "title": pid, "frame": "full",
                   "slots": [{"piece": pid, "size": "full", "data": self.PIECE_FIXTURES[pid]}]}
                  for i, pid in enumerate(contract_ids)]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        # requires 충족 픽스처이므로 검토요망(필수 누락)·미구현 경고가 없어야 한다.
        self.assertEqual([], rep.get("warnings") or [], f"전수 렌더 경고: {rep.get('warnings')}")
        self.assertNotIn("필수 필드 누락", html)
        self.assertNotIn("미구현", html)


class WireframeContractSmoke(unittest.TestCase):
    """W21 [3] 와이어프레임 결정기 계약 — wireframe.json 검증·병합·게이트 블록.

    오류 = 적용 중단(SSOT 안전) / 경고 = 표면화(R1·R2·R9·requires). 정본 = app/render/wireframe.py.
    """

    DECK = {"slides": [
        {"slide_id": "1", "title": "a", "fields": {"x": ["v"]}},
        {"slide_id": "2", "title": "b", "fields": {"x": ["v"]}},
    ]}

    @staticmethod
    def _wf(slides, selected_by="llm:test"):
        return {"schema_version": 1, "selected_by": selected_by, "slides": slides}

    def test_unknown_frame_or_piece_is_contract_error(self):
        """① 미정의 frame/piece·slide_id·selected_by 부재 = 오류(지어내기 차단, catalog_gap으로 유도)."""
        import wireframe
        wf = self._wf([
            {"slide_id": "1", "frame": "no_such_frame", "slots": [{"piece": "text_block", "binds": "x"}]},
            {"slide_id": "2", "frame": "full", "slots": [{"piece": "no_such_piece", "data": {"k": 1}}]},
            {"slide_id": "99", "frame": "full", "slots": [{"piece": "text_block", "binds": "x"}]},
        ], selected_by="")
        v = wireframe.validate(wf, self.DECK)
        joined = "\n".join(v["errors"])
        self.assertIn("no_such_frame", joined)
        self.assertIn("no_such_piece", joined)
        self.assertIn("99", joined)
        self.assertIn("selected_by", joined)

    def test_r9_variation_without_reason_is_surfaced(self):
        """② 같은 layout_group 내 frame 변주에 variation_reason 없음 → R9 경고(무사유 변주=AI티 신호)."""
        import wireframe
        wf = self._wf([
            {"slide_id": "1", "frame": "full", "layout_group": "g",
             "slots": [{"piece": "text_block", "binds": "x"}]},
            {"slide_id": "2", "frame": "split_v", "layout_group": "g",
             "slots": [{"piece": "text_block", "binds": "x"}]},
        ])
        v = wireframe.validate(wf, self.DECK)
        self.assertEqual([], v["errors"])
        self.assertTrue(any("R9" in w for w in v["warnings"]), f"R9 경고 없음: {v['warnings']}")
        # 사유를 선언하면 경고가 사라진다(의도의 기계 가시화).
        wf["slides"][1]["variation_reason"] = "전환 강조"
        v2 = wireframe.validate(wf, self.DECK)
        self.assertFalse(any("R9" in w for w in v2["warnings"]))

    def test_merge_and_gating_block(self):
        """③ 병합=frame/slots 이식+selected_by 각인, 게이트 블록=조합 통계·갭 표면화 키."""
        import copy
        import wireframe
        deck = copy.deepcopy(self.DECK)
        wf = self._wf([
            {"slide_id": "1", "message_type": "구조", "frame": "row_n", "layout_group": "g",
             "slots": [{"piece": "text_block", "binds": "x"}], "principles": ["3-2"],
             "knowledge_cards": ["카드A", "카드A"],
             "catalog_gap": ["chevron_flow"]},
        ])
        v = wireframe.validate(wf, deck)
        self.assertEqual([], v["errors"])
        self.assertEqual([{"slide_id": "1", "wanted": "chevron_flow"}], v["catalog_gap"])
        applied = wireframe.merge_into_deck(deck, wf)
        self.assertEqual(1, applied)
        s1 = deck["slides"][0]
        self.assertEqual("row_n", s1["frame"])
        self.assertEqual("llm:test", s1["wireframe_selected_by"])
        block = wireframe.gating_block(wf, v)
        for key in ("selected_by", "stats", "rule_warnings", "catalog_gap",
                    "applied_knowledge", "slides"):
            self.assertIn(key, block)
        self.assertEqual(["text_block"], block["slides"][0]["pieces"])
        self.assertEqual({
            "cards": ["카드A"],
            "slides_with_cards": 1,
            "slides_total": 1,
        }, block["applied_knowledge"])
        wf["slides"][0]["knowledge_cards"] = "카드A"
        v_bad = wireframe.validate(wf, deck)
        self.assertTrue(any("knowledge_cards" in e for e in v_bad["errors"]))


class RefineDesignSpecTests(unittest.TestCase):
    """W23 ④+ 디자인 고도화 목표 명세 계약 — design_spec.json 검증·형태 레퍼런스 수집·핸드오프.

    정본 = app/render/design_spec.py(wireframe.py의 자매). 오류=계약 위반(적용 중단) /
    경고=표면화 — wireframe과 동일 문법(결정 15·16·17).
    """

    DECK = {"slides": [
        {"slide_id": "1", "title": "a", "role": "전략", "key_message": "m1"},
        {"slide_id": "2", "title": "b", "role": "근거", "key_message": "m2"},
    ]}

    @staticmethod
    def _spec(slides, catalog_gap=None):
        return {"schema_version": 1, "run_id": "test", "generated_by": "llm:test",
                "slides": slides, "catalog_gap": catalog_gap or []}

    def test_missing_goal_is_contract_error(self):
        """① goal 없는 명세 = 계약 위반(의도 먼저)."""
        import design_spec
        spec = self._spec([
            {"slide_id": "1", "goal": "", "treatment": ["photo"], "image_kind": "mood"},
            {"slide_id": "2", "goal": "장표2 목표", "treatment": ["diagram"], "image_kind": "none"},
        ])
        v = design_spec.validate(spec, self.DECK)
        joined = "\n".join(v["errors"])
        self.assertTrue(("의도" in joined) or ("goal" in joined), f"goal 누락 오류 없음: {v['errors']}")

    def test_unknown_form_needs_piece_is_catalog_gap_error(self):
        """② form_needs에 미지 piece id → 오류에 catalog_gap 문구(지어내지 말라)."""
        import design_spec
        spec = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": ["diagram"], "image_kind": "conceptual",
             "form_needs": [{"kind": "piece", "id": "no_such_piece_xyz", "why": "형태 서술"}]},
            {"slide_id": "2", "goal": "목표2", "treatment": ["photo"], "image_kind": "none"},
        ])
        v = design_spec.validate(spec, self.DECK)
        joined = "\n".join(v["errors"])
        self.assertIn("no_such_piece_xyz", joined)
        self.assertIn("catalog_gap", joined)

    def test_image_kind_vocab_and_content_gap_surfaced(self):
        """③ image_kind 어휘 밖 = 오류 / content_gap 문자열 = 경고에 "내용 루프" 포함."""
        import design_spec
        spec = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": ["photo"], "image_kind": "portrait"},
            {"slide_id": "2", "goal": "목표2", "treatment": ["diagram"], "image_kind": "none",
             "content_gap": "근거 자산이 더 필요하다"},
        ])
        v = design_spec.validate(spec, self.DECK)
        joined_err = "\n".join(v["errors"])
        self.assertIn("portrait", joined_err)
        joined_warn = "\n".join(v["warnings"])
        self.assertIn("내용 루프", joined_warn)

    def test_collect_refs_copies_pieces_and_marks_frame_file_null(self):
        """④ collect_refs: 존재 piece 2종 + frame 1종 → design_refs/에 png 2장, frame은 file=null."""
        import design_spec
        spec = self._spec([
            {"slide_id": "1", "goal": "현위치 대비 목표 시각화", "treatment": ["diagram"],
             "image_kind": "conceptual",
             "form_needs": [
                 {"kind": "piece", "id": "matrix_2x2", "why": "현위치 대비 목표"},
                 {"kind": "piece", "id": "flow_arrow", "why": "단계 흐름"},
                 {"kind": "frame", "id": "hero_body", "why": "대형 강조"},
             ]},
        ])
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            manifest = design_spec.collect_refs(run, spec)
            refs_dir = run / "design_refs"
            self.assertTrue((refs_dir / "matrix_2x2.png").is_file())
            self.assertTrue((refs_dir / "flow_arrow.png").is_file())
            self.assertTrue((refs_dir / "refs_manifest.json").is_file())
            items = manifest["per_slide"]["1"]
            by_id = {it["id"]: it for it in items}
            self.assertEqual("design_refs/matrix_2x2.png", by_id["matrix_2x2"]["file"])
            self.assertEqual("design_refs/flow_arrow.png", by_id["flow_arrow"]["file"])
            self.assertIsNone(by_id["hero_body"]["file"])
            self.assertFalse(by_id["hero_body"]["gap"])

    def test_build_prompt_includes_real_ids_and_catalog_gap(self):
        """⑤ build_prompt: 어휘 표에 실제 frame/piece id 포함 + catalog_gap 문구 포함."""
        import design_spec
        prompt = design_spec.build_prompt(Path("."), self.DECK)
        self.assertIn("matrix_2x2", prompt)
        self.assertIn("hero_body", prompt)
        self.assertIn("catalog_gap", prompt)

    def test_build_handoff_contract_phrases_and_stage_order(self):
        """⑥ build_handoff: 계약 문구(불변·생성 금지·approve --ingest·stage9 --apply) 포함 +
        STAGE_ORDER에서 refine_bundle이 stage9_apply와 deck_review_bundle 사이."""
        import design_spec
        spec = self._spec([{"slide_id": "1", "goal": "목표", "treatment": ["photo"], "image_kind": "mood"}])
        manifest = {"per_slide": {}, "catalog_gap": []}
        handoff = design_spec.build_handoff(Path("."), spec, manifest)
        self.assertIn("불변", handoff)
        self.assertIn("생성 금지", handoff)
        self.assertIn("approve --ingest", handoff)
        self.assertIn("stage9 --apply", handoff)

        order = pipeline_state.STAGE_ORDER
        self.assertLess(order.index("stage9_apply"), order.index("refine_bundle"))
        self.assertLess(order.index("refine_bundle"), order.index("deck_review_bundle"))

    def test_none_without_reason_is_contract_error_reason_present_passes(self):
        """⑦ W27 P2 D5: image_kind=none인데 none_reason 없음 → 오류 / 있으면 통과."""
        import design_spec
        spec_no_reason = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": [], "image_kind": "none"},
            {"slide_id": "2", "goal": "목표2", "treatment": ["photo"], "image_kind": "mood"},
        ])
        v = design_spec.validate(spec_no_reason, self.DECK)
        joined = "\n".join(v["errors"])
        self.assertIn("none_reason", joined)

        spec_with_reason = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": [], "image_kind": "none",
             "none_reason": "텍스트 위계만으로 충분"},
            {"slide_id": "2", "goal": "목표2", "treatment": ["photo"], "image_kind": "mood"},
        ])
        v2 = design_spec.validate(spec_with_reason, self.DECK)
        self.assertEqual([], v2["errors"])

    def test_source_route_defaults_filled_and_bad_route_is_error(self):
        """⑧ W27 P2 D5: source_route 미지정 → 기본값(mood/conceptual→codex_gen,
        evidence→user_asset) 채움. 불량 값은 오류."""
        import design_spec
        spec = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": ["photo"], "image_kind": "mood"},
            {"slide_id": "2", "goal": "목표2", "treatment": ["diagram"], "image_kind": "evidence"},
        ])
        v = design_spec.validate(spec, self.DECK)
        self.assertEqual([], v["errors"])
        by_id = {s["slide_id"]: s for s in spec["slides"]}
        self.assertEqual("codex_gen", by_id["1"]["source_route"])
        self.assertEqual("user_asset", by_id["2"]["source_route"])

        bad_spec = self._spec([
            {"slide_id": "1", "goal": "목표", "treatment": ["photo"], "image_kind": "mood",
             "source_route": "no_such_route"},
        ])
        v_bad = design_spec.validate(bad_spec, self.DECK)
        joined = "\n".join(v_bad["errors"])
        self.assertIn("no_such_route", joined)
        self.assertIn("source_route", joined)


class PresetParitySmoke(unittest.TestCase):
    """W21-0 프리셋 6종 동급 참여 + 조각 병렬 리스트 어댑터(§6 이행 표, 결정 12 후속).

    preset = 이름 붙은 frame×piece 조합일 뿐 특권 없음 — compose.render_slide가 frame이 없고
    preset만 있을 때 presets.json에서 채워 넣는다. 어댑터는 조각 함수 내부 데이터 정규화로
    구현(새 조각 추가 없음 — ComposeEngineSmoke⑥이 계약-구현 완전성을 지킨다).
    """

    @staticmethod
    def _deck(slides):
        return {"meta": {"project": "preset smoke"}, "slides": slides}

    def _render(self, slides, td):
        out = Path(td) / "deck.html"
        rep = htmlgen.render_html(self._deck(slides), "core", out)
        return rep, out.read_text(encoding="utf-8")

    def test_preset_slide_renders_via_compose(self):
        """① frame 없이 preset만 있어도 조합 엔진으로 렌더된다(frame/slots를 presets.json에서 채움)."""
        slides = [{"slide_id": "1", "title": "t", "preset": "executive_summary",
                   "fields": {"main_claim": "주장", "supporting_points": ["a", "b"]}}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertIn("compose--hero_body", html)
        self.assertIn("주장", html)

    def test_org_roles_preset_parallel_list_adapter(self):
        """② p_org_table 어댑터: teams+roles(병렬 문자열 리스트) → {team, role} 3행. lead는 기존대로."""
        slides = [{"slide_id": "1", "title": "t", "preset": "org_roles", "fields": {
            "lead": "김PM",
            "teams": ["기획", "개발", "운영"],
            "roles": ["전략수립", "구현", "안정화"],
        }}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertEqual(3, html.count('<tr><th scope="row">'))
        self.assertIn("김PM", html)

    def test_org_table_nested_team_dicts_render_names_not_repr(self):
        """p_org_table: teams가 dict({name, roles[...]})이고 lead가 dict인 정본 형상 —
        dict repr('{')가 화면에 새지 않고 name·하위 roles·설명이 렌더된다. 팀 수 > roles 수여도
        min() 절단 없이 전 팀 표시(조용한 표시 소실 금지 — 5차 run slide 11 실측 결함)."""
        slides = [{"slide_id": "1", "title": "t", "preset": "org_roles", "fields": {
            "lead": {"name": "총괄 책임자", "description": "단일 창구"},
            "teams": [
                {"name": "상주 운영", "roles": ["선임 1명", "상담원 1명"]},
                {"name": "홍보 지원", "roles": ["콘텐츠 제작"]},
                {"name": "본사 지원", "roles": ["백업 3명"]},
            ],
            "roles": ["감독·지휘 책임", "투명성 보고"],
        }}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertNotIn("{&#x27;", html)   # dict repr 누출 없음(이스케이프된 형태 포함)
        self.assertNotIn("{'", html)
        self.assertEqual(3, html.count('<tr><th scope="row">'))  # 3팀 전부(절단 없음)
        for token in ("총괄 책임자", "단일 창구", "상주 운영", "본사 지원", "선임 1명"):
            self.assertIn(token, html)

    def test_case_card_parallel_lists_no_truncation(self):
        """p_case_card: cases 3 vs metrics·names 2(병렬 부족) — min() 절단으로 3번째 사례가
        조용히 사라지지 않고 전 건 표시된다(5차 run slide 2 실측 결함, org_table과 동일 원칙)."""
        slides = [{"slide_id": "1", "title": "t", "preset": "portfolio_cases", "fields": {
            "cases": ["A 위탁", "B 위탁", "C기관 상담 인력 위탁"],
            "metrics": ["m1", "m2"],
            "client_safe_names": ["A사", "B사"],
        }}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertEqual(3, html.count('<div class="cp-case">'))
        self.assertIn("C기관 상담 인력 위탁", html)

    def test_case_card_preset_parallel_list_adapter(self):
        """② p_case_card 어댑터: cases(문자열)+metrics+client_safe_names(병렬) → 카드 3장. dict 리스트 경로는 불변(⑥이 고정)."""
        slides = [{"slide_id": "1", "title": "t", "preset": "portfolio_cases", "fields": {
            "cases": ["A사 물류 자동화", "B사 CS 챗봇", "C사 재고 최적화"],
            "metrics": ["처리시간 40%↓", "응답시간 3배↑", "재고비용 15%↓"],
            "client_safe_names": ["A사(가명)", "B사(가명)", "C사(가명)"],
        }}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual([], rep.get("warnings") or [])
        self.assertEqual(3, html.count('class="cp-case"'))

    def test_roadmap_gantt_preset_parallel_list_adapter(self):
        """③+⑤-a p_timeline_gantt 어댑터: time_units+workstreams(cells) → period_labels/tasks 정규화.

        cells=[True, True, None] → start=1·end=2 → 3기간 중 2칸 점등. milestones는 이 어댑터 범위
        밖(무시). requires 정직성 검사는 조각 내부로 이관(requires=[]) — 병렬 입력에
        '필수 필드 누락' 오탐이 없어야 한다(⑤-a).
        """
        slides = [{"slide_id": "1", "title": "t", "preset": "roadmap_gantt", "fields": {
            "time_units": ["1월", "2월", "3월"],
            "workstreams": [{"label": "설계", "cells": [True, True, None]}],
        }}]
        with tempfile.TemporaryDirectory() as td:
            rep, html = self._render(slides, td)
        self.assertEqual(2, html.count("cp-gantt__on"))
        self.assertNotIn("필수 필드 누락", html)                    # ⑤-a: requires 오탐 없음
        self.assertEqual([], rep.get("warnings") or [])

    def test_roadmap_gantt_empty_input_surfaces_review(self):
        """⑤-b 무입력 gantt: 두 입력 형태 모두 없으면 빈 표 대신 검토요망 표면화(감춤 금지)."""
        slides = [{"slide_id": "1", "title": "t", "preset": "roadmap_gantt",
                   "fields": {"note": "일정 데이터 없음"}}]
        with tempfile.TemporaryDirectory() as td:
            _rep, html = self._render(slides, td)
        self.assertIn("검토요망", html)
        self.assertIn("timeline_gantt 입력 없음", html)
        self.assertNotIn("cp-gantt", html)                         # 빈 표를 조용히 그리지 않는다

    def test_wireframe_validate_accepts_preset_key(self):
        """④ wireframe.validate: 슬라이드 항목의 frame 대신 preset — 정의돼 있으면 통과, 없으면 오류."""
        import wireframe
        deck = {"slides": [{"slide_id": "1", "title": "a", "fields": {"lead": "PM"}}]}
        wf_ok = {"schema_version": 1, "selected_by": "llm:test",
                 "slides": [{"slide_id": "1", "preset": "org_roles"}]}
        v_ok = wireframe.validate(wf_ok, deck)
        self.assertEqual([], v_ok["errors"])
        self.assertIn("org_roles", v_ok["stats"]["presets_used"])

        wf_bad = {"schema_version": 1, "selected_by": "llm:test",
                  "slides": [{"slide_id": "1", "preset": "no_such_preset"}]}
        v_bad = wireframe.validate(wf_bad, deck)
        self.assertTrue(any("no_such_preset" in e for e in v_bad["errors"]), v_bad["errors"])


class BrandChromeTests(unittest.TestCase):
    """W22(통찰 14): 브랜드 크롬 — 로고=실자산 필수(생성 금지), 무브랜드 덱 바이트 불변,
    스킨 캐스케이드 해석 우선순위(CLI > design_brief.skin.skins > state rendered.skins)."""

    _PNG_1PX = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="

    @staticmethod
    def _deck():
        return {"meta": {"project": "P"}, "slides": [
            {"slide_id": 1, "role": "cover", "template_id": "cover",
             "title": "표지", "key_message": "핵심", "fields": {}},
            {"slide_id": 2, "role": "summary", "template_id": "card_grid",
             "title": "요약", "key_message": "메시지",
             "fields": {"main_claim": "주장", "supporting_points": ["a", "b"]}},
        ]}

    def _render(self, td: str, brand: dict | None):
        out = Path(td) / "deck.html"
        skins = [{"brand": brand}] if brand else None
        rep = htmlgen.render(self._deck(), catalog="core", skins=skins, out_path=out)
        return rep, out.read_text(encoding="utf-8")

    def test_brand_renders_with_real_asset(self):
        """① 실자산 로고 → 제안사 마크 전 장 + 클라이언트 마크 cover만 + data URI 임베드."""
        import base64 as _b64
        with tempfile.TemporaryDirectory() as td:
            logo = Path(td) / "logo.png"
            logo.write_bytes(_b64.b64decode(self._PNG_1PX))
            brand = {"client_name": "발주기관", "client_logo": str(logo),
                     "proposer_name": "제안사", "proposer_logo": str(logo)}
            rep, html = self._render(td, brand)
        # 마크업 마커로 센다(CSS(_BRAND_CSS)에도 클래스명이 등장하므로 bare 문자열 카운트 금지).
        proposer_mark = '<div class="slide__brand slide__brand--proposer">'
        client_mark = '<div class="slide__brand slide__brand--client">'
        self.assertEqual(2, html.count(proposer_mark), "제안사 마크는 전 장(2장)")
        # 클라이언트 마크는 cover 섹션에만: slide-1(cover)엔 있고 slide-2엔 없다.
        s1 = html.split('id="slide-1"')[1].split('id="slide-2"')[0]
        s2 = html.split('id="slide-2"')[1]
        self.assertIn(client_mark, s1)
        self.assertNotIn(client_mark, s2)
        self.assertIn("data:image/png;base64", html)
        self.assertEqual([], rep.get("warnings") or [], rep.get("warnings"))

    def test_brand_missing_logo_warns_not_fabricates(self):
        """② 선언된 로고 파일 없음 → 실자산 경고 + img 없음 + 이름 텍스트만(placeholder 금지)."""
        with tempfile.TemporaryDirectory() as td:
            brand = {"client_name": "예시여대", "client_logo": str(Path(td) / "없는로고.png"),
                     "proposer_name": None, "proposer_logo": None}
            rep, html = self._render(td, brand)
        self.assertTrue(any("실자산" in w for w in rep.get("warnings") or []), rep.get("warnings"))
        mark = html.split('<div class="slide__brand slide__brand--client">')[1].split("</div>")[0]
        self.assertNotIn("<img", mark)
        self.assertIn("예시여대", mark)

    def test_render_without_brand_byte_stable(self):
        """③ brand 없는 렌더에는 브랜드 마크업·CSS가 한 바이트도 없다."""
        with tempfile.TemporaryDirectory() as td:
            _rep, html = self._render(td, None)
        self.assertEqual(0, html.count("slide__brand"))

    def test_design_skins_resolution_priority(self):
        """④ _design_skins: CLI > design_brief.skin.skins, brand는 인라인 스킨으로 최후승."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            brief = {
                "schema_version": 1,
                "skin": {"pack": "core", "skins": ["quartz"]},
                "brand": {"client_name": "예시여대", "client_logo": None,
                          "proposer_name": None, "proposer_logo": None,
                          "placement": {"client": "cover", "proposer": "all"}},
            }
            (run / "design_brief.json").write_text(
                json.dumps(brief, ensure_ascii=False, indent=2), encoding="utf-8")

            skins = proposal_pipeline._design_skins(run, SimpleNamespace(skins=None))
            self.assertEqual("quartz", skins[0])
            self.assertIsInstance(skins[-1], dict)
            self.assertEqual("예시여대", skins[-1]["brand"]["client_name"])

            skins_cli = proposal_pipeline._design_skins(run, SimpleNamespace(skins="lecture-dark"))
            self.assertEqual("lecture-dark", skins_cli[0])
            self.assertIsInstance(skins_cli[-1], dict)
            self.assertEqual("예시여대", skins_cli[-1]["brand"]["client_name"])


class W24ProgressBarSmoke(unittest.TestCase):
    """W24(결정 11 ①③ + 목표조정 6b): you-are-here 진행 바 — 결정론 위치 판정 + status 상단 노출."""

    @staticmethod
    def _cp(decision: bool = False, design: bool = False):
        return {
            "start": {"cleared_at": "2026-07-13T00:00:00"},
            "decision": {"cleared_at": "2026-07-13T00:00:01" if decision else None},
            "design": {"cleared_at": "2026-07-13T00:00:02" if design else None},
        }

    def test_progress_position_deterministic_mapping(self):
        pos = pipeline_state.progress_position
        st = lambda *names: {n: {"at": "2026-07-13T00:00:00"} for n in names}
        # [1] 기록 없음 / 내용 만들기 중(render·message_map)
        self.assertEqual("1", pos({}, self._cp()))
        self.assertEqual("1", pos(st("start", "message_map", "render"), self._cp()))
        # [3] decision 청산 후 다음 미기록 / wireframe_* 기록
        self.assertEqual("3", pos(st("render"), self._cp(decision=True)))
        self.assertEqual("3", pos(st("render", "wireframe_bundle"), self._cp(decision=True)))
        # [4] stage9·design_brief 계열
        self.assertEqual("4", pos(st("render", "design_brief"), self._cp(decision=True)))
        self.assertEqual("4", pos(st("render", "stage9_apply"), self._cp(decision=True)))
        # [4+] refine_* 있고 approve 없음
        self.assertEqual("4+", pos(st("render", "stage9_apply", "refine_collect"), self._cp(decision=True)))
        # [5] approve 또는 design 체크포인트 청산
        self.assertEqual("5", pos(st("render", "refine_collect", "approve"), self._cp(decision=True)))
        self.assertEqual("5", pos(st("render"), self._cp(decision=True, design=True)))

    def test_format_status_starts_with_progress_bar(self):
        view = {"run_id": "r", "mode": None, "input": None, "has_state_file": False,
                "stages": {"stage9_bundle": {"at": "2026-07-13T00:00:00", "source": "recorded"}},
                "checkpoints": self._cp(decision=True), "message_map": None, "warnings": [],
                "next": {"kind": "command", "why": "w", "command": "c"}}
        out = pipeline_state.format_status(view)
        lines = out.splitlines()
        self.assertTrue(lines[0].startswith("공정 지도: [1]내용 만들기"))
        self.assertEqual("현재: ▶ [4] 디자인 입히기", lines[1])
        # resolve()가 계산한 progress가 있으면 그것을 우선한다
        view["progress"] = {"position": "4+", "label": "[4+] 디자인 고도화"}
        self.assertIn("현재: ▶ [4+] 디자인 고도화", pipeline_state.format_status(view))


class W25AnalysisDigestLocationSmoke(unittest.TestCase):
    """W25(NORTHSTAR 결정 13): 분석카드/프롬프트/digest를 workspace/로 통합.

    쓰기=새 위치만, 읽기=새 위치 우선 + 레거시(vendor) 폴백, start의 run/analysis 복제.
    """

    def test_save_card_writes_to_workspace_analysis(self):
        with tempfile.TemporaryDirectory() as td:
            analysis = Path(td) / "analysis"
            analyzer = SimpleNamespace(os=os, ANALYSIS_DIR=str(analysis))
            with mock.patch.object(server, "_load_analyzer", return_value=analyzer):
                saved = server.save_card("W25-CARD", "# 분석카드\nW25 내용")
            saved_path = Path(saved["path"])
            self.assertTrue(saved_path.is_file())
            self.assertEqual(analysis.resolve(), saved_path.parent.resolve())

    def test_find_analysis_file_falls_back_to_legacy(self):
        with tempfile.TemporaryDirectory() as td:
            new_dir = Path(td) / "new_analysis"
            legacy_dir = Path(td) / "legacy_analysis"
            new_dir.mkdir()
            legacy_dir.mkdir()
            legacy_card = legacy_dir / "W25-LEGACY_분석카드.md"
            legacy_card.write_text("# 레거시 카드", encoding="utf-8")
            with (
                mock.patch.object(server, "ANALYSIS", new_dir),
                mock.patch.object(server, "ANALYSIS_LEGACY", legacy_dir),
            ):
                found = server._find_analysis_file("W25-LEGACY_분석카드.md")
                self.assertEqual(legacy_card.resolve(), found.resolve())
                # 새 위치에 있으면 새 위치를 우선한다.
                new_card = new_dir / "W25-BOTH_분석카드.md"
                new_card.write_text("# 새 카드", encoding="utf-8")
                (legacy_dir / "W25-BOTH_분석카드.md").write_text("# 레거시 카드2", encoding="utf-8")
                found2 = server._find_analysis_file("W25-BOTH_분석카드.md")
                self.assertEqual(new_card.resolve(), found2.resolve())

    def test_start_copies_bid_analysis_assets_into_run(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            analysis_dir = root / "analysis"
            analysis_dir.mkdir()
            bid = "W25-START-001"
            (analysis_dir / f"{bid}_분석카드.md").write_text("# 카드", encoding="utf-8")
            (analysis_dir / f"{bid}_프롬프트.txt").write_text("프롬프트 본문", encoding="utf-8")
            run_dir = root / "run"
            run_dir.mkdir()
            with (
                mock.patch.object(proposal_pipeline, "ANALYSIS_DIR", analysis_dir),
                mock.patch.object(proposal_pipeline, "ANALYSIS_DIR_LEGACY", root / "no_such_legacy"),
            ):
                copied = proposal_pipeline._copy_bid_analysis_assets(run_dir, bid)
            self.assertEqual(2, len(copied))
            self.assertTrue((run_dir / "analysis" / f"{bid}_분석카드.md").is_file())
            self.assertTrue((run_dir / "analysis" / f"{bid}_프롬프트.txt").is_file())


class W26InstitutionResearchSmoke(unittest.TestCase):
    """W26(목표조정 8·9): 기관 공개 조사(문서 밖 근거) 검증·브랜드→스킨 변환·번들."""

    def test_validate_reports_missing_source_bad_hex_and_hook_source_warning(self):
        import institution_research as ir
        res = {
            "institution": "",
            "sources": [],
            "content_hooks": [{"claim": "출처 없는 훅", "use_in": "도입 직인용"}],
            "brand_tokens": {"colors": {"primary": "notahex"}},
        }
        v = ir.validate(res)
        joined_err = "\n".join(v["errors"])
        self.assertIn("institution", joined_err)
        self.assertIn("sources", joined_err)
        self.assertIn("hex", joined_err)
        self.assertTrue(any("source 없음" in w for w in v["warnings"]))

    def test_validate_passes_with_complete_hex_and_sourced_hook(self):
        import institution_research as ir
        res = {
            "institution": "테스트대학교",
            "sources": ["https://example.org/about"],
            "content_hooks": [{"claim": "훅", "use_in": "도입 직인용", "source": "https://example.org/news"}],
            "brand_tokens": {"colors": {"primary": "#123456"}},
        }
        v = ir.validate(res)
        self.assertEqual([], v["errors"])
        self.assertEqual([], v["warnings"])

    def test_to_skin_maps_primary_accent_and_brand(self):
        import institution_research as ir
        res = {
            "institution": "테스트대학교",
            "sources": ["https://example.org/about"],
            "brand_tokens": {
                "colors": {"primary": "#123456", "accent": "#ABCDEF"},
                "fonts": {"family": "Pretendard"},
                "logo": {"path": "logo.png"},
            },
        }
        skin = ir.to_skin(res, "test_inst")
        self.assertEqual("123456", skin["colors"]["navy"])
        self.assertEqual("ABCDEF", skin["colors"]["orange"])
        self.assertEqual("Pretendard", skin["fonts"]["family"])
        self.assertEqual("테스트대학교", skin["brand"]["client_name"])
        self.assertEqual("logo.png", skin["brand"]["client_logo"])
        self.assertEqual("test_inst", skin["_meta"]["name"])
        self.assertEqual("https://example.org/about", skin["_meta"]["provenance"])
        self.assertFalse(skin["_meta"]["self_contained"])

    def test_bundle_prompt_mentions_citation_and_source(self):
        import institution_research as ir
        prompt = ir.build_prompt("some_run", institution="테스트기관")
        self.assertIn("직인용", prompt)
        self.assertIn("출처", prompt)

    def test_apply_registers_skin_and_updates_design_brief(self):
        skin_id = f"w26_test_{uuid.uuid4().hex[:8]}"
        skin_path = proposal_pipeline.SKINS_DIR / f"{skin_id}.json"
        self.assertFalse(skin_path.exists(), "테스트 스킨 id 충돌 — 기존 파일이 이미 있음")
        try:
            with tempfile.TemporaryDirectory() as td:
                run = Path(td) / "run"
                run.mkdir()
                (run / "institution_research.json").write_text(
                    json.dumps({
                        "institution": "가상대학교",
                        "sources": ["https://example.org/about"],
                        "brand_tokens": {"colors": {"primary": "#123456"}},
                        "knowledge_used": {"cards": [], "web": []},  # ε패킷 안전장치①
                    }, ensure_ascii=False),
                    encoding="utf-8",
                )
                (run / "design_brief.json").write_text(
                    json.dumps({
                        "schema_version": 1,
                        "skin": {"pack": "core", "skins": []},
                        "brand": {"client_name": None, "client_logo": None},
                    }, ensure_ascii=False),
                    encoding="utf-8",
                )
                result = proposal_pipeline._research_apply(
                    run, SimpleNamespace(skin_id=skin_id),
                )
                self.assertEqual(skin_id, result["skin_id"])
                self.assertTrue(skin_path.is_file())
                skin = json.loads(skin_path.read_text(encoding="utf-8"))
                self.assertEqual("123456", skin["colors"]["navy"])

                brief = design_brief.load(run)
                self.assertIn(skin_id, brief["skin"]["skins"])
                self.assertEqual("가상대학교", brief["brand"]["client_name"])

                state = pipeline_state.load(run)
                self.assertIn("research_apply", state["stages"])
        finally:
            if skin_path.exists():
                skin_path.unlink()
            self.assertFalse(skin_path.exists())  # 스킨 레지스트리 오염 0 보장

    # ---- W31 마찰14: skin.value(계약 차용 소스) 승계 -----------------------------------------

    def test_apply_fills_empty_skin_value_but_preserves_existing(self):
        """빈 skin.value는 등록 스킨으로 채우고(승계), 이미 값이 있으면 건드리지 않는다."""
        skin_id = f"w31_f14_{uuid.uuid4().hex[:8]}"
        skin_path = proposal_pipeline.SKINS_DIR / f"{skin_id}.json"
        self.assertFalse(skin_path.exists())
        try:
            with tempfile.TemporaryDirectory() as td:
                run = Path(td) / "run"
                run.mkdir()
                (run / "institution_research.json").write_text(
                    json.dumps({
                        "institution": "가상기관",
                        "sources": ["https://example.org/about"],
                        "brand_tokens": {"colors": {"primary": "#654321"}},
                        "knowledge_used": {"cards": [], "web": []},  # ε패킷 안전장치①
                    }, ensure_ascii=False), encoding="utf-8")
                (run / "design_brief.json").write_text(
                    json.dumps({"schema_version": 1, "skin": {"pack": "core", "skins": []},
                                "brand": {}}, ensure_ascii=False), encoding="utf-8")
                proposal_pipeline._research_apply(run, SimpleNamespace(skin_id=skin_id))
                brief = design_brief.load(run)
                self.assertEqual(skin_id, brief["skin"]["value"])  # 빈 값 -> 승계

            # 이미 값이 있으면 보존(사용자 결정 우선).
            skin_id2 = f"w31_f14b_{uuid.uuid4().hex[:8]}"
            skin_path2 = proposal_pipeline.SKINS_DIR / f"{skin_id2}.json"
            with tempfile.TemporaryDirectory() as td2:
                run2 = Path(td2) / "run"
                run2.mkdir()
                (run2 / "institution_research.json").write_text(
                    json.dumps({
                        "institution": "가상기관2",
                        "sources": ["https://example.org/about"],
                        "brand_tokens": {"colors": {"primary": "#111111"}},
                        "knowledge_used": {"cards": [], "web": []},  # ε패킷 안전장치①
                    }, ensure_ascii=False), encoding="utf-8")
                (run2 / "design_brief.json").write_text(
                    json.dumps({"schema_version": 1,
                                "skin": {"value": "inkline", "pack": "core", "skins": []},
                                "brand": {}}, ensure_ascii=False), encoding="utf-8")
                proposal_pipeline._research_apply(run2, SimpleNamespace(skin_id=skin_id2))
                brief2 = design_brief.load(run2)
                self.assertEqual("inkline", brief2["skin"]["value"])  # 기존값 보존
                self.assertIn(skin_id2, brief2["skin"]["skins"])  # skins(렌더 체인)에는 여전히 추가
            if skin_path2.exists():
                skin_path2.unlink()
        finally:
            if skin_path.exists():
                skin_path.unlink()

    def test_apply_without_brief_records_applied_skin_for_later_lookup(self):
        """design_brief.json이 아직 없을 때: 등록 스킨을 institution_research.json에 남긴다
        (design_brief.build_default이 나중에 조회할 수 있게, 마찰14 ②)."""
        skin_id = f"w31_f14c_{uuid.uuid4().hex[:8]}"
        skin_path = proposal_pipeline.SKINS_DIR / f"{skin_id}.json"
        try:
            with tempfile.TemporaryDirectory() as td:
                run = Path(td) / "run"
                run.mkdir()
                (run / "institution_research.json").write_text(
                    json.dumps({
                        "institution": "가상기관3",
                        "sources": ["https://example.org/about"],
                        "brand_tokens": {"colors": {"primary": "#222222"}},
                        "knowledge_used": {"cards": [], "web": []},  # ε패킷 안전장치①
                    }, ensure_ascii=False), encoding="utf-8")
                proposal_pipeline._research_apply(run, SimpleNamespace(skin_id=skin_id))
                res = json.loads((run / "institution_research.json").read_text(encoding="utf-8"))
                self.assertEqual(skin_id, res["_applied_skin"]["skin_id"])

                # design_brief 기본값 생성이 이 등록 스킨을 skin.value 초안으로 조회한다.
                deck = {"slides": [{"slide_id": 1, "role": "표지", "title": "표지",
                                     "body": [], "key_message": "", "fields": {}}]}
                brief = design_brief.build_default(run, deck, skins_dir=proposal_pipeline.SKINS_DIR)
                self.assertEqual(skin_id, brief["skin"]["value"])
        finally:
            if skin_path.exists():
                skin_path.unlink()

    def test_build_default_ignores_applied_skin_if_file_missing(self):
        """등록 기록은 있는데 실제 skins/<id>.json이 없으면(지워짐 등) 초안으로 쓰지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            (run / "institution_research.json").write_text(
                json.dumps({"institution": "가상기관4", "_applied_skin": {"skin_id": "no-such-skin-xyz"}},
                           ensure_ascii=False), encoding="utf-8")
            deck = {"slides": [{"slide_id": 1, "role": "표지", "title": "표지",
                                 "body": [], "key_message": "", "fields": {}}]}
            brief = design_brief.build_default(run, deck, skins_dir=proposal_pipeline.SKINS_DIR)
            self.assertIsNone(brief["skin"]["value"])


class KnowledgeGapsStatusSmoke(unittest.TestCase):
    """어휘 갭 로그 처리상태 갱신 — find/update_status(멱등·사람 편집 보존·매칭 정밀)."""

    GAPS = [
        {"slide_id": "3", "wanted": "pillar_card 부재 — 병렬 주장"},
        {"slide_id": "5", "wanted": "pillar_card 부재 — 동일"},
    ]

    def test_update_status_only_matched_slide(self):
        """① --gap slide_id는 그 slide만 신설함으로 갱신, 나머지는 미해결 보존."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_gaps.record(run, self.GAPS, "wireframe")
            res = knowledge_gaps.update_status(run, "신설함", slide_id="3")
            self.assertEqual(1, res["updated"])
            data = json.loads((run / "knowledge_gaps.json").read_text(encoding="utf-8"))
            by_slide = {e["slide_id"]: e["status"] for e in data["entries"]}
            self.assertEqual("신설함", by_slide["3"])
            self.assertEqual(knowledge_gaps.DEFAULT_STATUS, by_slide["5"])

    def test_update_status_idempotent_and_no_match_is_noop(self):
        """② 재실행 멱등 + 매칭 0건이면 파일 무변경(조용한 신규 생성 없음)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_gaps.record(run, self.GAPS, "wireframe")
            knowledge_gaps.update_status(run, "신설함", slide_id="3")
            first = (run / "knowledge_gaps.json").read_text(encoding="utf-8")
            knowledge_gaps.update_status(run, "신설함", slide_id="3")
            self.assertEqual(first, (run / "knowledge_gaps.json").read_text(encoding="utf-8"))
            res = knowledge_gaps.update_status(run, "신설함", slide_id="999")
            self.assertEqual(0, res["updated"])

    def test_find_by_slide_and_need(self):
        """③ find는 slide_id·need로 회수(bundle --gap 경로). 없으면 빈 리스트."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_gaps.record(run, self.GAPS, "wireframe")
            hits = knowledge_gaps.find(run, slide_id="3")
            self.assertEqual(1, len(hits))
            self.assertIn("병렬 주장", hits[0]["need"])
            self.assertEqual([], knowledge_gaps.find(run, slide_id="42"))

    def test_update_requires_a_filter(self):
        """④ slide_id·need 둘 다 없으면 무엇을 고칠지 모르므로 ValueError(지어내지 않는다)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_gaps.record(run, self.GAPS, "wireframe")
            with self.assertRaises(ValueError):
                knowledge_gaps.update_status(run, "신설함")


class CurateSmoke(unittest.TestCase):
    """큐레이션 생애주기 — 라이브러리 스캔 · register(싱크백) · 참고자료 반입(파일+링크).

    지어내지 않는다: 없는 자산 register는 중단, 반입 스캔은 시스템 프리뷰·예시 주석을
    사용자 자료로 세지 않는다. 정본 = proposal_system/scripts/curate.py (DESIGN_ASSETS_LANE §5-④-③).
    """

    def test_scan_library_reads_real_skins_and_guides_no_fabrication(self):
        """① 라이브러리 스캔이 실제 skins/·config 가이드를 필수 키와 함께 반환(지어내지 않음)."""
        entries = curate.scan_library()
        ids = {e["id"] for e in entries}
        # 실제 스킨(현 repo) — 존재를 지어내지 않고 실측
        self.assertTrue({"quartz", "lecture-dark", "univ_sample"} <= ids, ids)
        for e in entries:
            for k in ("kind", "id", "source_path", "exists", "registered"):
                self.assertIn(k, e)
            self.assertIn(e["kind"], ("skin", "guide"))

    def test_register_copies_into_assets_and_is_idempotent(self):
        """② register = design-assets/로 복사 + 매니페스트 registered 반영, 재실행 멱등."""
        with tempfile.TemporaryDirectory() as td:
            base = Path(td)
            skins = base / "skins"
            skins.mkdir()
            (skins / "demo.json").write_text(
                json.dumps({"colors": {"navy": "112233"},
                            "_meta": {"name": "demo", "self_contained": True, "provenance": "unit"}},
                           ensure_ascii=False),
                encoding="utf-8",
            )
            assets = base / "design-assets"
            with (
                mock.patch.object(curate, "SKINS_DIR", skins),
                mock.patch.object(curate, "ASSETS_DIR", assets),
                mock.patch.object(curate, "ASSET_SKINS", assets / "skins"),
                mock.patch.object(curate, "ASSET_GUIDES", assets / "guides"),
                mock.patch.object(curate, "MANIFEST_JSON", assets / "curation_manifest.json"),
                mock.patch.object(curate, "MANIFEST_MD", assets / "curation_library.md"),
                mock.patch.object(curate, "scan_guides", lambda: []),
            ):
                res = curate.register("demo")
                self.assertEqual("demo", res["id"])
                self.assertTrue((assets / "skins" / "demo.json").is_file())
                mani = json.loads((assets / "curation_manifest.json").read_text(encoding="utf-8"))
                demo = next(e for e in mani["entries"] if e["id"] == "demo")
                self.assertTrue(demo["registered"])
                first = (assets / "skins" / "demo.json").read_text(encoding="utf-8")
                curate.register("demo")  # 멱등
                self.assertEqual(first, (assets / "skins" / "demo.json").read_text(encoding="utf-8"))

    def test_register_unknown_id_aborts_not_silent(self):
        """③ 없는 id register = ValueError(조용한 통과·지어내기 금지)."""
        with (
            mock.patch.object(curate, "scan_skins", lambda: []),
            mock.patch.object(curate, "scan_guides", lambda: []),
        ):
            with self.assertRaises(ValueError):
                curate.register("no_such_asset")

    def test_scan_intake_excludes_system_previews_and_example_links(self):
        """④ 반입 스캔: 사용자 파일만(시스템 프리뷰 제외) · 실제 링크만(주석 예시 제외)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            refs = run / curate.REFS_DIRNAME
            refs.mkdir()
            (refs / "moodboard.png").write_bytes(b"user")          # 사용자 반입 파일
            (refs / "pillar_card.png").write_bytes(b"system")      # 시스템 프리뷰(매니페스트 등재)
            (refs / curate.REFS_MANIFEST).write_text(
                json.dumps({"per_slide": {"3": [{"file": "design_refs/pillar_card.png"}]}}),
                encoding="utf-8",
            )
            (refs / curate.REFS_NOTE).write_text(
                "링크:\nhttps://real.example/board  실제 링크\n"
                "<!-- 예시: https://example.com/ignore-me -->\n",
                encoding="utf-8",
            )
            scanned = curate.scan_intake(run)
            self.assertEqual(["moodboard.png"], scanned["files"])
            self.assertEqual(["https://real.example/board"], scanned["links"])

    def test_open_intake_preserves_human_edits(self):
        """⑤ refs.md가 이미 있으면 덮지 않는다(사람이 붙인 링크 보존)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            note = curate.open_intake(run)
            note.write_text("https://my.link/keep\n", encoding="utf-8")
            curate.open_intake(run)  # 재호출
            self.assertEqual("https://my.link/keep\n", note.read_text(encoding="utf-8"))

    # -------------------------------------------------------------------
    # DF3(2026-07-24) — sync_master_assets: 확정 배경·장식을 design-assets/references/로 싱크백
    # (register()가 skin/guide 전용이라 커버 못 하는 마스터 자산용 최소 훅, §2b docstring 근거).
    # -------------------------------------------------------------------

    def test_sync_master_assets_requires_design_contract(self):
        """⑥ design_contract.json 자체가 없으면 ValueError(지어내지 않는다)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            with self.assertRaises(ValueError):
                curate.sync_master_assets(run)

    def test_sync_master_assets_requires_background_or_decor(self):
        """⑦ 계약은 있지만 background/decor_slots가 둘 다 없으면 ValueError(싱크백할 게 없음)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "design_contract.json").write_text(
                json.dumps({"chrome_contract": {"chrome": {"frame": {"image": None}}}}),
                encoding="utf-8",
            )
            with self.assertRaises(ValueError):
                curate.sync_master_assets(run)

    def test_sync_master_assets_copies_background_and_decor_missing_surfaces(self):
        """⑧ 배경+장식(하나는 원본 있음, 하나는 없음)을 references/<run명>/로 복사, 부재는 표면화만."""
        with tempfile.TemporaryDirectory() as td:
            base = Path(td)
            run = base / "runs" / "sample_run"
            run.mkdir(parents=True)
            (run / "bg.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            (run / "design_contract.json").write_text(json.dumps({
                "chrome_contract": {
                    "chrome": {"frame": {"image": "bg.png"}},
                    "decor_slots": [
                        {"id": "corner", "image": "corner.png", "anchor": "top-right", "width": 80},
                    ],
                }
            }), encoding="utf-8")
            assets = base / "design-assets"
            with mock.patch.object(curate, "ASSET_REFERENCES", assets / "references"):
                res = curate.sync_master_assets(run)
            self.assertEqual(1, len(res["copied"]))
            self.assertEqual("bg.png", Path(res["copied"][0]).name)
            self.assertEqual(["corner.png"], res["missing"])
            self.assertTrue((assets / "references" / "sample_run" / "bg.png").is_file())


class CheckpointAckApiSmoke(unittest.TestCase):
    """/api/ack는 검증된 run의 사람 관문 파일만 대시보드 provenance로 쓴다 (W27 P0)."""

    def test_ack_writes_dashboard_artifact(self):
        with tempfile.TemporaryDirectory() as td:
            runs = Path(td) / "runs"
            run = runs / "gen_demo"
            run.mkdir(parents=True)
            with mock.patch.object(server, "RUNS", runs):
                code, body = _post_json("/api/ack", {
                    "run_id": run.name, "gate": "design_refs", "decision": "confirm",
                })
            self.assertEqual(200, code)
            self.assertTrue(body["ok"])
            ack = pipeline_state.read_ack(run, "design_refs")
            self.assertEqual("confirm", ack["decision"])
            self.assertEqual("dashboard", ack["via"])

    def test_ack_rejects_path_traversal(self):
        with tempfile.TemporaryDirectory() as td:
            runs = Path(td) / "runs"
            runs.mkdir()
            with mock.patch.object(server, "RUNS", runs):
                code, _ = _post_json("/api/ack", {
                    "run_id": "../outside", "gate": "design", "decision": "confirm",
                })
            self.assertEqual(400, code)

    def test_ack_rejects_unknown_gate(self):
        with tempfile.TemporaryDirectory() as td:
            runs = Path(td) / "runs"
            run = runs / "gen_demo"
            run.mkdir(parents=True)
            with mock.patch.object(server, "RUNS", runs):
                code, _ = _post_json("/api/ack", {
                    "run_id": run.name, "gate": "research", "decision": "confirm",
                })
            self.assertEqual(400, code)
            self.assertFalse((run / pipeline_state.ACK_DIR).exists())

    def test_skip_only_allowed_for_skippable_gates(self):
        with tempfile.TemporaryDirectory() as td:
            runs = Path(td) / "runs"
            run = runs / "gen_demo"
            run.mkdir(parents=True)
            with mock.patch.object(server, "RUNS", runs):
                decision_code, _ = _post_json("/api/ack", {
                    "run_id": run.name, "gate": "decision", "decision": "skip",
                })
                skeleton_code, body = _post_json("/api/ack", {
                    "run_id": run.name, "gate": "skeleton_review", "decision": "skip",
                })
        self.assertEqual(400, decision_code)
        self.assertEqual(200, skeleton_code)
        self.assertEqual("skip", body["ack"]["decision"])


class ImageWorkorderSmoke(unittest.TestCase):
    """검토 대기 카드의 이미지 분업 지시서 요약 (W27 P2 — design_spec.json → 카드)."""

    def test_workorder_routes_and_none_count(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "design_spec.json").write_text(json.dumps({"slides": [
                {"slide_id": "s01", "image_kind": "evidence"},  # 기본값 → user_asset
                {"slide_id": "s02", "image_kind": "mood"},      # 기본값 → codex_gen
                {"slide_id": "s03", "image_kind": "conceptual", "source_route": "web_sample"},
                {"slide_id": "s04", "image_kind": "none", "none_reason": "텍스트 위계만으로 충분"},
            ]}, ensure_ascii=False), encoding="utf-8")
            wo = server._image_workorder(run)
            self.assertEqual(["s01"], wo["user_asset"])
            self.assertEqual(["s02"], wo["codex_gen"])
            self.assertEqual(["s03"], wo["web_sample"])
            self.assertEqual(1, wo["none"])

    def test_workorder_absent_when_no_spec(self):
        with tempfile.TemporaryDirectory() as td:
            self.assertIsNone(server._image_workorder(Path(td)))


class OptionalCheckpointsSmoke(unittest.TestCase):
    """research는 optional 유지, design_refs·design은 대시보드 ack 전속 human 관문 (W27 P0).

    정주행에서 "절차 없이 그대로 넘어가는" 문제의 해결: 소프트 정지와 달리 `go`를 다시 쳐도
    자동으로 안 지나간다(implied-clear 없음). 정본 = pipeline_state (DESIGN_ASSETS_LANE §5-④, 결정 2026-07-14).
    """

    @staticmethod
    def _stg(t):
        return {"at": f"2026-07-14T00:{t:02d}:00"}

    @staticmethod
    def _write_ack(run, gate, decision="confirm", via="dashboard", at=None):
        path = pipeline_state.ack_path(run, gate)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps({
            "gate": gate, "decision": decision,
            "at": at or "2999-01-01T00:00:00+09:00", "via": via,
        }), encoding="utf-8")
        return path

    def _bid_run(self, td, **cps):
        run = Path(td) / "run"
        run.mkdir(parents=True)
        st = pipeline_state._blank(run)
        st["mode"] = "direct"
        st["input"] = {"kind": "bid"}
        for k, v in cps.items():
            st["checkpoints"][k] = {"cleared_at": v}
        pipeline_state.save(run, st)
        return run

    def _refine_phase(self, run):
        """refine_collect까지 끝난 [4+] 상태로 조립(design_refs 관문 도달 조건)."""
        st = pipeline_state.load(run)
        st["stages"] = {n: self._stg(t) for n, t in [
            ("render", 2), ("review_resolve", 2), ("wireframe_bundle", 2), ("wireframe_apply", 3),
            ("design_brief", 3), ("stage9_bundle", 4), ("stage9_apply", 5),
            ("refine_bundle", 6), ("refine_collect", 7)]}
        pipeline_state.save(run, st)
        for f in ("deck.json", "design_brief.json", "design_overrides.json",
                  "design_spec.json", "wireframe.json"):
            (run / f).write_text("{}", encoding="utf-8")

    def _skeleton_phase(self, run):
        """message_map 수거와 skeleton 역제안까지 끝난 스토리라인 생성 전 상태."""
        st = pipeline_state.load(run)
        st["stages"].update({
            "message_map": self._stg(1),
            "skeleton": self._stg(2),
        })
        pipeline_state.save(run, st)
        (run / "message_map.json").write_text("{}", encoding="utf-8")
        (run / "skeleton.json").write_text("{}", encoding="utf-8")
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")

    def _wireframe_phase(self, run):
        """wireframe 적용까지 끝난 테마 입히기 전 상태."""
        st = pipeline_state.load(run)
        st["stages"].update({
            "render": self._stg(2),
            "review_resolve": self._stg(2),
            "wireframe_bundle": self._stg(3),
            "wireframe_apply": self._stg(4),
        })
        pipeline_state.save(run, st)
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        (run / "wireframe.json").write_text("{}", encoding="utf-8")
        (run / "gating_report.json").write_text("{}", encoding="utf-8")

    def test_research_offered_for_bid_not_brief_and_autosatisfied_by_file(self):
        """① research: bid+start통과면 선택 관문, brief면 안 뜸, institution_research.json 있으면 자동 충족."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00")
            n = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", n["kind"])
            self.assertEqual("research", n["checkpoint"])
            self.assertTrue(n.get("optional"))
            # 조사 파일이 생기면 자동 충족(다음 단계로)
            (run / "institution_research.json").write_text("{}", encoding="utf-8")
            self.assertNotEqual("research", pipeline_state.resolve(run)["next"].get("checkpoint"))
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"; run.mkdir(parents=True)
            st = pipeline_state._blank(run); st["mode"] = "direct"; st["input"] = {"kind": "brief"}
            st["checkpoints"]["start"] = {"cleared_at": "2026-07-14T00:00:00"}
            pipeline_state.save(run, st)
            self.assertNotEqual("research", pipeline_state.resolve(run)["next"].get("checkpoint"))

    def test_design_refs_gates_handoff_until_confirm(self):
        """② design_refs: refine_collect 후 선택 관문 → 통과(clear)해야 refine_handoff로."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            self._refine_phase(run)
            n = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", n["kind"])
            self.assertEqual("design_refs", n["checkpoint"])
            self.assertTrue(n.get("optional"))
            pipeline_state.clear_checkpoint(run, "design_refs")
            self.assertEqual("refine_handoff", pipeline_state.resolve(run)["next"].get("stage"))

    def test_optional_checkpoint_not_implied_by_downstream_stages(self):
        """③ 선택 관문은 산출물로 implied-clear 되지 않는다(에이전트가 밀고 지나가지 못함)."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            self._refine_phase(run)
            st = pipeline_state.load(run)
            st["stages"]["refine_handoff"] = self._stg(8)  # 하류 산출물이 있어도
            pipeline_state.save(run, st)
            eff = pipeline_state._effective_checkpoints(st["checkpoints"], st["stages"])
            self.assertIsNone(eff["design_refs"]["cleared_at"])  # design_refs는 여전히 미통과

    def test_human_commands_fail_without_ack(self):
        """① refine --handoff와 approve는 ack 없이는 main에서 비0 종료한다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            self._refine_phase(run)
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch.object(_sys, "argv", [
                    "proposal_pipeline.py", "refine", "--run", run.name, "--handoff",
                ]), mock.patch("sys.stderr", new=io.StringIO()):
                    self.assertNotEqual(0, proposal_pipeline.main())
                with mock.patch.object(_sys, "argv", [
                    "proposal_pipeline.py", "approve", "--run", run.name,
                ]), mock.patch("sys.stderr", new=io.StringIO()):
                    self.assertNotEqual(0, proposal_pipeline.main())

    def test_go_confirm_rejects_human_checkpoint(self):
        """② --confirm은 human 관문에 비0이며 안내를 stderr로 낸다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            self._refine_phase(run)
            err = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stderr", new=err):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=True, json=False))
            self.assertNotEqual(0, rc)
            self.assertIn("[GO 거부]", err.getvalue())
            self.assertIn("--confirm은 이 관문에 무효", err.getvalue())

        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(
                td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00"
            )
            self._wireframe_phase(run)
            err = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stderr", new=err):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=True, json=False))
            self.assertNotEqual(0, rc)
            self.assertIn("--confirm은 이 관문에 무효", err.getvalue())

    def test_dashboard_ack_is_consumed_and_cleared(self):
        """③ ack 배치 후 go가 통과하고 checkpoint cleared_at을 기록한다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            self._refine_phase(run)
            self._write_ack(run, "design_refs")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_go_refine_handoff", return_value=True), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            cleared = pipeline_state.load(run)["checkpoints"]["design_refs"]["cleared_at"]
            self.assertIsNotNone(cleared)
            self.assertIn("[GO] ack 확인: design_refs decision=confirm via=dashboard", out.getvalue())

        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(
                td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00"
            )
            self._wireframe_phase(run)
            with mock.patch.object(server, "RUNS", Path(td)):
                code, body = _post_json("/api/ack", {
                    "run_id": run.name,
                    "gate": "wireframe_review",
                    "decision": "skip",
                })
            self.assertEqual(200, code)
            self.assertEqual("skip", body["ack"]["decision"])
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_go_design_brief", return_value=True), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            # ⚠ 이 단언은 부하가 걸린 기계에서 간헐적으로 실패한다(W32 실측 ~4/20, 재현 조건
            # 미확정 — 배포판·원본 양쪽, chromium 유무 무관). 유력 지점은 go의 staleness 분기
            # `is_stale(run, gate, ack["at"])`: 감시 파일(wireframe.json)의 mtime이 ack보다
            # 새면 관문이 재무장되고 cleared_at이 안 찍힌다. 다음에 잡는 사람이 맨손으로
            # 시작하지 않게 실패 메시지에 판단 재료를 실어 둔다.
            _st = pipeline_state.load(run)
            self.assertIsNotNone(
                _st["checkpoints"]["wireframe_review"]["cleared_at"],
                f"ack={pipeline_state.read_ack(run, 'wireframe_review')} "
                f"cp={_st['checkpoints'].get('wireframe_review')} "
                f"wf_mtime={(run / 'wireframe.json').stat().st_mtime if (run / 'wireframe.json').exists() else None} "
                f"go_out={out.getvalue()[-400:]}"
            )
            self.assertIn(
                "[GO] ack 확인: wireframe_review decision=skip via=dashboard",
                out.getvalue(),
            )

    def test_non_dashboard_ack_is_invalid(self):
        """④ via가 dashboard가 아니면 ack 무효이며 waiting_human 상태를 유지한다.

        W31 리허설 마찰2: design_refs는 기본(standard) 프로파일에서 신호 없으면 자동 통과 대상이라,
        이 테스트의 관심사(무효 ack 거부)와 무관하게 통과해버릴 수 있다 — full 프로파일(전 관문 정지)
        로 고정해 gates 다이얼의 영향을 배제하고 원래 검증 대상만 본다.
        """
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00")
            gates.save_config(run, profile="full")
            self._refine_phase(run)
            self._write_ack(run, "design_refs", via="session")
            self.assertIsNone(pipeline_state.read_ack(run, "design_refs"))
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNone(pipeline_state.load(run)["checkpoints"]["design_refs"]["cleared_at"])
            self.assertIn("[GO 대기] waiting_human:design_refs", out.getvalue())

    def test_research_confirm_remains_optional(self):
        """⑤ research optional 관문은 기존처럼 go --confirm으로 통과한다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00")
            waiting = {"next": {"kind": "checkpoint", "checkpoint": "research"}}
            done = {"next": {"kind": "done"}}
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_run_state_view",
                                   side_effect=[waiting, done, done]), \
                 mock.patch.object(pipeline_state, "format_status", return_value="status"):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=True, json=False))
            self.assertEqual(0, rc)
            self.assertIsNotNone(pipeline_state.load(run)["checkpoints"]["research"]["cleared_at"])

    def test_skeleton_review_gates_only_before_storyline_and_accepts_skip_ack(self):
        """① skeleton 완료 뒤 관문, 기존 storyline은 면제, dashboard skip ack는 소비된다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", research="2026-07-14T00:00:01")
            self._skeleton_phase(run)
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt["kind"])
            self.assertEqual("skeleton_review", nxt["checkpoint"])
            self.assertTrue(nxt["optional"])
            self.assertTrue(nxt["human"])
            legacy = run / "storyline.json"
            legacy.write_text("{}", encoding="utf-8")
            legacy_next = pipeline_state.resolve(run, render_input=legacy)["next"]
            self.assertNotEqual("skeleton_review", legacy_next.get("checkpoint"))

        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00", research="2026-07-14T00:00:01")
            self._skeleton_phase(run)
            self._write_ack(run, "skeleton_review", decision="skip")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch("sys.stdout", new=out):
                    rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNotNone(pipeline_state.load(run)["checkpoints"]["skeleton_review"]["cleared_at"])

    def test_decision_rejects_confirm_and_consumes_dashboard_ack(self):
        """② decision은 human 관문이며 --confirm 거부, dashboard ack만 clearance를 만든다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2026-07-14T00:00:00")
            st = pipeline_state.load(run)
            st["stages"]["render"] = self._stg(2)
            pipeline_state.save(run, st)
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("decision", nxt["checkpoint"])
            self.assertIn(str(run / "storyline.json"), nxt["review"])
            self.assertIn(str(run / "message_map.json"), nxt["review"])
            err = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch("sys.stderr", new=err):
                    rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=True, json=False))
            self.assertNotEqual(0, rc)
            self.assertIn("[GO 거부]", err.getvalue())
            self._write_ack(run, "decision")
            waiting = {"next": nxt}
            done = {"next": {"kind": "done", "why": "done", "command": "done"}}
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch.object(proposal_pipeline, "_run_state_view", side_effect=[waiting, done, done]):
                    with mock.patch.object(pipeline_state, "format_status", return_value="status"):
                        with mock.patch("sys.stdout", new=out):
                            rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNotNone(pipeline_state.load(run)["checkpoints"]["decision"]["cleared_at"])

    def test_recorded_decision_rearms_and_rejects_old_ack_then_accepts_new_ack(self):
        """③ 감시 파일 수정은 recorded clearance와 낡은 ack를 무효화하고 새 ack로 해소된다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(td, start="2020-01-01T00:00:00", decision="2020-01-01T00:01:00")
            st = pipeline_state.load(run)
            st["stages"]["render"] = self._stg(2)
            pipeline_state.save(run, st)
            for filename in ("storyline.json", "message_map.json"):
                path = run / filename
                path.write_text("{}", encoding="utf-8")
                os.utime(path, (1_700_000_000, 1_700_000_000))
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("decision", nxt["checkpoint"])
            self.assertTrue(nxt["why"].startswith("[재무장]"))
            summary = server._run_status_summary(run)
            self.assertEqual("decision", summary["human_gate"])
            self.assertTrue(summary["human_label"].startswith("[재무장]"))
            self.assertFalse(summary["human_skippable"])
            old_cleared = pipeline_state.load(run)["checkpoints"]["decision"]["cleared_at"]
            self._write_ack(run, "decision", at="2021-01-01T00:00:00+09:00")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch.object(proposal_pipeline, "_go_resolution_skeleton"):
                    with mock.patch("sys.stdout", new=out):
                        rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertEqual(old_cleared, pipeline_state.load(run)["checkpoints"]["decision"]["cleared_at"])
            self.assertIn("ack가 산출물 변경보다 오래됨 - 재검토 필요", out.getvalue())
            self._write_ack(run, "decision", at="2024-01-01T00:00:00+09:00")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch.object(proposal_pipeline, "_go_review_resolve", return_value=True):
                    with mock.patch.object(proposal_pipeline, "_go_resolution_skeleton"):
                        with mock.patch("sys.stdout", new=out):
                            rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertNotEqual("decision", pipeline_state.resolve(run)["next"].get("checkpoint"))

    def test_implied_decision_clearance_is_not_rearmed(self):
        """④ 상태파일 없는 레거시 지문 clearance는 감시 산출물보다 오래돼도 재무장하지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            run.mkdir(exist_ok=True)
            for filename in ("manifest_render.json", "design_brief.json", "storyline.json", "message_map.json"):
                (run / filename).write_text("{}", encoding="utf-8")
            view = pipeline_state.resolve(run)
            self.assertEqual("inferred", view["checkpoints"]["decision"]["source"])
            self.assertNotEqual("decision", view["next"].get("checkpoint"))

    def test_approve_rejects_stale_design_clearance(self):
        """⑤ design clearance 뒤 deck.html 변경 시 approve는 새 사람 ack 전까지 비0이다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(
                td, start="2020-01-01T00:00:00", decision="2020-01-01T00:01:00",
                design="2020-01-01T00:02:00",
            )
            deck_html = run / "deck.html"
            deck_html.write_text("<html></html>", encoding="utf-8")
            (run / "deck.json").write_text("{}", encoding="utf-8")
            os.utime(deck_html, (1_700_000_000, 1_700_000_000))
            err = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                with mock.patch.object(_sys, "argv", ["proposal_pipeline.py", "approve", "--run", run.name]):
                    with mock.patch("sys.stderr", new=err):
                        rc = proposal_pipeline.main()
            self.assertNotEqual(0, rc)
            self.assertIn("재검토 필요", err.getvalue())

    def test_wireframe_review_surface_legacy_guard_and_rearm(self):
        """wireframe_apply 후 정지, 레거시 최초 면제, recorded clearance 뒤 갱신 재무장."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bid_run(
                td, start="2026-07-14T00:00:00", decision="2026-07-14T00:01:00"
            )
            self._wireframe_phase(run)
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt["kind"])
            self.assertEqual("wireframe_review", nxt["checkpoint"])
            self.assertTrue(nxt["optional"])
            self.assertTrue(nxt["human"])
            self.assertEqual(
                [str(run / "deck.html"), str(run / "wireframe.json"),
                 str(run / "gating_report.json")],
                nxt["review"],
            )

            st = pipeline_state.load(run)
            st["stages"]["design_brief"] = self._stg(5)
            pipeline_state.save(run, st)
            (run / "design_brief.json").write_text("{}", encoding="utf-8")
            self.assertNotEqual(
                "wireframe_review", pipeline_state.resolve(run)["next"].get("checkpoint")
            )

            pipeline_state.clear_checkpoint(run, "wireframe_review")
            os.utime(run / "wireframe.json", (2_000_000_000, 2_000_000_000))
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("wireframe_review", nxt["checkpoint"])
            self.assertTrue(nxt["why"].startswith("[재무장]"))
            self.assertIn("wireframe.json", nxt["why"])
            summary = server._run_status_summary(run)
            self.assertEqual("wireframe_review", summary["human_gate"])
            self.assertTrue(summary["human_label"].startswith("[재무장]"))
            self.assertTrue(summary["human_skippable"])


class KnowledgeRefsSmoke(unittest.TestCase):
    """W27 P1a 레퍼런스 시각 채널 — 지식 카드(examples/source 링크) → ig 폴더 jpg 결정론 수집.

    정본 = design_spec.py의 _resolve_card/_collect_knowledge_refs(D1·D2, CONTEXT/DESIGN_KNOWLEDGE_TEETH.md
    §2·§3 P1a). 실제 <개발 원본 전용 경로> 의존하지 않고 tempfile로 가짜 카드 루트를 구성한다.
    """

    @staticmethod
    def _write(path: Path, text: str) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text, encoding="utf-8")

    @classmethod
    def _fake_images_root(cls, root: Path) -> Path:
        """카드 루트 + grouped-images 폴더를 갖춘 가짜 reference_images_root 구성."""
        # 카드A(와이어프레임) — examples 링크 있음 → 예시 카드 경유(resolved_via=examples).
        cls._write(root / "cards" / "와이어프레임" / "카드A.md", """---
name: 카드A
claim: 카드A 원칙 주장
type: principle
layer: wireframe
apply: proactive
source: peedori_ ig_999999
examples: [예시A]
---

본문.

**조작적 정의**
- 규칙 1
- 규칙 2

[[다른카드]]
""")
        # 예시A — 카드A가 인용하는 example 카드(source ig_111111).
        cls._write(root / "cards" / "examples" / "예시A.md", """---
name: 예시A
proves: 예시A 증명
type: example
kind: before_after
source: peedori_ ig_111111
---

**교정 규칙**: 이렇게 고친다.

[[카드A]]
""")
        # 카드B(테마) — examples 빈 리스트 → 폴백 = 자기 source(ig_222222).
        cls._write(root / "cards" / "테마" / "카드B.md", """---
name: 카드B
claim: 카드B 원칙 주장
type: principle
layer: theme
apply: proactive
source: peedori_ ig_222222
examples: []
---

본문 설명.
""")
        # 카드C(와이어프레임) — ig_333333에 이미지 10장(6장 상한 검증용).
        cls._write(root / "cards" / "와이어프레임" / "카드C.md", """---
name: 카드C
claim: 카드C 원칙 주장
type: principle
layer: wireframe
apply: proactive
source: peedori_ ig_333333
examples: []
---

본문.
""")

        for i in range(1, 3):
            img = root / "grouped-images" / "ig_111111" / f"0{i}.jpg"
            img.parent.mkdir(parents=True, exist_ok=True)
            img.write_bytes(b"fake-jpg")
        img = root / "grouped-images" / "ig_222222" / "01.jpg"
        img.parent.mkdir(parents=True, exist_ok=True)
        img.write_bytes(b"fake-jpg")
        for i in range(1, 11):
            img = root / "grouped-images" / "ig_333333" / f"{i:02d}.jpg"
            img.parent.mkdir(parents=True, exist_ok=True)
            img.write_bytes(b"fake-jpg")
        return root

    def test_wireframe_prompt_indexes_wireframe_cards_only(self):
        """와이어프레임 카드의 name·claim만 색인하고 테마·본문은 주입하지 않는다."""
        import design_spec
        import wireframe

        with tempfile.TemporaryDirectory() as td:
            images_root = self._fake_images_root(Path(td) / "img_root")
            with mock.patch.object(
                design_spec, "_reference_images_root", return_value=images_root
            ):
                prompt = wireframe.build_prompt(WireframeContractSmoke.DECK)

        self.assertIn("## 디자인지식 카드 색인 (형태 원칙 - 와이어프레임 층 전용)", prompt)
        self.assertIn("- 카드A: 카드A 원칙 주장", prompt)
        self.assertIn("- 카드C: 카드C 원칙 주장", prompt)
        self.assertNotIn("카드B 원칙 주장", prompt)
        self.assertNotIn("예시A 증명", prompt)
        self.assertNotIn("ig_999999", prompt)
        self.assertNotIn("규칙 1", prompt)
        self.assertIn('"knowledge_cards"', prompt)

    def test_wireframe_prompt_surfaces_missing_knowledge_root(self):
        """reference_images_root 부재는 하드 실패 없이 고정 문구로 표면화한다."""
        import design_spec
        import wireframe

        with mock.patch.object(design_spec, "_reference_images_root", return_value=None):
            prompt = wireframe.build_prompt(WireframeContractSmoke.DECK)
        self.assertIn(
            "(디자인지식 색인 없음 - config knowledge.reference_images_root 미설정)",
            prompt,
        )

    def test_examples_link_resolves_to_example_card_ig_folder(self):
        """① 카드A(examples: [예시A]) → 예시A의 source ig_111111 이미지로 해석된다."""
        import design_spec
        with tempfile.TemporaryDirectory() as td:
            images_root = self._fake_images_root(Path(td) / "img_root")
            cards_root = images_root / "cards"
            resolved = design_spec._resolve_card(cards_root, images_root, "카드A")
            self.assertIsNone(resolved["gap"])
            self.assertEqual("examples", resolved["resolved_via"])
            self.assertEqual(2, len(resolved["images"]))
            self.assertTrue(all(p.parent.name == "ig_111111" for p in resolved["images"]))
            self.assertIn("규칙 1", resolved["watch_for"])

    def test_empty_examples_falls_back_to_own_source(self):
        """② examples: [] 카드(카드B) → 자기 source(ig_222222)로 폴백, resolved_via=source."""
        import design_spec
        with tempfile.TemporaryDirectory() as td:
            images_root = self._fake_images_root(Path(td) / "img_root")
            cards_root = images_root / "cards"
            resolved = design_spec._resolve_card(cards_root, images_root, "카드B")
            self.assertIsNone(resolved["gap"])
            self.assertEqual("source", resolved["resolved_via"])
            self.assertEqual(1, len(resolved["images"]))
            self.assertEqual("ig_222222", resolved["images"][0].parent.name)

    def test_unknown_card_slug_surfaces_gap_without_crashing(self):
        """③ 미지 카드 슬러그 → 예외 없이 gap 사유가 채워진다(지어내지 않음)."""
        import design_spec
        with tempfile.TemporaryDirectory() as td:
            images_root = self._fake_images_root(Path(td) / "img_root")
            cards_root = images_root / "cards"
            resolved = design_spec._resolve_card(cards_root, images_root, "존재하지않는카드")
            self.assertEqual([], resolved["images"])
            self.assertIsNotNone(resolved["gap"])
            self.assertIn("존재하지않는카드", resolved["gap"])

    def test_missing_reference_images_root_skips_without_breaking_existing_collect(self):
        """④ reference_images_root 부재 → 경고만 찍고 스킵, 기존 form_needs collect 결과는 불변."""
        import design_spec
        spec = {
            "schema_version": 1, "run_id": "test", "generated_by": "llm:test",
            "slides": [
                {"slide_id": "1", "goal": "목표", "treatment": ["diagram"], "image_kind": "conceptual",
                 "form_needs": [{"kind": "piece", "id": "matrix_2x2", "why": "현위치 대비 목표"}],
                 "knowledge_cards": ["카드A"]},
            ],
            "catalog_gap": [],
        }
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_spec, "_reference_images_root", return_value=None):
                manifest = design_spec.collect_refs(run, spec)
            self.assertEqual([], manifest["knowledge_refs"])
            self.assertTrue((run / "design_refs" / "matrix_2x2.png").is_file())
            items = manifest["per_slide"]["1"]
            self.assertEqual("matrix_2x2", items[0]["id"])

    def test_per_card_limit_and_idempotent_recopy(self):
        """⑤ 카드당 최대 6장 상한 + 재실행 시 중복 복사 없음(멱등)."""
        import design_spec
        spec = {
            "schema_version": 1, "run_id": "test", "generated_by": "llm:test",
            "slides": [
                {"slide_id": "1", "goal": "목표", "treatment": ["diagram"], "image_kind": "conceptual",
                 "knowledge_cards": ["카드C"]},
            ],
            "catalog_gap": [],
        }
        with tempfile.TemporaryDirectory() as td:
            images_root = self._fake_images_root(Path(td) / "img_root")
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_spec, "_reference_images_root", return_value=images_root):
                manifest1 = design_spec.collect_refs(run, spec)
                dest_dir = run / "design_refs" / "knowledge" / "카드C"
                files1 = sorted(dest_dir.glob("*.jpg"))
                self.assertEqual(6, len(files1))
                mtimes1 = {p.name: p.stat().st_mtime_ns for p in files1}

                manifest2 = design_spec.collect_refs(run, spec)
                files2 = sorted(dest_dir.glob("*.jpg"))
                self.assertEqual(6, len(files2))
                mtimes2 = {p.name: p.stat().st_mtime_ns for p in files2}
                self.assertEqual(mtimes1, mtimes2)  # 재복사 없었음(멱등)

            kr1 = manifest1["knowledge_refs"][0]
            self.assertEqual(6, len(kr1["images"]))
            self.assertIsNone(kr1["gap"])
            kr2 = manifest2["knowledge_refs"][0]
            self.assertEqual(kr1["images"], kr2["images"])


class JourneyFoldersSmoke(unittest.TestCase):
    """W31 R7(CONTEXT/JOURNEY.md) — 단계 폴더 여정(journey/). 회귀 방지:
    ①산출물이 생겼을 때만 폴더가 열림(수납처 선개방, 도달로 안 열림) ②storyline·wireframe
    가독 뷰(R1) 렌더 ③정본 변경 시 뷰 재생성(stale) ④폴더 안에 정본 사본이 없음(이중화 금지)
    ⑤A3 발주처 조사(디자인 계열)가 B1 게이트 전에도 07_테마확정을 선개방
    ⑥go_cmd가 실제로 sync를 호출한다."""

    def test_folder_opens_only_when_artifact_appears(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            report = journey_folders.sync(run)
            # 01/02는 run 존재만으로 항상 연다(A0/A1은 run 생성 전 단계 — 포인터만).
            self.assertIn("01", report["newly_opened"])
            self.assertIn("02", report["newly_opened"])
            folder05 = journey_folders.folder_path(run, "05")
            self.assertFalse(folder05.is_dir())  # storyline 없음 → 05는 아직 안 연다.

            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "제목", "message": "메시지",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            report2 = journey_folders.sync(run)
            self.assertIn("05", report2["newly_opened"])
            self.assertTrue(folder05.is_dir())
            self.assertTrue((folder05 / journey_folders.MANUAL_NAME).is_file())
            self.assertTrue((folder05 / journey_folders.STORYLINE_VIEW_NAME).is_file())

            # 멱등: 다시 sync해도 05는 "새로" 열리지 않는다(이미 열린 폴더는 계속 남는다).
            report3 = journey_folders.sync(run)
            self.assertNotIn("05", report3["newly_opened"])
            self.assertIn("05", report3["active"])

    def test_storyline_and_wireframe_readable_views_render_expected_fields(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T", "governing_message": "핵심"},
                "slides": [
                    {"n": 1, "section": "표지", "title": "표지 제목", "message": "표지 메시지",
                     "bullets": ["불릿A"], "example": True, "flag": ["검토요망: 근거 확인"]},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            (run / "wireframe.json").write_text(json.dumps({
                "schema_version": 1, "selected_by": "llm:test",
                "slides": [
                    {"slide_id": "1", "message_type": "서사", "frame": "hero_body", "rendition": "boxed",
                     "layout_group": "bookend",
                     "slots": [{"piece": "text_block", "size": "hero", "binds": "title"}],
                     "principles": ["표지 프레임"], "catalog_gap": []},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            journey_folders.sync(run)
            sview = (journey_folders.folder_path(run, "05")
                     / journey_folders.STORYLINE_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("자동 생성", sview)
            self.assertIn("표지 제목", sview)
            self.assertIn("[예시]", sview)
            self.assertIn("[검토요망: 근거 확인]", sview)
            self.assertIn("불릿A", sview)
            wview = (journey_folders.folder_path(run, "06")
                     / journey_folders.WIREFRAME_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("text_block", wview)
            self.assertIn("hero_body", wview)
            self.assertIn("표지 프레임", wview)

    def test_view_regenerates_when_canonical_source_changes(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            storyline_path = run / "storyline.json"
            storyline_path.write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "원본 제목", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            journey_folders.sync(run)
            view_path = journey_folders.folder_path(run, "05") / journey_folders.STORYLINE_VIEW_NAME
            self.assertIn("원본 제목", view_path.read_text(encoding="utf-8"))

            # 뷰를 인위적으로 과거로 돌려 "정본이 더 새롭다"를 확실히 한다(mtime 해상도 회피).
            past = view_path.stat().st_mtime - 10
            os.utime(view_path, (past, past))
            storyline_path.write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "수정된 제목", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            report = journey_folders.sync(run)
            self.assertIn(f"05/{journey_folders.STORYLINE_VIEW_NAME}", report["views_rendered"])
            refreshed = view_path.read_text(encoding="utf-8")
            self.assertIn("수정된 제목", refreshed)
            self.assertNotIn("원본 제목", refreshed)

    def test_no_canonical_duplication_inside_journey_folders(self):
        """이중화 금지(R7): journey/ 아래 어떤 폴더에도 정본 JSON의 사본이 없다 — 매뉴얼(포인터),
        '*_읽기.md' 파생 뷰, `회의체_메모.md`(사람 편집물), 그리고 W31 리허설 마찰5의
        `산출물.html`/`_전체여정.html`(클릭 링크+계보 파생 뷰 — 정본 사본이 아니라 링크뿐) 종류만
        존재한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            (run / "wireframe.json").write_text(json.dumps({
                "schema_version": 1, "selected_by": "x", "slides": [],
            }, ensure_ascii=False), encoding="utf-8")
            (run / "institution_research.json").write_text(json.dumps({
                "institution": "X", "brand_tokens": {"colors": {"primary": "#000"}},
            }, ensure_ascii=False), encoding="utf-8")
            journey_folders.sync(run)
            journey_root = journey_folders.journey_root(run)
            canonical_names = {"storyline.json", "wireframe.json", "institution_research.json",
                                "message_map.json", "design_brief.json"}
            found_any = False
            for path in journey_root.rglob("*"):
                if path.is_file():
                    found_any = True
                    self.assertNotIn(path.name, canonical_names,
                                      f"journey 폴더에 정본 사본 발견: {path}")
                    self.assertTrue(
                        path.name == journey_folders.MANUAL_NAME
                        or path.name == journey_folders.MEETING_NOTE_NAME
                        or path.name == journey_folders.OUTPUT_VIEW_NAME
                        or path.name == journey_folders.OVERVIEW_NAME
                        or path.name.endswith("_읽기.md")
                        or path.name.endswith("_읽기.html"),  # 마찰13: 도식 뷰(html)도 파생 뷰
                        f"예기치 않은 파일 종류(파생 뷰/매뉴얼/회의체 메모만 허용): {path}",
                    )
            self.assertTrue(found_any)

    def test_meeting_note_created_once_and_user_edits_survive_resync(self):
        """R4 유실 방지(2026-07-21 수정): `회의체_메모.md`는 없을 때만 최초 1회 생성되고,
        사람이 적은 내용은 이후 몇 번을 sync()해도(=매 go) 절대 덮어써지지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            report1 = journey_folders.sync(run)
            note_path = journey_folders.folder_path(run, "05") / journey_folders.MEETING_NOTE_NAME
            self.assertIn(f"05/{journey_folders.MEETING_NOTE_NAME}", report1["meeting_notes_created"])
            self.assertIn("사람 편집물", note_path.read_text(encoding="utf-8"))

            # 사람이 R4 칸을 채운다.
            user_text = (
                "> ✍️ 사람 편집물 — 자유롭게 수정, 시스템이 덮어쓰지 않음(이 폴더에서는 이 파일이 정본이다).\n\n"
                "# 회의체 메모 (R4)\n\n"
                "- 회의체 이름: 내용 동결 회의\n"
                "- 참석자: 김철수, 이영희\n"
                "- 판단 기준: 정성평가 배점 정합\n"
            )
            note_path.write_text(user_text, encoding="utf-8")

            # storyline.json을 바꿔 파생 뷰 재생성까지 겹쳐도 회의체 메모는 그대로여야 한다.
            storyline = json.loads((run / "storyline.json").read_text(encoding="utf-8"))
            storyline["slides"][0]["title"] = "제목 변경"
            (run / "storyline.json").write_text(json.dumps(storyline, ensure_ascii=False), encoding="utf-8")
            past = note_path.stat().st_mtime - 10
            os.utime(note_path, (past, past))

            report2 = journey_folders.sync(run)
            report3 = journey_folders.sync(run)  # 2회 이상 재실행(매 go)해도 유실 없음을 확인
            self.assertNotIn(f"05/{journey_folders.MEETING_NOTE_NAME}", report2["meeting_notes_created"])
            self.assertNotIn(f"05/{journey_folders.MEETING_NOTE_NAME}", report3["meeting_notes_created"])
            self.assertEqual(user_text, note_path.read_text(encoding="utf-8"))

            # 매뉴얼(포인터)에는 회의체_메모.md를 가리키는 안내만 있고 R4 칸 자체는 없다(유실 원인 제거 확인).
            manual_text = (journey_folders.folder_path(run, "05")
                           / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
            self.assertIn(journey_folders.MEETING_NOTE_NAME, manual_text)
            self.assertNotIn("[사용자 기입]", manual_text)

    def test_meeting_note_legacy_name_migrates_with_content(self):
        """마찰 10(2026-07-21 사용자): 구명 `회의체_메모.md`가 있으면 내용 보존한 채
        `회의록_메모.md`로 이관되고, 새 템플릿으로 덮어쓰지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            folder = journey_folders.folder_path(run, "05")
            folder.mkdir(parents=True, exist_ok=True)
            legacy = folder / "회의체_메모.md"
            legacy.write_text("사람이 적어둔 내용", encoding="utf-8")
            report = journey_folders.sync(run)
            new_path = folder / journey_folders.MEETING_NOTE_NAME
            self.assertEqual("회의록_메모.md", journey_folders.MEETING_NOTE_NAME)
            self.assertFalse(legacy.exists())
            self.assertEqual("사람이 적어둔 내용", new_path.read_text(encoding="utf-8"))
            # 이관은 신규 생성이 아니다 — meeting_notes_created에 잡히지 않는다.
            self.assertNotIn(f"05/{journey_folders.MEETING_NOTE_NAME}",
                             report["meeting_notes_created"])

    def test_institution_research_preopens_theme_folder_before_gate(self):
        """R7 수납처 선개방: institution_research.json이 생기면 B1(테마 확정) 게이트 전에도
        07_테마확정/이 즉시 열리고 브랜드 요약이 담긴다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self.assertFalse(journey_folders.folder_path(run, "07").is_dir())
            (run / "institution_research.json").write_text(json.dumps({
                "institution": "테스트대학교",
                "brand_tokens": {"colors": {"primary": "#004388", "accent": None},
                                  "fonts": {"family": None}, "logo": {"path": None}},
            }, ensure_ascii=False), encoding="utf-8")
            report = journey_folders.sync(run)
            self.assertIn("07", report["newly_opened"])
            brand_view = journey_folders.folder_path(run, "07") / journey_folders.BRAND_VIEW_NAME
            self.assertTrue(brand_view.is_file())
            self.assertIn("#004388", brand_view.read_text(encoding="utf-8"))

    def test_go_cmd_drives_journey_sync(self):
        """go 흐름 통합: go_cmd가 journey_folders.sync를 실제로 호출해 폴더를 만든다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            pipeline_state.record(run, "render")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertTrue(
                (run / "journey" / "05_내용동결" / journey_folders.STORYLINE_VIEW_NAME).is_file()
            )
            self.assertIn("[JOURNEY]", out.getvalue())


class DesignContractSmoke(unittest.TestCase):
    """W31 R-B(R2·R3·R5, CONTEXT/JOURNEY.md) — run별 디자인 계약(design_contract.json). 회귀 방지:
    ①생성(초안 병합·동결) ②차용(brief.skin.value) vs 중립(skins/_neutral.json) ③chrome/image
    2계약 분리 — image_contract에 canvas/chrome이 새지 않음(혼입 제거) ④design_contract 없는 run은
    종전 skin_path 폴백(파일럿·레거시 회귀 보존) ⑤B1 위치의 theme_confirm 게이트(선택·대시보드 ack)."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")

    # ---- ① 생성 + ② 차용 vs 중립 ----------------------------------------------------------

    def test_build_neutral_when_no_skin_value(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
        self.assertEqual("neutral", contract["meta"]["source"])
        self.assertIn(design_contract.NEUTRAL_SKIN_NAME, contract["meta"]["draft_path"])

    def test_build_borrows_named_skin_when_brief_has_value(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            brief = {"skin": {"value": "inkline"}}
            contract = design_contract.build(run, brief=brief, skins_dir=self.SKINS_DIR)
        self.assertEqual("inkline", contract["meta"]["source"])
        self.assertIn("inkline.json", contract["meta"]["draft_path"])
        # 차용본은 창고 스킨(inkline.json)의 결정값을 그대로 초안으로 가져온다.
        self.assertEqual("E8590C", contract["chrome_contract"]["colors"]["accent"])

    def test_build_falls_back_to_neutral_for_unknown_borrow_target(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            brief = {"skin": {"value": "no-such-skin-xyz"}}
            contract = design_contract.build(run, brief=brief, skins_dir=self.SKINS_DIR)
        self.assertEqual("neutral", contract["meta"]["source"])

    def test_all_warehouse_skins_are_borrowable(self):
        """안건 5 확정(2026-07-21, FEATURE_MAP_W31.md) — 창고(skins/*.json)의 모든 스킨은 입고 경로
        무관(졸업=curate 등록 · 외부 이식=add-skin)하게 차용 가능해야 한다: build가 예외 없이
        chrome/image 2계약으로 분리되고 image_contract에 canvas/chrome이 새지 않는다."""
        for skin_path in sorted(self.SKINS_DIR.glob("*.json")):
            skin_id = skin_path.stem
            with self.subTest(skin=skin_id):
                with tempfile.TemporaryDirectory() as td:
                    run = Path(td)
                    brief = {"skin": {"value": skin_id}}
                    contract = design_contract.build(run, brief=brief, skins_dir=self.SKINS_DIR)
                self.assertEqual(skin_id, contract["meta"]["source"])
                self.assertIsInstance(contract["chrome_contract"], dict)
                self.assertIsInstance(contract["image_contract"], dict)
                self.assertNotIn("canvas", contract["image_contract"])
                self.assertNotIn("chrome", contract["image_contract"])

    def test_go_executor_writes_contract_once_and_records_state(self):
        """proposal_pipeline._go_design_contract 배선 — 실제 executor가 파일을 만들고 idempotent."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            (run / "design_brief.json").write_text(
                json.dumps({"skin": {"value": "inkline"}}, ensure_ascii=False), encoding="utf-8")
            proposal_pipeline._go_design_contract(run, SimpleNamespace())
            self.assertTrue(design_contract.exists(run))
            first = design_contract.path(run).read_text(encoding="utf-8")
            state = pipeline_state.load(run)
            self.assertIn("design_contract", state["stages"])
            self.assertEqual("inkline", state["stages"]["design_contract"].get("source"))
            # 재실행해도 기존 파일 보존(사람 편집본일 수 있다 — design_brief와 동일 계약).
            proposal_pipeline._go_design_contract(run, SimpleNamespace())
            self.assertEqual(first, design_contract.path(run).read_text(encoding="utf-8"))

    # ---- ③ chrome/image 2계약 분리 — 혼입 제거 ---------------------------------------------

    def test_image_contract_excludes_canvas_and_chrome_keys(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract = design_contract.build(
                run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
        self.assertNotIn("canvas", contract["image_contract"])
        self.assertNotIn("chrome", contract["image_contract"])
        self.assertIn("canvas", contract["chrome_contract"])
        self.assertIn("chrome", contract["chrome_contract"])
        # 공용 토큰(색·폰트)은 조립·이미지 양쪽에 필요하므로 겹쳐도 된다(혼입과 다른 문제).
        self.assertIn("colors", contract["image_contract"])
        self.assertIn("colors", contract["chrome_contract"])

    def test_bundle_with_contract_omits_chrome_structure_from_prompt(self):
        """imagedeck.bundle이 design_contract를 쓰면 프롬프트에 canvas/chrome 구조 정보가 없다
        (R5 — 종전엔 'Design contract' JSON 블록이 header_h/footer_h까지 실었다). DF1(프롬프트
        다이어트, DECK_FIRST_DESIGN.md §2-①) 이후로는 그 블록 자체가 완전 스킨이어도 사라졌으므로
        — 이 테스트는 "블록이 canvas/chrome을 뺀 채 존재"가 아니라 "그런 구조 키가 프롬프트
        어디에도 리터럴로 새지 않는다"를 확인한다(R5 취지의 DF1판 강화 성립)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            contract = design_contract.build(
                run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            rep = imagedeck.bundle(run, None, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertTrue(rep["design_contract"])
        # DF1: 종전 "Design contract (skin ...)" JSON 덤프 블록 자체가 폐지됐다.
        self.assertNotIn("Design contract (skin", prompt)
        self.assertNotIn('"canvas"', prompt)
        self.assertNotIn('"chrome"', prompt)
        # image_contract(계약 레벨)는 여전히 canvas/chrome을 안 갖는다(R5 원 취지, 계약 구조 불변).
        self.assertNotIn("canvas", contract["image_contract"])
        self.assertNotIn("chrome", contract["image_contract"])

    def test_resolve_chrome_skin_uses_contract_chrome_only(self):
        """compose/pptx가 읽는 _resolve_chrome_skin이 chrome_contract만 반환(image_contract 아님)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            contract = design_contract.build(
                run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            resolved = imagedeck._resolve_chrome_skin(run, manifest)
        self.assertEqual(contract["chrome_contract"], resolved)

    # ---- ④ design_contract 없는 run — 폴백(파일럿·레거시 회귀 보존) ------------------------

    def test_bundle_without_contract_preserves_legacy_prompt_shape(self):
        """design_contract.json이 없는 폴백(skin_path 하나를 chrome/image 양쪽에 그대로 쓰는
        레거시·파일럿 경로)도 DF1(DECK_FIRST_DESIGN.md §3 DF1 행) 이후로는 계약 경로와 동일한
        수렴 프롬프트 문법을 쓴다 — 종전엔 이 경로만 chrome/canvas까지 그대로 실었지만(구조
        키가 새는 "버그(?)"), 그 원인이던 JSON 덤프 블록 자체가 폐지됐으므로 이 경로도 더 이상
        새지 않는다. 회귀 방지 대상은 "레거시 문법 보존"이 아니라 manifest.full_skin=True 기록."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            skin_path = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            rep = imagedeck.bundle(run, skin_path, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertFalse(rep["design_contract"])
        self.assertEqual("inkline", rep["skin"])
        self.assertTrue(rep["full_skin"])  # 하위호환 메타는 여전히 True(is_full_skin 판정 불가 관례)
        self.assertNotIn("Design contract (skin", prompt)
        self.assertNotIn('"chrome"', prompt)
        self.assertNotIn('"canvas"', prompt)
        self.assertIn("Art direction (본문 이미지는 틀에 끼워지는 본문", prompt)

    def test_imagedeck_cmd_bundle_defaults_to_neutral_not_inkline(self):
        """CLI 폴백(design_contract 없음·design_brief.skin.value 없음)의 기본값이 더 이상
        inkline이 아니라 _neutral이다(W29 자동폴백 폐기, 용어 정의 이행 과제)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            self._storyline(run)
            (run / "design_brief.json").write_text(json.dumps({"brand": {}}), encoding="utf-8")
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                proposal_pipeline.imagedeck_cmd(SimpleNamespace(
                    run=run.name, bundle=True, collect=False, compose=False,
                    skin=None, wireframe_mode="off", ref=None, ab=None))
            manifest = json.loads((run / imagedeck.MANIFEST_NAME).read_text(encoding="utf-8"))
        self.assertEqual("_neutral", manifest["skin"])

    # ---- ⑤ theme_confirm 게이트(B1) — 선택·대시보드 ack 가능 -------------------------------

    def test_theme_confirm_is_optional_and_skippable(self):
        self.assertIn("theme_confirm", pipeline_state.OPTIONAL_CHECKPOINTS)
        self.assertIn("theme_confirm", pipeline_state.HUMAN_CHECKPOINTS)
        self.assertIn("theme_confirm", server.SKIPPABLE_ACK_GATES)

    def test_gate_requires_contract_then_theme_confirm_before_imagedeck_bundle(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            pipeline_state.set_render_route(run, "image_infographic", "off")  # 뼈대 스킵(D13) — 이 테스트의 관심사 아님
            self._storyline(run)
            (run / "design_brief.json").write_text(json.dumps({"brand": {}}), encoding="utf-8")
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            pipeline_state.record(run, "render")
            pipeline_state.record(run, "review_resolve")
            pipeline_state.record(run, "design_brief")
            pipeline_state.clear_checkpoint(run, "decision")
            # design_contract.json이 아직 없다 -> next는 그 단계를 만들라는 command.
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("command", nxt["kind"])
            self.assertEqual("design_contract", nxt.get("stage"))
            # 계약이 생기면 -> theme_confirm 관문(선택·human).
            contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            pipeline_state.record(run, "design_contract", source="neutral")
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt["kind"])
            self.assertEqual("theme_confirm", nxt["checkpoint"])
            self.assertTrue(nxt.get("optional"))
            # ack(대시보드 confirm 또는 --confirm 스킵과 동치) -> 이미지 라우트 기본이므로 imagedeck_bundle.
            pipeline_state.clear_checkpoint(run, "theme_confirm")
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("imagedeck_bundle", nxt.get("stage"))

    def test_theme_gate_not_required_for_run_already_past_it(self):
        """레거시/파일럿 run(design_contract 도입 전에 이미 imagedeck_bundle까지 진행)은
        design_contract·theme_confirm을 소급 요구하지 않는다(회귀 보존)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            pipeline_state.set_render_route(run, "image_infographic", "off")  # 뼈대 스킵(D13) — 이 테스트의 관심사 아님
            self._storyline(run)
            (run / "design_brief.json").write_text(json.dumps({"brand": {}}), encoding="utf-8")
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            pipeline_state.record(run, "render")
            pipeline_state.record(run, "review_resolve")
            pipeline_state.record(run, "design_brief")
            pipeline_state.record(run, "imagedeck_bundle")  # 이미 하류로 진행된 레거시 run
            pipeline_state.clear_checkpoint(run, "decision")
            nxt = pipeline_state.resolve(run)["next"]
            self.assertNotEqual("design_contract", nxt.get("stage"))
            self.assertNotEqual("theme_confirm", nxt.get("checkpoint"))
            self.assertFalse(design_contract.exists(run))  # 파일을 새로 만들라고도 하지 않는다.

    def test_go_cmd_creates_contract_and_stops_at_theme_confirm(self):
        """go 통합: design_brief 확정 뒤 첫 `go`가 계약을 만들고, theme_confirm에서 멈춘다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            pipeline_state.set_render_route(run, "image_infographic", "off")  # 뼈대 스킵(D13) — 이 테스트의 관심사 아님
            self._storyline(run)
            (run / "design_brief.json").write_text(
                json.dumps({"skin": {"value": "inkline"}, "brand": {}}), encoding="utf-8")
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            pipeline_state.record(run, "render")
            pipeline_state.record(run, "review_resolve")
            pipeline_state.record(run, "design_brief")
            pipeline_state.clear_checkpoint(run, "decision")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertTrue(design_contract.exists(run))
            self.assertIn("[DESIGN_CONTRACT]", out.getvalue())
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("theme_confirm", nxt.get("checkpoint"))
            # 07_테마확정 파생 뷰도 같은 go 호출에서 열린다(journey_folders 배선).
            view_path = (run / "journey" / journey_folders.FOLDERS["07"]
                         / journey_folders.DESIGN_CONTRACT_VIEW_NAME)
            self.assertTrue(view_path.is_file())
            self.assertIn("chrome_contract", view_path.read_text(encoding="utf-8"))

    # ---- W31 마찰15: 차용 = 중립 위 딥머지(대체 아님) ---------------------------------------

    def test_partial_skin_borrow_inherits_neutral_structure(self):
        """부분 스킨(colors/brand/_meta뿐인 institution_research --apply 산출)을 차용해도
        canvas/export/chrome 등 구조 키는 중립에서 상속된다 — 마찰15 재현의 핵심 회귀 방지."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            brief = {"skin": {"value": "예시공공기관"}}  # 창고 실물 부분 스킨(읽기 전용 사용)
            contract = design_contract.build(run, brief=brief, skins_dir=self.SKINS_DIR)
        self.assertEqual("예시공공기관", contract["meta"]["source"])
        self.assertTrue(contract["meta"]["merged_with_neutral"])
        chrome = contract["chrome_contract"]
        self.assertIn("canvas", chrome)
        self.assertIn("export", chrome)
        self.assertIn("chrome", chrome)
        self.assertEqual(1920, chrome["canvas"]["width"])  # 중립 구조 그대로 상속
        # 부분 스킨이 준 브랜드색은 살아있다(구조 상속과 공존 — 대체가 아니라 병합).
        self.assertEqual("005DAB", chrome["colors"]["navy"])
        self.assertEqual("0099DB", chrome["colors"]["orange"])
        # 중립의 다른 색 키(accent 등)도 사라지지 않는다(부분 스킨엔 없던 키).
        self.assertIn("accent", chrome["colors"])

    def test_full_skin_borrow_equivalent_to_direct_skin_file(self):
        """완전 스킨(inkline)을 딥머지해도 결과가 원본 스킨과 실질 동일하다(회귀 없음) —
        중립에만 있는 키가 추가될 수 있으나(무해), 겹치는 키의 값은 전부 차용 쪽이 이긴다."""
        inkline_raw = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract = design_contract.build(
                run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
        chrome = contract["chrome_contract"]
        self.assertEqual(inkline_raw["colors"], chrome["colors"])
        self.assertEqual(inkline_raw["canvas"], chrome["canvas"])
        self.assertEqual(inkline_raw["chrome"], chrome["chrome"])
        self.assertEqual(inkline_raw["export"], chrome["export"])
        self.assertEqual(inkline_raw["typography"], chrome["typography"])

    def test_validate_structure_raises_human_message_when_structure_missing(self):
        broken = {"chrome_contract": {"colors": {"accent": "000000"}}, "image_contract": {}}
        with self.assertRaises(design_contract.DesignContractError) as ctx:
            design_contract.validate_structure(broken)
        msg = str(ctx.exception)
        self.assertIn("canvas", msg)
        self.assertIn("export", msg)
        self.assertIn("중립 병합", msg)

    def test_build_raises_when_run_overrides_strip_required_structure_key(self):
        """run 조정 채널이 구조 키를 지우면(예: canvas를 null로) 동결이 사람 말 오류로 막힌다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / design_contract.RUN_OVERRIDES_NAME).write_text(
                json.dumps({"canvas": None}), encoding="utf-8")
            with self.assertRaises(design_contract.DesignContractError) as ctx:
                design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
        self.assertIn("canvas", str(ctx.exception))

    def test_imagedeck_bundle_survives_partial_skin_borrow(self):
        """마찰15 재현의 역검증: 부분 스킨 차용 계약으로 imagedeck.bundle을 돌려도 크래시하지
        않는다(이전엔 canvas_dims의 처리 안 된 traceback으로 죽었다)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            contract = design_contract.build(
                run, brief={"skin": {"value": "예시공공기관"}}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")  # 크래시하면 여기서 예외
        self.assertTrue(manifest["design_contract"])
        self.assertEqual("예시공공기관", manifest["skin"])

    # ---- W31 마찰14 ④: 재동결 통로(design_contract.prev.json + theme_confirm 재무장) -------

    def _frozen_run(self, td, skin_value="inkline"):
        run = Path(td)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        pipeline_state.set_render_route(run, "image_infographic", "off")
        self._storyline(run)
        (run / "design_brief.json").write_text(
            json.dumps({"skin": {"value": skin_value}, "brand": {}}), encoding="utf-8")
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        pipeline_state.record(run, "render")
        pipeline_state.record(run, "review_resolve")
        pipeline_state.record(run, "design_brief")
        pipeline_state.clear_checkpoint(run, "decision")
        contract = design_contract.build(run, brief=design_brief.load(run) or {}, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        pipeline_state.record(run, "design_contract", source=(contract.get("meta") or {}).get("source"))
        return run

    def test_go_warns_when_brief_is_newer_than_frozen_contract(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            # 계약보다 브리핑을 미래 시각으로 만들어(결정론 — 기존 재무장 테스트와 동일 관례) stale을 재현한다.
            os.utime(run / "design_brief.json", (2_000_000_000, 2_000_000_000))
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
        self.assertIn("[GO 경고]", out.getvalue())
        self.assertIn("--refreeze-contract", out.getvalue())

    def test_go_no_warning_when_contract_still_fresh(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
        self.assertNotIn("[GO 경고]", out.getvalue())

    def test_refreeze_contract_preserves_prev_and_rearms_theme_confirm(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td, skin_value="inkline")
            pipeline_state.clear_checkpoint(run, "theme_confirm")  # 이미 확정된 것으로 가정
            old_contract_text = design_contract.path(run).read_text(encoding="utf-8")
            # 사람이 브리핑을 고쳐 다른 스킨을 차용하기로 한다.
            brief = design_brief.load(run)
            brief["skin"]["value"] = "_neutral"
            design_brief.save(run, brief)
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                proposal_pipeline.go_cmd(SimpleNamespace(
                    run=run.name, confirm=False, json=False, refreeze_contract=True))
            self.assertIn("[REFREEZE]", out.getvalue())
            prev_path = run / "design_contract.prev.json"
            self.assertTrue(prev_path.is_file())
            self.assertEqual(old_contract_text, prev_path.read_text(encoding="utf-8"))
            new_contract = design_contract.load(run)
            self.assertEqual("_neutral", new_contract["meta"]["source"])
            # theme_confirm이 재동결된 계약(새 mtime)을 보고 다음 판정에서 재무장된다.
            nxt = pipeline_state.resolve(run)["next"]
            self.assertEqual("checkpoint", nxt.get("kind"))
            self.assertEqual("theme_confirm", nxt.get("checkpoint"))
            self.assertIn("재무장", nxt.get("why", ""))


class PromptReconstructionSmoke(unittest.TestCase):
    """DF1(프롬프트 다이어트, CONTEXT/DECK_FIRST_DESIGN.md §2-①·§3) — imagedeck bundle 프롬프트
    조립: 종전 "제약 수준 = 차용 수준"(완전 스킨만 계약 전체 JSON 덤프, W31 β1 마찰19) 두 갈래를
    하나로 수렴했다. full_skin 값과 무관하게 항상 같은 문법(_design_section) — 이웃 크롬 브리핑
    + 브랜드 고지(있으면) + 분량/넘침 행동 한계(있으면) + 색·배경 자유 명시 — 이 실리고,
    테마 정의(colors·color_roles·typography·layout·components)는 프롬프트에서 완전히 빠졌다
    (디자인 정본은 이제 HTML 틀뿐). full_skin 필드 자체는 계약/매니페스트 메타로만 남는다."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")

    def _prompt02(self, run, brief):
        self._storyline(run)
        contract = design_contract.build(run, brief=brief, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        imagedeck.bundle(run, None, wireframe_mode="off")
        return contract, (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")

    def test_neutral_contract_gets_art_direction_not_json_dump(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract, prompt = self._prompt02(run, {})
        self.assertFalse(design_contract.is_full_skin(contract))
        self.assertIn("Art direction (본문 이미지는 틀에 끼워지는 본문", prompt)
        self.assertIn("너는 이 디자인의 일부다", prompt)
        self.assertNotIn("Design contract (skin", prompt)
        # 테마 정의(색 역할·컴포넌트·거버넌스 키)가 프롬프트에 통째로 새지 않는다.
        self.assertNotIn("color_roles", prompt)
        self.assertNotIn("generation_rules", prompt)

    def test_full_skin_contract_no_longer_dumps_theme_spec(self):
        """DF1 핵심 회귀 방지: 완전 스킨(inkline)을 차용해도 더 이상 계약 전체 JSON을 프롬프트에
        싣지 않는다(종전 "제약 수준 = 차용 수준" 분기 폐지)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract, prompt = self._prompt02(run, {"skin": {"value": "inkline"}})
        self.assertTrue(design_contract.is_full_skin(contract))
        self.assertIn("Art direction (본문 이미지는 틀에 끼워지는 본문", prompt)
        self.assertIn("너는 이 디자인의 일부다", prompt)
        self.assertNotIn("Design contract (skin - 값 변경 금지):", prompt)
        self.assertNotIn("color_roles", prompt)
        self.assertNotIn("generation_rules", prompt)
        # inkline 고유 강조색 hex가 프롬프트에 나열되지 않는다(테마 사양 나열 0건).
        inkline = json.loads((self.SKINS_DIR / "inkline.json").read_text(encoding="utf-8"))
        accent = (inkline.get("colors") or {}).get("accent")
        if accent:
            self.assertNotIn(accent, prompt)

    def test_full_and_neutral_paths_converge_to_identical_prompt_grammar(self):
        """DF1 검증 기준: full_skin true/false 두 경로의 프롬프트가 동일 문법으로 수렴됐는지
        실제 조립 결과로 대조한다(추정이 아니라 실측) — 이 스토리라인/스킨 조합은 문법뿐 아니라
        바이트까지 동일하다(둘 다 중립 위 딥머지 구조라 content_limits/overflow_policy도 같다)."""
        with tempfile.TemporaryDirectory() as td1, tempfile.TemporaryDirectory() as td2:
            _, neutral_prompt = self._prompt02(Path(td1), {})
            _, full_prompt = self._prompt02(Path(td2), {"skin": {"value": "inkline"}})
        self.assertEqual(neutral_prompt, full_prompt)

    def test_constraints_drop_color_opinions_keep_behavior_rules(self):
        """색 취향(테마)은 빠지고, 오버플로 대응 같은 행동 규칙은 완전/축소 구분 없이 유지된다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            _, prompt = self._prompt02(run, {})
        self.assertIn("no font shrinking", prompt)  # 오버플로 A/B 분할 지시 — 행동 규칙이라 존치
        self.assertNotIn("no olive or green", prompt)  # 색 취향(테마) — DF1로 제거
        self.assertNotIn("no dark full background", prompt)  # 색 취향(테마) — DF1로 제거
        self.assertIn("no invented facts", prompt)
        self.assertIn("no photorealism, no 3D, no watermark", prompt)
        self.assertIn("하한 22px까지 축소를 허용", prompt)  # 중립 typography.minimum_body_size=22
        self.assertIn("색·서브 컬러·배경 연출·컬러 아이콘은 자유다", prompt)

    def test_content_limits_and_overflow_kept_as_behavior_rule(self):
        """content_limits/overflow_policy는 "디자인"이 아니라 "행동 한계"라 DF1에서도 남는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            _, prompt = self._prompt02(run, {})
        self.assertIn("분량/넘침 한계(행동 규칙 - 디자인 아님)", prompt)
        self.assertIn("content_limits", prompt)
        self.assertIn("overflow_policy", prompt)

    def test_partial_skin_contract_shows_brand_notice(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract, prompt = self._prompt02(run, {"skin": {"value": "예시공공기관"}})
        self.assertFalse(design_contract.is_full_skin(contract))  # self_contained=false(부분 스킨)
        self.assertIn("이 기관의 브랜드 색: #005DAB·#0099DB", prompt)
        self.assertIn("쓰라/마라 강제 없음", prompt)

    def test_neutral_contract_has_no_brand_notice(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            _, prompt = self._prompt02(run, {})
        self.assertNotIn("이 기관의 브랜드 색", prompt)

    def test_full_skin_flag_recorded_in_contract_and_manifest(self):
        """DF1 이후에도 full_skin 필드 자체는 계약/매니페스트에 하위호환으로 남는다(소비만 수렴
        — is_full_skin 호출부·describe_for_view 등은 그대로 이 값을 읽는다). 다만 이 값이 더 이상
        프롬프트 문법을 가르지 않는다(위 convergence 테스트로 별도 확인)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            contract = design_contract.build(
                run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
        self.assertTrue(manifest["full_skin"])
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertTrue(slide2["full_skin"])

    def test_legacy_fallback_without_contract_uses_unified_prompt(self):
        """design_contract.json이 없는 폴백 경로(파일럿·레거시)도 DF1 수렴 문법을 그대로 쓴다 —
        회귀 방지 대상은 "문법 차이"가 아니라 full_skin=True 기록 자체(is_full_skin 판정 불가 시
        관례 — 계약 기반 full_skin=false 경로와도 문법은 동일하다)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            skin_path = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin_path, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertTrue(manifest["full_skin"])
        self.assertIn("Art direction (본문 이미지는 틀에 끼워지는 본문", prompt)
        self.assertNotIn("Design contract (skin - 값 변경 금지):", prompt)


class ReferenceTiersSmoke(unittest.TestCase):
    """W31 β1(리허설 마찰20) — 레퍼런스 3계층 조회(장별 > 전체 > 시드) + 문단 유무 분기."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문A", "template_id": "strategy_pillars", "fields": {}},
            {"n": 3, "title": "본문B", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")

    def _bundle(self, run):
        self._storyline(run)
        contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        return imagedeck.bundle(run, None, wireframe_mode="off")

    def test_seed_used_when_nothing_in_run(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._bundle(run)
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("Reference roles:", prompt)
        self.assertIn("(기본값 - 교체 가능)", prompt)
        self.assertIn("original_style_reference.png", prompt)

    def test_global_ref_overrides_seed(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gdir = run / "imagedeck_refs" / "global"
            gdir.mkdir(parents=True)
            (gdir / "brand_style.png").write_bytes(b"\x89PNG\r\n")
            self._bundle(run)
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("brand_style.png", prompt)
        self.assertNotIn("(기본값 - 교체 가능)", prompt)
        self.assertNotIn("original_style_reference.png", prompt)

    def test_slide_specific_ref_overrides_global(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gdir = run / "imagedeck_refs" / "global"
            gdir.mkdir(parents=True)
            (gdir / "brand_style.png").write_bytes(b"\x89PNG\r\n")
            sdir = run / "imagedeck_refs" / "slides" / "02"
            sdir.mkdir(parents=True)
            (sdir / "slide2_only.png").write_bytes(b"\x89PNG\r\n")
            self._bundle(run)
            prompt2 = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            prompt3 = (run / imagedeck.PROMPTS_DIR / "03.md").read_text(encoding="utf-8")
        self.assertIn("slide2_only.png", prompt2)
        self.assertNotIn("brand_style.png", prompt2)   # 장별이 전체를 이긴다(02)
        self.assertIn("brand_style.png", prompt3)       # 03은 장별이 없어 전체로 폴백
        self.assertNotIn("slide2_only.png", prompt3)

    def test_no_references_anywhere_removes_paragraph(self):
        """시드 폴더까지 비면 'Reference roles' 문단 자체가 사라지고 자립 지시로 대체된다."""
        with tempfile.TemporaryDirectory() as td, tempfile.TemporaryDirectory() as empty_seed:
            run = Path(td)
            with mock.patch.object(imagedeck, "SEED_REFS_DIR", Path(empty_seed)):
                self._bundle(run)
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertNotIn("Reference roles:", prompt)
        self.assertIn("자립 지시(레퍼런스 없음)", prompt)
        self.assertIn("art direction", prompt)

    def test_cli_ref_overrides_everything(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gdir = run / "imagedeck_refs" / "global"
            gdir.mkdir(parents=True)
            (gdir / "brand_style.png").write_bytes(b"\x89PNG\r\n")
            self._storyline(run)
            contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            manifest = imagedeck.bundle(run, None, wireframe_mode="off", refs=["explicit_cli_ref.png"])
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("explicit_cli_ref.png", prompt)
        self.assertNotIn("brand_style.png", prompt)
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("cli", slide2["references_source"])


class DeckPreviewSmoke(unittest.TestCase):
    """DF4(CONTEXT/DECK_FIRST_DESIGN.md §2-③·§3 DF4 행) — 덱 프리뷰(틀+배경, 본문 비움) 렌더 +
    4계층 레퍼런스 조회(slide > global > deck_preview > seed). 완료 조건 ①②③:
    ①4계층 순서 정확(slide/global 최우선 불변 - deck_preview는 그 아래·seed보다 위)
    ②deck_preview 존재 시 bundle 프롬프트에 자동 첨부 ③계약 미동결 시 사람 말 오류.

    조회 순서·자동 첨부·계약 게이트는 가짜 PNG를 deck_preview/에 직접 심어 playwright 없이도
    검증한다(앵커 지시) - 실제 렌더(html_to_png 왕복)만 playwright 가용 시로 skipTest한다."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문A", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")

    def _frozen_run(self, run):
        self._storyline(run)
        contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        return contract

    def _plant_preview_png(self, run, cls="content"):
        d = run / "imagedeck_refs" / imagedeck.DECK_PREVIEW_DIRNAME
        d.mkdir(parents=True, exist_ok=True)
        (d / f"{cls}.png").write_bytes(b"\x89PNG\r\n")

    # -- 완료 조건 ①: 4계층 조회 순서 ----------------------------------------

    def test_deck_preview_used_when_no_slide_or_global_refs(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            self._plant_preview_png(run, "content")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("deck_preview", slide2["references_source"])
        self.assertIn("content.png", prompt)
        self.assertIn("deck preview - 이 이미지가 끼워질 완성 틀", prompt)

    def test_global_ref_still_overrides_deck_preview(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            self._plant_preview_png(run, "content")
            gdir = run / "imagedeck_refs" / "global"
            gdir.mkdir(parents=True)
            (gdir / "brand.png").write_bytes(b"\x89PNG\r\n")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("global", slide2["references_source"])
        self.assertIn("brand.png", prompt)
        self.assertNotIn("content.png", prompt)

    def test_slide_specific_ref_still_overrides_deck_preview(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            self._plant_preview_png(run, "content")
            sdir = run / "imagedeck_refs" / "slides" / "02"
            sdir.mkdir(parents=True)
            (sdir / "only.png").write_bytes(b"\x89PNG\r\n")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("slide", slide2["references_source"])

    def test_deck_preview_still_overrides_seed(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            self._plant_preview_png(run, "content")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("deck_preview", slide2["references_source"])

    def test_deck_preview_is_class_specific_not_cross_attached(self):
        """deck_preview/full_image.png만 있고 content.png는 없으면 content 클래스 장에는 이 층이
        붙지 않고(다른 클래스의 틀을 섞어 붙이면 오히려 혼선) 시드로 폴백한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            self._plant_preview_png(run, "full_image")  # content가 아니라 full_image만 존재
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)  # 본문A = content 클래스
        self.assertEqual("seed", slide2["references_source"])

    def test_no_references_anywhere_still_falls_back_to_freestanding(self):
        """4계층 전부 비면(deck_preview 포함) 기존 문법대로 자립 지시로 대체된다(회귀 없음)."""
        with tempfile.TemporaryDirectory() as td, tempfile.TemporaryDirectory() as empty_seed:
            run = Path(td)
            with mock.patch.object(imagedeck, "SEED_REFS_DIR", Path(empty_seed)):
                self._frozen_run(run)
                imagedeck.bundle(run, None, wireframe_mode="off")
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("자립 지시(레퍼런스 없음)", prompt)

    # -- 완료 조건 ③: 계약 미동결/클래스 미선언 시 사람 말 오류 ----------------

    def test_preview_requires_frozen_contract(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            with self.assertRaises(imagedeck.ImagedeckError) as ctx:
                imagedeck.render_deck_preview(run)
        self.assertIn("design_contract.json", str(ctx.exception))

    def test_preview_requires_slide_classes(self):
        """slide_classes 없는(레거시/부분) 계약은 프리뷰 대상 클래스가 없다는 사람 말 오류."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            contract = {"schema_version": 1, "meta": {"source": "x"},
                        "chrome_contract": {"canvas": {}, "export": {"width": 100, "height": 100}},
                        "image_contract": {}}
            design_contract.save(run, contract)
            with self.assertRaises(imagedeck.ImagedeckError) as ctx:
                imagedeck.render_deck_preview(run)
        self.assertIn("slide_classes", str(ctx.exception))

    # -- 실제 렌더(html_to_png 왕복) - playwright 가용 시에만 -------------------

    def test_render_deck_preview_produces_real_pngs_per_class(self):
        import rasterize
        if not rasterize.available():
            self.skipTest(rasterize.unavailable_reason())
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            rep = imagedeck.render_deck_preview(run)
            classes = {r["class"] for r in rep["rendered"]}
            self.assertEqual({"content", "full_image"}, classes)  # image=none(cover/toc/divider) 제외
            self.assertEqual({"cover", "toc", "divider"}, set(rep["skipped_classes"]))
            for r in rep["rendered"]:
                p = Path(r["out"])
                self.assertTrue(p.is_file())
                self.assertGreater(p.stat().st_size, 0)

    def test_render_deck_preview_then_bundle_auto_attaches_real_file(self):
        """완료 조건 ②의 실제 렌더 왕복판: 진짜 html_to_png로 구운 프리뷰가 bundle에 자동 첨부."""
        import rasterize
        if not rasterize.available():
            self.skipTest(rasterize.unavailable_reason())
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._frozen_run(run)
            imagedeck.render_deck_preview(run)
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        self.assertEqual("deck_preview", slide2["references_source"])
        self.assertIn("deck_preview", prompt)
        self.assertIn("content.png", prompt)


class RefreezeDownstreamNoticeSmoke(unittest.TestCase):
    """W31 β1(리허설 마찰18) — 재동결 후 하류(imagedeck 등) 재적용 안내 + status stale 표면화."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}}]},
            ensure_ascii=False), encoding="utf-8")

    def _frozen_run(self, td):
        run = Path(td)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        pipeline_state.set_render_route(run, "image_infographic", "off")
        self._storyline(run)
        (run / "design_brief.json").write_text(
            json.dumps({"skin": {"value": "inkline"}, "brand": {}}), encoding="utf-8")
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        pipeline_state.record(run, "render")
        pipeline_state.record(run, "review_resolve")
        pipeline_state.record(run, "design_brief")
        contract = design_contract.build(run, brief=design_brief.load(run) or {}, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        pipeline_state.record(run, "design_contract", source=(contract.get("meta") or {}).get("source"))
        return run

    def test_refreeze_notices_downstream_when_imagedeck_bundle_recorded(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            imagedeck.bundle(run, None, wireframe_mode="off")
            pipeline_state.record(run, "imagedeck_bundle")
            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                proposal_pipeline._refreeze_contract(run)
        self.assertIn("이미 이미지 단계를 지났다", out.getvalue())
        self.assertIn("imagedeck --bundle 재실행", out.getvalue())

    def test_refreeze_silent_when_no_downstream_recorded(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                proposal_pipeline._refreeze_contract(run)
        self.assertNotIn("이미 이미지 단계를 지났다", out.getvalue())
        self.assertNotIn("하류 단계", out.getvalue())

    def test_status_flags_contract_newer_than_manifest(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            imagedeck.bundle(run, None, wireframe_mode="off")
            manifest_path = run / imagedeck.MANIFEST_NAME
            contract_path = design_contract.path(run)
            os.utime(manifest_path, (1_000_000_000, 1_000_000_000))
            os.utime(contract_path, (1_000_000_100, 1_000_000_100))
            warnings = pipeline_state.resolve(run)["warnings"]
        self.assertTrue(any("[하류 stale]" in w for w in warnings))

    def test_status_silent_when_manifest_is_fresh(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._frozen_run(td)
            imagedeck.bundle(run, None, wireframe_mode="off")
            # 번들이 계약보다 나중(정상 순서) — 재번들 필요 없음.
            warnings = pipeline_state.resolve(run)["warnings"]
        self.assertFalse(any("[하류 stale]" in w for w in warnings))


class PromptOverlayDesignSignalSmoke(unittest.TestCase):
    """W31 R-C(R6·R9, CONTEXT/JOURNEY.md) — 프롬프트 오버레이(사람 편집물) + review_badges
    저점수 장 디자인 신호 주입. 회귀 방지:
    ①오버레이(imagedeck_prompts_local/NN.md)가 있으면 프롬프트 말미에 명확한 구획으로 병합
    ②재번들해도(오버레이 원본은 사람 쪽 폴더에 있으니) 매번 다시 붙는다
    ③오버레이가 없으면 프롬프트에 아무 영향이 없다(구획 자체가 생기지 않는다)
    ④저점수(밋밋·발산추천) 장에는 배경이미지·디자인지식 적용 권장 신호가 자동 주입된다
    ⑤충실 장에는 아무것도 주입되지 않는다(과잉 장식 방지)
    ⑥채점 데이터(design_brief.json) 자체가 없는 run은 크래시도, 강제 주입도 없이 조용히 생략
    ⑦오버레이 폴더는 시스템이 빈 폴더+README만 만들고 NN.md는 절대 만들거나 지우지 않는다."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "전략 개요", "template_id": "strategy_pillars", "fields": {}},
            {"n": 3, "title": "일정", "template_id": "process_steps", "fields": {}},
        ]}, ensure_ascii=False), encoding="utf-8")

    def _design_brief(self, run, verdicts):
        rhythm = [{"slide_id": sid, "verdict": v} for sid, v in verdicts.items()]
        (run / "design_brief.json").write_text(
            json.dumps({"page_rhythm": {"slides": rhythm}}, ensure_ascii=False), encoding="utf-8")

    def _skin(self):
        return imagedeck.resolve_skin("inkline", self.SKINS_DIR)

    def _slide(self, manifest, n):
        return next(s for s in manifest["slides"] if s.get("n") == n)

    def test_bundle_opens_overlay_dir_with_readme_only(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            local = run / imagedeck.PROMPTS_LOCAL_DIR
            self.assertTrue(local.is_dir())
            self.assertTrue((local / "README.md").is_file())
            self.assertEqual(["README.md"], [p.name for p in local.iterdir()])

    def test_overlay_merges_into_generated_prompt(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            (run / imagedeck.PROMPTS_LOCAL_DIR / "02.md").write_text(
                "추가 지시: 청록색 배경을 강조해줘.", encoding="utf-8")
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertIn(
                "## 사람 추가 지시(오버레이 — imagedeck_prompts_local/02.md)", prompt)
            self.assertIn("청록색 배경을 강조해줘", prompt)
            self.assertTrue(self._slide(manifest, 2)["overlay_merged"])

    def test_overlay_survives_repeated_rebundle(self):
        """재번들(prompts_dir의 이전 *.md는 매번 삭제 후 재생성)해도 오버레이 원본은 사람 쪽
        폴더에 그대로 있으니, 병합 결과가 매번 다시 나타난다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            overlay_path = run / imagedeck.PROMPTS_LOCAL_DIR / "02.md"
            overlay_path.write_text("이 지시는 살아남아야 한다.", encoding="utf-8")
            for _ in range(3):
                imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertIn("이 지시는 살아남아야 한다", prompt)
            # 오버레이 원본 자체는 시스템이 절대 건드리지 않는다(정본은 사람 쪽).
            self.assertEqual("이 지시는 살아남아야 한다.", overlay_path.read_text(encoding="utf-8"))

    def test_no_overlay_means_no_prompt_effect(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertNotIn("사람 추가 지시", prompt)
            self.assertFalse(self._slide(manifest, 2)["overlay_merged"])

    def test_low_score_verdicts_inject_design_signal_block(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            self._design_brief(run, {1: "충실", 2: "발산추천", 3: "밋밋"})
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            p2 = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            p3 = (run / imagedeck.PROMPTS_DIR / "03.md").read_text(encoding="utf-8")
            self.assertIn("디자인 지표(review_badges 채점", p2)
            self.assertIn("발산추천", p2)
            self.assertIn("배경 이미지 생성을 적극 권장", p2)
            self.assertIn("디자인 지표(review_badges 채점", p3)
            self.assertIn("밋밋", p3)
            self.assertTrue(self._slide(manifest, 2)["design_signal_injected"])
            self.assertTrue(self._slide(manifest, 3)["design_signal_injected"])
            # 표지(html 전용 - 프롬프트 없음)도 지표 뷰용으로 verdict는 기록되지만 주입은 없다.
            cover = self._slide(manifest, 1)
            self.assertEqual("충실", cover["design_verdict"])
            self.assertFalse(cover["design_signal_injected"])
            self.assertFalse(cover["overlay_merged"])

    def test_substantial_verdict_gets_no_injection(self):
        """충실(thin 아님) 장은 과잉 장식 방지를 위해 신호를 주입하지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            self._design_brief(run, {1: "충실", 2: "충실", 3: "충실"})
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            p2 = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            p3 = (run / imagedeck.PROMPTS_DIR / "03.md").read_text(encoding="utf-8")
            self.assertNotIn("디자인 지표(review_badges 채점", p2)
            self.assertNotIn("디자인 지표(review_badges 채점", p3)
            self.assertFalse(self._slide(manifest, 2)["design_signal_injected"])
            self.assertFalse(self._slide(manifest, 3)["design_signal_injected"])

    def test_missing_score_data_gracefully_skipped(self):
        """design_brief.json 자체가 없는 run(파일럿·직접 bundle 호출)은 채점 신호를
        강제하지 않고 조용히 생략한다(크래시 없음)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            self.assertFalse((run / "design_brief.json").exists())
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            for s in manifest["slides"]:
                self.assertIsNone(s.get("design_verdict"))
                self.assertFalse(s.get("design_signal_injected"))
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertNotIn("디자인 지표(review_badges 채점", prompt)


class JourneyDesignMetricsViewSmoke(unittest.TestCase):
    """W31 R-C(R9, CONTEXT/JOURNEY.md) — 08_프롬프트확인/디자인지표_읽기.md 파생 뷰.
    정본=imagedeck_manifest.json(파생 뷰 규칙 준수) — 저점수/충실/미채점 3가지 표기,
    매뉴얼(포인터)의 오버레이 현황 안내까지 회귀 방지."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _storyline(self, run):
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {}},
            {"n": 2, "title": "전략 개요", "template_id": "strategy_pillars", "fields": {}},
        ]}, ensure_ascii=False), encoding="utf-8")

    def _skin(self):
        return imagedeck.resolve_skin("inkline", self.SKINS_DIR)

    def test_metrics_view_renders_verdict_table_and_manual_pointer(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            (run / "design_brief.json").write_text(json.dumps({"page_rhythm": {"slides": [
                {"slide_id": 1, "verdict": "충실"}, {"slide_id": 2, "verdict": "밋밋"}]}},
                ensure_ascii=False), encoding="utf-8")
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            journey_folders.sync(run)
            folder08 = journey_folders.folder_path(run, "08")
            view = (folder08 / journey_folders.DESIGN_METRICS_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("밋밋", view)
            self.assertIn("배경이미지 생성 권장", view)
            self.assertIn("주입 없음(충실", view)
            manual = (folder08 / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
            self.assertIn(journey_folders.DESIGN_METRICS_VIEW_NAME, manual)
            self.assertIn(journey_folders.PROMPTS_LOCAL_DIRNAME, manual)

    def test_metrics_view_reports_no_scoring_when_design_brief_absent(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            journey_folders.sync(run)
            view = (journey_folders.folder_path(run, "08")
                    / journey_folders.DESIGN_METRICS_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("채점 미실시", view)
            self.assertIn("채점 데이터 없음", view)

    def test_manual_overlay_pointer_reflects_existing_overlay_files(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._storyline(run)
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            (run / imagedeck.PROMPTS_LOCAL_DIR / "02.md").write_text("추가 지시", encoding="utf-8")
            journey_folders.sync(run)
            manual = (journey_folders.folder_path(run, "08")
                      / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
            self.assertIn("02.md", manual)
            self.assertIn("오버레이 존재 현황", manual)


class GatesProfileSmoke(unittest.TestCase):
    """W31 리허설 마찰2(관문 다이얼 + 조건부 승격, CONTEXT/REHEARSAL_FRICTIONS_W31.md #2) 회귀 방지.

    검증 대상: ①프로파일별 정지/통과 목록 ②조건부 승격(자동 대상도 신호 나쁘면 재정지)
    ③신호 부재 시 보수적 기본값(스킵 가능=통과·비스킵=정지) ④auto 기록이 사람 confirm과 구분됨
    ⑤go --gates 중도 변경 지속 ⑥go_cmd 통합(사람 ack 없이도 auto_pass 관문을 조용히 통과).
    """

    @staticmethod
    def _stg(t):
        return {"at": f"2026-07-22T00:{t:02d}:00"}

    def _wireframe_review_run(self, td, *, bad_signal: bool) -> Path:
        """wireframe_apply 완료 · decision 청산 · wireframe_review 미청산 — [3]->[4] 경계.
        (OptionalCheckpointsSmoke._wireframe_phase와 동형 — 이 클래스 전속 사본.)"""
        run = Path(td) / "run"
        run.mkdir(parents=True)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        pipeline_state.clear_checkpoint(run, "decision")
        st = pipeline_state.load(run)
        st["stages"].update({
            "render": self._stg(2), "review_resolve": self._stg(2),
            "wireframe_bundle": self._stg(3), "wireframe_apply": self._stg(4),
        })
        pipeline_state.save(run, st)
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        (run / "wireframe.json").write_text("{}", encoding="utf-8")
        counts = ({"발산추천": 2, "밋밋": 1, "충실": 0} if bad_signal
                  else {"발산추천": 0, "밋밋": 0, "충실": 3})
        (run / "gating_report.json").write_text(json.dumps({
            "review_needed_total": 4 if bad_signal else 0,
            "review_badges": {"counts": counts},
        }, ensure_ascii=False), encoding="utf-8")
        return run

    # ① 프로파일별 정지/통과 목록 -------------------------------------------------
    def test_profile_defaults_match_dial_design(self):
        # full: 전 관문 정지(현행 편집 모드와 동일).
        for gid in gates.GATE_IDS:
            self.assertEqual("stop", gates.PROFILE_DEFAULTS["full"][gid])
        # standard(기본): 회의 관문(theme_confirm·비스킵 2종)만 정지, 나머지 3종은 자동 대상.
        std = gates.PROFILE_DEFAULTS["standard"]
        for gid in ("design_refs", "skeleton_review", "wireframe_review"):
            self.assertEqual("auto", std[gid])
        for gid in ("theme_confirm", "imagedeck_prompt_ack", "imagedeck_ack"):
            self.assertEqual("stop", std[gid])
        # express: 스킵 가능 4종 + 비스킵 2종 전부 '자동' 분류(단, 비스킵은 decide()가 조건부로 되돌린다).
        for gid in gates.GATE_IDS:
            self.assertEqual("auto", gates.PROFILE_DEFAULTS["express"][gid])

    def test_no_gates_json_defaults_to_standard(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            cfg = gates.load_config(run)
        self.assertEqual("standard", cfg["profile"])
        self.assertEqual("default", cfg["source"])
        self.assertEqual({}, cfg["overrides"])

    def test_full_profile_forces_stop_even_with_clean_signal(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)  # 신호 깨끗
            gates.save_config(run, profile="full")
            d = gates.decide(run, "wireframe_review")
        self.assertEqual("stop", d["action"])
        self.assertEqual("stop", d["mode"])  # 조건부 승격이 아니라 프로파일 자체가 정지

    # ② 조건부 승격 ---------------------------------------------------------------
    def test_conditional_promotion_stops_auto_gate_on_bad_signal(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=True)  # standard 기본(auto) + 신호 나쁨
            d = gates.decide(run, "wireframe_review")
        self.assertEqual("auto", d["mode"])       # 프로파일상 자동 통과 대상이었으나
        self.assertEqual("stop", d["action"])     # 신호가 나빠 재정지(조건부 승격)
        self.assertTrue(d["signal"]["bad"])
        self.assertIn("검토요망", d["reason"])

    def test_clean_signal_auto_passes_under_standard(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)
            d = gates.decide(run, "wireframe_review")
        self.assertEqual("auto_pass", d["action"])

    # ③ 신호 부재 — 보수적 기본값 ---------------------------------------------------
    def test_missing_signal_permissive_for_skippable_gate(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)  # gating_report.json 자체가 없다(스켈레톤 검토는 render 이전 지점)
            d = gates.decide(run, "skeleton_review")
        self.assertFalse(d["signal"]["available"])
        self.assertEqual("auto_pass", d["action"])

    def test_missing_signal_conservative_for_non_skippable_gate(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gates.save_config(run, profile="express")  # 가장 관대한 프로파일이어도
            d = gates.decide(run, "imagedeck_prompt_ack")  # imagedeck_manifest.json 없음
        self.assertFalse(d["signal"]["available"])
        self.assertEqual("stop", d["action"])  # 비스킵 2종은 신호 없으면 보수적으로 정지

    def test_imagedeck_collect_failure_escalates_ack_gate(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gates.save_config(run, profile="express")
            (run / "imagedeck_collect.json").write_text(json.dumps({
                "slides": [{"n": 1, "status": "ok"}, {"n": 2, "status": "px_mismatch"}],
            }, ensure_ascii=False), encoding="utf-8")
            d = gates.decide(run, "imagedeck_ack")
        self.assertEqual("stop", d["action"])
        self.assertIn("수거 검증 불합격", d["reason"])

    # ④ auto 기록이 사람 confirm과 구분됨 -------------------------------------------
    def test_auto_ack_is_distinguishable_from_dashboard_ack(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)
            d = gates.decide(run, "wireframe_review")
            gates.write_auto_ack(run, "wireframe_review", d)
            # read_ack(대시보드 전용 필터)는 auto 기록을 사람 ack로 인정하지 않는다.
            self.assertIsNone(pipeline_state.read_ack(run, "wireframe_review"))
            # 표시 전용 read_any_ack는 via='auto'·decision='auto'로 구분해서 보여준다.
            peek = pipeline_state.read_any_ack(run, "wireframe_review")
        self.assertEqual("auto", peek["via"])
        self.assertEqual("auto", peek["decision"])
        self.assertIn("reason", peek)

    # ⑤ 중도 변경 지속 -------------------------------------------------------------
    def test_go_gates_flag_persists_profile_mid_course(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir(parents=True)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                proposal_pipeline.go_cmd(SimpleNamespace(
                    run=run.name, confirm=False, json=False, gates="express"))
            self.assertEqual("express", gates.load_config(run)["profile"])

    # ⑥ go_cmd 통합 — 사람 ack 없이 auto_pass 관문을 조용히 통과 ------------------------
    def test_go_cmd_auto_passes_wireframe_review_when_signal_clean(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_go_design_brief", return_value=True), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNotNone(
                pipeline_state.load(run)["checkpoints"]["wireframe_review"]["cleared_at"]
            )
            ack = pipeline_state.read_any_ack(run, "wireframe_review")
        self.assertEqual("auto", ack["via"])
        self.assertIn("[GO] 자동 통과", out.getvalue())
        # 사람 ack로는 인정되지 않는다(read_ack은 여전히 None) — confirm과 구분 유지.
        self.assertIsNone(pipeline_state.read_ack(run, "wireframe_review"))

    def test_go_cmd_stops_at_wireframe_review_when_signal_bad(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=True)
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNone(
                pipeline_state.load(run)["checkpoints"]["wireframe_review"]["cleared_at"]
            )
        self.assertIn("waiting_human:wireframe_review", out.getvalue())
        self.assertIn("신호 나쁨", out.getvalue())

    def test_go_cmd_full_profile_stops_wireframe_review_regardless_of_signal(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)  # 신호는 깨끗해도
            gates.save_config(run, profile="full")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIsNone(
                pipeline_state.load(run)["checkpoints"]["wireframe_review"]["cleared_at"]
            )
        self.assertIn("waiting_human:wireframe_review", out.getvalue())

    def test_dashboard_ack_still_wins_over_auto_pass(self):
        """사람이 이미 대시보드에서 ack했다면(비록 gates가 auto_pass 판정이어도) 그 ack가 소비된다
        (기록에는 via=dashboard로 남아 auto와 절대 섞이지 않는다) — W27 ack 기계 호환 불변."""
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td, bad_signal=False)
            path = pipeline_state.ack_path(run, "wireframe_review")
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(json.dumps({
                "gate": "wireframe_review", "decision": "confirm",
                "at": "2999-01-01T00:00:00+09:00", "via": "dashboard",
            }), encoding="utf-8")
            out = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_go_design_brief", return_value=True), \
                 mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc)
            self.assertIn("via=dashboard", out.getvalue())
            ack = pipeline_state.read_any_ack(run, "wireframe_review")
        self.assertEqual("dashboard", ack["via"])


class JourneyCheckSmoke(unittest.TestCase):
    """W31 리허설 마찰4(REHEARSAL_FRICTIONS_W31.md #4) — journey 폴더의 검토_체크.md 채널.

    검증: ①발급(체크박스·라운드 토큰·사람 전속 문구) ②멱등(같은 라운드 재작성 안 함)
    ③수거(체크+토큰 일치 → ack) ④토큰 불일치(구 라운드) 무시 ⑤재무장 시 새 토큰 미체크 재발급
    ⑥대시보드 ack가 이미 있으면 덮어쓰지 않음("먼저 온 쪽이 이김") ⑦pipeline_state.read_ack가
    via="journey_check"도 사람 ack로 인정 ⑧go_cmd 통합 ⑨Claude 대리 금지 문구(파일 + CLAUDE.md).
    """

    def _stg(self, t):
        return {"at": f"2026-07-22T01:{t:02d}:00"}

    def _wireframe_review_run(self, td) -> Path:
        """wireframe_apply 완료 · decision 청산 · wireframe_review 미청산 — [3]->[4] 경계.
        관문 다이얼을 full로 고정해 gates.py 자동 통과가 이 테스트에 개입하지 않게 한다."""
        run = Path(td) / "run"
        run.mkdir(parents=True)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        gates.save_config(run, profile="full")
        pipeline_state.clear_checkpoint(run, "decision")
        st = pipeline_state.load(run)
        st["stages"].update({
            "render": self._stg(2), "review_resolve": self._stg(2),
            "wireframe_bundle": self._stg(3), "wireframe_apply": self._stg(4),
        })
        pipeline_state.save(run, st)
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        (run / "wireframe.json").write_text("{}", encoding="utf-8")
        return run

    # ① 발급 -----------------------------------------------------------------
    def test_issue_creates_form_with_checkboxes_token_and_human_only_header(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            path = journey_check.issue(run, "wireframe_review")
            self.assertIsNotNone(path)
            text = path.read_text(encoding="utf-8")
            self.assertIn("사람 전속", text)
            self.assertIn("Claude", text)
            self.assertIn("- [ ] 검토 완료", text)
            self.assertIn("건너뛰기", text)  # wireframe_review는 스킵 가능 관문
            self.assertIsNotNone(journey_check._extract_token(text))

    def test_non_skippable_gate_omits_skip_checkbox(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "12")
            folder.mkdir(parents=True)
            path = journey_check.issue(run, "design")
            text = path.read_text(encoding="utf-8")
            self.assertIn("- [ ] 검토 완료", text)
            self.assertNotIn("- [ ] 건너뛰기", text)
            self.assertIn("건너뛸 수 없음", text)

    def test_no_folder_mapping_or_unopened_folder_yields_no_issue(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            # design_refs는 journey 폴더 매핑 자체가 없다(기존 갭 — 대시보드 채널만).
            self.assertIsNone(journey_check.issue(run, "design_refs"))
            # wireframe_review는 매핑은 있으나 06 폴더가 아직 열리지 않았다.
            self.assertIsNone(journey_check.issue(run, "wireframe_review"))

    def test_issue_idempotent_within_same_round_preserves_human_check(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            path.write_text(
                path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료"),
                encoding="utf-8",
            )
            again = journey_check.issue(run, "wireframe_review")
            self.assertIsNone(again)  # 같은 라운드 — 재작성 안 함
            self.assertIn("- [x] 검토 완료", path.read_text(encoding="utf-8"))

    # ② 수거 -------------------------------------------------------------------
    def test_collect_ack_confirms_when_checked_and_token_matches(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            path.write_text(
                path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료"),
                encoding="utf-8",
            )
            written = journey_check.collect_ack(run, "wireframe_review")
            self.assertIsNotNone(written)
            ack = pipeline_state.read_ack(run, "wireframe_review")
            self.assertEqual("confirm", ack["decision"])
            self.assertEqual("journey_check", ack["via"])

    def test_collect_ack_skip_only_for_skippable_gate(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            path.write_text(
                path.read_text(encoding="utf-8").replace("- [ ] 건너뛰기(스킵 가능 관문만 표기)",
                                                          "- [x] 건너뛰기(스킵 가능 관문만 표기)"),
                encoding="utf-8",
            )
            journey_check.collect_ack(run, "wireframe_review")
            ack = pipeline_state.read_ack(run, "wireframe_review")
            self.assertEqual("skip", ack["decision"])

    # ③ 토큰 불일치 무시 ----------------------------------------------------------
    def test_collect_ack_ignores_stale_round_token(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            text = path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료")
            text = journey_check._TOKEN_RE.sub("<!-- round-token: 000000000000 -->", text)
            path.write_text(text, encoding="utf-8")
            self.assertIsNone(journey_check.collect_ack(run, "wireframe_review"))
            self.assertIsNone(pipeline_state.read_ack(run, "wireframe_review"))

    # ④ 재무장 — 새 토큰 미체크 재발급 -----------------------------------------------
    def test_reissue_on_rearm_gives_fresh_unchecked_token(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            wf = run / "wireframe.json"
            wf.write_text("{}", encoding="utf-8")
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            old_token = journey_check._extract_token(path.read_text(encoding="utf-8"))
            path.write_text(
                path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료"),
                encoding="utf-8",
            )
            self.assertIsNotNone(journey_check.collect_ack(run, "wireframe_review"))  # 이번 라운드 소비·기록
            # 감시 산출물 변경(재무장 조건) — mtime을 명확히 미래로 옮겨 토큰이 확실히 달라지게 한다.
            future = wf.stat().st_mtime + 100
            wf.write_text('{"changed": true}', encoding="utf-8")
            os.utime(wf, (future, future))
            reissued = journey_check.issue(run, "wireframe_review")
            self.assertIsNotNone(reissued)
            new_text = path.read_text(encoding="utf-8")
            new_token = journey_check._extract_token(new_text)
            self.assertNotEqual(old_token, new_token)
            self.assertIn("- [ ] 검토 완료", new_text)  # 미체크로 재발급
            # 이전 체크는 유실이 아니라 이미 ack json에 기록되어 있다.
            self.assertEqual("journey_check", pipeline_state.read_any_ack(run, "wireframe_review")["via"])

    # ⑤ 대시보드 등가 — 먼저 온 쪽이 이김 --------------------------------------------
    def test_dashboard_ack_wins_and_is_not_overwritten(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            ack_path = pipeline_state.ack_path(run, "wireframe_review")
            ack_path.parent.mkdir(parents=True, exist_ok=True)
            ack_path.write_text(json.dumps({
                "gate": "wireframe_review", "decision": "confirm",
                "at": "2999-01-01T00:00:00+09:00", "via": "dashboard",
            }), encoding="utf-8")
            journey_check.issue(run, "wireframe_review")
            path = folder / journey_check.CHECK_NAME
            path.write_text(
                path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료"),
                encoding="utf-8",
            )
            self.assertIsNone(journey_check.collect_ack(run, "wireframe_review"))
            self.assertEqual("dashboard", pipeline_state.read_ack(run, "wireframe_review")["via"])

    def test_pipeline_state_read_ack_accepts_journey_check_via(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            ack_path = pipeline_state.ack_path(run, "decision")
            ack_path.parent.mkdir(parents=True, exist_ok=True)
            ack_path.write_text(json.dumps({
                "gate": "decision", "decision": "confirm",
                "at": "2999-01-01T00:00:00+09:00", "via": "journey_check",
            }), encoding="utf-8")
            ack = pipeline_state.read_ack(run, "decision")
        self.assertIsNotNone(ack)
        self.assertEqual("journey_check", ack["via"])

    # ⑥ go_cmd 통합 --------------------------------------------------------------
    def test_go_cmd_issues_then_consumes_journey_check(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td)
            out1 = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch("sys.stdout", new=out1):
                rc1 = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc1)
            self.assertIn("waiting_human:wireframe_review", out1.getvalue())
            check_path = journey_folders.folder_path(run, "06") / journey_check.CHECK_NAME
            self.assertTrue(check_path.is_file())

            check_path.write_text(
                check_path.read_text(encoding="utf-8").replace("- [ ] 검토 완료", "- [x] 검토 완료"),
                encoding="utf-8",
            )
            out2 = io.StringIO()
            with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run), \
                 mock.patch.object(proposal_pipeline, "_go_design_brief", return_value=True), \
                 mock.patch("sys.stdout", new=out2):
                rc2 = proposal_pipeline.go_cmd(SimpleNamespace(run=run.name, confirm=False, json=False))
            self.assertEqual(0, rc2)
            self.assertIn("via=journey_check", out2.getvalue())
            self.assertIsNotNone(
                pipeline_state.load(run)["checkpoints"]["wireframe_review"]["cleared_at"]
            )

    # ⑦ Claude 대리 금지 문구 ------------------------------------------------------
    @_requires_claude_md
    def test_claude_prohibition_text_present_in_form_and_claude_md(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            folder = journey_folders.folder_path(run, "06")
            folder.mkdir(parents=True)
            path = journey_check.issue(run, "wireframe_review")
            text = path.read_text(encoding="utf-8")
        self.assertIn("Claude", text)
        self.assertIn("대신", text)
        claude_md = (REPO_ROOT / "CLAUDE.md").read_text(encoding="utf-8")
        self.assertIn("검토_체크.md", claude_md)
        self.assertIn("Claude", claude_md)
        self.assertIn("대리", claude_md)


class JourneyOutputViewSmoke(unittest.TestCase):
    """W31 리허설 마찰5(REHEARSAL_FRICTIONS_W31.md #5) — 산출물.html(클릭 링크+계보) + _전체여정.html."""

    def test_output_view_lists_artifact_link_and_lineage(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            st = pipeline_state.load(run)
            st["stages"].update({
                "message_map": {"at": "2026-07-22T01:00:00"},
                "skeleton": {"at": "2026-07-22T02:00:00"},
                "storyline_bundle": {"at": "2026-07-22T03:00:00"},
            })
            pipeline_state.save(run, st)
            (run / "message_map.json").write_text("{}", encoding="utf-8")
            (run / "skeleton.json").write_text("{}", encoding="utf-8")
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            report = journey_folders.sync(run)
            view_path = journey_folders.folder_path(run, "04") / journey_folders.OUTPUT_VIEW_NAME
            self.assertTrue(view_path.is_file())
            self.assertIn(f"04/{journey_folders.OUTPUT_VIEW_NAME}", report["views_rendered"])
            html = view_path.read_text(encoding="utf-8")
            self.assertIn("storyline.json", html)  # 상대 링크(href)
            self.assertIn("2026-07-22T03:00:00", html)  # 완료 시각
            self.assertIn("메시지맵", html)  # 계보 사슬에 상위 단계 등장
            self.assertIn("자동 생성", html)
            # 이중화 금지: 정본 JSON 내용(사본)이 아니라 링크·계보 텍스트만.
            self.assertNotIn('"n": 1', html)

    def test_output_view_absent_when_folder_open_but_catalog_artifact_not_yet_produced(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "wireframe_prompt").mkdir(parents=True)
            (run / "wireframe_prompt" / "prompt.md").write_text("prompt", encoding="utf-8")
            journey_folders.sync(run)
            folder06 = journey_folders.folder_path(run, "06")
            self.assertTrue(folder06.is_dir())  # 06은 열렸다(프롬프트 존재)
            view_path = folder06 / journey_folders.OUTPUT_VIEW_NAME
            self.assertFalse(view_path.is_file())  # wireframe.json(정본)이 아직 없어 산출물이 없다

    def test_overview_maps_steps_and_marks_current_position(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            report = journey_folders.sync(run)
            overview_path = journey_folders.journey_root(run) / journey_folders.OVERVIEW_NAME
            self.assertTrue(overview_path.is_file())
            self.assertIn(journey_folders.OVERVIEW_NAME, report["views_rendered"])
            html = overview_path.read_text(encoding="utf-8")
            self.assertIn("04_내용만들기", html)
            self.assertIn("14_마무리", html)
            self.assertIn("현재 위치", html)


class StatusPendingAckSmoke(unittest.TestCase):
    """W31 리허설 마찰7(REHEARSAL_FRICTIONS_W31.md #7) — status의 미소비 ack 표시.

    상태 판정(cleared_at) 자체는 바뀌지 않는다 — go만이 소비한다(불변). status는 표시만 덧붙인다.
    """

    def _wireframe_review_run(self, td) -> Path:
        run = Path(td) / "run"
        run.mkdir(parents=True)
        pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
        pipeline_state.clear_checkpoint(run, "decision")
        st = pipeline_state.load(run)
        st["stages"].update({
            "render": {"at": "2026-07-22T01:00:00"},
            "review_resolve": {"at": "2026-07-22T01:00:00"},
            "wireframe_bundle": {"at": "2026-07-22T01:00:00"},
            "wireframe_apply": {"at": "2026-07-22T01:00:00"},
        })
        pipeline_state.save(run, st)
        (run / "deck.html").write_text("<html></html>", encoding="utf-8")
        (run / "wireframe.json").write_text("{}", encoding="utf-8")
        return run

    def test_status_shows_pending_ack_before_go_consumes(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td)
            ack_path = pipeline_state.ack_path(run, "wireframe_review")
            ack_path.parent.mkdir(parents=True, exist_ok=True)
            ack_path.write_text(json.dumps({
                "gate": "wireframe_review", "decision": "confirm",
                "at": "2999-01-01T00:00:00+09:00", "via": "dashboard",
            }), encoding="utf-8")
            view = pipeline_state.resolve(run)
            text = pipeline_state.format_status(view)
            self.assertIn("waiting_human:wireframe_review", text)
            self.assertIn("ack 있음(대시보드/폴더 체크)", text)
            self.assertIn("via=dashboard", text)
            self.assertIsNone(
                pipeline_state.load(run)["checkpoints"]["wireframe_review"]["cleared_at"]
            )  # 상태 판정 로직 불변 — 표시만 추가됐다.

    def test_status_silent_when_no_ack_yet(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._wireframe_review_run(td)
            view = pipeline_state.resolve(run)
            text = pipeline_state.format_status(view)
            self.assertIn("waiting_human:wireframe_review", text)
            self.assertNotIn("ack 있음", text)


class ReviewBadgesContentUnitsSmoke(unittest.TestCase):
    """W31 리허설 마찰 16 — 채점 내용 단위 현대화(불릿+채워진 필드 리프, 예시 제외).
    구 기준은 불릿만 세어 필드형 덱 전 장이 자동 '얇음'(22/22 실측)이었다."""

    def _rb(self):
        import review_badges
        return review_badges

    def _slide(self, **kw):
        base = {"slide_id": 1, "role": "process", "title": "프로세스",
                "key_message": "이 장의 핵심 메시지는 충분히 길다 스무자 넘김",
                "body": [], "fields": {}, "review_needed": []}
        base.update(kw)
        return base

    def test_field_rich_slide_without_bullets_is_substantial(self):
        rb = self._rb()
        deck = {"slides": [self._slide(fields={"steps": [
            {"label": "착수", "desc": "요구 분석"},
            {"label": "설계", "desc": "구조 확정"},
        ]})]}
        badge = rb.compute_review_badges(deck)["badges"][0]
        self.assertEqual("충실", badge["verdict"])
        self.assertGreaterEqual(badge["signals"]["content_units"], 3)

    def test_example_marked_fields_do_not_count_as_content(self):
        rb = self._rb()
        deck = {"slides": [self._slide(fields={
            "cases": [{"client": "A사", "desc": "실적", "is_example": True}],
            "note": "[예시] 교체 필요",
        })]}
        badge = rb.compute_review_badges(deck)["badges"][0]
        self.assertEqual(0, badge["signals"]["field_values"])
        self.assertNotEqual("충실", badge["verdict"])

    def test_bullet_era_slide_still_substantial(self):
        rb = self._rb()
        deck = {"slides": [self._slide(body=["하나", "둘", "셋"])]}
        self.assertEqual("충실", rb.compute_review_badges(deck)["badges"][0]["verdict"])

    def test_empty_slide_still_thin(self):
        rb = self._rb()
        deck = {"slides": [self._slide(key_message="짧다")]}
        self.assertNotEqual("충실", rb.compute_review_badges(deck)["badges"][0]["verdict"])


class W31EmphasisHeroAxisBSmoke(unittest.TestCase):
    """W31 R9 2축화(리허설 마찰17, 2026-07-21 확정) — 디자인 강조(emphasis=hero) 축B.

    강조는 구조 결정이라 A5(내용 동결)에서 확정하고, 하류(뼈대·이미지)는 소비만 한다. 회귀 방지:
    ①hero 후보 결정론 유도(표지·간지·결론 + 핵심 전략 축 지지 장, 축 자체도 상한 있음)
    ②storyline_읽기.md에 후보/확정 딱지 + A5 고지 문구
    ③05 매뉴얼에 '디자인 강조 장 확정' 결정 항목 고지
    ④emphasis가 adapt_storyline→deck.json까지 전달되고 스키마를 통과
    ⑤imagedeck bundle이 emphasis 장에 축B 강조 블록을 주입(축A 문구와 섞이지 않는다)
    ⑥디자인지표_읽기.md가 장별로 [보완 축A/강조 축B/그대로]를 표기
    ⑦review_badges가 emphasis 장을 thin 판정에서 제외(충실 강제 + signals.emphasis=True)
    ⑧emphasis 없는 storyline은 adapt_storyline 출력에 키 자체가 생기지 않는다(회귀 없음)."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _skin(self):
        return imagedeck.resolve_skin("inkline", self.SKINS_DIR)

    def test_hero_candidates_include_frame_roles_and_capped_axis_reps(self):
        doc = {"meta": {"project": "T"}, "slides": [
            {"n": 1, "section": "표지", "title": "표지"},
            {"n": 2, "section": "목차", "title": "목차"},
            {"n": 3, "section": "전략", "title": "축1-1", "supports_axis": "axis1"},
            {"n": 4, "section": "데이터", "title": "축1-2", "supports_axis": "axis1"},
            {"n": 5, "section": "전략", "title": "축2-1", "supports_axis": "axis2"},
            {"n": 6, "section": "전략", "title": "축3-1", "supports_axis": "axis3"},
            {"n": 7, "section": "전략", "title": "축4-1", "supports_axis": "axis4"},
            {"n": 8, "section": "마무리", "title": "마무리"},
        ]}
        cands = journey_folders.compute_hero_candidates(doc)
        self.assertIn(1, cands)      # 표지
        self.assertIn(8, cands)      # 결론/마무리
        self.assertNotIn(2, cands)   # 목차는 후보 아님
        self.assertIn(3, cands)      # axis1 대표(첫 등장)
        self.assertNotIn(4, cands)   # 같은 축 두번째 슬라이드는 후보 아님(축당 1장)
        self.assertTrue({3, 5, 6}.issubset(cands))
        self.assertNotIn(7, cands)   # 축 후보 자체도 상한(_HERO_AXIS_CAP=3) 밖

    def test_storyline_view_shows_candidate_and_confirmed_hero_tags_plus_intro_notice(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [
                    {"n": 1, "section": "표지", "title": "표지", "message": "m", "bullets": []},
                    {"n": 2, "section": "전략", "title": "핵심전략", "message": "m", "bullets": [],
                     "supports_axis": "axis1", "emphasis": "hero"},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            journey_folders.sync(run)
            view = (journey_folders.folder_path(run, "05")
                    / journey_folders.STORYLINE_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("이 단계에서 디자인 강조(hero) 장을 확정한다", view)
            self.assertIn("⭐ 강조 후보", view)   # 표지(1) = 후보
            self.assertIn("⭐ 강조 확정", view)   # 2 = emphasis 확정

    def test_manual_05_mentions_emphasis_decision(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [{"n": 1, "section": "표지", "title": "T", "message": "m",
                            "bullets": [], "flag": []}],
            }, ensure_ascii=False), encoding="utf-8")
            journey_folders.sync(run)
            manual = (journey_folders.folder_path(run, "05")
                      / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
            self.assertIn("디자인 강조 장 확정", manual)
            self.assertIn("emphasis", manual)

    def test_adapt_storyline_carries_emphasis_and_schema_accepts_hero(self):
        doc = {"meta": {"project": "T"}, "slides": [
            {"n": 1, "section": "전략", "title": "강조장", "message": "m",
             "bullets": ["짧게"], "emphasis": "hero"},
        ]}
        deck = _adapt.adapt_storyline(doc, pack="core")
        slide = deck["slides"][0]
        self.assertEqual("hero", slide["emphasis"])
        import slide_model
        self.assertEqual([], slide_model.validate(deck, "slide_model"))

    def test_adapt_storyline_omits_emphasis_key_when_absent(self):
        """회귀: emphasis 없는 storyline은 종전과 동일하게 키 자체가 생기지 않는다."""
        doc = {"meta": {"project": "T"}, "slides": [
            {"n": 1, "section": "전략", "title": "일반장", "message": "m", "bullets": ["b"]},
        ]}
        deck = _adapt.adapt_storyline(doc, pack="core")
        self.assertNotIn("emphasis", deck["slides"][0])

    def test_review_badges_forces_substantial_and_excludes_from_low_score_counts(self):
        import review_badges
        thin_hero = {"slide_id": 1, "role": "전략", "title": "강조",
                     "key_message": "짧다", "body": [], "fields": {}, "review_needed": [],
                     "emphasis": "hero"}
        report = review_badges.compute_review_badges({"slides": [thin_hero]})
        badge = report["badges"][0]
        self.assertEqual("충실", badge["verdict"])
        self.assertTrue(badge["signals"]["emphasis"])
        self.assertEqual(0, report["counts"]["밋밋"])
        self.assertEqual(0, report["counts"]["발산추천"])
        self.assertEqual(1, report["counts"]["충실"])

    def test_review_badges_signals_omit_emphasis_key_when_not_hero(self):
        import review_badges
        slide = {"slide_id": 1, "role": "전략", "title": "일반",
                 "key_message": "충분히 긴 메시지라 스무자를 넘긴다", "body": ["a", "b", "c"],
                 "fields": {}, "review_needed": []}
        badge = review_badges.compute_review_badges({"slides": [slide]})["badges"][0]
        self.assertNotIn("emphasis", badge["signals"])

    def test_wireframe_prompt_recommends_hero_body_for_emphasis_slide_only(self):
        import wireframe
        deck = {"slides": [
            {"slide_id": 1, "role": "전략", "title": "강조장", "key_message": "m",
             "fields": {}, "emphasis": "hero"},
            {"slide_id": 2, "role": "일정", "title": "보통장", "key_message": "m", "fields": {}},
        ]}
        prompt = wireframe.build_prompt(deck)
        self.assertIn("emphasis=hero", prompt)
        self.assertIn("hero_body", prompt)
        idx2 = prompt.index("slide 2")
        self.assertNotIn("emphasis=hero", prompt[idx2:])

    def test_imagedeck_bundle_injects_axis_b_block_distinct_from_axis_a(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
                {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
                {"n": 2, "title": "핵심 전략", "template_id": "strategy_pillars", "fields": {},
                 "emphasis": "hero"},
                {"n": 3, "title": "일정", "template_id": "process_steps", "fields": {}},
            ]}, ensure_ascii=False), encoding="utf-8")
            # emphasis 장은 review_badges가 충실을 강제하므로, 그 실제 상황을 그대로 재현
            # (design_brief에도 충실로 기록됨) — 축A 문구가 섞이지 않는지가 이 테스트의 핵심.
            (run / "design_brief.json").write_text(json.dumps({"page_rhythm": {"slides": [
                {"slide_id": 1, "verdict": "충실"},
                {"slide_id": 2, "verdict": "충실"},
                {"slide_id": 3, "verdict": "밋밋"},
            ]}}, ensure_ascii=False), encoding="utf-8")
            manifest = imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            p2 = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertIn("강조 축B", p2)
            self.assertIn("챕터-간지-여백", p2)
            self.assertIn("강조색-하나-나머지-무채색", p2)
            self.assertNotIn("축A", p2)
            slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
            self.assertTrue(slide2["emphasis_signal_injected"])
            self.assertFalse(slide2["design_signal_injected"])  # 충실 처리라 축A는 안 뜬다

            p3 = (run / imagedeck.PROMPTS_DIR / "03.md").read_text(encoding="utf-8")
            self.assertIn("축A 보완", p3)
            self.assertNotIn("강조 축B", p3)

    def test_design_metrics_view_marks_axis_a_axis_b_and_untouched(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
                {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {}},
                {"n": 2, "title": "핵심 전략", "template_id": "strategy_pillars", "fields": {},
                 "emphasis": "hero"},
                {"n": 3, "title": "일정", "template_id": "process_steps", "fields": {}},
            ]}, ensure_ascii=False), encoding="utf-8")
            (run / "design_brief.json").write_text(json.dumps({"page_rhythm": {"slides": [
                {"slide_id": 1, "verdict": "충실"},
                {"slide_id": 2, "verdict": "충실"},
                {"slide_id": 3, "verdict": "밋밋"},
            ]}}, ensure_ascii=False), encoding="utf-8")
            imagedeck.bundle(run, self._skin(), wireframe_mode="off")
            journey_folders.sync(run)
            view = (journey_folders.folder_path(run, "08")
                    / journey_folders.DESIGN_METRICS_VIEW_NAME).read_text(encoding="utf-8")
            self.assertIn("강조 축B", view)
            self.assertIn("보완 축A", view)
            self.assertIn("그대로", view)


class BannedHouseTermsSurfaceGuard(unittest.TestCase):
    """W31 P2 재발 방지 장치(CLEANUP_W31_PLAN §P2) — 사용자 대면 표면에 하우스 용어가 새면 빨간불.
    실제 사고 재현 경로: htmlgen r_data()가 모든 렌더 CSS에 house_a-* 클래스를 무조건 삽입하던
    누출을 P2에서 수리(7e7d8c0 직전) — 이 가드가 그 회귀를 영구 차단한다. 역사 문서·창고 스킨·
    출처 주석은 대상 아님(렌더 산출물만 검사)."""

    BANNED = ("house_a", "house_b", "하우스B", "quartz", "예시스튜디오", "quartz_guide")

    def _storyline(self):
        return {"meta": {"project": "가드"}, "slides": [
            {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": []},
            {"n": 2, "section": "데이터", "title": "시장", "message": "규모",
             "template_id": "data_interpretation", "example": True,
             "fields": {"metric": "m", "comparison": [{"label": "A", "value": 70},
                                                      {"label": "B", "value": 30}],
                        "interpretation": ["예시 해석"]}},
        ]}

    def test_core_render_surface_has_no_house_terms(self):
        deck = _adapt.adapt_storyline(self._storyline(), pack="core")
        _bind.bind_deck(deck, _ingest._load_templates("core"))
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            htmlgen.render_html(deck, "core", out)
            html = out.read_text(encoding="utf-8").lower()
        for term in self.BANNED:
            self.assertNotIn(term.lower(), html,
                             f"사용자 대면 렌더 산출물에 하우스 용어 누출: {term}")


class W31CompanyWarehouseSmoke(unittest.TestCase):
    """W31 리허설 마찰6 — 제안사(자사) 프로필 창고(`proposal_system/companies/<id>/`).

    발주처 조사(institution_research)와 대칭인 자사 사실 소스. 여기서는 모듈 자체(스키마
    검증·병합·gaps 축적)와 run 투입 배선(start --company·enrich 채움·fictional 가드)을
    검증한다. 실제 창고(`proposal_system/companies/`)를 건드리지 않도록 `company._COMPANIES_ROOT`
    를 임시 디렉터리로 패치한다."""

    def _isolated_root(self, td):
        root = Path(td) / "companies"
        return mock.patch.object(company, "_COMPANIES_ROOT", root)

    def test_ensure_scaffold_creates_assets_intake_readmes_and_list_reads_it_back(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            company.ensure_scaffold("acme")
            self.assertTrue((company.assets_dir("acme") / "README.md").is_file())
            self.assertTrue((company.intake_dir("acme") / "README.md").is_file())
            self.assertEqual([], company.list_companies())  # profile.json 없으면 표에 안 뜬다

            company.save("acme", {
                "schema_version": 1, "fictional": False,
                "overview": {"name": {"value": "에이씨엠", "source": "test"}},
                "track_records": [{"client": "X", "description": "d", "source": "test"}],
            })
            rows = company.list_companies()
            self.assertEqual(1, len(rows))
            self.assertEqual("acme", rows[0]["id"])
            self.assertEqual("에이씨엠", rows[0]["name"])
            self.assertEqual(1, rows[0]["track_records"])
            self.assertIn("에이씨엠", company.format_list(rows))

    def test_validate_detects_missing_source_and_bad_fictional_type(self):
        profile = {
            "fictional": "yes",  # bool 아님 → 오류
            "overview": {"name": {"value": "회사"}},  # source 없음 → 오류
            "strengths": [{"value": "강점"}],  # source 없음 → 오류
            "track_records": [{"client": "A", "description": "d", "source": "ok"}],
        }
        v = company.validate(profile)
        joined = "\n".join(v["errors"])
        self.assertIn("fictional", joined)
        self.assertIn("overview.name", joined)
        self.assertIn("strengths[0]", joined)

    def test_validate_passes_when_all_items_sourced(self):
        profile = {
            "fictional": True,
            "overview": {"name": {"value": "회사", "source": "s"}},
            "strengths": [{"value": "강점", "source": "s"}],
        }
        v = company.validate(profile)
        self.assertEqual([], v["errors"])

    def test_merge_profile_appends_new_and_updates_existing_without_duplicating(self):
        existing = {
            "fictional": True,
            "overview": {"name": {"value": "회사", "source": "seed"}},
            "track_records": [{"client": "A", "description": "d1", "metric": "m1", "source": "seed"}],
        }
        incoming = {
            "track_records": [
                {"client": "A", "description": "d1", "metric": "m1-갱신", "source": "인테이크"},  # 갱신
                {"client": "B", "description": "d2", "source": "인테이크"},  # 신규
            ],
        }
        merged, diff = company.merge_profile(existing, incoming)
        self.assertEqual(2, len(merged["track_records"]))  # 중복 추가 아님(항목 추가·갱신)
        rec_a = next(r for r in merged["track_records"] if r["client"] == "A")
        self.assertEqual("m1-갱신", rec_a["metric"])
        self.assertEqual(1, diff["added"]["track_records"])
        self.assertEqual(1, diff["updated"]["track_records"])
        # 기존 overview는 incoming이 건드리지 않았으니 그대로 보존.
        self.assertEqual("회사", merged["overview"]["name"]["value"])

    def test_apply_cmd_rejects_missing_source_and_writes_nothing(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            candidate = Path(td) / "candidate.json"
            candidate.write_text(json.dumps({
                "overview": {"name": {"value": "무출처회사"}},  # source 없음
            }, ensure_ascii=False), encoding="utf-8")
            with self.assertRaises(proposal_pipeline.PipelineInputError):
                proposal_pipeline._company_apply_cmd("newco", str(candidate))
            self.assertFalse(company.exists("newco"))

    def test_apply_cmd_merges_and_updates_gaps_for_missing_sections(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            candidate = Path(td) / "candidate.json"
            candidate.write_text(json.dumps({
                "fictional": False,
                "overview": {"name": {"value": "신규회사", "source": "회사소개서.pdf"}},
            }, ensure_ascii=False), encoding="utf-8")
            result = proposal_pipeline._company_apply_cmd("newco", str(candidate))
            self.assertTrue(company.exists("newco"))
            profile = company.load("newco")
            self.assertEqual("신규회사", profile["overview"]["name"]["value"])
            self.assertTrue(company.gaps_path("newco").is_file())
            gaps_text = company.gaps_path("newco").read_text(encoding="utf-8")
            self.assertIn("특장점", gaps_text)  # strengths 비어있음이 구조 스캔으로 표면화
            self.assertTrue(result["gaps"])

    def test_gaps_append_is_deterministic_and_dedupes(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            added1 = company.append_gaps("dupco", ["같은 항목", "다른 항목"])
            added2 = company.append_gaps("dupco", ["같은 항목"])  # 이미 있음 → 스킵
            self.assertEqual(2, added1)
            self.assertEqual(0, added2)
            text = company.gaps_path("dupco").read_text(encoding="utf-8")
            self.assertEqual(1, text.count("같은 항목"))

    def test_start_company_records_selection_and_survives_unknown_id(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            root = Path(td)
            runs = root / "workspace" / "runs"
            runs.mkdir(parents=True)
            args = SimpleNamespace(bid="TEST-W31-COMPANY-001", brief=None, mode="direct",
                                   run_name=None, selected_by=None, company="ghost-co")
            buf = io.StringIO()
            with mock.patch.object(proposal_pipeline, "REPO_ROOT", root), \
                 mock.patch.object(proposal_pipeline, "RUNS", runs), \
                 mock.patch("sys.stdout", buf):
                proposal_pipeline.start_cmd(args)
            run = next(runs.iterdir())
            sel = company.load_selection(run)
            self.assertEqual("ghost-co", sel["company_id"])
            self.assertIn("profile.json이 아직 없다", buf.getvalue())  # 미등록 id도 차단 없이 기록+안내

    def test_unselected_run_has_no_selection_and_enrich_untouched(self):
        """미선택 run 회귀 불변: company_profile 인자가 아예 없던 시절과 바이트 동일해야 한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            self.assertIsNone(company.load_selection(run))
            deck = {"meta": {"pack": "core"}, "slides": [
                {"slide_id": "s1", "template_id": "portfolio_cases", "role": "company", "fields": {}},
            ]}
            without_arg = _enrich.enrich_deck(json.loads(json.dumps(deck)), {}, "")
            with_none = _enrich.enrich_deck(json.loads(json.dumps(deck)), {}, "", company_profile=None)
            self.assertEqual(without_arg, with_none)
            self.assertNotIn("example", without_arg["slides"][0])  # 채워질 근거가 없으니 그대로 검토요망

    def test_enrich_fills_company_cases_and_org_from_profile(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            company.save("gaon-systems", {
                "fictional": True,
                "overview": {"name": {"value": "가온", "source": "s"}},
                "track_records": [{"client": "A교육청", "description": "구축", "metric": "성과",
                                    "source": "s"}],
                "organization": {
                    "lead": {"name": "김도현", "description": "PM", "source": "s"},
                    "teams": [{"name": "개발팀", "roles": ["백엔드"], "source": "s"}],
                },
            })
            profile = company.load("gaon-systems")
            deck = {"meta": {"pack": "core"}, "slides": [
                {"slide_id": "s1", "template_id": "portfolio_cases", "role": "company", "fields": {}},
                {"slide_id": "s2", "template_id": "org_roles", "fields": {}},
            ]}
            out = _enrich.enrich_deck(deck, {}, "", company_profile=profile)
            case_slide = out["slides"][0]
            self.assertEqual("A교육청", case_slide["fields"]["cases"][0]["client"])
            self.assertTrue(case_slide["fields"]["cases"][0]["is_example"])  # fictional=True 전파
            self.assertTrue(case_slide["example"])
            self.assertIn(_bind.EXAMPLE_REVIEW_TAG, case_slide["review_needed"])
            org_slide = out["slides"][1]
            self.assertEqual("김도현", org_slide["fields"]["lead"]["name"])
            self.assertEqual("개발팀", org_slide["fields"]["teams"][0]["name"])
            self.assertTrue(org_slide["example"])

    def test_enrich_non_fictional_company_does_not_force_example_tag(self):
        profile = {
            "fictional": False,
            "overview": {"name": {"value": "실사", "source": "s"}},
            "track_records": [{"client": "실고객", "description": "실사례", "source": "s"}],
        }
        deck = {"meta": {"pack": "core"}, "slides": [
            {"slide_id": "s1", "template_id": "portfolio_cases", "role": "company", "fields": {}},
        ]}
        out = _enrich.enrich_deck(deck, {}, "", company_profile=profile)
        self.assertNotIn("example", out["slides"][0])
        self.assertEqual("실고객", out["slides"][0]["fields"]["cases"][0]["client"])

    def test_storyline_and_message_map_bundles_attach_company_summary_when_selected(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            company.save("gaon-systems", {
                "fictional": True,
                "overview": {"name": {"value": "가온시스템즈", "source": "s"}},
            })
            run = Path(td) / "run"
            run.mkdir()
            (run / "brief.md").write_text("# 브리프\n샘플 요구사항", encoding="utf-8")
            company.save_selection(run, "gaon-systems")

            out = proposal_pipeline.bundle_storyline_from_brief(
                run, SimpleNamespace(pack="core"),
            )
            text = out.read_text(encoding="utf-8")
            self.assertIn("가온시스템즈", text)
            self.assertIn("제안사(자사) 프로필", text)

            mm_out = proposal_pipeline.bundle_message_map(run, SimpleNamespace())
            mm_text = mm_out.read_text(encoding="utf-8")
            self.assertIn("가온시스템즈", mm_text)

    def test_collect_company_gap_entries_uses_real_korean_role_and_template_id(self):
        """실측 계약: adapt_storyline은 role을 section 원문("제안업체" 등)으로 남긴다 —
        skeleton.py의 영문 role 태그("company")는 deck.json까지 오지 않는다(회귀 방지)."""
        deck = {"slides": [
            {"slide_id": 1, "role": "제안업체", "template_id": "portfolio_cases",
             "title": "Ⅱ. 제안업체", "review_needed": ["[예시 데이터] 실데이터로 교체 필요"]},
            {"slide_id": 2, "role": "전략", "template_id": "strategy_pillars",
             "title": "무관 장표", "review_needed": ["다른 검토요망"]},
            {"slide_id": 3, "role": "조직", "template_id": "org_roles",
             "title": "인력 구성", "review_needed": []},
        ]}
        entries = proposal_pipeline._collect_company_gap_entries(deck, "demo-run")
        self.assertEqual(1, len(entries))
        self.assertIn("Ⅱ. 제안업체", entries[0])
        self.assertIn("demo-run", entries[0])

    def test_render_run_accumulates_gaps_for_selected_company_when_review_needed_remains(self):
        with tempfile.TemporaryDirectory() as td, self._isolated_root(td):
            company.save("gaon-systems", {
                "fictional": True,
                "overview": {"name": {"value": "가온", "source": "s"}},
                # track_records 비움 → portfolio_cases를 못 채워 review_needed가 남는다.
            })
            runs_root = Path(td) / "workspace" / "runs"
            runs_root.mkdir(parents=True)
            run = runs_root / "run"
            run.mkdir()
            company.save_selection(run, "gaon-systems")
            storyline = {
                "meta": {"project": "가온 데모"},
                "slides": [
                    {"n": 1, "section": "표지", "title": "표지", "message": "M", "bullets": []},
                    {"n": 2, "section": "제안업체", "title": "Ⅱ. 제안업체", "message": "M",
                     "template_id": "portfolio_cases", "bullets": ["소개"]},
                ],
            }
            sl_path = run / "storyline.json"
            sl_path.write_text(json.dumps(storyline, ensure_ascii=False), encoding="utf-8")
            args = SimpleNamespace(
                stage6=None, storyline=str(sl_path), run_dir=str(run), stage7=None, stage8=None,
                project=None, pack="core", pattern_sets=None, pptx=False, pptx_mode="native",
                skins=None, analysis=None, rfp=None, anonymize_config=None, json=False,
            )
            with mock.patch.object(proposal_pipeline, "RUNS", runs_root):
                proposal_pipeline.render_run(args)
            gaps_path = company.gaps_path("gaon-systems")
            self.assertTrue(gaps_path.is_file())
            self.assertIn("제안업체", gaps_path.read_text(encoding="utf-8"))

    def test_bundles_unchanged_bytes_when_no_company_selected(self):
        """미선택이면 회사 블록 자체가 안 생겨 바이트 불변(institution_research와 같은 계약)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            (run / "brief.md").write_text("# 브리프\n샘플", encoding="utf-8")
            out1 = proposal_pipeline.bundle_storyline_from_brief(run, SimpleNamespace(pack="core"))
            text1 = out1.read_text(encoding="utf-8")
            self.assertNotIn("제안사(자사) 프로필", text1)


class MasterDesignSmoke(unittest.TestCase):
    """W31 R10 v2(β2, 리허설 마찰26) — 덱 마스터 디자인 공정(imagedeck --master-bundle/--master-apply).

    정본: CONTEXT/JOURNEY.md R10절(v2) + CONTEXT/REHEARSAL_FRICTIONS_W31.md 26행."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def test_master_bundle_prompt_has_five_input_elements_and_multi_candidate_instruction(self):
        """복합 입력함 5요소(발주처·자사·주제·레퍼런스·디자인지식 pull) + 복수 후보안 지시."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            rep = imagedeck.master_bundle(run, skins_dir=self.SKINS_DIR)
            text = Path(rep["prompt"]).read_text(encoding="utf-8")
        self.assertIn("(1) 발주처 브랜드", text)
        self.assertIn("(2) 자사 아이덴티티", text)
        self.assertIn("(3) 주제", text)
        self.assertIn("(4) 레퍼런스", text)
        self.assertIn("(5) 디자인지식 pull", text)
        self.assertIn("공유뇌 `ref/디자인지식/`을 능동 조회하라", text)
        self.assertIn("Claude 디자인 스킬", text)
        self.assertIn("복수 후보안", text)
        self.assertIn("발주처축", text)
        self.assertIn("자사축", text)
        self.assertIn("주제축", text)
        self.assertIn("master_design.json", text)
        self.assertIn('"density": "standard | spacious | dense"', text)

    def test_master_bundle_works_without_content_design_first_route(self):
        """storyline.json이 전혀 없는 run에서도 동작한다(디자인 선행 루트, 내용 의존 금지)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self.assertFalse((run / "storyline.json").is_file())
            rep = imagedeck.master_bundle(run, skins_dir=self.SKINS_DIR)
            text = Path(rep["prompt"]).read_text(encoding="utf-8")
        self.assertFalse(rep["content_present"])
        self.assertIn("내용(storyline) 유무와 무관하게 동작한다", text)
        self.assertIn("표지 + 샘플 본문", text)

    def test_master_apply_rejects_invalid_schema(self):
        """look 없음·density 오타·assets 파일 부재는 오류(적용 중단), chosen_axis 자유값은 경고만."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            doc = {"look": "", "density": "huge", "assets": ["missing.png"], "chosen_axis": "자유서술축"}
            result = imagedeck.validate_master_design(doc, run)
        self.assertTrue(result["errors"])
        joined = " ".join(result["errors"])
        self.assertIn("look", joined)
        self.assertIn("density", joined)
        self.assertIn("missing.png", joined)
        self.assertTrue(result["warnings"])  # chosen_axis 관례 밖은 경고만(차단 아님)

    def _make_master_doc(self, run: Path, density="spacious"):
        asset = run / "hero.png"
        asset.write_bytes(b"\x89PNG\r\n\x1a\n")
        return {
            "look": "따뜻한 아이보리 바탕에 손그림풍 플랫 선화, 강조색 하나만 채도 높게",
            "density": density,
            "assets": ["hero.png"],
            "chosen_axis": "발주처축",
            "sources": ["institution_research.json"],
            "knowledge_used": {"cards": [], "web": []},  # ε패킷 안전장치① — 생략하면 적용 차단
        }

    def test_master_apply_records_art_direction_and_density_prev_preserved(self):
        """계약에 art_direction·density 기록 — 기존 계약이 있었으면 prev로 보존(재동결 문법)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            old_contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, old_contract)
            old_text = design_contract.path(run).read_text(encoding="utf-8")

            doc = self._make_master_doc(run)
            doc_path = run / imagedeck.MASTER_DESIGN_NAME
            doc_path.write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline._master_apply_cmd(run, SimpleNamespace(master_file=None))
            self.assertEqual(0, rc)

            prev_path = run / "design_contract.prev.json"
            self.assertTrue(prev_path.is_file())
            self.assertEqual(old_text, prev_path.read_text(encoding="utf-8"))

            new_contract = design_contract.load(run)
            self.assertEqual("spacious", new_contract["density"])
            self.assertEqual(doc["look"], new_contract["art_direction"]["look"])
            self.assertEqual("발주처축", new_contract["art_direction"]["chosen_axis"])
            self.assertIn("[MASTER APPLY]", out.getvalue())

    def test_master_apply_registers_series_reference_and_bundle_auto_includes(self):
        """확정 시안이 imagedeck_refs/global/에 등록되고, 이후 imagedeck.bundle이 자동 동봉한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            doc = self._make_master_doc(run, density="standard")
            doc_path = run / imagedeck.MASTER_DESIGN_NAME
            doc_path.write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                proposal_pipeline._master_apply_cmd(run, SimpleNamespace(master_file=None))

            registered = run / "imagedeck_refs" / "global" / "hero.png"
            self.assertTrue(registered.is_file())

            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T"},
                "slides": [
                    {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
                    {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            self.assertEqual("global", manifest["slides"][1]["references_source"])
            prompt02 = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
            self.assertIn("hero.png", prompt02)
            self.assertIn("확정 룩(마스터 시안, 07_테마확정)", prompt02)

    # -------------------------------------------------------------------
    # DF3(2026-07-24, CONTEXT/DECK_FIRST_DESIGN.md §2-②·§3) — 자산 생성 절차 편입:
    # 전 장 공유 배경 PNG·코너 장식을 --master-bundle/--master-apply에 편입.
    # -------------------------------------------------------------------

    def _png(self, path, w, h):
        import struct
        import zlib
        def ch(t, d):
            c = t + d
            return struct.pack(">I", len(d)) + c + struct.pack(">I", zlib.crc32(c) & 0xffffffff)
        path.parent.mkdir(parents=True, exist_ok=True)
        raw = b"".join(b"\x00" + b"\xff" * (w * 3) for _ in range(h))
        path.write_bytes(b"\x89PNG\r\n\x1a\n" + ch(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
                         + ch(b"IDAT", zlib.compress(raw)) + ch(b"IEND", b""))

    def test_master_bundle_prompt_includes_background_and_decor_instructions(self):
        """DF3 ①: 마스터 시안 지시문에 전 장 공유 배경 PNG·코너 장식 생성 항목 + 스키마 키."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            rep = imagedeck.master_bundle(run, skins_dir=self.SKINS_DIR)
            text = Path(rep["prompt"]).read_text(encoding="utf-8")
        self.assertIn("전 장 공유 배경 PNG", text)
        self.assertIn("저대비", text)
        self.assertIn("코너 장식 후보", text)
        self.assertIn("투명 배경", text)
        self.assertNotIn("SVG", text)  # 생성 루트=Codex PNG 우선, SVG 지시 금지(사용자 정책)
        self.assertIn('"background":', text)
        self.assertIn('"decor_slots":', text)
        self.assertIn("chrome_contract.chrome.frame.image", text)
        self.assertIn("chrome_contract.decor_slots", text)

    def test_validate_master_design_background_and_decor_slots(self):
        """DF3 ②: background/decor_slots — 실재 검증(assets와 동일 수준), 미지정은 검사 자체가 없다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            # 미지정 = 하위호환(검사 없음, 오류 0)
            base = {"look": "무드", "density": "standard", "assets": []}
            result = imagedeck.validate_master_design(base, run)
            self.assertEqual([], result["errors"])

            # background 파일 없음 = 오류
            doc = dict(base, background="missing_bg.png")
            result = imagedeck.validate_master_design(doc, run)
            self.assertTrue(any("background" in e and "missing_bg.png" in e for e in result["errors"]))

            # decor_slots: anchor 미지원·image 없음·width 없음 = 각각 오류
            doc = dict(base, decor_slots=[{"id": "bad", "anchor": "middle", "image": "", "width": None}])
            result = imagedeck.validate_master_design(doc, run)
            joined = " ".join(result["errors"])
            self.assertIn("anchor 미지원", joined)
            self.assertIn("image 경로 없음", joined)
            self.assertIn("width 없음", joined)

            # decor_slots: image 파일 부재 = 오류
            doc = dict(base, decor_slots=[{"id": "ghost", "anchor": "top-left", "image": "nope.png", "width": 80}])
            result = imagedeck.validate_master_design(doc, run)
            self.assertTrue(any("ghost" in e and "이미지 없음" in e for e in result["errors"]))

            # 전부 유효(실재 파일) = 오류 없음
            (run / "bg.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            (run / "corner.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            doc = dict(base, background="bg.png",
                      decor_slots=[{"id": "corner", "anchor": "top-right", "image": "corner.png", "width": 100}])
            result = imagedeck.validate_master_design(doc, run)
            self.assertEqual([], result["errors"])

    def test_master_apply_freezes_background_and_decor_into_chrome_contract(self):
        """DF3 ③: background/decor_slots가 imagedeck/design_assets/로 복사 + chrome_contract에 동결."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "bg_src").mkdir()
            (run / "bg_src" / "hero_bg.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            (run / "corner_src").mkdir()
            (run / "corner_src" / "corner.png").write_bytes(b"\x89PNG\r\n\x1a\n")
            doc = self._make_master_doc(run)
            doc["background"] = "bg_src/hero_bg.png"
            doc["decor_slots"] = [{"id": "corner", "anchor": "top-right", "image": "corner_src/corner.png",
                                   "width": 90, "opacity": 0.8}]
            doc_path = run / imagedeck.MASTER_DESIGN_NAME
            doc_path.write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                rc = proposal_pipeline._master_apply_cmd(run, SimpleNamespace(master_file=None))
            self.assertEqual(0, rc)

            contract = design_contract.load(run)
            chrome = contract["chrome_contract"]
            bg_rel = chrome["chrome"]["frame"]["image"]
            self.assertTrue(bg_rel.startswith("imagedeck/design_assets/"))
            self.assertTrue((run / bg_rel).is_file())
            self.assertEqual(1, len(chrome["decor_slots"]))
            decor_rel = chrome["decor_slots"][0]["image"]
            self.assertTrue(decor_rel.startswith("imagedeck/design_assets/"))
            self.assertTrue((run / decor_rel).is_file())
            self.assertEqual("top-right", chrome["decor_slots"][0]["anchor"])
            self.assertEqual(90, chrome["decor_slots"][0]["width"])
            self.assertIn("배경 동결", out.getvalue())
            self.assertIn("장식 슬롯 동결", out.getvalue())
            self.assertIn("design-assets 싱크백", out.getvalue())
            self.assertIn("curate --sync-master", out.getvalue())

    def test_master_apply_without_background_decor_is_backward_compatible(self):
        """DF3 ④: background/decor_slots 미지정 = imagedeck/design_assets/ 안 만들고 계약도 손대지 않는다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            doc = self._make_master_doc(run)
            self.assertNotIn("background", doc)
            self.assertNotIn("decor_slots", doc)
            doc_path = run / imagedeck.MASTER_DESIGN_NAME
            doc_path.write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                proposal_pipeline._master_apply_cmd(run, SimpleNamespace(master_file=None))

            self.assertFalse((run / imagedeck.MASTER_ASSETS_DIR).exists())
            contract = design_contract.load(run)
            chrome = contract["chrome_contract"]
            self.assertIsNone(chrome["chrome"]["frame"]["image"])
            self.assertNotIn("decor_slots", chrome)
            self.assertNotIn("배경 동결", out.getvalue())
            self.assertNotIn("장식 슬롯 동결", out.getvalue())
            self.assertNotIn("싱크백", out.getvalue())

            # 계약 그 자체는 background/decor_slots 없는 일반 build()와 바이트 동일해야 한다
            # (art_direction/density만 추가되는 것은 기존 마스터 시안 동작 그대로 - 회귀 아님).
            plain = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            self.assertEqual(plain["chrome_contract"], chrome)

    def test_master_apply_end_to_end_compose_renders_background_and_decor(self):
        """DF3 실측: master-apply로 동결한 배경·장식이 compose(HTML)에 실제로 반영된다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._png(run / "hero_bg.png", 40, 20)
            self._png(run / "corner.png", 10, 10)
            doc = self._make_master_doc(run)
            doc["background"] = "hero_bg.png"
            doc["decor_slots"] = [{"id": "corner", "anchor": "bottom-right", "image": "corner.png", "width": 60}]
            doc_path = run / imagedeck.MASTER_DESIGN_NAME
            doc_path.write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")

            out = io.StringIO()
            with mock.patch("sys.stdout", new=out):
                proposal_pipeline._master_apply_cmd(run, SimpleNamespace(master_file=None))

            (run / "storyline.json").write_text(json.dumps({
                "meta": {"project": "T프로젝트"},
                "slides": [
                    {"n": 1, "title": "T프로젝트 제안서", "template_id": "cover_slide", "fields": {}},
                    {"n": 2, "title": "본문", "template_id": "strategy_pillars", "fields": {}},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            manifest = imagedeck.bundle(run, None, wireframe_mode="off")
            for s in manifest["slides"]:
                if s.get("render") == "html":
                    continue
                px = s["expected_px"]
                self._png(run / "imagedeck" / "slides" / s["out_name"], px["w"], px["h"])
            rep = imagedeck.compose(run)
            self.assertEqual([], rep["warnings"])
            html = (run / "deck.images.html").read_text(encoding="utf-8")
            contract = design_contract.load(run)
            bg_rel = contract["chrome_contract"]["chrome"]["frame"]["image"]
            self.assertIn(f"background-image:url({(run / bg_rel).resolve().as_uri()})", html)
            self.assertIn('class="decor"', html)

    def test_density_band_warning_surfaces_when_nonstandard_and_storyline_exists(self):
        """density != standard 이고 storyline이 이미 있으면 go/status 경고가 표면화된다(차단 아님)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({"meta": {}, "slides": []}), encoding="utf-8")
            (run / "design_contract.json").write_text(json.dumps({"density": "dense"}), encoding="utf-8")
            warnings = pipeline_state._density_band_warning(run)
            self.assertTrue(warnings)
            self.assertIn("[밀도 비표준]", warnings[0])
            self.assertIn("A5 부분 재생성 권장", warnings[0])

            (run / "design_contract.json").write_text(json.dumps({"density": "standard"}), encoding="utf-8")
            self.assertEqual([], pipeline_state._density_band_warning(run))

            # storyline이 없으면(디자인 선행, 아직 내용 없음) 비표준이어도 침묵(경고 대상 자체가 없다).
            (run / "storyline.json").unlink()
            (run / "design_contract.json").write_text(json.dumps({"density": "dense"}), encoding="utf-8")
            self.assertEqual([], pipeline_state._density_band_warning(run))

    def test_express_profile_shortens_master_route_note(self):
        """gates 프로파일 express면 07 매뉴얼의 마스터 안내가 축약된다(권장 문구만)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            standard_note = journey_folders._master_route_note(run)
            self.assertIn("imagedeck --master-bundle", standard_note)

            gates.save_config(run, profile="express")
            express_note = journey_folders._master_route_note(run)
            self.assertIn("express", express_note)
            self.assertLess(len(express_note), len(standard_note))
            self.assertNotIn("imagedeck --master-bundle", express_note)

    def test_full_skin_contract_skips_master_route_note(self):
        """차용 스킨이 이미 완전 스펙(full_skin)이면 '차용본이 곧 마스터' 안내로 생략을 권한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            contract = design_contract.build(run, brief={"skin": {"value": "inkline"}}, skins_dir=self.SKINS_DIR)
            self.assertTrue(design_contract.is_full_skin(contract))
            design_contract.save(run, contract)
            note = journey_folders._master_route_note(run)
        self.assertIn("차용본이 곧 마스터", note)
        self.assertIn("생략해도 된다", note)

    def test_master_design_block_wired_into_message_map_and_storyline_bundles(self):
        """확정 룩이 있으면 message_map·storyline 핸드오프에 동봉된다(company_profile_block과 같은 문법)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            (run / "brief.md").write_text("# 브리프\n샘플", encoding="utf-8")
            contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            contract["density"] = "dense"
            contract["art_direction"] = {"look": "차분한 네이비 바탕에 얇은 라인 아이콘", "chosen_axis": "주제축"}
            design_contract.save(run, contract)

            mm_out = proposal_pipeline.bundle_message_map(run, SimpleNamespace())
            mm_text = mm_out.read_text(encoding="utf-8")
            self.assertIn("확정 룩(마스터 시안, 07_테마확정 · R10 v2)", mm_text)
            self.assertIn("차분한 네이비 바탕에 얇은 라인 아이콘", mm_text)
            self.assertIn("밀도 비표준", mm_text)

            sl_out = proposal_pipeline.bundle_storyline_from_brief(run, SimpleNamespace(pack="core"))
            sl_text = sl_out.read_text(encoding="utf-8")
            self.assertIn("확정 룩(마스터 시안, 07_테마확정 · R10 v2)", sl_text)

    def test_bundles_unchanged_bytes_when_no_master_design(self):
        """마스터 시안 미확정 run은 블록 자체가 안 생겨 바이트 불변(institution_research와 같은 계약)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            (run / "brief.md").write_text("# 브리프\n샘플", encoding="utf-8")
            out1 = proposal_pipeline.bundle_storyline_from_brief(run, SimpleNamespace(pack="core"))
            text1 = out1.read_text(encoding="utf-8")
            self.assertNotIn("확정 룩(마스터 시안", text1)


# =============================================================================
# W31 γ패킷 — 리허설 마찰 22~25 (오버플로 종단 처리)
# =============================================================================

class BandViolationSignalSmoke(unittest.TestCase):
    """마찰22: 분량 밴드 위반이 gates 신호·핸드오프 하드 제약·경고 위험도로 표면화되는지."""

    def test_band_violation_escalates_prompt_ack_gate_even_under_express(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gates.save_config(run, profile="express")  # 가장 관대한 프로파일이어도
            (run / "gating_report.json").write_text(json.dumps({
                "length_rhythm": {"slides": 5, "band_violations": [
                    {"slide_id": 2, "words": 90, "band": [20, 70], "kind": "over"},
                    {"slide_id": 3, "words": 95, "band": [20, 70], "kind": "over"},
                    {"slide_id": 4, "words": 100, "band": [20, 70], "kind": "over"},
                ]},
            }, ensure_ascii=False), encoding="utf-8")
            d = gates.decide(run, "imagedeck_prompt_ack")
        self.assertEqual("stop", d["action"])
        self.assertIn("분량 밴드 위반", d["reason"])
        self.assertIn("마찰22", d["reason"])

    def test_two_violations_below_threshold_do_not_escalate(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gates.save_config(run, profile="express")
            (run / "gating_report.json").write_text(json.dumps({
                "length_rhythm": {"slides": 5, "band_violations": [
                    {"slide_id": 2, "words": 90, "band": [20, 70], "kind": "over"},
                    {"slide_id": 3, "words": 95, "band": [20, 70], "kind": "over"},
                ]},
            }, ensure_ascii=False), encoding="utf-8")
            d = gates.decide(run, "imagedeck_prompt_ack")
        self.assertEqual("auto_pass", d["action"])  # 2건 < 임계 3 — 재정지 안 함

    def test_length_band_is_hard_constraint_in_storyline_handoff_prompt(self):
        block = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn("length_band", block)
        self.assertIn("하드 제약", block)
        self.assertIn("밴드 초과 금지", block)

    def test_length_rhythm_warning_mentions_image_stage_overflow_risk(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "gating_report.json").write_text(json.dumps({
                "length_rhythm": {"slides": 5, "band_violations": [
                    {"slide_id": 3, "words": 200, "band": [20, 70], "kind": "over"},
                ]},
            }, ensure_ascii=False), encoding="utf-8")
            warns = pipeline_state._length_rhythm_warnings(run)
        joined = " ".join(warns)
        self.assertIn("분량 밴드 위반", joined)
        self.assertIn("오버플로", joined)
        self.assertIn("08_프롬프트확인", joined)

    def test_design_metrics_view_shows_band_violations_prominently(self):
        doc = {"slides": [{"n": 1, "out_name": "01_x.png"}]}
        band_violations = [{"slide_id": 2, "words": 90, "band": [20, 70], "kind": "over"}]
        view = journey_folders.render_design_metrics_view(doc, band_violations=band_violations)
        self.assertIn("분량 밴드 위반 1장", view)
        self.assertIn("마찰22", view)
        self.assertIn("마찰23", view)


class OverflowSplitSmoke(unittest.TestCase):
    """마찰23: 하한 폰트로도 수용 불가 확정 장의 결정론 사전 분할(bundle) — 발생/비발생/분할 포기."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def _run(self, td, *, big_bullets=False, unsplittable=False):
        run = Path(td)
        slides = [
            {"n": 1, "title": "정상장", "template_id": "executive_summary",
             "message": "짧은 메시지", "bullets": ["근거1", "근거2"], "fields": {}},
        ]
        if big_bullets:
            # 1792x784 본문·min_font 22px 기준 용량 추정치(~1488자)의 확실한 초과(2000자, 20개 항목
            # → 분할 가능한 리스트).
            slides.append({
                "n": 2, "title": "과다분량장", "template_id": "executive_summary",
                "message": "m", "bullets": ["다" * 100 for _ in range(20)], "fields": {},
            })
        if unsplittable:
            # 용량은 초과하지만 분할 가능한 리스트가 전혀 없다(fields는 스칼라 문자열, bullets 비었음).
            slides.append({
                "n": 3, "title": "분할불가장", "template_id": "executive_summary",
                "message": "m", "bullets": [], "fields": {"summary": "다" * 2200},
            })
        (run / "storyline.json").write_text(json.dumps({
            "meta": {"project": "T"}, "slides": slides,
        }, ensure_ascii=False), encoding="utf-8")
        return run

    def test_overflowing_slide_with_splittable_bullets_is_split_into_ab(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td, big_bullets=True)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            splits = manifest.get("overflow_splits") or []
            variants = [s for s in manifest["slides"] if s["n"] == 2]
            pa = (run / next(v for v in variants if v["variant"] == "A")["prompt_file"]).read_text(encoding="utf-8")
            pb = (run / next(v for v in variants if v["variant"] == "B")["prompt_file"]).read_text(encoding="utf-8")
        self.assertEqual(1, len(splits))
        self.assertEqual(2, splits[0]["n"])
        self.assertEqual(2, len(variants))
        self.assertEqual({"A", "B"}, {v["variant"] for v in variants})
        self.assertTrue(all(v.get("content_split") for v in variants))
        self.assertIn("사전 분할", pa)
        self.assertNotEqual(pa, pb)

    def test_normal_slide_is_not_split(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td, big_bullets=True)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
        n1 = [s for s in manifest["slides"] if s["n"] == 1]
        self.assertEqual(1, len(n1))
        self.assertIsNone(n1[0]["variant"])
        self.assertNotIn("content_split", n1[0])

    def test_unsplittable_overflow_gives_up_with_warning(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._run(td, unsplittable=True)
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
        n3 = [s for s in manifest["slides"] if s["n"] == 3]
        self.assertEqual(1, len(n3))
        self.assertIsNone(n3[0]["variant"])
        self.assertTrue(n3[0].get("overflow_split_skipped"))
        self.assertIn("분할", n3[0].get("overflow_split_skip_reason", ""))
        self.assertEqual([], manifest.get("overflow_splits"))


class PixelHeuristicsSmoke(unittest.TestCase):
    """마찰24: collect 결정론 픽셀 휴리스틱(PIL) — warn/clean/PIL 부재 강등 + scaffold 자동 생성."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    def test_dense_edge_ink_flags_warn(self):
        from PIL import Image, ImageDraw
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "edge.png"
            img = Image.new("RGB", (400, 300), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            draw.rectangle([0, 0, 399, imagedeck.EDGE_BAND_PX - 1], fill=(0, 0, 0))  # 상단 띠 전체를 잉크로
            img.save(p)
            rep = imagedeck.pixel_heuristics(p)
        self.assertTrue(rep["available"])
        self.assertIn("edge_ink_high", rep["flags"])

    def test_clean_white_image_has_no_flags(self):
        from PIL import Image
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "clean.png"
            Image.new("RGB", (400, 300), (255, 255, 255)).save(p)
            rep = imagedeck.pixel_heuristics(p)
        self.assertTrue(rep["available"])
        self.assertEqual([], rep["flags"])

    def test_gracefully_degrades_when_pil_unavailable(self):
        with mock.patch.object(imagedeck, "_pil_available", return_value=False):
            rep = imagedeck.pixel_heuristics(Path("does-not-matter.png"))
        self.assertEqual({"available": False}, rep)

    def _png(self, path, w, h):
        import struct
        import zlib
        def ch(t, d):
            c = t + d
            return struct.pack(">I", len(d)) + c + struct.pack(">I", zlib.crc32(c) & 0xffffffff)
        path.parent.mkdir(parents=True, exist_ok=True)
        raw = b"".join(b"\x00" + b"\xff" * (w * 3) for _ in range(h))
        path.write_bytes(b"\x89PNG\r\n\x1a\n" + ch(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
                         + ch(b"IDAT", zlib.compress(raw)) + ch(b"IEND", b""))

    def test_collect_auto_generates_review_scaffold_without_explicit_flag(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
                {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
                {"n": 2, "title": "본문", "template_id": "process_steps", "fields": {}},
            ]}, ensure_ascii=False), encoding="utf-8")
            skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
            manifest = imagedeck.bundle(run, skin, wireframe_mode="off")
            for s in manifest["slides"]:
                if s.get("render") == "html":
                    continue
                px = s["expected_px"]
                self._png(run / "imagedeck" / "slides" / s["out_name"], px["w"], px["h"])
            rep = imagedeck.collect(run)
            review_path = run / "imagedeck_review.md"
            self.assertTrue(review_path.is_file())  # --review-scaffold 없이도 자동 생성
        self.assertIn("pixel_warnings", rep)
        self.assertIn("pixel_heuristics_available", rep)


class ChromeTitleOverflowSmoke(unittest.TestCase):
    """마찰25: 크롬(HTML) 상단(제목) auto-fit 단계적 축소 + 결정론 경고 신호."""

    def _deck(self, title):
        fields = {"main_claim": "핵심 주장", "supporting_points": ["근거1", "근거2"]}
        slide = {"slide_id": 1, "role": "제안개요", "template_id": "executive_summary",
                 "title": title, "key_message": "메시지", "body": [], "fields": fields}
        return {"meta": {"project": "테스트 프로젝트"}, "slides": [slide]}

    def test_long_title_gets_shrink_class_and_css_injected(self):
        long_title = "가" * 45   # TITLE_LEN_SM(40) 초과
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            htmlgen.render_html(self._deck(long_title), "core", out)
            html = out.read_text(encoding="utf-8")
        self.assertIn("title--sm", html)
        self.assertIn(".slide__title.title--sm", html)  # CSS 규칙도 함께 실려야 실제로 축소된다

    def test_short_title_keeps_bytes_unchanged_no_fit_class(self):
        """회귀(기대값 갱신 없이 그대로 유지돼야 함): 짧은 제목뿐인 덱은 auto-fit 클래스/CSS가
        전혀 섞이지 않는다 — 이 패킷 이전 바이트와 동일해야 한다(같은 입력 → 같은 출력 결정론)."""
        with tempfile.TemporaryDirectory() as td:
            out1 = Path(td) / "a.html"
            out2 = Path(td) / "b.html"
            htmlgen.render_html(self._deck("짧은 제목"), "core", out1)
            htmlgen.render_html(self._deck("짧은 제목"), "core", out2)
            html1 = out1.read_text(encoding="utf-8")
            html2 = out2.read_text(encoding="utf-8")
        self.assertNotIn("title--sm", html1)
        self.assertNotIn("title--md", html1)
        self.assertEqual(html1, html2)  # 결정론 — 같은 입력은 항상 같은 바이트

    def test_design_checks_flags_title_overflow_risk(self):
        import design_checks
        html = ('<html><body><section id="slide-1" class="slide">'
                '<h2 class="slide__title">' + ("가" * 45) + '</h2></section></body></html>')
        dc = design_checks.compute_design_checks(html)
        self.assertIn("title_overflow_risk", dc["slides"][0]["flags"])
        self.assertEqual(1, dc["summary"]["title_overflow_risk"])
        self.assertEqual("warn", dc["status"])


class W31ArchiveRoundtripSmoke(unittest.TestCase):
    """W31 리허설 마찰9 — run 보관소 왕복(`archive --run/--restore/--list`, `YYYY-MM_한글명`).

    실 `workspace/runs`·`workspace/archive`·`dashboard/last_search.json`을 건드리지 않도록
    `archive.ARCHIVE_ROOT`·`archive.LAST_SEARCH_PATH`·`proposal_pipeline.RUNS`를 임시
    디렉터리로 패치한다(company.py W31 창고 테스트와 동일 격리 관례)."""

    def _isolated(self, td):
        root = Path(td)
        runs = root / "workspace" / "runs"
        runs.mkdir(parents=True)
        arc_root = root / "workspace" / "archive"
        last_search = root / "no_last_search.json"  # 존재하지 않음 — last_search 폴백 비활성(격리)
        return runs, arc_root, last_search

    def _make_run(self, runs_root, run_id, *, bid=None):
        run = runs_root / run_id
        run.mkdir(parents=True)
        pipeline_state.init(
            run, mode="direct",
            input_kind="bid" if bid else "brief", input_ref=bid or "brief.md",
            selected_by="user", bid=bid,
        )
        return run

    def _mark_completed(self, run):
        (run / "approval.json").write_text(
            json.dumps({"status": "approved", "timestamp": "2026-07-20T10:00:00"}, ensure_ascii=False),
            encoding="utf-8",
        )

    def test_archive_moves_run_and_records_meta(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = self._make_run(runs, "gen_TEST-001")
            self._mark_completed(run)
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                rep = archive.archive_run(run, name="테스트기관-홍보용역")
            self.assertFalse(run.exists())  # 원 위치에서 사라짐(이동 — 복사 아님)
            dest = rep["dest"]
            self.assertTrue(dest.is_dir())
            self.assertTrue((dest / "approval.json").is_file())  # 자립 — 산출물 그대로 딸려옴
            self.assertRegex(dest.name, r"^\d{4}-\d{2}_테스트기관-홍보용역$")
            meta = json.loads((dest / archive.META_NAME).read_text(encoding="utf-8"))
            self.assertEqual("gen_TEST-001", meta["run_id"])
            self.assertEqual("테스트기관-홍보용역", meta["korean_name"])
            self.assertEqual("user_specified", meta["name_source"])

    def test_infer_name_prefers_analysis_card_over_fallback(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = self._make_run(runs, "gen_R26TEST-000", bid="R26TEST-000")
            analysis_dir = run.parent.parent / "analysis"
            analysis_dir.mkdir(parents=True)
            (analysis_dir / "R26TEST-000_분석카드.md").write_text(
                "# 📋 입찰 분석카드 — 가상기관 테스트 홍보 용역\n\n"
                "| 구분 | 내용 |\n|---|---|\n"
                "| 발주처(공고/수요기관) | 가상기관 (테스트용, 가상) |\n",
                encoding="utf-8",
            )
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                name, source = archive.infer_name(run)
            self.assertEqual("분석카드", source)
            self.assertEqual("가상기관-홍보용역", name)

    def test_infer_name_falls_back_to_machine_id_when_nothing_found(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = self._make_run(runs, "gen_TEST-NOINFO")
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                name, source = archive.infer_name(run)
            self.assertEqual("gen_TEST-NOINFO", name)
            self.assertEqual("fallback_id", source)

    def test_restore_round_trip_is_lossless(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = self._make_run(runs, "gen_TEST-RT-001")
            (run / "deck.json").write_text(json.dumps({"marker": "원본유지"}, ensure_ascii=False),
                                           encoding="utf-8")
            self._mark_completed(run)
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                rep = archive.archive_run(run, name="복귀테스트")
                restored = archive.restore(rep["folder"], runs_root=runs)
            dest = restored["restored_to"]
            self.assertEqual(runs / "gen_TEST-RT-001", dest)
            self.assertTrue(dest.is_dir())
            deck = json.loads((dest / "deck.json").read_text(encoding="utf-8"))
            self.assertEqual("원본유지", deck["marker"])
            self.assertTrue((dest / "approval.json").is_file())
            self.assertTrue((dest / "pipeline_state.json").is_file())

    def test_archive_folder_name_conflict_gets_dash2_suffix(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run1 = self._make_run(runs, "gen_TEST-C1")
            run2 = self._make_run(runs, "gen_TEST-C2")
            self._mark_completed(run1)
            self._mark_completed(run2)
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                rep1 = archive.archive_run(run1, name="같은이름")
                rep2 = archive.archive_run(run2, name="같은이름")
            self.assertNotEqual(rep1["folder"], rep2["folder"])
            self.assertTrue(rep2["folder"].endswith("-2"))

    def test_restore_rejects_when_active_run_id_already_exists(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = self._make_run(runs, "gen_TEST-DUP")
            self._mark_completed(run)
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                rep = archive.archive_run(run, name="충돌")
                (runs / "gen_TEST-DUP").mkdir()  # 활성으로 재생성(복귀 충돌 유발) — 사람 말 오류 기대
                with self.assertRaises(archive.ArchiveError):
                    archive.restore(rep["folder"], runs_root=runs)
            # 거부됐으니 보관본은 그대로(파괴 없음)
            self.assertTrue((arc_root / rep["folder"] / "approval.json").is_file())

    def test_completed_run_detection_and_list_table(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run_active = self._make_run(runs, "gen_TEST-ACTIVE-DONE")
            self._mark_completed(run_active)
            self._make_run(runs, "gen_TEST-INCOMPLETE")  # approval.json 없음 — 목록에 안 뜬다
            run_to_archive = self._make_run(runs, "gen_TEST-TO-ARCHIVE")
            self._mark_completed(run_to_archive)
            self.assertTrue(archive.is_completed(run_active))
            self.assertFalse(archive.is_completed(runs / "gen_TEST-INCOMPLETE"))
            with mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                archive.archive_run(run_to_archive, name="보관됨")
                text = archive.format_list(runs, arc_root)
            self.assertIn("gen_TEST-ACTIVE-DONE", text)
            self.assertNotIn("gen_TEST-INCOMPLETE", text)
            self.assertIn("보관됨", text)

    def test_start_threshold_hint_fires_at_3_and_silent_below(self):
        with tempfile.TemporaryDirectory() as td:
            runs, _arc_root, _last_search = self._isolated(td)
            self.assertIsNone(archive.start_threshold_hint(runs))  # 0건 — 침묵
            for i in range(2):
                run = self._make_run(runs, f"gen_TEST-TH-{i}")
                self._mark_completed(run)
            self.assertIsNone(archive.start_threshold_hint(runs))  # 2건 — 아직 임계 미만, 침묵
            run3 = self._make_run(runs, "gen_TEST-TH-2")
            self._mark_completed(run3)
            hint = archive.start_threshold_hint(runs)  # 3건 — 알림
            self.assertIsNotNone(hint)
            self.assertIn("3", hint)
            self.assertIn("archive --list", hint)

    def test_ship_prints_archive_hint_after_success(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run = runs / "gen_TEST-SHIP-HINT"
            run.mkdir(parents=True)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            (run / "deck.json").write_text(json.dumps({"meta": {}, "slides": []}), encoding="utf-8")
            (run / "deck.html").write_text("<html></html>", encoding="utf-8")
            pipeline_state.record(run, "render", pack="core", skins=None)
            pipeline_state.clear_checkpoint(run, "design")  # 사람 관문 통과 기록(대시보드 ack 대체)

            args = SimpleNamespace(run=str(run), pptx=False, pptx_mode="native",
                                   ingest=None, source=None)
            buf = io.StringIO()
            with mock.patch.object(proposal_pipeline, "RUNS", runs), \
                 mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search), \
                 mock.patch("sys.stdout", buf):
                rc = proposal_pipeline.ship_cmd(args)
            self.assertEqual(0, rc)
            out = buf.getvalue()
            self.assertIn("보관: archive --run gen_TEST-SHIP-HINT", out)
            self.assertIn("한글명 제안", out)

    def test_status_run_prints_archive_hint_only_when_completed(self):
        with tempfile.TemporaryDirectory() as td:
            runs, arc_root, last_search = self._isolated(td)
            run_done = self._make_run(runs, "gen_TEST-STATUS-DONE")
            self._mark_completed(run_done)
            run_open = self._make_run(runs, "gen_TEST-STATUS-OPEN")

            with mock.patch.object(proposal_pipeline, "RUNS", runs), \
                 mock.patch.object(archive, "ARCHIVE_ROOT", arc_root), \
                 mock.patch.object(archive, "LAST_SEARCH_PATH", last_search):
                buf_done = io.StringIO()
                with mock.patch("sys.stdout", buf_done):
                    proposal_pipeline.status_run_cmd(SimpleNamespace(run=str(run_done), json=False))
                buf_open = io.StringIO()
                with mock.patch("sys.stdout", buf_open):
                    proposal_pipeline.status_run_cmd(SimpleNamespace(run=str(run_open), json=False))
            self.assertIn("보관: archive --run gen_TEST-STATUS-DONE", buf_done.getvalue())
            self.assertNotIn("보관:", buf_open.getvalue())  # 미완료 run은 침묵


class KCPacketSmoke(unittest.TestCase):
    """KC 패킷(2026-07-24 확정) — 지식 체크 3+1 모델: ① 기획 입구 · ③ 산출 출구 · 검토_체크.md 연동.

    ②(디자인 입구)는 기존 W27 디자인지식 pull 그대로 — 이 패킷에서 손대지 않는다(대상 밖).
    관문 프로파일 연동: standard(기본)=①+③, express=①·③ 모두 생략, full=③에 장별 샘플링
    심화 안내 추가.
    """

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    # ① message_map 핸드오프 — 기획 입구 pull ---------------------------------

    def test_message_map_handoff_includes_knowledge_pull_by_default(self):
        """gates.json 미설정 = standard 프로파일 — pull 요구 문구가 기본으로 동봉된다.

        ε패킷(2026-07-23): 경로가 vault 재편(ref/기획지식/메시지설계·ref/기획지식/경험설계)을
        반영하고, 지시문 조립은 하드코딩이 아니라 pipeline.config.json의 knowledge_stages 표
        (knowledge_ledger 모듈)를 소비한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "brief.md").write_text("# 브리프\n샘플", encoding="utf-8")
            out = proposal_pipeline.bundle_message_map(run, SimpleNamespace())
            text = out.read_text(encoding="utf-8")
        self.assertIn("지식 pull", text)
        self.assertIn("ref/기획지식/메시지설계/", text)
        self.assertIn("ref/기획지식/경험설계/", text)
        self.assertIn("obsidian_search", text)
        # 안전장치①(2026-07-23): 보고 의무는 프로파일과 무관하게 항상 동봉된다.
        self.assertIn("knowledge_used", text)

    def test_message_map_handoff_omits_knowledge_pull_under_express(self):
        """express는 'vault를 조회하라'는 pull 넛지만 생략한다 — 보고 의무(안전장치①)는
        express에서도 남는다(2026-07-23 확정: 자동 모드에서도 보고 없이 진행 금지)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "brief.md").write_text("# 브리프\n샘플", encoding="utf-8")
            gates.save_config(run, profile="express")
            out = proposal_pipeline.bundle_message_map(run, SimpleNamespace())
            text = out.read_text(encoding="utf-8")
        self.assertNotIn("[지식 pull", text)
        self.assertNotIn("ref/기획지식/경험설계/", text)
        self.assertIn("knowledge_used", text)  # 보고 의무는 express에서도 남는다

    def test_message_map_knowledge_pull_text_helper(self):
        """ε패킷: 시그니처가 (run, profile)로 바뀌었고, express도 더 이상 None이 아니다
        (보고 의무는 항상 동봉 — 안전장치①)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self.assertIsNotNone(message_map.knowledge_pull_text(run, "express"))
            self.assertNotIn("[지식 pull", message_map.knowledge_pull_text(run, "express"))
            self.assertIsNotNone(message_map.knowledge_pull_text(run, "standard"))
            self.assertIsNotNone(message_map.knowledge_pull_text(run, "full"))
            self.assertIsNotNone(message_map.knowledge_pull_text(run, None))

    # ③ imagedeck_review.md — 산출 출구 지식 대조 ------------------------------

    def _bundled_run(self, td):
        run = Path(td)
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문", "template_id": "process_steps", "fields": {}},
        ]}, ensure_ascii=False), encoding="utf-8")
        skin = imagedeck.resolve_skin("inkline", self.SKINS_DIR)
        imagedeck.bundle(run, skin, wireframe_mode="off")
        return run

    def test_review_scaffold_includes_knowledge_check_by_default(self):
        """ε패킷(2026-07-23): 폴더 목록이 하드코딩이 아니라 config knowledge_stages.imagedeck_review
        표(기획지식/경험설계 — 테마 카드 대조는 master_design 단계로 이동)를 소비한다."""
        with tempfile.TemporaryDirectory() as td:
            run = self._bundled_run(td)
            imagedeck.review_scaffold(run)
            text = (run / "imagedeck_review.md").read_text(encoding="utf-8")
        self.assertIn("지식 대조", text)
        self.assertIn("ref/기획지식/경험설계/", text)
        self.assertIn("지식 사용 기록", text)  # ε패킷 원장 수거용 기입란
        self.assertNotIn("장별 샘플링 심화", text)  # standard엔 없음(full 전용)

    def test_review_scaffold_omits_knowledge_check_under_express(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._bundled_run(td)
            gates.save_config(run, profile="express")
            imagedeck.review_scaffold(run)
            text = (run / "imagedeck_review.md").read_text(encoding="utf-8")
        self.assertNotIn("지식 대조", text)

    def test_review_scaffold_full_profile_adds_sampling_guidance(self):
        with tempfile.TemporaryDirectory() as td:
            run = self._bundled_run(td)
            gates.save_config(run, profile="full")
            imagedeck.review_scaffold(run)
            text = (run / "imagedeck_review.md").read_text(encoding="utf-8")
        self.assertIn("지식 대조", text)
        self.assertIn("장별 샘플링 심화", text)

    # 검토_체크.md — ①·③ 지식 체크 수행 확인 항목 ------------------------------

    def test_journey_check_decision_gate_includes_kc_item_by_default(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            journey_folders.folder_path(run, "05").mkdir(parents=True)
            path = journey_check.issue(run, "decision")
            text = path.read_text(encoding="utf-8")
        self.assertIn("지식 체크 확인 — ① 기획 입구", text)

    def test_journey_check_imagedeck_ack_gate_includes_kc_item_by_default(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            journey_folders.folder_path(run, "11").mkdir(parents=True)
            path = journey_check.issue(run, "imagedeck_ack")
            text = path.read_text(encoding="utf-8")
        self.assertIn("지식 체크 확인 — ③ 산출 출구", text)

    def test_journey_check_omits_kc_item_under_express(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            gates.save_config(run, profile="express")
            journey_folders.folder_path(run, "05").mkdir(parents=True)
            path = journey_check.issue(run, "decision")
            text = path.read_text(encoding="utf-8")
        self.assertNotIn("지식 체크 확인", text)

    def test_unrelated_gate_never_gets_kc_item(self):
        """skeleton_review는 KC 매핑 밖(①·③ 대상 아님) — 프로파일과 무관하게 문구가 없다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            journey_folders.folder_path(run, "04").mkdir(parents=True)
            path = journey_check.issue(run, "skeleton_review")
            text = path.read_text(encoding="utf-8")
        self.assertNotIn("지식 체크 확인", text)


class EpsilonPacketSmoke(unittest.TestCase):
    """ε패킷(2026-07-23 확정) — 지식 소비 체계: config 1점화·원장·지시 오버레이·안전장치 3중.

    KCPacketSmoke(①·③ 지점 hook)와 겹치지 않게, 여기는 (1) config 표 소비의 일반성(가짜
    config 픽스처로 실vault 무의존), (2) research 신설 pull, (3) 원장 기록·파생 뷰, (4) 안전장치
    ①(보고 누락 차단)·②(웹 사용 조건부 재정지, 비허용 단계 웹 오류)·③(상시 표면화), (5) 지식_지시.md
    오버레이 동봉을 검증한다."""

    def _map(self, *, cards=None, web=None):
        return {
            "governing_message": "발주처를 위한 통합 소통체계",
            "strategy_axes": [{
                "id": "axis1", "message": "하위 메시지 1",
                "evidence_slots": [{"type": "데이터", "desc": "근거 1", "status": "filled", "source": "출처"}],
            }],
            "knowledge_used": {"cards": cards or [], "web": web or []},
        }

    # (1) config 1점화 — 가짜 config 픽스처로 검증(실vault·실pipeline.config.json 무의존) --------

    def test_stage_table_reads_from_config_fixture_not_hardcoded(self):
        fixture = {"knowledge_stages": {
            "storyline": {"pull": ["가짜지식/픽스처층"], "web_search": True},
        }}
        table = knowledge_ledger.load_stage_table(fixture)
        self.assertEqual(["가짜지식/픽스처층"], knowledge_ledger.pull_folders("storyline", table))
        self.assertTrue(knowledge_ledger.web_allowed("storyline", table))
        # 실제 config(web_search=false)와 다르다는 것 자체가 "하드코딩이 아니라 표를 읽는다"는 증거.
        self.assertFalse(knowledge_ledger.web_allowed("storyline"))

    def test_unconfigured_stage_falls_back_to_empty_safely(self):
        """설정에 없는 단계는 크래시 없이 빈 pull/False로 우아 폴백(배포판·픽스처 무의존)."""
        table = knowledge_ledger.load_stage_table({"knowledge_stages": {}})
        self.assertEqual([], knowledge_ledger.pull_folders("message_map", table))
        self.assertFalse(knowledge_ledger.web_allowed("message_map", table))

    # (2) research 신설 pull — 사용자 확정: 기존엔 없었다 -----------------------------------------

    def test_research_bundle_now_includes_knowledge_pull(self):
        import institution_research as ir_mod
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            prompt = ir_mod.build_prompt(run, institution="테스트기관", profile="standard")
        self.assertIn("지식 pull", prompt)
        self.assertIn("ref/기획지식/메시지설계/", prompt)
        self.assertIn("ref/기획지식/경험설계/", prompt)
        self.assertIn("knowledge_used", prompt)

    # (3) 원장 기록 + 파생 뷰 ------------------------------------------------------------------

    def test_message_map_collect_records_ledger_and_derived_view(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir(parents=True)
            journey_folders.folder_path(run, "05").mkdir(parents=True)  # 05_내용동결(decision)
            doc = self._map(cards=["여정-단계-접점"])
            message_map.map_path(run).write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")
            buf = io.StringIO()
            with mock.patch("sys.stdout", buf):
                proposal_pipeline._go_message_map_collect(run, SimpleNamespace())
            ledger = knowledge_ledger.load_ledger(run)
            view_path = journey_folders.folder_path(run, "05") / knowledge_ledger.LEDGER_VIEW_NAME
            self.assertEqual(["여정-단계-접점"], ledger["stages"]["message_map"]["cards"])
            self.assertTrue(view_path.is_file())
            self.assertIn("여정-단계-접점", view_path.read_text(encoding="utf-8"))
            self.assertIn("지식: 카드 1", buf.getvalue())  # 안전장치③: go 출력에도 상시 표면화

    # (4-①) 보고 누락 = 수거 검증 실패(자동 모드에서도) -------------------------------------------

    def test_missing_knowledge_used_block_blocks_collect(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            doc = self._map()
            del doc["knowledge_used"]  # 조용한 생략 시뮬레이션
            run.mkdir(parents=True)
            message_map.map_path(run).write_text(json.dumps(doc, ensure_ascii=False), encoding="utf-8")
            with self.assertRaises(proposal_pipeline.PipelineInputError) as ctx:
                proposal_pipeline._go_message_map_collect(run, SimpleNamespace())
        self.assertIn("knowledge_used", str(ctx.exception))

    def test_validate_knowledge_used_rejects_missing_block_directly(self):
        errors, _ = knowledge_ledger.validate_knowledge_used({"a": 1}, "message_map")
        self.assertTrue(any("knowledge_used" in e for e in errors), msg=str(errors))

    # (4-②) 웹 사용 → 사람 관문 조건부 재정지(auto/express여도) -----------------------------------

    def test_web_usage_forces_conditional_gate_stop(self):
        """wireframe 단계(wireframe_review 게이트)를 쓴다 — theme_confirm은 standard 프로파일에서
        이미 항상 stop이라(PROFILE_DEFAULTS) 웹 신호의 효과를 보여줄 수 없다. wireframe_review는
        standard에서 기본 auto(자동 통과 대상)라, 웹 사용 조건부 재정지가 실제로 갈리는 지점이다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_ledger.record(run, "wireframe", {"cards": [], "web": [
                {"url": "https://example.com/ref", "purpose": "레퍼런스 무드보드"}
            ]})
            decision = gates.decide(run, "wireframe_review")  # wireframe → wireframe_review 매핑
        self.assertEqual("stop", decision["action"])
        self.assertTrue(any("웹" in r for r in (decision["signal"] or {}).get("reasons", [])),
                        msg=str(decision))

    def test_no_web_usage_does_not_force_stop(self):
        """웹 원장이 비어 있으면(또는 없으면) 이 신호는 기존 게이트 판정을 전혀 바꾸지 않는다
        (회귀 방지 — base 신호가 없을 때의 기존 보수적/관대 분기가 그대로 살아 있어야 한다).
        wireframe_review는 standard 프로파일 기본이 auto이고, gating_report.json도 없어 base
        신호가 unavailable → 스킵 가능 관문이라 auto_pass가 유지돼야 한다."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            knowledge_ledger.record(run, "wireframe", {"cards": ["카드A"], "web": []})
            decision = gates.decide(run, "wireframe_review")
        self.assertNotEqual("stop", decision["action"])

    # (4-② 대칭) web_search=false 단계에서 web 항목 → 검증 오류 -----------------------------------

    def test_web_item_on_disallowed_stage_is_validation_error(self):
        doc = {"knowledge_used": {"cards": [], "web": [
            {"url": "https://example.com", "purpose": "금지된 사용"}
        ]}}
        errors, _ = knowledge_ledger.validate_knowledge_used(doc, "storyline")  # storyline web_search=false
        self.assertTrue(any("web_search=false" in e for e in errors), msg=str(errors))

    # (5) 지식_지시.md 오버레이 동봉 -------------------------------------------------------------

    def test_knowledge_instruction_overlay_attached_to_handoff(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            journey_folders.folder_path(run, "05").mkdir(parents=True)
            path = knowledge_ledger.ensure_overlay(run, "message_map")
            text = path.read_text(encoding="utf-8")
            text = text.replace(
                "## [추가로 참고]\n(추가로 조회했으면 하는 지식 카드·주제를 한 줄씩 적어라. 없으면 비워둔다.)",
                "## [추가로 참고]\n- 픽스처-카드-슬러그",
            )
            path.write_text(text, encoding="utf-8")
            block = knowledge_ledger.handoff_block(run, "message_map", "standard")
        self.assertIn("사용자 지식 지시", block)
        self.assertIn("픽스처-카드-슬러그", block)
        # 2026-07-23 사용자 확정: 웹 허용 단계는 "무단 수행"이 아니라 추천→동의 흐름.
        self.assertIn("사용자에게 웹 조사를 추천", block)
        self.assertIn("무단 수행 금지", block)

    def test_overlay_issued_only_once_and_never_overwritten(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            journey_folders.folder_path(run, "05").mkdir(parents=True)
            path1 = knowledge_ledger.ensure_overlay(run, "message_map")
            path1.write_text(path1.read_text(encoding="utf-8") + "\n사람이 적은 메모\n", encoding="utf-8")
            path2 = knowledge_ledger.ensure_overlay(run, "message_map")  # 두 번째 발급 시도
            self.assertEqual(path1, path2)
            self.assertIn("사람이 적은 메모", path1.read_text(encoding="utf-8"))

    # (5 안전장치③) status/go 상시 표면화 --------------------------------------------------------

    def test_status_surfaces_knowledge_line(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            knowledge_ledger.record(run, "message_map", {"cards": ["카드A"], "web": []})
            view = pipeline_state.resolve(run)
            text = pipeline_state.format_status(view)
        self.assertTrue(any("카드 1" in ln for ln in view["knowledge"]), msg=str(view["knowledge"]))
        self.assertIn("지식 사용 (ε패킷 원장", text)

    def test_stage_with_no_ledger_entry_is_silent(self):
        """기록 없는 run은 강제로 아무것도 표면화하지 않는다(강제 아님 — 조용함이 정상)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="brief", input_ref="x")
            view = pipeline_state.resolve(run)
        self.assertEqual([], view["knowledge"])


class DesignKnowledgeCardCarrySmoke(unittest.TestCase):
    """δ패킷 — A6(뼈대 결정기, wireframe.json slides[].knowledge_cards)가 이미 고른 디자인지식
    카드를 imagedeck 장별 프롬프트에 결정론 운반(본문 발췌 + 실물 이미지). 선택은 LLM(pull),
    운반은 코드(결정론) — design_spec.py의 R9 축B 카드 인용과 같은 문법(기계 무차별 주입 아님).

    정본 = proposal_system/scripts/design_knowledge_cards.py. 실제 vault(<개발 원본 전용 경로>)·knowhow에
    의존하지 않고 tempfile로 가짜 vault/knowhow를 구성한다(배포판·CI 안전)."""

    SKINS_DIR = Path(__file__).resolve().parent.parent / "skins"

    @staticmethod
    def _write(path: Path, text: str) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text, encoding="utf-8")

    def _fake_vault(self, root: Path) -> Path:
        # 카드A(와이어프레임) — 자기 ref_images 없음, examples 1홉으로 이미지 2장 확보.
        self._write(root / "와이어프레임" / "카드A.md", """---
name: 카드A
claim: 카드A 원칙 주장
type: principle
layer: wireframe
examples: [예시A]
---

본문 설명.

**조작적 정의**
- 규칙 1
- 규칙 2

[[다른카드]]
""")
        self._write(root / "examples" / "예시A.md", """---
name: 예시A
proves: 예시A 증명
type: example
source: peedori_ ig_111111
ref_images:
  - design://insta/peedori/ig_111111/01.jpg   # before
  - design://insta/peedori/ig_111111/02.jpg   # after
---

**교정 규칙**: 이렇게 고친다.
""")
        # 카드B(와이어프레임) — 자기 ref_images 직접 보유(1장).
        self._write(root / "와이어프레임" / "카드B.md", """---
name: 카드B
claim: 카드B 원칙 주장
layer: wireframe
ref_images:
  - design://insta/peedori/ig_222222/01.jpg
---

본문.

**조작적 정의**
- 규칙 X
""")
        # 카드C(테마) — 조작적 정의가 600자 상한을 넘는 긴 카드(절단 검증용).
        long_line = "가" * 700
        self._write(root / "테마" / "카드C.md", f"""---
name: 카드C
claim: 카드C 원칙 주장
layer: theme
---

**조작적 정의**
{long_line}
""")
        # 카드D(와이어프레임) — 이미지 5장(장당 상한 4 초과 검증용).
        d_images = "\n".join(
            f"  - design://insta/peedori/ig_444444/{i:02d}.jpg" for i in range(1, 6)
        )
        self._write(root / "와이어프레임" / "카드D.md", f"""---
name: 카드D
claim: 카드D 원칙 주장
layer: wireframe
ref_images:
{d_images}
---

본문.
""")
        return root

    def _fake_knowhow(self, root: Path) -> Path:
        groups = {
            "ig_111111": ["01.jpg", "02.jpg"],
            "ig_222222": ["01.jpg"],
            "ig_444444": [f"{i:02d}.jpg" for i in range(1, 6)],
        }
        for ig, names in groups.items():
            for name in names:
                p = root / "design" / "insta" / "peedori" / ig / name
                p.parent.mkdir(parents=True, exist_ok=True)
                p.write_bytes(b"fake-jpg")
        return root

    def _storyline(self, run: Path) -> None:
        (run / "storyline.json").write_text(json.dumps({"meta": {"project": "T"}, "slides": [
            {"n": 1, "title": "표지", "template_id": "cover_slide", "fields": {"project_title": "T"}},
            {"n": 2, "title": "본문A", "template_id": "strategy_pillars", "fields": {}},
        ]}, ensure_ascii=False), encoding="utf-8")

    def _wireframe(self, run: Path, cards: list) -> None:
        (run / "wireframe.json").write_text(json.dumps({
            "schema_version": 1, "selected_by": "llm:test",
            "slides": [
                {"slide_id": "2", "message_type": "구조", "frame": "flow_seq", "rendition": "boxed",
                 "slots": [{"piece": "flow_arrow", "size": "wide", "binds": "*"}],
                 "knowledge_cards": cards, "catalog_gap": []},
            ],
        }, ensure_ascii=False), encoding="utf-8")

    def _bundle(self, run: Path, cards: list, wireframe_mode: str = "on"):
        self._storyline(run)
        self._wireframe(run, cards)
        contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
        design_contract.save(run, contract)
        return imagedeck.bundle(run, None, wireframe_mode=wireframe_mode)

    # -- ① 운반 절 존재 + claim·조작적 정의 발췌 -----------------------------

    def test_carry_section_present_with_claim_and_excerpt(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                self._bundle(run, ["카드A"])
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn(design_knowledge_cards.CARRY_HEADER, prompt)
        self.assertIn("카드A", prompt)
        self.assertIn("카드A 원칙 주장", prompt)
        self.assertIn("규칙 1", prompt)
        self.assertIn("규칙 2", prompt)

    # -- ② 카드 미발견 표기(조용한 실패 금지) --------------------------------

    def test_missing_card_marked_not_silently_dropped(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                self._bundle(run, ["존재하지않는카드"])
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("존재하지않는카드", prompt)
        self.assertIn("[카드 미발견]", prompt)

    # -- ③ 카드당 상한 절단 ---------------------------------------------------

    def test_excerpt_truncated_at_cap(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            card = design_knowledge_cards.load_card("카드C", vault=vault)
        self.assertTrue(card["found"])
        self.assertLessEqual(
            len(card["excerpt"]), design_knowledge_cards.CARD_EXCERPT_LIMIT + len(" ...(절단)")
        )
        self.assertIn("...(절단)", card["excerpt"])

    # -- ④ ref_images design:// 해석(카드 자기 보유분) ------------------------

    def test_own_ref_images_resolve_to_real_files(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            card = design_knowledge_cards.load_card("카드B", vault=vault)
            self.assertTrue(card["found"])
            self.assertEqual(1, len(card["images"]))
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                carry = design_knowledge_cards.carry_knowledge(["카드B"])
            self.assertEqual(1, len(carry["images"]))
            resolved_path = Path(carry["images"][0]["path"])
            self.assertTrue(resolved_path.is_file())
            self.assertEqual("ig_222222", resolved_path.parent.name)

    # -- ⑤ examples 1홉 ref_images 해석 ---------------------------------------

    def test_examples_one_hop_ref_images_resolved(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                carry = design_knowledge_cards.carry_knowledge(["카드A"])
            self.assertEqual(2, len(carry["images"]))
            self.assertTrue(all(Path(img["path"]).is_file() for img in carry["images"]))
            self.assertEqual("예시A", carry["images"][0]["via"])

    # -- ⑥ 장당 이미지 상한(프롬프트 비대 방지) -------------------------------

    def test_image_cap_per_slide(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                carry = design_knowledge_cards.carry_knowledge(["카드D"])  # 카드D는 이미지 5장
        self.assertEqual(design_knowledge_cards.IMAGE_LIMIT_PER_SLIDE, len(carry["images"]))
        self.assertEqual([], carry["images_gap"])  # 상한 초과분은 gap이 아니라 조용히 절단

    # -- ⑦ config/vault 부재 우아 생략(크래시 금지) ---------------------------

    def test_missing_vault_degrades_gracefully_without_crash(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=None), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=None):
                self._bundle(run, ["아무카드"])  # 크래시하면 여기서 예외
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("[카드 미발견]", prompt)  # vault 없음 -> 못 찾은 카드로 정직하게 표면화

    def test_missing_pipeline_config_section_degrades_gracefully(self):
        """design_knowledge 절 자체가 config에 없는 배포판 시나리오 — 크래시 없이 빈 결과."""
        with tempfile.TemporaryDirectory() as td:
            fake_config = Path(td) / "pipeline.config.json"
            fake_config.write_text(json.dumps({"paths": {}}), encoding="utf-8")
            with mock.patch.object(design_knowledge_cards, "PIPELINE_CONFIG_PATH", fake_config):
                self.assertIsNone(design_knowledge_cards._vault_dir())
                self.assertIsNone(design_knowledge_cards._knowhow_root())
                carry = design_knowledge_cards.carry_knowledge(["카드A"])
        self.assertEqual(["카드A"], carry["missing"])
        self.assertEqual([], carry["images"])

    # -- ⑧ manifest 기록(journey 08 요약 소스) --------------------------------

    def test_manifest_records_knowledge_carried_summary(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                manifest = self._bundle(run, ["카드A", "존재하지않는카드"])
        slide2 = next(s for s in manifest["slides"] if s.get("n") == 2)
        kc = slide2["knowledge_carried"]
        self.assertEqual(1, kc["cards"])
        self.assertEqual(2, kc["images"])
        self.assertIn("존재하지않는카드", kc["missing"])

    def test_html_only_slide_records_knowledge_carried_none(self):
        """cover 등 html 전용 장은 프롬프트 자체가 없어 운반 대상이 없다(None — 미측정과 구분)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            manifest = self._bundle(run, [])
        slide1 = next(s for s in manifest["slides"] if s.get("n") == 1)
        self.assertEqual("html", slide1.get("render"))
        self.assertIsNone(slide1["knowledge_carried"])

    # -- ⑨ 구조 레퍼런스 — β1 3계층과 같은 지위로 Reference roles에 합류 -----

    def test_structure_reference_appended_to_reference_roles_list(self):
        with tempfile.TemporaryDirectory() as td:
            vault = self._fake_vault(Path(td) / "vault")
            knowhow = self._fake_knowhow(Path(td) / "knowhow")
            run = Path(td) / "run"
            run.mkdir()
            with mock.patch.object(design_knowledge_cards, "_vault_dir", return_value=vault), \
                 mock.patch.object(design_knowledge_cards, "_knowhow_root", return_value=knowhow):
                self._bundle(run, ["카드B"])
                prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertIn("Reference roles:", prompt)
        self.assertIn("구조 레퍼런스", prompt)
        self.assertIn("카드B", prompt)

    # -- ⑩ 회귀 없음: knowledge_cards 미인용 run은 기존 문법 그대로 -----------

    def test_no_knowledge_cards_field_leaves_prompt_unchanged(self):
        """wireframe.json 슬라이드에 knowledge_cards 필드 자체가 없으면(기존 run) 카드 운반
        절이 생기지 않는다 - vault/knowhow에 손도 대지 않는다(회귀 없음)."""
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            self._storyline(run)
            (run / "wireframe.json").write_text(json.dumps({
                "schema_version": 1, "selected_by": "llm:test",
                "slides": [
                    {"slide_id": "2", "message_type": "구조", "frame": "flow_seq", "rendition": "boxed",
                     "slots": [{"piece": "flow_arrow", "size": "wide", "binds": "*"}], "catalog_gap": []},
                ],
            }, ensure_ascii=False), encoding="utf-8")
            contract = design_contract.build(run, brief={}, skins_dir=self.SKINS_DIR)
            design_contract.save(run, contract)
            imagedeck.bundle(run, None, wireframe_mode="on")
            prompt = (run / imagedeck.PROMPTS_DIR / "02.md").read_text(encoding="utf-8")
        self.assertNotIn(design_knowledge_cards.CARRY_HEADER, prompt)
        self.assertNotIn("구조 레퍼런스", prompt)


# ---------------------------------------------------------------------------
# W32-A 패킷 — 시연 1차 마찰 28·29·31 (산출물 오염 차단: 내부 어휘·원시 데이터의 청중 노출)
# ---------------------------------------------------------------------------

class W32FieldsShapeCoerceSmoke(unittest.TestCase):
    """마찰28: 문자열 기대 자리에 들어온 객체를 관용 코어스 + 경고 표면화(원시 dict 노출 금지)."""

    def _deck(self, cases):
        slide = {"slide_id": 1, "role": "portfolio", "template_id": "portfolio_cases",
                 "title": "유사 수행 사례", "key_message": "메시지", "body": [],
                 "fields": {"cases": cases}}
        return {"meta": {"project": "테스트 프로젝트"}, "slides": [slide]}

    def test_dict_case_is_coerced_not_dumped_raw(self):
        cases = [{"name": "복지관 파일럿", "description": "어르신 30명 대상"}]
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html(self._deck(cases), "core", out)
            html = out.read_text(encoding="utf-8")
        self.assertNotIn("{&#x27;name&#x27;", html)   # 원시 dict가 조판되면 안 된다
        self.assertNotIn("&#x27;description&#x27;", html)
        self.assertIn("복지관 파일럿", html)           # 라벨은 살아야 한다
        self.assertIn("어르신 30명 대상", html)        # 상세도 살아야 한다
        self.assertTrue(any("shape 불일치" in w for w in rep["warnings"]),
                        f"코어스는 경고로 표면화돼야 한다: {rep['warnings']}")

    def test_string_cases_unchanged_and_no_warning(self):
        """회귀: 정상 shape(문자열 배열)은 종전과 동일 바이트 + 경고 0."""
        cases = ["복지관 파일럿: 어르신 30명 대상", "지자체 확산: 3개 시군"]
        with tempfile.TemporaryDirectory() as td:
            out1, out2 = Path(td) / "a.html", Path(td) / "b.html"
            rep = htmlgen.render_html(self._deck(cases), "core", out1)
            htmlgen.render_html(self._deck(cases), "core", out2)
            html1 = out1.read_text(encoding="utf-8")
            html2 = out2.read_text(encoding="utf-8")
        self.assertEqual(html1, html2)
        self.assertEqual([], [w for w in rep["warnings"] if "shape 불일치" in w])

    def test_coercion_covers_plain_bullets_too(self):
        """⒞ 전수 점검: cases 전용이 아니라 _esc를 타는 모든 문자열 자리가 덮인다."""
        slide = {"slide_id": 1, "role": "content", "title": "본문", "key_message": "",
                 "body": [{"label": "항목", "text": "설명"}], "fields": {}}
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html({"meta": {"project": "P"}, "slides": [slide]}, "core", out)
            html = out.read_text(encoding="utf-8")
        self.assertNotIn("{&#x27;label&#x27;", html)
        self.assertIn("항목: 설명", html)
        self.assertTrue(any("shape 불일치" in w for w in rep["warnings"]))

    def _steps_slide(self, title):
        return {"slide_id": 1, "role": "프로세스", "template_id": "process_steps",
                "title": title, "key_message": "", "body": [],
                "fields": {"steps": [{"name": "착수", "description": "킥오프 회의"}],
                           "outputs": ["착수보고서"]}}

    def test_layouts_core_fields_already_safe(self):
        """⒞ 전수 점검 결과의 고정: agenda·process_steps는 htmlgen REGISTRY가 아니라 layouts_core가
        렌더하는데, 그쪽은 처음부터 _display/_pick로 dict를 자체 처리하고 있었다(= fields는 이미 안전).
        이 사실을 테스트로 못 박아 둔다 — 나중에 _display가 걷히면 여기서 잡힌다."""
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            htmlgen.render_html({"meta": {"project": "P"}, "slides": [self._steps_slide("추진 절차")]},
                                "core", out)
            html = out.read_text(encoding="utf-8")
        self.assertNotIn("{&#x27;name&#x27;", html)
        self.assertIn("착수", html)
        self.assertIn("킥오프 회의", html)

    def test_layouts_core_head_title_is_coerced(self):
        """layouts_core에 실제로 남아 있던 구멍: 이 모듈은 자체 _head를 갖고 있어 제목·역할·메시지가
        _esc로 직행한다 — dict 제목이면 원시 dict가 크롬에 노출됐다. 코어스 + 경고까지 확인한다
        (레이아웃 모듈은 bare name 로드라 기록 인스턴스 통일도 함께 검증된다)."""
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html(
                {"meta": {"project": "P"},
                 "slides": [self._steps_slide({"name": "추진 절차", "description": "4단계"})]},
                "core", out)
            html = out.read_text(encoding="utf-8")
        self.assertNotIn("{&#x27;name&#x27;", html)
        self.assertIn("추진 절차", html)
        self.assertTrue(any("shape 불일치" in w for w in rep["warnings"]),
                        f"layouts_core 경로의 코어스도 경고로 표면화돼야 한다: {rep['warnings']}")

    def test_prompt_documents_portfolio_cases_shape(self):
        """⒜ 상류 예방: 프롬프트가 portfolio_cases shape를 문자열 배열로 명시한다."""
        import storyline_prompt
        text = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn("portfolio_cases", text)
        self.assertIn("라벨: 설명", text)
        self.assertIn("객체를 넣지 말 것", text)


class W32InternalMarkLeakSmoke(unittest.TestCase):
    """마찰29: message 내 "(axisN 지지)" 내부 표기 — 프롬프트 지시 제거 + 수거 경고."""

    def test_prompt_no_longer_asks_axis_mark_in_message(self):
        src = Path(proposal_pipeline.__file__).read_text(encoding="utf-8")
        self.assertNotIn('그 축의 id를 함께 명시하라', src)
        self.assertIn("message 본문에 축 id를 적지 마라", src)

    def test_collect_surfaces_leftover_marks(self):
        import contextlib
        doc = {"slides": [{"n": 3, "message": "안전한 돌봄을 만든다 (axis1 지지)"},
                          {"n": 4, "message": "정상 메시지"}]}
        buf = io.StringIO()
        with contextlib.redirect_stdout(buf):
            proposal_pipeline._surface_internal_marks(doc)
        out = buf.getvalue()
        self.assertIn("내부 표기", out)
        self.assertIn("장 3", out)
        self.assertNotIn("장 4", out)

    def test_clean_storyline_is_silent(self):
        import contextlib
        buf = io.StringIO()
        with contextlib.redirect_stdout(buf):
            proposal_pipeline._surface_internal_marks({"slides": [{"n": 1, "message": "정상"}]})
        self.assertEqual("", buf.getvalue())


class W32AgendaBadgeChannelSmoke(unittest.TestCase):
    """마찰31: "해소수단 미배정" 배지를 장표에서 걷고 검수 채널(warnings)로 보낸다."""

    def test_unassigned_relief_not_drawn_on_slide(self):
        import compose
        html = compose.p_agenda({"items": [{"title": "사업 이해", "relief": "원인"},
                                           {"title": "추진 전략"}]})
        self.assertNotIn("해소수단 미배정", html)      # 청중 노출 금지
        self.assertNotIn("cp-agenda__badge--none", html)
        self.assertIn("원인", html)                     # 신고된 relief는 종전대로 표기
        self.assertIn("추진 전략", html)                # 항목 자체는 그대로 그린다

    def test_unassigned_relief_goes_to_review_channel_not_warnings(self):
        """검수 노트는 warnings가 아니다 — relief는 선택 필드라 미신고 자체는 결함이 아니고,
        warnings에 섞으면 "warnings=0 = 무결" 계약이 깨진다(전수 렌더 테스트가 그 계약에 의존)."""
        import compose
        warnings = []
        slide = {"slide_id": 7, "frame": "full", "title": "목차", "role": "목차",
                 "slots": [{"piece": "agenda", "data": {"items": [{"title": "사업 이해"}]}}]}
        html = compose.render_slide(slide, warnings)
        notes = compose.drain_notes()
        self.assertNotIn("해소수단 미배정", html)
        self.assertEqual([], warnings)
        self.assertTrue(any("해소수단 미신고" in n for n in notes),
                        f"검수 신호는 별도 채널로 가야 한다: {notes}")
        self.assertEqual([], compose.drain_notes())   # 수거 후 비워진다(다음 슬라이드로 누수 없음)

    def test_render_report_carries_review_notes(self):
        slide = {"slide_id": 1, "frame": "full", "title": "목차", "role": "목차",
                 "slots": [{"piece": "agenda", "data": {"items": [{"title": "사업 이해"}]}}]}
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "deck.html"
            rep = htmlgen.render_html({"meta": {"project": "P"}, "slides": [slide]}, "core", out)
        self.assertEqual([], rep["warnings"])
        self.assertTrue(any("해소수단 미신고" in n for n in rep.get("review_notes") or []),
                        f"리포트에 검수 노트가 실려야 한다: {rep}")


# ---------------------------------------------------------------------------
# W32-D 패킷 — 시연 1차 마찰 30 (지식 원장의 가시성 공백 + 범위 config + 강등)
# ---------------------------------------------------------------------------

class W32FormIntentSmoke(unittest.TestCase):
    """마찰36(기능 제안·가역): form_intent/art_note 통로 — 계약 등재·지식 공급·뼈대 주입·
    프롬프트 명시 구획·작업용 키 유출 차단. 되돌림 조항은 REHEARSAL_FRICTIONS_W31.md 행 36."""

    def test_contract_registers_optional_fields(self):
        import storyline_prompt
        text = storyline_prompt.STORYLINE_SCHEMA_BLOCK
        self.assertIn("form_intent", text)
        self.assertIn("art_note", text)
        # 사슬: 형태 의도 → template_id → required_fields (fields 미기입 공격이 이 줄의 몫)
        self.assertIn("그에 맞는 template_id를 고르고", text)
        # emphasis(사람 전속)와 성격 구분 — LLM 직접 기입 허용을 명시해야 실제로 채워진다
        self.assertIn("직접 적어도 됩니다", text)

    def test_storyline_pull_includes_wireframe_knowledge(self):
        """⑴ 어휘 공급 — 지식이 없으면 형태를 생각할 재료 자체가 프롬프트에 안 들어온다."""
        import knowledge_ledger as kl
        self.assertIn("디자인지식/와이어프레임", kl.pull_folders("storyline"))

    def _deck(self, extra=None):
        s = {"slide_id": 1, "role": "전략", "template_id": "strategy_pillars",
             "title": "추진 전략", "key_message": "메시지", "body": ["항목"],
             "fields": {"pillars": ["A", "B", "C"]}}
        s.update(extra or {})
        return {"meta": {"project": "P"}, "slides": [s]}

    def test_wireframe_prompt_injects_form_intent(self):
        import wireframe
        text = wireframe.build_prompt(self._deck({"form_intent": "3단계 순차 흐름"}))
        self.assertIn("3단계 순차 흐름", text)
        self.assertIn("그대로 베끼지 말 것", text)   # 결정기 독립 판단 계약(오염 신호 가드)

    def test_wireframe_prompt_byte_identical_without_form_intent(self):
        """회귀: form_intent 없는 덱은 종전 프롬프트와 바이트 동일(emphasis 훅과 같은 방식)."""
        import wireframe
        text = wireframe.build_prompt(self._deck())
        self.assertNotIn("form_intent", text)

    def test_image_prompt_filters_working_keys_and_carries_art_note(self):
        """⑷⑸: 작업용 키는 덤프에서 빠지고, art_note는 명시 구획으로 실린다."""
        import imagedeck
        slide = {"n": 3, "title": "장", "bullets": ["내용"],
                 "visual": "구명 형태 메모", "form_intent": "좌우 대비",
                 "art_note": "숫자가 주인공", "evidence": "작업 메모",
                 "supports_axis": "axis1", "deck_class": "content"}
        public = {k: v for k, v in slide.items() if k not in imagedeck._PROMPT_WORKING_KEYS}
        self.assertEqual({"n", "title", "bullets"}, set(public))
        self.assertIn("art_note_block", imagedeck._PROMPT_TEMPLATE)
        self.assertIn("{storyline_json}", imagedeck._PROMPT_TEMPLATE)

    def test_adapter_carries_form_intent_to_deck(self):
        """종단 사슬의 병목: 결정기는 deck.json을 읽으므로 어댑터(allowlist)가 form_intent를
        운반해야 한다 — 여기가 빠지면 storyline에 적혀도 조용히 끊긴다(유닛 우회로는 안 잡히는
        구멍, 실제로 첫 구현에서 놓쳤다). art_note는 미운반이 정답(imagedeck이 storyline 직독)."""
        doc = {"slides": [{"n": 1, "section": "전략", "title": "T", "message": "M",
                           "bullets": ["b"], "form_intent": "좌우 대비",
                           "art_note": "숫자가 주인공"}]}
        deck = _adapt.adapt_storyline(doc, project="P", pack="core")
        s = deck["slides"][0]
        self.assertEqual("좌우 대비", s.get("form_intent"))
        self.assertNotIn("art_note", s)

    def test_material_file_exists_with_split_fields(self):
        """실물 스펙 재료(강의 덱 60장 form/art 분리본) — 없으면 skip(외부 경로)."""
        p = Path(r"<개발 원본 전용 경로>")
        if not p.is_file():
            self.skipTest("재료 파일 부재(외부 경로)")
        d = json.loads(p.read_text(encoding="utf-8"))
        slides = d.get("slides") or []
        self.assertEqual(60, len(slides))
        self.assertEqual(60, sum(1 for s in slides if s.get("form_intent")))


class W32BakedKnowledgeSmoke(unittest.TestCase):
    """마찰30 본체: 코드에 구워진 지식도 원장 뷰에 병기(pull과 구분 표기)."""

    def _run_with_deck(self, pieces):
        td = tempfile.mkdtemp()
        run = Path(td)
        slides = [{"slide_id": i, "slots": [{"piece": p}]} for i, p in enumerate(pieces, 1)]
        (run / "deck.json").write_text(
            json.dumps({"meta": {"project": "P", "pack": "core"}, "slides": slides},
                       ensure_ascii=False), encoding="utf-8")
        return run

    def test_baked_section_lists_used_pieces_with_source(self):
        import knowledge_ledger as kl
        run = self._run_with_deck(["agenda", "pillar_card"])
        view = kl.render_ledger_view(run, {"stages": {}})
        self.assertIn("## 구운 지식", view)   # 낱말이 아니라 절 제목으로 판정(머리말 안내와 구분)
        self.assertIn("`agenda`", view)
        # agenda 조각의 원전 카드가 추적 가능해야 한다(마찰31 배지의 출처)
        self.assertIn("목차는-상대의-두려움-목록이다", view)

    def test_pull_and_baked_are_distinguished(self):
        import knowledge_ledger as kl
        run = self._run_with_deck(["agenda"])
        view = kl.render_ledger_view(run, {"stages": {}})
        self.assertIn("pull 아님", view)   # 채널 구분이 명시돼야 한다(설계상 분리)

    def test_no_pieces_no_section(self):
        """회귀: 조각 조합 덱이 아니면(레거시 template 경로) 구운 지식 절도 없다."""
        import knowledge_ledger as kl
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "deck.json").write_text(
                json.dumps({"meta": {"pack": "core"}, "slides": [{"slide_id": 1}]}),
                encoding="utf-8")
            self.assertNotIn("## 구운 지식", kl.render_ledger_view(run, {"stages": {}}))

    def test_no_deck_is_silent(self):
        import knowledge_ledger as kl
        with tempfile.TemporaryDirectory() as td:
            self.assertEqual([], kl._baked_knowledge_summary(Path(td)))


class W32KnowledgeScopeSmoke(unittest.TestCase):
    """마찰30 별건①: 기획지식 루트 카드가 '1차 범위 밖'으로 오탐되지 않는다(config 확정)."""

    def test_planning_root_is_in_scope(self):
        import knowledge_ledger as kl
        for stage in ("research", "message_map", "storyline", "imagedeck_review"):
            with self.subTest(stage=stage):
                self.assertIn("기획지식", kl.pull_folders(stage))

    def test_prefix_match_covers_subfolders(self):
        """루트를 넣으면 하위 폴더도 범위 안(판정이 prefix 매칭) — 이 확장이 의도된 것임을 고정."""
        import knowledge_ledger as kl
        folders = kl.pull_folders("storyline")
        for folder in ("기획지식", "기획지식/메시지설계", "기획지식/examples"):
            with self.subTest(folder=folder):
                self.assertTrue(any(folder == f or folder.startswith(f + "/") for f in folders))


class W32UnverifiedCardDemotionSmoke(unittest.TestCase):
    """마찰30 별건②: vault 실물이 없는 카드는 강등 표기 + 표면화(조용한 통과 금지)."""

    def test_unverified_card_marked_demoted_in_view(self):
        import knowledge_ledger as kl
        with tempfile.TemporaryDirectory() as td:
            ledger = {"stages": {"wireframe": {
                "cards": ["없는-카드", "있는-카드"], "web": [],
                "card_origins": {"있는-카드": {"folder": "디자인지식/와이어프레임", "in_scope": True},
                                 "없는-카드": {"folder": None, "in_scope": None}}}}}
            view = kl.render_ledger_view(Path(td), ledger)
        self.assertIn("강등", view)
        self.assertIn("유령노트", view)          # 실측 원인 1위를 안내에 남긴다
        self.assertNotIn("있는-카드 ⛔", view)   # 확인된 카드는 강등되지 않는다

    def test_unverified_cards_helper(self):
        import knowledge_ledger as kl
        entry = {"cards": ["a", "b"], "card_origins": {"a": {"folder": "기획지식"}, "b": {"folder": None}}}
        self.assertEqual(["b"], kl._unverified_cards(entry))
        self.assertEqual([], kl._unverified_cards({"cards": [], "card_origins": {}}))


# ---------------------------------------------------------------------------
# W32-C 패킷 — 시연 1차 마찰 35·27 (검토 표면 사용성)
# ---------------------------------------------------------------------------

class W32ViewerSharedSmoke(unittest.TestCase):
    """마찰35: deck.images.html에 화면 맞춤 축소 + 키보드 넘김 — 뷰어 스크립트 공용화."""

    def test_deck_html_nav_bytes_unchanged(self):
        """회귀: deck.html의 내비 블록은 viewer 모듈로 옮겼을 뿐 문자열이 같아야 한다."""
        import viewer
        self.assertIn("addEventListener('keydown'", viewer.NAV_JS)
        self.assertIn("현재 슬라이드 자동 스냅", viewer.NAV_JS)
        self.assertEqual(viewer.NAV_JS, htmlgen._NAV_JS)

    def test_deck_stage_inlined_and_close_token_escaped(self):
        """deck-stage를 인라인한다. **닫는 script 토큰 이스케이프가 핵심** — 원본 주석에 그 토큰이
        있어 그대로 넣으면 HTML 파서가 스크립트를 거기서 끊는다(조용히 깨지는 함정, 실제로 밟았다)."""
        import viewer
        js = viewer.deck_stage_script()
        self.assertTrue(js, "벤더 컴포넌트를 읽어야 한다")
        self.assertIn("customElements.define('deck-stage'", js)
        self.assertNotIn("</" + "script", js)        # 조기 종료 유발 토큰이 남으면 안 된다
        self.assertIn("<\\/" + "script", js)         # 이스케이프된 형태로 보존
        # (`"<\/"`로 쓰면 파이썬이 SyntaxWarning을 뱉는다 — 설치 직후 화면에 경고가 찍히면
        #  통과인데도 실패처럼 보인다. 뜻은 같고 경고만 없앤 표기.)

    def test_stage_wrap_and_sync_wiring(self):
        import viewer
        self.assertIn('width="1920"', viewer.stage_open(1920, 1080))
        close = viewer.stage_close()
        self.assertIn("</deck-stage>", close)
        self.assertIn("BroadcastChannel", close)     # 두 창 동기화 주 채널
        self.assertIn("addEventListener('storage'", close)   # 폴백(브라우저 편차 보험)
        self.assertIn("no-rail", close)              # 전면 창은 레일 끔

    def test_compose_html_carries_stage(self):
        """실제 산출물(deck.images.html)에 컴포넌트가 실린다 — 두 표면의 기능 불일치 해소."""
        import imagedeck
        so, sc, css = imagedeck._stage_wrap(1920, 1080)
        self.assertIn("<deck-stage", so)
        self.assertIn("customElements.define('deck-stage'", sc)
        self.assertIn("deck-stage:not(:defined)", css)   # FOUC 차단

    def test_rasterize_path_has_no_stage(self):
        """회귀: PNG 래스터라이즈 경로는 deck-stage를 쓰면 안 된다 —
        활성 장만 보이고 화면에 맞춰 축소돼 실물 px가 어긋난다."""
        import imagedeck
        self.assertEqual(("", "", ""), imagedeck._stage_wrap(1920, 1080, enabled=False))


class W32PartialProduceStatusSmoke(unittest.TestCase):
    """마찰27: status가 부분 생산 진행률을 표시(판정 로직 불변 — 표시만)."""

    def _run_with(self, n_done, n_total):
        td = tempfile.mkdtemp()
        run = Path(td)
        slides = [{"n": i, "out_name": f"{i:02d}_s.png", "render": "image"}
                  for i in range(1, n_total + 1)]
        (run / "imagedeck_manifest.json").write_text(
            json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8")
        sd = run / "imagedeck" / "slides"
        sd.mkdir(parents=True)
        for s in slides[:n_done]:
            (sd / s["out_name"]).write_bytes(b"x")
        return run

    def test_progress_counts_only_image_slides(self):
        run = self._run_with(5, 20)
        self.assertEqual((5, 20), pipeline_state.imagedeck_progress(run))
        self.assertFalse(pipeline_state.imagedeck_images_present(run))

    def test_no_manifest_is_zero(self):
        with tempfile.TemporaryDirectory() as td:
            self.assertEqual((0, 0), pipeline_state.imagedeck_progress(Path(td)))

    def test_html_only_slides_excluded(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            (run / "imagedeck_manifest.json").write_text(json.dumps({"slides": [
                {"n": 1, "out_name": "01.png", "render": "html"},
                {"n": 2, "out_name": "02.png", "render": "image"}]}), encoding="utf-8")
            (run / "imagedeck" / "slides").mkdir(parents=True)
            self.assertEqual((0, 1), pipeline_state.imagedeck_progress(run))


# ---------------------------------------------------------------------------
# W32-B 패킷 — 시연 1차 마찰 32·34 (데모 진행을 막던 실패 모드)
# ---------------------------------------------------------------------------

class W32CheckboxTypoSmoke(unittest.TestCase):
    """마찰32: 사람 손 편집 채널은 사람의 오타를 전제한다 — `[x ]` 같은 공백 변형 수용."""

    def test_whitespace_variants_are_accepted(self):
        import journey_check
        for text in ("- [x] 검토 완료", "- [X] 검토 완료", "- [x ] 검토 완료",
                     "- [ x] 검토 완료", "- [  X ] 검토 완료"):
            with self.subTest(text=text):
                self.assertEqual((True, False), journey_check._extract_checks(text))

    def test_unchecked_stays_unchecked(self):
        """회귀: 빈 체크는 종전대로 미체크다(관대함이 무조건 통과가 되면 안 된다)."""
        import journey_check
        self.assertEqual((False, False), journey_check._extract_checks("- [ ] 검토 완료"))
        self.assertEqual((False, False), journey_check._extract_checks("- [] 검토 완료"))
        self.assertEqual((False, True), journey_check._extract_checks(
            "- [ ] 검토 완료\n- [x ] 건너뛰기(스킵 가능 관문만 표기)"))

    def test_unknown_mark_is_surfaced_not_ignored(self):
        """⒝ 조용한 무시 제거: 수용 범위 밖 표기는 사람에게 알린다."""
        import journey_check
        marks = journey_check._unknown_marks("- [v] 검토 완료\n- [o] 건너뛰기")
        self.assertEqual(2, len(marks))
        self.assertIn("[v] 검토 완료", marks)
        self.assertEqual([], journey_check._unknown_marks("- [x] 검토 완료\n- [ ] 건너뛰기"))


class W32RejectedArtifactSmoke(unittest.TestCase):
    """마찰34: px-FAIL 산출물을 `.rejected.png`로 개명 — 증거 보존 + skip 고착 해소."""

    def _manifest(self, run):
        (run / "imagedeck").mkdir(parents=True, exist_ok=True)
        (run / "imagedeck" / "slides").mkdir(parents=True, exist_ok=True)
        manifest = {"gen_canvas": {"w": 1792, "h": 784},
                    "slides": [{"n": 3, "out_name": "03_test.png", "render": "image",
                                "expected_px": {"w": 1792, "h": 784},
                                "prompt_file": "imagedeck_prompts/03.md"}]}
        (run / "imagedeck_manifest.json").write_text(
            json.dumps(manifest, ensure_ascii=False), encoding="utf-8")
        (run / "imagedeck_prompts").mkdir(parents=True, exist_ok=True)
        (run / "imagedeck_prompts" / "03.md").write_text("프롬프트", encoding="utf-8")
        return manifest

    def _bad_png(self, path):
        """기대 px와 다른 최소 PNG(휴리스틱 없이 png_dims가 읽을 수 있으면 충분)."""
        try:
            from PIL import Image
        except ImportError:
            self.skipTest("Pillow 없음")
        Image.new("RGB", (1896, 830), "white").save(path)

    def test_failed_slide_is_renamed_and_reproduced_next_run(self):
        import imagedeck
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._manifest(run)
            out = run / "imagedeck" / "slides" / "03_test.png"

            def bad_runner(prompt, meta):
                self._bad_png(out)
                return ""

            rep1 = imagedeck.produce(run, bad_runner)
            self.assertEqual(["03_test.png"], rep1["failed"])
            self.assertFalse(out.exists(), "불량본은 정본 이름으로 남으면 안 된다(skip 고착)")
            self.assertTrue((run / "imagedeck" / "slides" / "03_test.rejected.png").exists(),
                            "증거본은 보존돼야 한다")

            # 재실행 = skip이 아니라 재위임(자기 안내와 일치) + 증거본은 누적 보존
            calls = []

            def bad_runner2(prompt, meta):
                calls.append(meta["slide"])
                self._bad_png(out)
                return ""

            rep2 = imagedeck.produce(run, bad_runner2)
            self.assertEqual([3], calls, f"재실행은 실패분을 다시 위임해야 한다: {rep2}")
            self.assertEqual([], rep2["skipped"])
            self.assertTrue((run / "imagedeck" / "slides" / "03_test.rejected2.png").exists())

    def test_good_slide_still_skips_on_rerun(self):
        """회귀: 정상 산출물은 종전대로 skip(재실행 안전 설계를 깨지 않는다)."""
        import imagedeck
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._manifest(run)
            out = run / "imagedeck" / "slides" / "03_test.png"
            try:
                from PIL import Image
            except ImportError:
                self.skipTest("Pillow 없음")
            Image.new("RGB", (1792, 784), "white").save(out)
            rep = imagedeck.produce(run, lambda p, m: self.fail("정상본은 재위임하면 안 된다"))
            self.assertEqual(["03_test.png"], rep["skipped"])

    def test_rejected_is_excluded_from_counts(self):
        import imagedeck
        self.assertTrue(imagedeck.is_rejected(Path("03_a.rejected.png")))
        self.assertTrue(imagedeck.is_rejected(Path("03_a.rejected2.png")))
        self.assertFalse(imagedeck.is_rejected(Path("03_a.png")))


class W32ManualProduceRouteSmoke(unittest.TestCase):
    """W32 수동 생산 루트 — codex/agy CLI 없는 사용자용 (감지 → 여정 09 가이드 → adopt 수거).

    회귀 방지: ①감지가 실측(shutil.which)이다 ②가이드가 여정 09 폴더에 자기완결로 생성된다
    ③adopt가 파일명(장 번호 시작)으로만 매칭·PNG 변환·정확 px 리사이즈·개명 배치한다(순서 추측
    금지 - unmatched 보고) ④produce는 codex 미감지 시 실패가 아니라 가이드 생성으로 전환한다
    ⑤journey 09 매뉴얼이 가이드 존재 시 수동 루트 안내로 바뀐다."""

    def _manifest(self, run: Path, slides: list[dict]) -> None:
        (run / "imagedeck_prompts").mkdir(parents=True, exist_ok=True)
        for s in slides:
            if s.get("render") != "html":
                pf = f"imagedeck_prompts/{int(s['n']):02d}{s.get('variant') or ''}.md"
                s.setdefault("prompt_file", pf)
                (run / pf).write_text(f"# 프롬프트 장 {s['n']}", encoding="utf-8")
        (run / imagedeck.MANIFEST_NAME).write_text(json.dumps({
            "schema_version": 1, "gen_canvas": {"w": 64, "h": 32}, "slides": slides,
        }, ensure_ascii=False), encoding="utf-8")

    def test_detect_producers_uses_which(self):
        with mock.patch("imagedeck.shutil.which", return_value=None):
            det = imagedeck.detect_producers()
        self.assertEqual({"codex": None, "agy": None}, det)
        with mock.patch("imagedeck.shutil.which", side_effect=lambda n: f"/bin/{n}"):
            det = imagedeck.detect_producers()
        self.assertEqual("/bin/codex", det["codex"])
        self.assertEqual("/bin/agy", det["agy"])

    def test_manual_guide_lands_in_journey_09_and_is_self_contained(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            self._manifest(run, [
                {"n": 1, "out_name": "01_표지.png", "expected_px": {"w": 64, "h": 32}},
                {"n": 2, "out_name": "02_본문.png", "expected_px": {"w": 64, "h": 32}},
                {"n": 3, "render": "html", "deck_class": "toc"},
            ])
            guide = imagedeck.write_manual_guide(run)
            self.assertEqual(run / "journey" / "09_이미지생산" / imagedeck.MANUAL_GUIDE_NAME, guide)
            text = guide.read_text(encoding="utf-8")
        # 자기완결: 최종 파일명·기대 px·adopt 명령·레퍼런스 첨부 지시가 다 있어야 한다.
        self.assertIn("01_표지.png", text)
        self.assertIn("64x32", text)
        self.assertIn("--adopt", text)
        self.assertIn("Reference roles", text)
        self.assertIn("HTML 전용 장 1개", text)
        self.assertIn("없음", text)  # 상태 열 실측(아직 이미지 없음)

    def test_manual_guide_requires_bundle_first(self):
        with tempfile.TemporaryDirectory() as td:
            with self.assertRaises(imagedeck.ImagedeckError):
                imagedeck.write_manual_guide(Path(td))

    @unittest.skipUnless(imagedeck._pil_available(), "Pillow 미설치 - adopt는 Pillow 필수")
    def test_adopt_matches_converts_resizes_and_reports(self):
        from PIL import Image
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            self._manifest(run, [
                {"n": 1, "out_name": "01_표지.png", "expected_px": {"w": 64, "h": 32}},
                {"n": 5, "variant": "A", "out_name": "05A_실험.png", "expected_px": {"w": 40, "h": 20}},
            ])
            src = Path(td) / "downloads"
            src.mkdir()
            # 크기·포맷이 제각각인 다운로드 산출물: 01=jpg 대형, 05A=png, 무관 파일 1개.
            Image.new("RGB", (200, 100), "white").save(src / "1.jpg")
            Image.new("RGB", (30, 30), "black").save(src / "05A.png")
            (src / "image (3).png").write_bytes(b"not-a-real-target")
            rep = imagedeck.adopt(run, src)
            self.assertEqual(["01_표지.png", "05A_실험.png"], sorted(rep["adopted"]))
            self.assertEqual(["image (3).png"], rep["unmatched"])  # 순서 추측 배치 금지
            self.assertEqual([], rep["missing"])
            # 계약 이행 실측: PNG + 정확 px.
            self.assertEqual((64, 32), imagedeck.png_dims(run / imagedeck.SLIDES_DIR / "01_표지.png"))
            self.assertEqual((40, 20), imagedeck.png_dims(run / imagedeck.SLIDES_DIR / "05A_실험.png"))
            # 재실행 = 교체(replaced) — 사람이 폴더에 넣은 명시 행동이라 skip이 아니다.
            rep2 = imagedeck.adopt(run, src)
            self.assertEqual([], rep2["adopted"])
            self.assertEqual(["01_표지.png", "05A_실험.png"], sorted(rep2["replaced"]))
            # adopt가 가이드 상태 열을 최신화한다.
            text = (run / "journey" / "09_이미지생산" / imagedeck.MANUAL_GUIDE_NAME).read_text(encoding="utf-8")
        self.assertIn("| OK |", text)

    def test_adopt_without_pillow_raises_clear_error(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td) / "run"
            run.mkdir()
            self._manifest(run, [{"n": 1, "out_name": "01_x.png", "expected_px": {"w": 8, "h": 8}}])
            src = Path(td) / "dl"
            src.mkdir()
            with mock.patch.dict(_sys.modules, {"PIL": None}):
                with self.assertRaises(imagedeck.ImagedeckError) as ctx:
                    imagedeck.adopt(run, src)
        self.assertIn("Pillow", str(ctx.exception))

    def test_produce_without_codex_switches_to_manual_route(self):
        # CLI 층: produce가 실패(1)가 아니라 가이드 생성 + 종료 0으로 **경로 전환**한다.
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            pipeline_state.init(run, mode="direct", input_kind="bid", input_ref="T")
            self._manifest(run, [{"n": 1, "out_name": "01_x.png", "expected_px": {"w": 8, "h": 8}}])
            args = SimpleNamespace(run=str(run), bundle=False, produce=True, collect=False,
                                   compose=False, only="", timeout=1)
            with mock.patch("imagedeck.shutil.which", return_value=None):
                with mock.patch.object(proposal_pipeline, "_render_run_dir", return_value=run):
                    rc = proposal_pipeline.imagedeck_cmd(args)
            self.assertEqual(0, rc)
            guide = run / "journey" / "09_이미지생산" / imagedeck.MANUAL_GUIDE_NAME
            self.assertTrue(guide.is_file())
            # 전환 사실이 state에 남는다(문서 고고학 방지 - status가 답하는 관례).
            state = pipeline_state.load(run)
            self.assertIn("imagedeck_manual_guide", state.get("stages") or {})

    def test_journey_09_manual_switches_wording_when_guide_exists(self):
        with tempfile.TemporaryDirectory() as td:
            run = Path(td)
            slides = run / "imagedeck" / "slides"
            slides.mkdir(parents=True)
            (slides / "01_x.png").write_bytes(b"png")
            folder = journey_folders.folder_path(run, "09")
            journey_folders._populate_09(run, folder, {})
            base = (folder / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
            self.assertNotIn("수동 루트", base)
            self.assertIn("W32", base)  # 기본 매뉴얼도 수동 전환 존재를 고지
            (folder / journey_folders.MANUAL_GUIDE_NAME_09).write_text("가이드", encoding="utf-8")
            journey_folders._populate_09(run, folder, {})
            manual = (folder / journey_folders.MANUAL_NAME).read_text(encoding="utf-8")
        self.assertIn("수동 루트", manual)
        self.assertIn("--adopt", manual)


class BrowserDependencyProbeSmoke(unittest.TestCase):
    """W32 마찰37 회귀 방지 — chromium 부재를 '가용'으로 오판하면 스모크가 빨간불로 끝난다.

    사고 경로: requirements.txt가 playwright **패키지**를 설치하므로, 판정이 `import playwright`
    뿐이면 브라우저 바이너리를 안 받은 사람(README가 '선택'이라 적어 둔 그대로)에게
    available()=True를 돌려주고 launch에서 터진다. 강의 35장(설치 테이크)이 이 화면을 찍는다.
    """

    def test_probe_reports_missing_chromium_as_unavailable(self):
        import rasterize
        done = SimpleNamespace(returncode=0, stdout=b"0")   # 자식 프로세스: 실행 파일 없음
        with mock.patch.object(rasterize, "_BROWSER_PROBE", None), \
             mock.patch.object(rasterize.subprocess, "run", return_value=done):
            ok, reason = rasterize.probe()
        self.assertFalse(ok)
        self.assertIn("chromium", reason)
        self.assertIn(rasterize.CHROMIUM_INSTALL_HINT, reason)   # 다음 한 줄을 사람 말로 준다

    def test_probe_reports_present_chromium_as_available(self):
        import rasterize
        done = SimpleNamespace(returncode=0, stdout=b"1")
        with mock.patch.object(rasterize, "_BROWSER_PROBE", None), \
             mock.patch.object(rasterize.subprocess, "run", return_value=done):
            ok, _reason = rasterize.probe()
            self.assertTrue(ok)
            self.assertEqual("", rasterize.unavailable_reason())   # 가용하면 사유 없음

    def test_layout_probe_reuses_rasterize_verdict_no_second_judgement(self):
        """layout_probe는 두 번째 판정을 하지 않는다 — rasterize가 False면 무조건 False."""
        import layout_probe, rasterize
        with mock.patch.object(rasterize, "available", return_value=False):
            self.assertFalse(layout_probe.available())

    def test_probe_is_silent_on_stderr(self):
        """playwright 드라이버의 종료 잔여 로그를 삼킨다 — 통과인데 실패처럼 보이면 안 된다."""
        import rasterize
        captured = {}

        def fake_run(cmd, **kw):
            captured.update(kw)
            return SimpleNamespace(returncode=0, stdout=b"1")

        with mock.patch.object(rasterize, "_BROWSER_PROBE", None), \
             mock.patch.object(rasterize.subprocess, "run", side_effect=fake_run):
            rasterize.probe()
        self.assertEqual(rasterize.subprocess.DEVNULL, captured.get("stderr"))


if __name__ == "__main__":
    unittest.main(verbosity=2)
